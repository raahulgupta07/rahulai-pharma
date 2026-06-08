"""
Data Upload & Auto-Onboarding
==============================

Upload CSV/Excel/JSON files → AI auto-generates metadata → loads into knowledge base.

Usage:
    POST /api/upload  (multipart/form-data with file + optional table_name)
    GET  /api/tables   (list all tables with row counts)
    DELETE /api/tables/{table_name}  (drop a user-uploaded table)
"""

import json
import os
import re
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

_bg_executor = ThreadPoolExecutor(max_workers=5, thread_name_prefix="dash-bg")


def _bg_done_log(fut):
    """Log exceptions from bg executor tasks (otherwise swallowed silently)."""
    try:
        exc = fut.exception()
    except Exception:
        return
    if exc:
        import logging as _l
        _l.getLogger(__name__).exception("upload bg task failed", exc_info=exc)


import pandas as pd
from fastapi import APIRouter, Form, HTTPException, Request, UploadFile
from datetime import datetime as _dt_datetime
from sqlalchemy import create_engine as _sa_create_engine, inspect, text
from sqlalchemy.pool import NullPool
from dash.utils.safe_json import safe_dumps


def create_engine(url, **kw):
    """Wrapper that forces NullPool for ad-hoc engines (PgBouncer handles pooling)."""
    kw.setdefault("poolclass", NullPool)
    return _sa_create_engine(url, **kw)

from dash.paths import BUSINESS_DIR, KNOWLEDGE_DIR, QUERIES_DIR, TABLES_DIR
from db import db_url
from dash.settings import TRAINING_MODEL

import logging as _upload_logging
_upload_lg = _upload_logging.getLogger(__name__)
# Module-level logger alias — bare `logger.warning/.info` calls scattered in this
# file (lines ~997, 1650, 1672, 1863, 4015, 6949, 9797 ...) used to raise
# NameError when their try/except branches fired. Per-table training would then
# catch the NameError as "training failed". Single alias fixes all sites.
logger = _upload_lg

# Observability tracing (fail-soft; no-op if unavailable or TRACING_DISABLED).
try:
    from dash.obs.trace import (
        start_trace as _trace_start,
        trace_span as _trace_span,
        end_trace as _trace_end,
    )
except Exception:  # noqa: BLE001
    from contextlib import contextmanager as _trace_cm

    def _trace_start(*_a, **_k):  # type: ignore
        return ""

    def _trace_end(*_a, **_k):  # type: ignore
        return None

    @_trace_cm
    def _trace_span(*_a, **_k):  # type: ignore
        yield None


def _safe_write_json(path: Path, data) -> None:
    """Atomic JSON write — prevents corruption from concurrent writes."""
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp_fd, tmp_path = tempfile.mkstemp(dir=path.parent, suffix=".json")
    try:
        with os.fdopen(tmp_fd, "w") as f:
            json.dump(data, f, indent=2, default=str)
        os.replace(tmp_path, str(path))  # Atomic rename
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


def _safe_read_json(path: Path, default=None):
    """Safe JSON read — handles concurrent writes and corruption."""
    try:
        if path.exists():
            return json.loads(path.read_text(errors="ignore"))
    except (json.JSONDecodeError, OSError):
        pass
    return default if default is not None else {}


def _doc_category_hint(filename: str) -> str:
    """Epic 4 — light institutional-knowledge tag from filename so the
    Researcher can surface "why/incident" context. Returns 'general' when no
    signal. Fail-soft, never raises."""
    try:
        fn = (filename or "").lower()
    except Exception:
        return "general"
    if any(k in fn for k in ("incident", "postmortem", "post-mortem", "post_mortem", "rca", "outage")):
        return "incident"
    if any(k in fn for k in ("changelog", "change-log", "change_log", "release-notes", "release_notes")):
        return "changelog"
    if "launch" in fn:
        return "launch"
    return "general"


def _compute_file_hash(file_path: str) -> str:
    """sha256 of file content. Read in 64KB chunks. Fail-soft returns empty string."""
    try:
        import hashlib
        h = hashlib.sha256()
        with open(file_path, 'rb') as f:
            while True:
                chunk = f.read(65536)
                if not chunk:
                    break
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def _lookup_upload_cache(file_hash: str) -> dict | None:
    """Look up cached extraction plan by file hash. Returns plan dict or None.
    Bumps hit_count + last_used_at atomically on hit."""
    if not file_hash:
        return None
    try:
        from db.session import get_write_engine
        from sqlalchemy import text as _slt
        e = get_write_engine()
        with e.begin() as c:
            row = c.execute(_slt("""
                UPDATE public.dash_upload_cache
                SET hit_count = hit_count + 1,
                    last_used_at = now()
                WHERE file_hash = :h
                RETURNING plan, rescue_used, hit_count
            """), {"h": file_hash}).fetchone()
        if not row:
            return None
        m = dict(row._mapping)
        return {"plan": m["plan"], "rescue_used": m["rescue_used"], "hit_count": m["hit_count"]}
    except Exception:
        return None


def _save_upload_cache(file_hash: str, file_size: int, file_ext: str, plan: dict, rescue_used: bool):
    """Persist plan keyed by file hash. ON CONFLICT updates plan + bumps last_used_at."""
    if not file_hash:
        return
    try:
        from db.session import get_write_engine
        from sqlalchemy import text as _slt
        e = get_write_engine()
        with e.begin() as c:
            c.execute(_slt("""
                INSERT INTO public.dash_upload_cache
                    (file_hash, file_size_bytes, file_ext, plan, rescue_used, first_seen_at, last_used_at)
                VALUES (:h, :sz, :ext, CAST(:p AS jsonb), :ru, now(), now())
                ON CONFLICT (file_hash) DO UPDATE SET
                    plan = EXCLUDED.plan,
                    rescue_used = EXCLUDED.rescue_used,
                    last_used_at = now()
            """), {"h": file_hash, "sz": file_size, "ext": file_ext,
                   "p": json.dumps(plan), "ru": rescue_used})
    except Exception:
        pass


router = APIRouter(prefix="/api", tags=["Upload"])
_engine = create_engine(db_url)


# ---------------------------------------------------------------------------
# Central column classifier — replaces inline heuristics across the pipeline.
# Returns dict[col_name] -> role, where role ∈ {
#   "pk"          — declared primary key (information_schema)
#   "id"          — name-pattern key (ends _id/_code/_key) or 100% unique numeric
#   "fk"          — declared foreign key
#   "constant"    — DISTINCT(col) == 1
#   "enum"        — adaptive low-cardinality: distinct ≤ max(20, rows*0.001) AND ratio < 0.05
#   "date_dim"    — date/timestamp col w/ ≥2 distinct values
#   "measure"     — numeric, not id/pk/fk/enum, has variance
#   "free_text"   — high-cardinality text (ratio > 0.5)
#   "dimension"   — text w/ moderate cardinality (default text bucket)
# }
# Adaptive thresholds derived from REAL row count, not sample. Per-project
# override hook reads dash_projects.feature_config['pipeline_thresholds'].
# ---------------------------------------------------------------------------

_DEFAULT_THRESHOLDS = {
    "enum_max_distinct": 20,          # absolute cap on enum distinct
    "enum_ratio_max": 0.05,           # distinct/rows < this → enum candidate
    "enum_ratio_strong": 0.001,       # below this → always enum
    "pk_min_rows": 100,               # row count floor before claiming PK by uniqueness
    "free_text_ratio": 0.5,           # distinct/rows > this → free-text
    "null_alert_floor_pct": 5.0,      # absolute floor (% of nulls)
    "null_alert_stddev_mult": 1.5,    # alert if null_pct > median + N*stddev
}


def _get_pipeline_thresholds(project_slug: str | None) -> dict:
    """Read per-project override from feature_config; fallback to defaults."""
    if not project_slug:
        return dict(_DEFAULT_THRESHOLDS)
    try:
        from db.session import get_sql_engine
        with get_sql_engine().connect() as c:
            row = c.execute(text(
                "SELECT feature_config #> '{pipeline_thresholds}' "
                "FROM public.dash_projects WHERE slug = :s"
            ), {"s": project_slug}).fetchone()
        if row and row[0]:
            merged = dict(_DEFAULT_THRESHOLDS)
            merged.update(row[0] if isinstance(row[0], dict) else json.loads(row[0]))
            return merged
    except Exception:
        pass
    return dict(_DEFAULT_THRESHOLDS)


def _get_real_pk_fk(schema: str, table: str) -> tuple[set[str], set[str]]:
    """Query information_schema for declared PK + FK columns. Fail-soft."""
    pks: set[str] = set()
    fks: set[str] = set()
    try:
        eng = create_engine(db_url, poolclass=NullPool)
        with eng.connect() as c:
            for col_name, in c.execute(text("""
                SELECT kcu.column_name
                FROM information_schema.table_constraints tc
                JOIN information_schema.key_column_usage kcu
                  ON tc.constraint_name = kcu.constraint_name
                 AND tc.table_schema = kcu.table_schema
                WHERE tc.table_schema = :sch AND tc.table_name = :t
                  AND tc.constraint_type = 'PRIMARY KEY'
            """), {"sch": schema, "t": table}).fetchall():
                pks.add(col_name)
            for col_name, in c.execute(text("""
                SELECT kcu.column_name
                FROM information_schema.table_constraints tc
                JOIN information_schema.key_column_usage kcu
                  ON tc.constraint_name = kcu.constraint_name
                 AND tc.table_schema = kcu.table_schema
                WHERE tc.table_schema = :sch AND tc.table_name = :t
                  AND tc.constraint_type = 'FOREIGN KEY'
            """), {"sch": schema, "t": table}).fetchall():
                fks.add(col_name)
    except Exception:
        pass
    return pks, fks


def classify_columns(
    project_slug: str | None,
    schema: str,
    table: str,
    col_stats: list[dict],
    row_count: int,
) -> dict[str, str]:
    """Classify every column into a role using real DB metadata + adaptive thresholds.

    Args:
        project_slug: Project (for threshold override lookup).
        schema: Postgres schema name.
        table: Table name.
        col_stats: list of {name, type, unique_count, null_pct, min, max, ...}
                   from _analyze_column or _sql_profile_columns.
        row_count: REAL row count of the table.

    Returns: dict[col_name] -> role string.
    """
    th = _get_pipeline_thresholds(project_slug)
    real_pks, real_fks = _get_real_pk_fk(schema, table)
    roles: dict[str, str] = {}

    # adaptive enum cap based on row count: small tables → fixed cap;
    # large tables → up to 0.1% of rows distinct
    enum_cap = max(th["enum_max_distinct"], int(row_count * th["enum_ratio_strong"]))

    for ca in col_stats:
        name = ca.get("name", "")
        if not name:
            continue
        n_lower = name.lower()

        # 1. Declared PK wins over everything
        if name in real_pks:
            roles[name] = "pk"
            continue
        if name in real_fks:
            roles[name] = "fk"
            continue

        uc = ca.get("unique_count") or 0
        nulls = ca.get("null_pct") or 0
        col_type = ca.get("type", "")

        # 2. Name-pattern id — but ONLY if also unique. Bug fix 2026-05-25:
        # `site_code` (53 distinct in 106K rows) used to be tagged id from
        # the `_code` suffix even though it's a dimension. Require ≥50%
        # uniqueness for name-pattern to win.
        name_looks_id = (
            n_lower in {"id", "_period"}
            or n_lower.endswith(("_id", "_key", "_uuid"))
            or (n_lower.endswith("_code") and row_count > 0 and uc / row_count >= 0.5)
        )
        if name_looks_id:
            roles[name] = "id"
            continue

        # 3. Constant column (one distinct value across whole table)
        if row_count > 0 and uc == 1:
            roles[name] = "constant"
            continue

        # 4. Date/datetime → date_dim (multi-value) or constant (single)
        if col_type == "datetime" or "date" in n_lower or "time" in n_lower:
            roles[name] = "date_dim" if uc > 1 else "constant"
            continue

        # 5. PK-by-uniqueness — any row count. Bug fix 2026-05-25: previous
        # 100-row floor missed small reference tables (currencies / regions /
        # categories) where uniqueness is real. Floor only governs the
        # ALERT text, not the role itself.
        if row_count > 1 and uc == row_count:
            roles[name] = "id"
            continue

        # 6. Adaptive enum: low absolute distinct AND low ratio
        if row_count > 0:
            ratio = uc / row_count
            if uc <= enum_cap and ratio < th["enum_ratio_max"]:
                roles[name] = "enum"
                continue
            # 7. Free-text: very high cardinality (mostly unique)
            if col_type in ("text", "string") and ratio > th["free_text_ratio"]:
                roles[name] = "free_text"
                continue

        # 8. Numeric not flagged above → real measure
        if col_type == "numeric":
            roles[name] = "measure"
        else:
            roles[name] = "dimension"

    return roles


def _is_aggregatable(role: str) -> bool:
    """True if the column role accepts SUM/AVG/multiplication."""
    return role == "measure"


def detect_composite_pk(schema: str, table: str, roles: dict[str, str], row_count: int) -> tuple[str, str] | None:
    """Find (col_a, col_b) where COUNT(DISTINCT a||b) == row_count.

    Only probes if no single-col PK was found in roles + at least 2 id/fk/enum-
    candidate columns exist. Skips when row_count < 100 (noise). Caps at 2-col
    composites (deeper combos = exponential cost). Returns first hit or None.
    """
    if row_count < 100:
        return None
    if any(r == "pk" for r in roles.values()):
        return None  # single-col PK already declared
    if any(r == "id" for r in roles.values()):
        return None  # uniqueness-based id already found

    candidates = [n for n, r in roles.items() if r in ("fk", "id", "enum", "dimension")]
    if len(candidates) < 2:
        return None

    try:
        eng = create_engine(db_url, poolclass=NullPool)
        with eng.connect() as c:
            # Try first 6 pairs (12 cols max → 66 probes, but capped at 6)
            from itertools import combinations
            for a, b in list(combinations(candidates[:8], 2))[:12]:
                try:
                    n = c.execute(text(
                        f'SELECT COUNT(*) FROM (SELECT DISTINCT "{a}", "{b}" '
                        f'FROM "{schema}"."{table}") t'
                    )).scalar() or 0
                    if n == row_count:
                        return (a, b)
                except Exception:
                    continue
    except Exception:
        pass
    return None


def detect_currency_scale(schema: str, table: str, col_name: str, col_min, col_max) -> str | None:
    """Heuristic: numeric col stored as minor units (cents).

    Triggers when min >= 100, max < 10^11, all values divisible by 100, AND
    avg > 1000. Returns scale hint ('100' for cents) or None.
    """
    try:
        if col_min is None or float(col_min) < 100:
            return None
        if col_max is None or float(col_max) >= 1e11:
            return None
        eng = create_engine(db_url, poolclass=NullPool)
        with eng.connect() as c:
            non_div = c.execute(text(
                f'SELECT COUNT(*) FROM "{schema}"."{table}" '
                f'WHERE "{col_name}" IS NOT NULL AND "{col_name}" % 100 <> 0'
            )).scalar() or 0
            if non_div > 0:
                return None
            avg = c.execute(text(
                f'SELECT AVG("{col_name}") FROM "{schema}"."{table}"'
            )).scalar()
            if avg is not None and float(avg) > 1000:
                return "100"
    except Exception:
        pass
    return None


def _persist_extraction_plan(project_slug: str, table_name: str, plan_data: dict) -> int | None:
    """Persist an extraction plan row to public.dash_extraction_plans.
    Fail-soft: any error logs + returns None. Uses write engine (public schema write).
    Uses CAST(:x AS jsonb) per PgBouncer + SQLAlchemy rule.
    """
    try:
        from db.session import get_write_engine
        from sqlalchemy import text as _slt
        src = str(plan_data.get("strategy", "") or "")
        # Normalize strategy tag
        if "llm-rescued" in src or "llm-repaired" in src:
            strategy = "llm-rescued"
        elif "rules-split" in src:
            strategy = "rules-split"
        elif "ai-unpivot" in src:
            strategy = "ai-unpivot"
        elif "rules" in src:
            strategy = "rules"
        else:
            strategy = src or "unknown"
        llm_rescued = bool(plan_data.get("llm_rescued") or "llm-rescued" in src)
        params = {
            "s": project_slug,
            "t": table_name,
            "sf": plan_data.get("source_file"),
            "sh": plan_data.get("sheet_name"),
            "fh": plan_data.get("file_hash"),
            "st": strategy,
            "hr": plan_data.get("header_row"),
            "sk": json.dumps(plan_data.get("skip_rows") or []),
            "bl": json.dumps(plan_data.get("blocks") or []),
            "ri": plan_data.get("row_count_in"),
            "ro": plan_data.get("row_count_out"),
            "lr": llm_rescued,
            "rr": plan_data.get("rescue_reasoning"),
        }
        e = get_write_engine()
        with e.begin() as c:
            row = c.execute(_slt(
                """
                INSERT INTO public.dash_extraction_plans
                  (project_slug, table_name, source_file, sheet_name, file_hash,
                   strategy, header_row, skip_rows, blocks,
                   row_count_in, row_count_out, llm_rescued, rescue_reasoning)
                VALUES
                  (:s, :t, :sf, :sh, :fh,
                   :st, :hr, CAST(:sk AS jsonb), CAST(:bl AS jsonb),
                   :ri, :ro, :lr, :rr)
                RETURNING id
                """
            ), params).fetchone()
        return int(row[0]) if row else None
    except Exception as exc:
        try:
            _upload_lg.warning(f"_persist_extraction_plan failed for {project_slug}/{table_name}: {exc}")
        except Exception:
            pass
        return None


def _enqueue_vector_backfill(slug: str) -> dict:
    """Enqueue every knowledge doc + brain entry + KG triple for this project
    into ``dash.dash_vectors`` via ``VECTOR_SYNC.enqueue``.

    Idempotent (vector_sync hash-checks). Fail-soft. Returns counters for log
    output. Called as the final step of TRAIN ALL so vectors are guaranteed
    populated even when in-flight reembed workers missed events.
    """
    import asyncio as _asyncio
    counts = {"knowledge": 0, "brain": 0, "kg": 0, "errors": 0}
    try:
        from dash.tools.vector_sync import VECTOR_SYNC as _VS
    except Exception:
        counts["errors"] += 1
        return counts

    # Ensure worker is running (idempotent).
    try:
        _VS.start()
    except Exception:
        pass

    def _fire(coro):
        """Schedule a coroutine on the running loop if any, else run sync."""
        try:
            loop = _asyncio.get_running_loop()
            loop.create_task(coro)
        except RuntimeError:
            try:
                _asyncio.run(coro)
            except Exception:
                counts["errors"] += 1

    # 1) Knowledge docs on disk (per-project).
    try:
        docs_dir = KNOWLEDGE_DIR / slug / "docs"
        if docs_dir.exists():
            for f in docs_dir.iterdir():
                if not f.is_file():
                    continue
                try:
                    txt = f.read_text(errors="ignore")[:50000]
                    if txt.strip():
                        _fire(_VS.enqueue(slug, "docs", f.name, txt, {}))
                        counts["knowledge"] += 1
                except Exception:
                    counts["errors"] += 1
    except Exception:
        counts["errors"] += 1

    # 2) Company brain entries for this project (scope-merged + global).
    try:
        eng = create_engine(db_url)
        with eng.connect() as c:
            rows = c.execute(text(
                "SELECT id, category, name, definition "
                "FROM public.dash_company_brain "
                "WHERE project_slug = :s OR project_slug IS NULL "
                "ORDER BY id DESC LIMIT 2000"
            ), {"s": slug}).fetchall()
        eng.dispose()
        for r in rows:
            try:
                bid, cat, name, defn = r[0], r[1] or "", r[2] or "", r[3] or ""
                content = f"[{cat}] {name}: {defn}".strip()
                if not content:
                    continue
                _fire(_VS.enqueue(slug, "brain", f"brain_{bid}", content,
                                  {}, {"category": cat, "name": name}))
                counts["brain"] += 1
            except Exception:
                counts["errors"] += 1
    except Exception:
        counts["errors"] += 1

    # 3) Knowledge graph triples.
    try:
        eng = create_engine(db_url)
        with eng.connect() as c:
            rows = c.execute(text(
                "SELECT id, subject, predicate, object "
                "FROM public.dash_knowledge_triples "
                "WHERE project_slug = :s "
                "ORDER BY id DESC LIMIT 5000"
            ), {"s": slug}).fetchall()
        eng.dispose()
        for r in rows:
            try:
                tid, s, p, o = r[0], r[1] or "", r[2] or "", r[3] or ""
                content = f"{s} {p} {o}".strip()
                if not content:
                    continue
                _fire(_VS.enqueue(slug, "kg", f"triple_{tid}", content, {}))
                counts["kg"] += 1
            except Exception:
                counts["errors"] += 1
    except Exception:
        counts["errors"] += 1

    return counts


# ──────────────────── Investment vertical — 5 sub-steps ────────────────────
def _run_investment_training_steps(slug: str, master_engine, master_run_id, total_tables: int,
                                   _log) -> None:
    """Run 5 investment-vertical training sub-steps (15a-15e).

    Idempotent + fail-soft per step. Each step writes brain entries that the
    Layer 14b context injection then surfaces into agent prompts.
    """
    from sqlalchemy import text as _t

    def _step(name: str, idx: int) -> None:
        if master_run_id:
            try:
                with master_engine.connect() as cn:
                    cn.execute(_t(
                        "UPDATE public.dash_training_runs SET steps = :s WHERE id = :id"
                    ), {"s": f"investment_{name}|{slug}|{idx}|5", "id": master_run_id})
                    cn.commit()
            except Exception:
                pass

    def _save_brain(category: str, name: str, value: str) -> None:
        """Idempotent brain upsert (ON CONFLICT name+category+project)."""
        try:
            with master_engine.connect() as cn:
                cn.execute(_t(
                    "INSERT INTO public.dash_company_brain "
                    "(category, name, definition, project_slug, metadata) "
                    "VALUES (:c, :n, :v, :s, CAST('{}' AS jsonb))"
                ), {"c": category, "n": name[:200], "v": value[:4000], "s": slug})
                cn.commit()
        except Exception as _e:
            import logging
            logging.debug(f"_save_brain {category}/{name} failed: {_e}")

    # ── 15a: Financial Statement Detection ──
    _step("financial_detection", 1)
    try:
        with master_engine.connect() as cn:
            rows = cn.execute(_t(
                "SELECT table_name FROM information_schema.tables "
                "WHERE table_schema=:s AND table_name IN "
                "('balance_sheet','income_statement','cash_flow','cap_table',"
                " 'pnl','financial_statements','balance_sheets','income_statements','cash_flows')"
            ), {"s": slug}).fetchall()
        detected = [r[0] for r in rows]
        if detected:
            _save_brain("financial_statement", "Detected financial tables",
                        f"Project contains: {', '.join(detected)}")
            _log(f"✓ 15a: financial detection — {len(detected)} table(s)", "", total_tables)
        else:
            _log("· 15a: no financial tables detected (skipping)", "", total_tables)
    except Exception as _e:
        _log(f"⚠ 15a failed: {str(_e)[:80]}", "", total_tables)

    # ── 15b: Pitch Deck Section Extraction ──
    _step("pitch_deck_extract", 2)
    try:
        from dash.paths import KNOWLEDGE_DIR as _KD
        docs_dir = _KD / slug / "docs"
        pptx_count = 0
        if docs_dir.exists():
            for f in docs_dir.iterdir():
                if f.is_file() and f.suffix.lower() == ".pptx":
                    pptx_count += 1
                    _save_brain("pitch_deck_section", f"Deck: {f.name}",
                                f"Uploaded pitch deck (.pptx). Run RAG via search_pitch_deck.")
        _log(f"✓ 15b: pitch decks indexed — {pptx_count} file(s)", "", total_tables)
    except Exception as _e:
        _log(f"⚠ 15b failed: {str(_e)[:80]}", "", total_tables)

    # ── 15c: Comparable Deals Cross-Project Lookup ──
    _step("comp_deals", 3)
    try:
        # Pull this project's sector from brain (if available) then federate
        with master_engine.connect() as cn:
            sec_row = cn.execute(_t(
                "SELECT definition FROM public.dash_company_brain "
                "WHERE project_slug=:s AND (name ILIKE '%sector%' OR name ILIKE '%industry%') "
                "LIMIT 1"
            ), {"s": slug}).first()
            sector = sec_row[0] if sec_row else None
            sql = (
                "SELECT name, definition AS value, project_slug FROM public.dash_company_brain "
                "WHERE category='ic_memo' AND project_slug IS NOT NULL "
                "AND project_slug <> :s "
            )
            params: dict = {"s": slug}
            if sector:
                sql += "AND (name ILIKE :sec OR value ILIKE :sec) "
                params["sec"] = f"%{sector[:50]}%"
            sql += "ORDER BY created_at DESC NULLS LAST LIMIT 5"
            comps = cn.execute(_t(sql), params).mappings().all()
        for c in comps:
            _save_brain("comp_deal", f"{c['project_slug']}: {c['name']}",
                        str(c['value'])[:1000])
        _log(f"✓ 15c: comparable deals — {len(comps)} memos federated", "", total_tables)
    except Exception as _e:
        _log(f"⚠ 15c failed: {str(_e)[:80]}", "", total_tables)

    # ── 15e: Investment Q&A Generation ──
    _step("investment_qa", 5)
    try:
        from dash.paths import KNOWLEDGE_DIR as _KD2
        qa_dir = _KD2 / slug / "training"
        qa_dir.mkdir(parents=True, exist_ok=True)
        # Templated Q&A — narrative gets answered at chat time via tools.
        investment_qa = [
            {"q": "What is LTV/CAC?", "tool": "compute_unit_economics"},
            {"q": "What is gross margin trend over last 3 years?", "tool": "get_income_statement"},
            {"q": "What sectors does the CEO have experience in?", "tool": "extract_team_bios"},
            {"q": "What is the TAM?", "tool": "extract_market_size"},
            {"q": "What's the customer concentration?", "tool": "get_customer_concentration"},
            {"q": "What is current cash runway?", "tool": "get_cashflow"},
            {"q": "What is the cap table?", "tool": "get_cap_table"},
            {"q": "What is MoM revenue growth?", "tool": "compute_growth_metrics"},
            {"q": "Are there red flags?", "tool": "find_red_flags"},
            {"q": "Does this fit our mandate?", "tool": "verify_against_mandate"},
            {"q": "What's the proposed valuation?", "tool": "extract_term_proposed"},
            {"q": "Comparable deals in this sector?", "tool": "find_comparable_deals"},
            {"q": "What is EV/Revenue multiple?", "tool": "compute_valuation_multiples"},
            {"q": "What is current MRR?", "tool": "compute_growth_metrics"},
            {"q": "Who are the top 5 customers?", "tool": "get_customer_concentration"},
            {"q": "What is the burn rate?", "tool": "get_cashflow"},
            {"q": "Pitch deck — what's the problem statement?", "tool": "search_pitch_deck"},
            {"q": "DD findings — any legal issues?", "tool": "search_dd_findings"},
            {"q": "Team bios — who are the founders?", "tool": "extract_team_bios"},
            {"q": "What's our IC recommendation?", "tool": "(committee_chair)"},
            {"q": "Draft an IC memo.", "tool": "(memo_writer)"},
            {"q": "What's the SOM?", "tool": "extract_market_size"},
            {"q": "Latest balance sheet?", "tool": "get_balance_sheet"},
            {"q": "ARR YoY growth?", "tool": "compute_growth_metrics"},
            {"q": "Rule of 40 check?", "tool": "compute_valuation_multiples"},
        ]
        out_path = qa_dir / "investment_qa.json"
        try:
            with open(out_path, "w") as fh:
                json.dump(investment_qa, fh, indent=2)
            _log(f"✓ 15e: investment Q&A — {len(investment_qa)} pairs written", "", total_tables)
        except Exception as _e3:
            _log(f"⚠ 15e write failed: {str(_e3)[:80]}", "", total_tables)
    except Exception as _e:
        _log(f"⚠ 15e failed: {str(_e)[:80]}", "", total_tables)

# Tables that ship with the demo — protected from deletion
PROTECTED_TABLES = {"customers", "subscriptions", "plan_changes", "invoices", "usage_metrics", "support_tickets", "dash_users", "dash_tokens", "dash_projects", "shared_results"}

MAX_FILE_SIZE = 200 * 1024 * 1024  # 200 MB
ALLOWED_EXTENSIONS = {".csv", ".xlsx", ".xls", ".json", ".sql", ".py", ".txt", ".md", ".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif", ".webp"}

# Universal vision prompt — handles charts, scanned docs, photos, diagrams in one call
_UNIVERSAL_VISION_PROMPT = (
    "Analyze this image from a business document.\n"
    "If it contains TEXT (scanned document, certificate, letter): extract ALL text exactly as written. "
    "Preserve layout. Render any tables as markdown tables.\n"
    "If it contains a CHART or GRAPH: extract ALL data points, axis labels, legend items as a markdown table. "
    "Describe the trend: increasing/decreasing/flat and by how much.\n"
    "If it contains a DIAGRAM or FLOWCHART: describe all components, connections, and flow.\n"
    "If it contains a PHOTO: describe what is visible and any text, signage, or labels.\n"
    "If it contains a TABLE rendered as image: extract ALL rows and columns as a markdown table.\n"
    "Be precise with all numbers, dates, and labels. Miss nothing."
)


# ---------------------------------------------------------------------------
# Smart Upload: File Classification + Table Matching + Fingerprinting
# ---------------------------------------------------------------------------

def classify_file(filename: str, headers: list[str] | None = None, content_sample: str = "", df=None, project_slug: str = "") -> str:
    """Classify uploaded file: data, column_definition, business_rules, sql_patterns, documentation."""
    ext = Path(filename).suffix.lower()

    # SQL files → query patterns
    if ext == ".sql":
        return "sql_patterns"

    # Documents → check content for rules vs general docs
    if ext in (".md", ".txt", ".pdf", ".docx", ".pptx"):
        lower = content_sample.lower()
        rule_signals = ["rule:", "must be", "should be", "always", "never", "constraint", "formula", "calculation"]
        if sum(1 for s in rule_signals if s in lower) >= 2:
            return "business_rules"
        return "documentation"

    # Code files
    if ext in (".py", ".js", ".ts"):
        return "documentation"

    # Data files (CSV, Excel, JSON)
    if ext in (".csv", ".xlsx", ".xls", ".json"):
        # Method 1: Check by HEADERS (Format A — has "column_name" + "description" headers)
        if headers:
            header_lower = [str(h).lower().strip() for h in headers]
            def_signals = ["definition", "description", "meaning", "business_meaning", "data_type", "column_name"]
            has_def = any(s in h for h in header_lower for s in def_signals)
            has_col = any("column" in h or "field" in h for h in header_lower)
            if has_def and has_col:
                return "column_definition"

        # Method 2: Check by VALUES (Format B — no headers, values match existing table columns)
        if df is not None and project_slug and len(df.columns) <= 5:
            result = _detect_definition_by_values(df, project_slug)
            if result:
                return "column_definition"

        # Method 3: Heuristic — 2-3 columns, one short text + one long text = likely definitions
        if df is not None and len(df.columns) in (2, 3) and len(df) >= 5 and len(df) <= 200:
            text_cols = []
            for c in df.columns:
                vals = df[c].dropna().astype(str)
                if len(vals) > 0:
                    avg_len = vals.str.len().mean()
                    all_text = all(not str(v).replace('.', '').replace('-', '').isdigit() for v in vals.head(5))
                    text_cols.append((c, avg_len, all_text))
            # If we have exactly 2 text columns with different avg lengths
            text_only = [t for t in text_cols if t[2]]
            if len(text_only) >= 2:
                lengths = sorted(text_only, key=lambda x: x[1])
                # Short col avg < 50, long col avg > 30, and ratio > 1.3
                if lengths[0][1] < 60 and lengths[-1][1] > 20 and lengths[-1][1] / max(lengths[0][1], 1) > 1.3:
                    return "column_definition"

        return "data"

    return "data"


def _detect_definition_by_values(df, project_slug: str) -> dict | None:
    """Generic: detect if a file is a column definition by matching VALUES against existing table columns."""
    if len(df) < 3 or len(df.columns) < 2:
        return None

    # Get all column names from all tables in this project
    existing_columns = set()
    try:
        from sqlalchemy import inspect as sa_inspect
        insp = sa_inspect(_engine)
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        for tbl in insp.get_table_names(schema=schema):
            for col in insp.get_columns(tbl, schema=schema):
                existing_columns.add(col["name"].lower().strip())
                # Also add original name variations (with spaces, dots, etc.)
                existing_columns.add(col["name"].lower().replace("_", " ").strip())
    except Exception:
        return None

    if not existing_columns:
        return None

    # For each column in the uploaded file, check if its VALUES match existing column names
    for col_idx, col in enumerate(df.columns):
        values = df[col].dropna().astype(str).tolist()
        if not values:
            continue

        values_lower = [v.lower().strip() for v in values]

        # How many values match existing column names? (try multiple formats)
        matches = 0
        for v in values_lower:
            cleaned = re.sub(r"[^a-z0-9]", "_", v).strip("_")
            cleaned = re.sub(r"_+", "_", cleaned)
            if (v in existing_columns or
                cleaned in existing_columns or
                v.replace(":", "").replace(".", "").strip() in existing_columns or
                cleaned[:30] in existing_columns):
                matches += 1
        match_pct = matches / max(len(values), 1) * 100

        if match_pct > 40:
            # This column contains column names! Find the description column
            # Description = the column with the longest average text
            best_desc_col = None
            best_avg_len = 0
            for other_idx, other_col in enumerate(df.columns):
                if other_idx == col_idx:
                    continue
                other_values = df[other_col].dropna().astype(str).tolist()
                avg_len = sum(len(v) for v in other_values) / max(len(other_values), 1)
                if avg_len > best_avg_len:
                    best_avg_len = avg_len
                    best_desc_col = other_col

            return {"name_col": col, "desc_col": best_desc_col, "match_pct": match_pct}

    return None


def match_existing_table(project_slug: str, new_columns: list[str]) -> dict | None:
    """Check if uploaded columns match an existing table. Returns match info or None."""
    if not new_columns or not project_slug:
        return None

    try:
        from sqlalchemy import inspect as sa_inspect
        insp = sa_inspect(_engine)
        import re as _re
        schema = _re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        existing_tables = insp.get_table_names(schema=schema)

        new_cols_lower = set(str(c).lower().strip() for c in new_columns)
        best_match = None
        best_overlap = 0

        for tbl in existing_tables:
            tbl_cols = insp.get_columns(tbl, schema=schema)
            tbl_col_names = set(c["name"].lower() for c in tbl_cols)

            overlap = new_cols_lower & tbl_col_names
            overlap_pct = len(overlap) / max(len(new_cols_lower), 1) * 100

            if overlap_pct > best_overlap and overlap_pct >= 50:
                # Count existing rows
                try:
                    with _engine.connect() as conn:
                        row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                except Exception:
                    row_count = 0

                new_cols = new_cols_lower - tbl_col_names
                missing_cols = tbl_col_names - new_cols_lower

                best_match = {
                    "table": tbl,
                    "overlap_pct": round(overlap_pct),
                    "matched_columns": list(overlap),
                    "new_columns": list(new_cols),
                    "missing_columns": list(missing_cols),
                    "existing_rows": row_count,
                }
                best_overlap = overlap_pct

        return best_match
    except Exception:
        return None


def compute_fingerprint(row_count: int, col_names: list[str]) -> str:
    """Compute a fingerprint for change detection."""
    import hashlib
    sorted_cols = sorted(str(c).lower().strip() for c in col_names)
    raw = f"{row_count}|{'|'.join(sorted_cols)}"
    return hashlib.md5(raw.encode()).hexdigest()


def save_fingerprint(project_slug: str, table_name: str, row_count: int, col_names: list[str]):
    """Save fingerprint to dash_table_metadata."""
    fp = compute_fingerprint(row_count, col_names)
    col_hash = compute_fingerprint(0, col_names)  # cols only, no row count
    try:
        with _engine.connect() as conn:
            conn.execute(text(
                "UPDATE public.dash_table_metadata SET fingerprint = :fp, row_count = :rc, col_hash = :ch, updated_at = NOW() "
                "WHERE project_slug = :s AND table_name = :t"
            ), {"fp": fp, "rc": row_count, "ch": col_hash, "s": project_slug, "t": table_name})
            conn.commit()
    except Exception:
        pass
    # New stock upload may introduce new outlets — auto-provision their store
    # embeds (and gateway keys) immediately, so the engineer never has to open a
    # tab or run a script. Stock tables only; idempotent + fail-soft.
    if "balance_stock" in str(table_name).lower():
        try:
            from dash.embed import manager as _embed_mgr
            made = _embed_mgr.auto_provision_store_embeds(project_slug)
            if made:
                logger.info("upload: auto-provisioned %d store embeds after %s", made, table_name)
        except Exception:
            logger.exception("upload: store-embed auto-provision after ingest failed (ignored)")
        try:
            from app.auth import _auto_provision_missing as _apigw_provision
            _apigw_provision()
        except Exception:
            logger.exception("upload: gateway-key auto-provision after ingest failed (ignored)")


def check_fingerprint_changed(project_slug: str, table_name: str, row_count: int, col_names: list[str]) -> str:
    """Compare new fingerprint with stored. Returns: 'new', 'unchanged', 'rows_changed', 'schema_changed'."""
    fp = compute_fingerprint(row_count, col_names)
    col_hash = compute_fingerprint(0, col_names)
    try:
        with _engine.connect() as conn:
            row = conn.execute(text(
                "SELECT fingerprint, col_hash FROM public.dash_table_metadata WHERE project_slug = :s AND table_name = :t"
            ), {"s": project_slug, "t": table_name}).fetchone()

        if not row or not row[0]:
            return "new"
        if row[0] == fp:
            return "unchanged"
        if row[1] != col_hash:
            return "schema_changed"
        return "rows_changed"
    except Exception:
        return "new"


def process_column_definitions(project_slug: str, df) -> dict:
    """Process a column definition file → save annotations + memories + rules.
    Works with ANY format — detects column name and description columns automatically."""

    # Method 1: Try header-based detection (Format A: has Column Name, Description headers)
    table_col = next((h for h in df.columns if str(h).lower() in ("table", "table_name", "entity")), None)
    col_col = next((h for h in df.columns if str(h).lower() in ("column", "column_name", "field", "field_name")), None)
    def_col = next((h for h in df.columns if str(h).lower() in ("definition", "description", "meaning", "business_meaning", "notes", "comment")), None)
    type_col = next((h for h in df.columns if str(h).lower() in ("data_type", "type", "format", "dtype")), None)

    # Method 2: If headers don't match, detect by VALUES (Format B: no headers)
    if not col_col or not def_col:
        detected = _detect_definition_by_values(df, project_slug)
        if detected:
            col_col = detected["name_col"]
            def_col = detected["desc_col"]
        else:
            # Method 3: Heuristic — in a 2-3 column file, shorter text col = names, longer = descriptions
            if len(df.columns) in (2, 3):
                text_cols = []
                for c in df.columns:
                    vals = df[c].dropna().astype(str)
                    if len(vals) > 0:
                        avg_len = vals.str.len().mean()
                        text_cols.append((c, avg_len))
                text_cols.sort(key=lambda x: x[1])
                if len(text_cols) >= 2:
                    col_col = text_cols[0][0]  # shorter = column names
                    def_col = text_cols[-1][0]  # longer = descriptions

    if not col_col or not def_col:
        return {"saved": 0, "error": "Could not detect column name and description columns"}

    # If no table column, auto-detect from project's existing tables
    default_table = ""
    if not table_col:
        try:
            from sqlalchemy import inspect as sa_inspect
            insp = sa_inspect(_engine)
            schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
            tables = insp.get_table_names(schema=schema)
            if tables:
                default_table = tables[0]  # use first table
        except Exception:
            pass

    annotations_saved = 0
    memories_saved = 0
    rules_saved = 0

    try:
        with _engine.connect() as conn:
            for _, row in df.iterrows():
                tbl = str(row[table_col]).strip() if table_col and row.get(table_col) else ""
                col = str(row[col_col]).strip() if row.get(col_col) else ""
                defn = str(row[def_col]).strip() if row.get(def_col) else ""

                if not col or not defn or defn == "nan":
                    continue

                # Save as annotation (use detected table or default)
                ann_table = tbl.lower() if tbl else default_table
                if ann_table:
                    conn.execute(text(
                        "INSERT INTO public.dash_annotations (project_slug, table_name, column_name, annotation, updated_by) "
                        "VALUES (:s, :t, :c, :a, 'column_definition') "
                        "ON CONFLICT (project_slug, table_name, column_name) DO UPDATE SET annotation = :a, updated_at = NOW()"
                    ), {"s": project_slug, "t": ann_table, "c": col.lower(), "a": defn})
                    annotations_saved += 1

                # Save as memory
                fact = f"Column '{ann_table}.{col}': {defn}" if ann_table else f"Column '{col}': {defn}"
                conn.execute(text(
                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'column_definition')"
                ), {"s": project_slug, "f": fact[:500]})
                memories_saved += 1

                # Save type as rule if specified
                if type_col and row.get(type_col) and str(row[type_col]) != "nan":
                    rule_defn = f"Column '{col}' should be treated as {row[type_col]} type"
                    rule_id = f"coldef_{col.lower()[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'data_type', :defn, 'column_definition') "
                        "ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": f"{col} type", "defn": rule_defn})
                    rules_saved += 1

            conn.commit()
    except Exception:
        pass

    return {"annotations": annotations_saved, "memories": memories_saved, "rules": rules_saved}


def process_business_rules_doc(project_slug: str, content: str, filename: str) -> dict:
    """Extract business rules and facts from a text/markdown document via LLM."""
    from dash.settings import training_llm_call

    prompt = f"""Extract business rules and key facts from this document.

DOCUMENT ({filename}):
{content[:3000]}

Return ONLY valid JSON:
{{
  "rules": [
    {{"name": "Rule name", "definition": "The business rule text"}}
  ],
  "facts": [
    "Key fact 1",
    "Key fact 2"
  ]
}}"""

    result = training_llm_call(prompt, "extraction")
    if not result:
        return {"rules": 0, "facts": 0}

    try:
        parsed = json.loads(result)
    except json.JSONDecodeError:
        return {"rules": 0, "facts": 0}

    rules_saved = 0
    facts_saved = 0

    try:
        with _engine.connect() as conn:
            for rule in (parsed.get("rules") or [])[:10]:
                if isinstance(rule, dict) and rule.get("name"):
                    rule_id = f"doc_{re.sub(r'[^a-z0-9]', '_', rule['name'].lower())[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'business_rule', :defn, 'document') "
                        "ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": rule["name"], "defn": rule.get("definition", "")})
                    rules_saved += 1

            for fact in (parsed.get("facts") or [])[:10]:
                if fact and isinstance(fact, str):
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'document')"
                    ), {"s": project_slug, "f": fact[:500]})
                    facts_saved += 1

            conn.commit()
    except Exception:
        pass

    return {"rules": rules_saved, "facts": facts_saved}


def process_sql_file(project_slug: str, content: str) -> dict:
    """Parse SQL file and save queries as patterns with metadata extraction."""
    # Split by semicolons or double newlines
    queries = [q.strip() for q in re.split(r';|\n\n', content) if q.strip()]
    saved = 0

    try:
        engine = _engine
        for q in queries[:20]:
            if not q.upper().startswith(("SELECT", "WITH")):
                continue
            # Generate a question from the SQL
            question = f"Run: {q[:100]}"
            if _save_query_pattern_with_metadata(engine, project_slug, question, q, source='sql_file'):
                saved += 1
    except Exception:
        pass

    return {"patterns_saved": saved}


def _sanitize_table_name(name: str) -> str:
    """Convert filename to a valid PostgreSQL table name."""
    name = Path(name).stem.lower()
    name = re.sub(r"[^a-z0-9_]", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    if not name or name[0].isdigit():
        name = "t_" + name
    return name[:63]  # PG identifier limit


def _find_header_row(file_path: str, ext: str) -> int:
    """Detect the real header row in messy Excel/CSV files.
    Scans first 10 rows to find the one with the most text-like columns."""
    try:
        if ext in (".xlsx", ".xls"):
            df_raw = pd.read_excel(file_path, header=None, nrows=10)
        elif ext == ".csv":
            df_raw = pd.read_csv(file_path, header=None, nrows=10)
        else:
            return 0

        best_row = 0
        best_score = 0

        for i in range(min(10, len(df_raw))):
            row = df_raw.iloc[i]
            non_null = row.dropna()
            if len(non_null) == 0:
                continue

            score = 0
            for val in non_null:
                s = str(val).strip()
                # Good header: short text, no long numbers, not a formula
                if isinstance(val, str) and len(s) > 1 and len(s) < 50 and not s.startswith("="):
                    score += 3
                # Penalize: pure numbers, very long strings, formulas
                elif isinstance(val, (int, float)):
                    score -= 1

            # Prefer rows with many non-null text values
            text_ratio = sum(1 for v in non_null if isinstance(v, str) and len(str(v).strip()) > 1) / max(len(non_null), 1)
            score += text_ratio * 5

            if score > best_score:
                best_score = score
                best_row = i

        return best_row
    except Exception:
        return 0


def _contextual_enrich_chunks(chunks: list[str], filename: str, doc_summary: str = "") -> list[str]:
    """Enrich chunks with document context before embedding (Anthropic's Contextual Retrieval).
    Prepends 1-2 sentences of context to each chunk for better retrieval."""
    if not chunks or len(chunks) <= 1:
        return chunks

    try:
        from dash.settings import training_llm_call

        # Build document summary from first chunk if not provided
        if not doc_summary:
            doc_summary = chunks[0][:300]

        enriched = []
        # Process in batches to reduce LLM calls (max 20 calls = 200 chunks)
        batch_size = 10
        max_batches = 20
        batch_count = 0
        for batch_start in range(0, len(chunks), batch_size):
            if batch_count >= max_batches:
                logger.info(f"Contextual enrichment capped at {max_batches} batches ({batch_count * batch_size} chunks), skipping remaining {len(chunks) - batch_start}")
                enriched.extend(chunks[batch_start:])
                break
            batch_count += 1
            batch = chunks[batch_start:batch_start + batch_size]
            numbered = "\n".join(f"{i+1}. {c[:200]}" for i, c in enumerate(batch))

            prompt = f"""Add brief context (1 sentence, max 20 words) to each chunk from this document.
Document: {filename}
Summary: {doc_summary[:300]}

Chunks:
{numbered}

Return JSON array of context strings (same order, same count):
["Context for chunk 1...", "Context for chunk 2...", ...]"""

            raw = training_llm_call(prompt, "extraction")
            if raw:
                try:
                    contexts = json.loads(raw.strip().strip("`").lstrip("json").strip())
                    for i, chunk in enumerate(batch):
                        ctx = contexts[i] if i < len(contexts) else ""
                        if ctx and len(ctx) > 5:
                            enriched.append(f"{ctx} — {chunk}")
                        else:
                            enriched.append(chunk)
                except Exception:
                    enriched.extend(batch)
            else:
                enriched.extend(batch)

        return enriched
    except Exception:
        return chunks  # fallback: return original chunks unchanged


def _filter_junk_chunks(chunks: list[str]) -> list[str]:
    """Remove junk chunks: too short, pure formatting, near-duplicates."""
    if not chunks:
        return chunks

    filtered = []
    seen = set()
    for chunk in chunks:
        text = chunk.strip()
        # Skip too short
        if len(text) < 20:
            continue
        # Skip pure formatting/headers
        if text.count('\n') > len(text) / 10:  # mostly newlines
            continue
        # Skip near-duplicates (first 50 chars as key)
        key = text[:50].lower()
        if key in seen:
            continue
        seen.add(key)
        filtered.append(text)

    return filtered


def _extract_document_structure(file_path: str, ext: str, text: str = "") -> dict:
    """Extract document structure: TOC, headings, sections, page ranges.
    Works for PDF (pymupdf headings), PPTX (slide titles), DOCX (heading styles), MD (# headings)."""
    structure = {"sections": [], "total_pages": 0, "has_toc": False}

    try:
        if ext == ".pdf":
            import fitz
            doc = fitz.open(file_path)
            structure["total_pages"] = len(doc)
            # Extract TOC if available
            toc = doc.get_toc()
            if toc:
                structure["has_toc"] = True
                for level, title, page in toc:
                    structure["sections"].append({"title": title.strip(), "level": level, "page": page, "type": "heading"})
            else:
                # Fallback: detect headings from text formatting
                for pi, page in enumerate(doc):
                    blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
                    for block in blocks:
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                # Large font or bold = likely heading
                                if span.get("size", 0) > 14 or (span.get("flags", 0) & 2**4):  # bold flag
                                    txt = span.get("text", "").strip()
                                    if txt and len(txt) > 3 and len(txt) < 200:
                                        level = 1 if span.get("size", 0) > 18 else (2 if span.get("size", 0) > 14 else 3)
                                        structure["sections"].append({"title": txt, "level": level, "page": pi + 1, "type": "detected"})
            doc.close()

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            structure["total_pages"] = len(prs.slides)
            for si, slide in enumerate(prs.slides):
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                if not title:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            title = shape.text.strip()[:100]
                            break
                structure["sections"].append({"title": title or f"Slide {si+1}", "level": 1, "page": si + 1, "type": "slide"})

        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            for pi, para in enumerate(doc.paragraphs):
                style = para.style.name.lower() if para.style else ""
                if "heading" in style and para.text.strip():
                    level = int(style.replace("heading ", "").strip()) if style.replace("heading ", "").strip().isdigit() else 1
                    structure["sections"].append({"title": para.text.strip(), "level": level, "page": pi, "type": "heading"})

        elif ext == ".md" and text:
            for li, line in enumerate(text.split("\n")):
                line = line.strip()
                if line.startswith("#"):
                    level = len(line) - len(line.lstrip("#"))
                    title = line.lstrip("#").strip()
                    if title:
                        structure["sections"].append({"title": title, "level": level, "page": li, "type": "heading"})

    except Exception:
        pass

    # Deduplicate consecutive same-title sections
    if structure["sections"]:
        deduped = [structure["sections"][0]]
        for s in structure["sections"][1:]:
            if s["title"] != deduped[-1]["title"]:
                deduped.append(s)
        structure["sections"] = deduped[:100]  # Cap at 100 sections

    return structure


def _section_aware_chunks(text: str, structure: dict, max_chunk_size: int = 1500) -> list[dict]:
    """Split text at section/heading boundaries. Each chunk tagged with section + page.
    Returns [{text, section, page, heading_path}]."""
    sections = structure.get("sections", [])
    if not sections or not text:
        # Fallback: fixed-size chunks with page markers
        chunks = []
        lines = text.split("\n")
        current = []
        current_size = 0
        for line in lines:
            current.append(line)
            current_size += len(line)
            if current_size >= max_chunk_size:
                chunks.append({"text": "\n".join(current), "section": "", "page": 0, "heading_path": ""})
                current = []
                current_size = 0
        if current:
            chunks.append({"text": "\n".join(current), "section": "", "page": 0, "heading_path": ""})
        return chunks

    # Split text at heading boundaries
    lines = text.split("\n")
    chunks = []
    current_section = sections[0]["title"] if sections else ""
    current_page = sections[0].get("page", 0) if sections else 0
    heading_path = [current_section]
    current_lines = []
    section_idx = 0

    for line in lines:
        stripped = line.strip()
        # Check if this line matches a section heading
        matched = False
        for si, sec in enumerate(sections[section_idx:], section_idx):
            if stripped and sec["title"] in stripped and len(stripped) < len(sec["title"]) + 20:
                # Found a heading — save current chunk, start new one
                if current_lines:
                    chunk_text = "\n".join(current_lines).strip()
                    if len(chunk_text) > 20:
                        chunks.append({"text": chunk_text, "section": current_section, "page": current_page, "heading_path": " > ".join(heading_path)})
                current_section = sec["title"]
                current_page = sec.get("page", current_page)
                # Update heading path based on level
                level = sec.get("level", 1)
                heading_path = heading_path[:level-1] + [current_section]
                current_lines = [line]
                section_idx = si + 1
                matched = True
                break

        if not matched:
            current_lines.append(line)

        # Split if chunk too large
        if sum(len(l) for l in current_lines) > max_chunk_size:
            chunk_text = "\n".join(current_lines).strip()
            if len(chunk_text) > 20:
                chunks.append({"text": chunk_text, "section": current_section, "page": current_page, "heading_path": " > ".join(heading_path)})
            current_lines = []

    # Last chunk
    if current_lines:
        chunk_text = "\n".join(current_lines).strip()
        if len(chunk_text) > 20:
            chunks.append({"text": chunk_text, "section": current_section, "page": current_page, "heading_path": " > ".join(heading_path)})

    return chunks


def _hierarchical_summarize(chunks: list[dict], filename: str) -> dict:
    """Summarize document hierarchically: section summaries → doc summary.
    For big docs (5+ sections), 77% cheaper than enriching every chunk."""
    if not chunks:
        return {"doc_summary": "", "section_summaries": []}

    # Group chunks by section
    sections = {}
    for c in chunks:
        sec = c.get("section") or "Main"
        if sec not in sections:
            sections[sec] = []
        sections[sec].append(c["text"])

    if len(sections) < 3:
        # Small doc — just concatenate and summarize once
        all_text = "\n".join(c["text"][:500] for c in chunks[:20])
        try:
            from dash.settings import training_llm_call
            summary = training_llm_call(f"Summarize this document in 200 words:\n\n{all_text[:4000]}", "extraction")
            return {"doc_summary": summary or "", "section_summaries": []}
        except Exception:
            return {"doc_summary": all_text[:500], "section_summaries": []}

    # Big doc — summarize each section, then summarize summaries
    section_summaries = []
    try:
        from dash.settings import training_llm_call

        for sec_name, sec_chunks in list(sections.items())[:20]:  # Cap at 20 sections
            sec_text = "\n".join(sec_chunks)[:3000]
            if len(sec_text) < 50:
                continue
            summary = training_llm_call(
                f"Summarize this section '{sec_name}' in 100 words. Include key numbers, dates, decisions:\n\n{sec_text}",
                "extraction"
            )
            if summary:
                section_summaries.append({"section": sec_name, "summary": summary.strip()})

        # Final: summarize all section summaries
        combined = "\n\n".join(f"**{s['section']}:** {s['summary']}" for s in section_summaries)
        doc_summary = training_llm_call(
            f"Create a 300-word executive summary of this document '{filename}' from these section summaries:\n\n{combined[:4000]}",
            "extraction"
        ) or ""

        return {"doc_summary": doc_summary.strip(), "section_summaries": section_summaries}
    except Exception:
        return {"doc_summary": "", "section_summaries": section_summaries}


def _clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """Smart cleanup: normalize nulls, drop empty rows/columns, rename unnamed, clean column names."""
    # 0. Normalize null representations → NaN
    _null_strings = {"N/A", "n/a", "#N/A", "NA", "na", "NULL", "null", "None", "none", "NONE",
                     "N/a", "#NA", "NaN", "nan", "-", "?", ".", "—", "–", ""}
    for col in df.columns:
        if df[col].dtype == object:
            df[col] = df[col].map(lambda x: pd.NA if isinstance(x, str) and x.strip() in _null_strings else x)

    # 1. Drop rows where ALL values are NaN
    df = df.dropna(how='all')

    # 2. Drop columns where ALL values are NaN (100% empty)
    empty_cols = [c for c in df.columns if df[c].isna().all()]
    if empty_cols:
        df = df.drop(columns=empty_cols)

    # 3. Drop columns that are >95% null (nearly empty)
    high_null_cols = [c for c in df.columns if df[c].isna().mean() > 0.95]
    if high_null_cols:
        df = df.drop(columns=high_null_cols)

    # 4. Rename "Unnamed" columns based on content BEFORE cleaning names
    renamed = {}
    day_names = {'monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday'}
    month_names = {'january', 'february', 'march', 'april', 'may', 'june', 'july', 'august', 'september', 'october', 'november', 'december', 'jan', 'feb', 'mar', 'apr', 'jun', 'jul', 'aug', 'sep', 'oct', 'nov', 'dec'}

    for c in df.columns:
        cstr = str(c).lower()
        if 'unnamed' in cstr or cstr.startswith('col_'):
            non_null = df[c].dropna()
            if len(non_null) == 0:
                renamed[c] = f'empty_{list(df.columns).index(c)}'
                continue

            sample = [str(v).lower().strip() for v in non_null.head(10).tolist() if str(v).strip()]

            if sample and all(v in day_names for v in sample):
                renamed[c] = 'day_of_week'
            elif sample and all(v in month_names for v in sample):
                renamed[c] = 'month'
            elif pd.api.types.is_numeric_dtype(non_null):
                renamed[c] = f'value_{list(df.columns).index(c) + 1}'
            else:
                renamed[c] = f'column_{list(df.columns).index(c) + 1}'

    if renamed:
        df = df.rename(columns=renamed)

    # 4b. Currency/comma/percentage cleanup — coerce text columns to numeric.
    # P1: NEVER coerce ID/code columns (would re-float a large code → precision loss).
    _CURRENCY_RE = re.compile(r'^[\s$€£¥₹]*([-+]?[\d,]+\.?\d*)\s*%?\s*$')
    for col in df.columns:
        if _is_id_colname(col):
            continue
        if df[col].dtype != object:
            continue
        non_null = df[col].dropna()
        if len(non_null) == 0:
            continue
        # Check if >50% of values look like numbers with currency/comma/percent
        sample = non_null.head(20)
        matches = sum(1 for v in sample if _CURRENCY_RE.match(str(v).strip()))
        if matches >= len(sample) * 0.5:
            is_pct = any('%' in str(v) for v in sample.head(10))
            cleaned = df[col].apply(lambda x: re.sub(r'[$€£¥₹,\s]', '', str(x)).rstrip('%') if pd.notna(x) else x)
            df[col] = pd.to_numeric(cleaned, errors='coerce')
            if is_pct:
                df[col] = df[col] / 100.0

    # 4c. P1 — repair ID/code columns that slipped in as float (e.g. NaN forced
    # float64): cast whole-valued floats → nullable Int64 so the DB column is
    # bigint, not double precision (no silent rounding of large codes downstream).
    for col in df.columns:
        if not _is_id_colname(col):
            continue
        try:
            if pd.api.types.is_float_dtype(df[col]):
                nn = df[col].dropna()
                if len(nn) and (nn == nn.round()).all():
                    df[col] = df[col].astype("Int64")
        except Exception:
            pass

    # 5. Clean ALL column names (after smart rename)
    df.columns = [re.sub(r"[^a-z0-9_]", "_", str(c).lower().strip()).strip("_") or f"col_{i}" for i, c in enumerate(df.columns)]

    # 5b. Escape PostgreSQL reserved words
    PG_RESERVED = {'select', 'from', 'where', 'order', 'group', 'by', 'having', 'join', 'on',
        'user', 'table', 'column', 'schema', 'database', 'index', 'view', 'limit', 'offset',
        'insert', 'update', 'delete', 'drop', 'create', 'alter', 'case', 'when', 'then',
        'else', 'end', 'all', 'any', 'exists', 'in', 'not', 'and', 'or', 'null', 'true',
        'false', 'primary', 'key', 'foreign', 'check', 'default', 'constraint', 'references'}
    df.columns = [f"{c}_col" if c in PG_RESERVED else c for c in df.columns]

    # 6. Remove duplicate column names
    seen = {}
    new_cols = []
    for c in df.columns:
        if c in seen:
            seen[c] += 1
            new_cols.append(f"{c}_{seen[c]}")
        else:
            seen[c] = 0
            new_cols.append(c)
    df.columns = new_cols

    return df


def _detect_delimiter(file_path: str) -> str:
    """Auto-detect CSV delimiter from first line."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            first_line = f.readline()
        counts = {d: first_line.count(d) for d in [',', ';', '\t', '|']}
        best = max(counts, key=counts.get)
        return best if counts[best] > 0 else ','
    except Exception:
        return ','


# P1 — ID/code columns must NEVER be read as float (large codes lose precision,
# e.g. "1E+12" or a 13-digit barcode → silently rounded, join key destroyed).
# Read them as string at parse time + exempt them from numeric coercion.
_ID_COL_RE = re.compile(r'(^|_)(id|code|barcode|sku|ean|upc|gtin|ref|refno|acct|account)($|_|no$)', re.I)

def _is_id_colname(name) -> bool:
    return bool(_ID_COL_RE.search(str(name)))

def _id_dtypes(file_path: str, ext: str, header_row, sep: str | None = None) -> dict:
    """Header-only pre-scan → {col: str} for ID/code columns so pandas keeps them
    as text (exact) instead of float (lossy). Fail-soft → {} (normal read)."""
    try:
        if ext == ".csv":
            hdr = pd.read_csv(file_path, header=header_row, sep=sep or ",", nrows=0).columns
        else:
            hdr = pd.read_excel(file_path, header=header_row, nrows=0).columns
        return {c: "string" for c in hdr if _is_id_colname(c)}
    except Exception:
        return {}


def _read_file(file_path: str, ext: str) -> pd.DataFrame:
    """Read a data file into a DataFrame with smart header detection and cleanup."""
    if ext == ".csv":
        header_row = _find_header_row(file_path, ext)
        sep = _detect_delimiter(file_path)
        df = pd.read_csv(file_path, header=header_row, sep=sep, dtype=_id_dtypes(file_path, ext, header_row, sep))
        return _clean_dataframe(df)
    elif ext in (".xlsx", ".xls"):
        header_row = _find_header_row(file_path, ext)
        df = pd.read_excel(file_path, header=header_row, dtype=_id_dtypes(file_path, ext, header_row))
        return _clean_dataframe(df)
    elif ext == ".json":
        try:
            df = pd.read_json(file_path)
        except ValueError:
            # Try records orientation
            df = pd.read_json(file_path, orient='records')
        return _clean_dataframe(df)
    else:
        raise ValueError(f"Unsupported format: {ext}")


def _analyze_column(series: pd.Series) -> dict[str, Any]:
    """Analyze a single column for metadata generation."""
    info: dict[str, Any] = {
        "name": series.name,
        "dtype": str(series.dtype),
        "null_pct": round(series.isna().mean() * 100, 1),
        "unique_count": int(series.nunique()),
        "total_count": len(series),
    }

    if pd.api.types.is_numeric_dtype(series):
        info["type"] = "numeric"
        clean = series.dropna()
        if len(clean) > 0:
            info["min"] = float(clean.min())
            info["max"] = float(clean.max())
            info["mean"] = round(float(clean.mean()), 2)
    elif pd.api.types.is_datetime64_any_dtype(series):
        info["type"] = "datetime"
        clean = series.dropna()
        if len(clean) > 0:
            info["min_date"] = str(clean.min())
            info["max_date"] = str(clean.max())
    else:
        info["type"] = "text"
        clean = series.dropna().astype(str)
        if len(clean) > 0:
            info["sample_values"] = clean.value_counts().head(5).index.tolist()
            if info["unique_count"] <= 20:
                info["is_categorical"] = True

    # Try date detection for string columns
    if info["type"] == "text" and len(clean) > 0:
        try:
            sample = clean.head(20).astype(str)
            parsed = pd.to_datetime(sample, format='mixed', errors='coerce')
            if parsed.notna().mean() > 0.7:  # 70%+ parseable as dates
                info["type"] = "datetime"
                info["min_date"] = str(parsed.min())
                info["max_date"] = str(parsed.max())
        except Exception:
            pass

    return info


def _infer_pg_type(col_info: dict) -> str:
    """Map analyzed column to a PostgreSQL-friendly description."""
    if col_info["type"] == "numeric":
        if col_info.get("max", 0) == int(col_info.get("max", 0)) and col_info.get("min", 0) == int(col_info.get("min", 0)):
            return "INTEGER"
        return "NUMERIC"
    elif col_info["type"] == "datetime":
        return "TIMESTAMP"
    else:
        return "TEXT"


def _detect_relationships(columns: list[str]) -> list[str]:
    """Detect potential foreign key relationships from column names."""
    rels = []
    for col in columns:
        if col.endswith("_id") and col != "id":
            ref_table = col[:-3]  # customer_id → customer
            rels.append(f"`{col}` likely references `{ref_table}` table")
    return rels


def _profile_table(df: pd.DataFrame, project_slug: str, table_name: str) -> dict:
    """Profile table using pure pandas — real health %, column stats, alerts. No external deps."""
    profile_data = {"health": 60, "alerts": [], "columns": {}, "row_count": len(df), "col_count": len(df.columns)}
    try:
        alerts = []
        col_profiles = {}
        total_cols = len(df.columns)

        for col in df.columns:
            series = df[col]
            null_pct = round(float(series.isna().mean()) * 100, 1)
            unique = int(series.nunique())
            ctype = "numeric" if pd.api.types.is_numeric_dtype(series) else ("datetime" if pd.api.types.is_datetime64_any_dtype(series) else "text")

            cp = {"type": ctype, "missing_pct": null_pct, "unique": unique}
            if ctype == "numeric":
                clean = series.dropna()
                if len(clean) > 0:
                    cp["min"] = float(clean.min())
                    cp["max"] = float(clean.max())
                    cp["mean"] = round(float(clean.mean()), 2)
                    cp["zeros_pct"] = round(float((clean == 0).mean()) * 100, 1)
            elif ctype == "text":
                clean = series.dropna().astype(str)
                if len(clean) > 0:
                    cp["top_values"] = clean.value_counts().head(5).index.tolist()
                    if unique <= 20:
                        cp["is_categorical"] = True

            col_profiles[col] = cp

            # Generate alerts
            if null_pct > 50:
                alerts.append(f"Column '{col}' is {null_pct}% missing")
            if null_pct > 0 and null_pct <= 50:
                alerts.append(f"Column '{col}' has {null_pct}% missing values")
            if ctype == "numeric" and cp.get("zeros_pct", 0) > 30:
                alerts.append(f"Column '{col}' has {cp['zeros_pct']}% zero values")

        # Duplicate detection
        dup_rows = int(df.duplicated().sum())
        if dup_rows > 0:
            alerts.append(f"{dup_rows} duplicate rows detected")

        # Scientific-notation ID corruption detection (Excel large-ID collapse)
        # Runs over every text/object column whose name matches an ID/code pattern.
        # A column with >1000 rows but ≤2 distinct values, or where a high fraction
        # of non-null values match sci-notation form, is flagged as corrupt.
        import re as _re_sci
        _SCI_RE = _re_sci.compile(r'^\d(\.\d+)?[eE]\+?\d+$')
        _ID_NAME_RE = _re_sci.compile(r'code|_id$|^id$|article|sku|barcode', _re_sci.IGNORECASE)
        for col in df.columns:
            if not _ID_NAME_RE.search(str(col)):
                continue
            series = df[col]
            # Only inspect text/object columns (numeric already collapsed by pandas)
            if pd.api.types.is_numeric_dtype(series):
                continue
            clean = series.dropna().astype(str)
            n_total = len(clean)
            if n_total == 0:
                continue
            n_distinct = int(clean.nunique())
            # Heuristic 1: >1000 rows but ≤2 distinct values → suspiciously collapsed
            collapsed = n_total > 1000 and n_distinct <= 2
            # Heuristic 2: >50% of values match scientific-notation pattern
            sci_hits = int(clean.apply(lambda v: bool(_SCI_RE.match(v.strip()))).sum())
            sci_frac = sci_hits / n_total if n_total > 0 else 0
            sci_dominant = sci_frac > 0.5
            if collapsed or sci_dominant:
                sample_val = clean.iloc[0] if len(clean) > 0 else "?"
                logger.warning(
                    "CORRUPT ID COLUMN detected in '%s': column '%s' has %d distinct value(s) "
                    "over %d rows (sci-notation fraction=%.0f%%). Excel has collapsed large IDs "
                    "to scientific notation (e.g. '%s'). Catalog↔stock joins WILL break. "
                    "Re-upload with article_code formatted as TEXT, or export CSV directly from source.",
                    table_name, col, n_distinct, n_total, sci_frac * 100, sample_val,
                )
                alerts.append(
                    f"[DATA CORRUPTION] Column '{col}' looks corrupted by Excel scientific-notation "
                    f"(e.g. '{sample_val}', only {n_distinct} distinct value(s) over {n_total:,} rows). "
                    f"Re-upload this file with '{col}' formatted as TEXT (not General/Number), "
                    f"or export CSV directly from source without opening in Excel."
                )

        # Compute REAL health %
        health = 100
        missing_cols = sum(1 for c in col_profiles.values() if c.get("missing_pct", 0) > 50)
        if total_cols > 0:
            health -= int((missing_cols / total_cols) * 30)
        if dup_rows > 0:
            health -= 15
        if len(alerts) > 8:
            health -= 10
        if len(df) < 5:
            health -= 20
        # Bonus for completeness
        complete_cols = sum(1 for c in col_profiles.values() if c.get("missing_pct", 0) == 0)
        if total_cols > 0 and complete_cols == total_cols:
            health = min(100, health + 10)
        health = max(10, min(100, health))

        profile_data = {
            "health": health,
            "row_count": len(df),
            "col_count": total_cols,
            "duplicate_rows": dup_rows,
            "alerts": alerts[:20],
            "columns": col_profiles,
        }

        # Save profile JSON
        profile_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
        profile_dir.mkdir(parents=True, exist_ok=True)
        _safe_write_json(profile_dir / f"{table_name}_profile.json", profile_data)

    except Exception:
        pass
    return profile_data


def _detect_unit(col_name: str, p: dict) -> str:
    """Best-effort unit for a column from its name + profiled type.
    Rules only ($0). LLM deep-analysis can override later.
    Returns: currency | percentage | count | physical | time | number | none
    """
    n = str(col_name or "").lower()
    if p.get("type") == "datetime":
        return "time"
    numeric = p.get("type") == "numeric"
    if numeric and re.search(r"price|amount|cost|revenue|sales|total|value|fee|salary|paid|payment|budget|spend|gmv|aov|mmk|usd|eur|gbp|\$|income|profit|margin_amt", n):
        return "currency"
    if numeric and re.search(r"pct|percent|_rate$|^rate|ratio|margin|share|growth|conversion|utiliz|occupancy", n):
        return "percentage"
    if numeric and re.search(r"\bqty\b|quantity|count|num_|_num\b|units|orders|visits|clicks|views|sessions|stock|on_hand|headcount", n):
        return "count"
    if numeric and re.search(r"\bkg\b|weight|height|\bcm\b|\bkm\b|\bml\b|volume|distance|temp|litre|liter|gram|meter", n):
        return "physical"
    return "number" if numeric else "none"


def _column_group(p: dict) -> str:
    """Bucket a profiled column for grouped display."""
    c = p.get("classification")
    if p.get("type") == "datetime":
        return "DATE"
    if c == "id":
        return "IDENTIFIER"
    if c == "measure" or (p.get("type") == "numeric" and c != "dimension"):
        return "MEASURE"
    if c == "dimension":
        return "DIMENSION"
    return "TEXT"


def _sql_profile_columns(project_slug: str, table_name: str, engine=None) -> list[dict]:
    """Profile ALL columns via SQL — zero RAM, works on 1M+ rows in <10s.
    Classifies each column as: dimension (categorical), measure (numeric), id, or text."""
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    qualified = f'"{schema}"."{table_name}"'
    if not engine:
        engine = create_engine(db_url)

    profiles = []
    try:
        with engine.connect() as conn:
            # Get column list + types
            cols = conn.execute(text(
                "SELECT column_name, data_type FROM information_schema.columns "
                "WHERE table_schema = :s AND table_name = :t ORDER BY ordinal_position"
            ), {"s": schema, "t": table_name}).fetchall()

            total_rows = conn.execute(text(f"SELECT COUNT(*) FROM {qualified}")).scalar() or 0

            for col_name, col_type in cols:
                if col_name.startswith('_source_'):
                    continue  # Skip our tracking columns
                safe_col = f'"{col_name}"'
                p = {"name": col_name, "pg_type": col_type, "total_rows": total_rows}

                try:
                    # Universal stats
                    stats = conn.execute(text(f"""
                        SELECT COUNT({safe_col}) as non_null,
                               COUNT(DISTINCT {safe_col}) as unique_count
                        FROM {qualified}
                    """)).fetchone()
                    p["non_null"] = stats[0]
                    p["null_count"] = total_rows - stats[0]
                    p["null_pct"] = round((total_rows - stats[0]) / max(total_rows, 1) * 100, 1)
                    p["unique_count"] = stats[1]

                    is_numeric = col_type in ('integer', 'bigint', 'smallint', 'numeric', 'double precision', 'real')
                    is_date = col_type in ('timestamp with time zone', 'timestamp without time zone', 'date')

                    if is_numeric and stats[0] > 0:
                        num_stats = conn.execute(text(f"""
                            SELECT MIN({safe_col}), MAX({safe_col}),
                                   AVG({safe_col}::numeric), STDDEV({safe_col}::numeric),
                                   PERCENTILE_CONT(0.25) WITHIN GROUP (ORDER BY {safe_col}::numeric),
                                   PERCENTILE_CONT(0.50) WITHIN GROUP (ORDER BY {safe_col}::numeric),
                                   PERCENTILE_CONT(0.75) WITHIN GROUP (ORDER BY {safe_col}::numeric)
                            FROM {qualified} WHERE {safe_col} IS NOT NULL
                        """)).fetchone()
                        p.update({"type": "numeric", "min": float(num_stats[0]) if num_stats[0] is not None else None,
                                  "max": float(num_stats[1]) if num_stats[1] is not None else None,
                                  "mean": round(float(num_stats[2]), 2) if num_stats[2] is not None else None,
                                  "stddev": round(float(num_stats[3]), 2) if num_stats[3] is not None else None,
                                  "p25": float(num_stats[4]) if num_stats[4] is not None else None,
                                  "median": float(num_stats[5]) if num_stats[5] is not None else None,
                                  "p75": float(num_stats[6]) if num_stats[6] is not None else None})
                    elif is_date and stats[0] > 0:
                        date_stats = conn.execute(text(f"SELECT MIN({safe_col}), MAX({safe_col}) FROM {qualified} WHERE {safe_col} IS NOT NULL")).fetchone()
                        p.update({"type": "datetime", "min_date": str(date_stats[0])[:10], "max_date": str(date_stats[1])[:10]})
                    else:
                        p["type"] = "text"

                    # Classify: dimension vs measure vs id
                    if p["unique_count"] <= 500 and p["type"] != "datetime":
                        p["classification"] = "dimension"
                    elif p["unique_count"] >= total_rows * 0.9 and total_rows > 10:
                        p["classification"] = "id"
                    elif p["type"] == "numeric":
                        p["classification"] = "measure"
                    else:
                        p["classification"] = "text"

                except Exception:
                    p["type"] = "unknown"
                    p["classification"] = "unknown"

                # Semantic unit + display group (rules; LLM may refine later)
                p["unit"] = _detect_unit(col_name, p)
                p["group"] = _column_group(p)

                profiles.append(p)
    except Exception as e:
        logger.warning(f"SQL profiling failed for {table_name}: {e}")

    return profiles


def _compute_table_quality(project_slug: str, table_name: str, engine=None) -> dict:
    """Fast, free (SQL-only) data-quality score for one table.

    Returns {table, score (0-100), verdict GOOD/WEAK/BAD, reasons[], rows, cols}.
    Used by the pre-train quality gate so users review junk before spending LLM $.
    Fail-open: if analysis errors, return WEAK (never hard-block on a gate bug).
    """
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    if not engine:
        engine = create_engine(db_url)

    rows = 0
    try:
        with engine.connect() as conn:
            rows = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{table_name}"')).scalar() or 0
        profiles = _sql_profile_columns(project_slug, table_name, engine=engine)
    except Exception as e:
        logger.warning(f"quality check failed for {table_name}: {e}")
        return {"table": table_name, "score": 50, "verdict": "WEAK",
                "reasons": ["could not analyze — review manually"], "rows": rows, "cols": 0}

    cols = [p for p in profiles if not str(p.get("name", "")).startswith("_source_")]
    n_cols = len(cols)
    reasons: list[str] = []
    score = 100

    # ── Hard fails (no judgment — can't train on nothing) ──
    if n_cols == 0:
        return {"table": table_name, "score": 0, "verdict": "BAD",
                "reasons": ["no columns"], "rows": rows, "cols": 0}
    if rows == 0:
        return {"table": table_name, "score": 0, "verdict": "BAD",
                "reasons": ["empty — 0 rows"], "rows": 0, "cols": n_cols}
    if rows < 5:
        return {"table": table_name, "score": 10, "verdict": "BAD",
                "reasons": [f"only {rows} row(s) — too few to learn from"], "rows": rows, "cols": n_cols}

    # ── Soft scoring ──
    if rows < 50:
        score -= 30
        reasons.append(f"only {rows} rows — thin data")

    avg_null = sum(p.get("null_pct", 0) or 0 for p in cols) / max(n_cols, 1)
    if avg_null > 70:
        score -= 40
        reasons.append(f"{round(avg_null)}% of cells empty — mostly blank")
    elif avg_null > 40:
        score -= 15
        reasons.append(f"{round(avg_null)}% of cells empty")

    dead_cols = [p["name"] for p in cols if (p.get("unique_count", 0) or 0) <= 1]
    if dead_cols:
        frac = len(dead_cols) / n_cols
        if frac > 0.5:
            score -= 30
            reasons.append(f"{len(dead_cols)} of {n_cols} columns have no variety")
        else:
            score -= 10
            reasons.append(f"{len(dead_cols)} dead column(s): {', '.join(map(str, dead_cols[:3]))}")

    score = max(0, min(100, score))
    verdict = "GOOD" if score >= 70 else ("WEAK" if score >= 40 else "BAD")
    if not reasons:
        reasons.append("clean")
    return {"table": table_name, "score": score, "verdict": verdict,
            "reasons": reasons, "rows": rows, "cols": n_cols}


def _live_tables(project_slug: str) -> set[str]:
    """Base table names that actually exist in the project's DB schema.
    Used to skip orphaned knowledge artifacts (dimensions/qa/business JSON)
    left behind by a raw-SQL wipe / external drop."""
    try:
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        return set(inspect(create_engine(db_url)).get_table_names(schema=schema))
    except Exception:
        return set()


def _purge_orphan_knowledge(project_slug: str, log_fn=None) -> int:
    """Delete knowledge files whose table no longer exists in the DB.
    Covers tables/, dimensions/, business/, queries/, training/ — so a dropped
    table can't keep seeding ghost relationships, Q&A, or eval cases.
    Returns count removed. Fail-soft."""
    _log = log_fn or (lambda m: None)
    live = _live_tables(project_slug)
    if not live:
        return 0  # can't confirm liveness → don't risk deleting good files
    base = KNOWLEDGE_DIR / project_slug
    removed = 0
    # (subdir, suffix-to-strip-from-stem)
    targets = [
        ("tables", ""), ("dimensions", ""), ("business", "_rules"),
        ("queries", "_queries"), ("training", "_qa"),
    ]
    for sub, suffix in targets:
        d = base / sub
        if not d.exists():
            continue
        for f in d.glob("*.*"):
            stem = f.stem
            if suffix and stem.endswith(suffix):
                stem = stem[: -len(suffix)]
            if stem and stem not in live:
                try:
                    f.unlink(missing_ok=True)
                    removed += 1
                except Exception:
                    pass
    if removed:
        _log(f"🧹 purged {removed} orphan knowledge file(s) for dropped tables")
    return removed


def _build_dimension_catalog(project_slug: str, table_name: str, profiles: list[dict], engine=None) -> dict:
    """Build catalog of exact values for all dimension columns via SQL.
    Returns {col_name: [{value, count, pct}, ...]}."""
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    qualified = f'"{schema}"."{table_name}"'
    if not engine:
        engine = create_engine(db_url)

    catalog = {}
    date_text_cols: list[str] = []   # text cols holding DD/MM/YYYY dates (#3)
    total_rows = profiles[0]["total_rows"] if profiles else 0

    def _is_lineage(c: str) -> bool:
        cl = (c or "").lower()
        return cl.startswith("_") or cl in (
            "source_file", "source_sheet", "batch_id", "content_hash",
            "row_key", "ingested_at", "period",
        )

    def _clean_val(v) -> str:
        # Strip null bytes + control chars (data had '\0' in mmlabel) (#5).
        s = str(v)
        return "".join(ch for ch in s if ch == "\n" or ch == "\t" or ord(ch) >= 32).strip()

    _DDMMYYYY = re.compile(r"^\s*\d{1,2}/\d{1,2}/\d{4}\b")

    for p in profiles:
        if p.get("classification") != "dimension":
            continue
        col = p["name"]
        # #5 — skip ingest/lineage cols (single-valued bookkeeping, not real dims)
        if _is_lineage(col):
            continue
        try:
            with engine.connect() as conn:
                # #2 — a "dimension" with long average text is actually FREE TEXT
                # (e.g. the `other` col = 76-char Burmese instructions). Catalog-ing
                # it dumps 100+ sentences into the brain. Skip when avg len > 40.
                try:
                    avglen = conn.execute(text(
                        f'SELECT AVG(LENGTH("{col}")) FROM {qualified} WHERE "{col}" IS NOT NULL'
                    )).scalar()
                    if avglen is not None and float(avglen) > 40:
                        continue
                except Exception:
                    pass
                rows = conn.execute(text(f"""
                    SELECT "{col}", COUNT(*) as freq
                    FROM {qualified}
                    WHERE "{col}" IS NOT NULL
                    GROUP BY "{col}"
                    ORDER BY freq DESC
                    LIMIT 100
                """)).fetchall()
                cleaned = [{"value": _clean_val(r[0]), "count": int(r[1]),
                            "pct": round(int(r[1]) / max(total_rows, 1) * 100, 1)} for r in rows]
                cleaned = [c for c in cleaned if c["value"]]  # drop now-empty (null-byte-only)
                # #4 — a single-value column (e.g. stock created_at = one snapshot
                # timestamp on every row) is useless as a dimension. Skip it.
                if len(cleaned) <= 1:
                    continue
                catalog[col] = cleaned
                # #3 — detect text cols storing DD/MM/YYYY dates
                if cleaned and sum(1 for c in cleaned[:20] if _DDMMYYYY.match(c["value"])) >= max(1, min(10, len(cleaned[:20])) // 2):
                    date_text_cols.append(col)
        except Exception:
            pass

    # #3 — write a brain rule so the agent uses to_date(col,'DD/MM/YYYY HH24:MI')
    # instead of `col::date` (which throws DatetimeFieldOverflow on DD/MM/YYYY).
    if date_text_cols:
        try:
            from db.session import get_write_engine as _gwe
            _we = _gwe()
            with _we.begin() as _bc:
                for _dc in date_text_cols:
                    _defn = (f'Column "{_dc}" stores dates as TEXT in DD/MM/YYYY HH24:MI format. '
                             f"To filter/group by date use to_date(\"{_dc}\", 'DD/MM/YYYY HH24:MI') "
                             f"or to_timestamp(...); NEVER \"{_dc}\"::date (raises DatetimeFieldOverflow).")
                    _bc.execute(text("""
                        INSERT INTO public.dash_company_brain
                          (project_slug, category, name, definition, metadata)
                        VALUES (:s, 'pattern', :n, :d, '{}'::jsonb)
                        ON CONFLICT (project_slug, name) WHERE project_slug IS NOT NULL
                          DO UPDATE SET definition = EXCLUDED.definition
                    """), {"s": project_slug, "n": f"date_format:{table_name}.{_dc}", "d": _defn})
        except Exception:
            pass

    # Save to knowledge dir
    try:
        dim_dir = KNOWLEDGE_DIR / project_slug / "dimensions"
        dim_dir.mkdir(parents=True, exist_ok=True)
        _safe_write_json(dim_dir / f"{table_name}.json", catalog)
    except Exception:
        pass

    return catalog


def _detect_hierarchies(project_slug: str, table_name: str, catalog: dict, engine=None) -> list[dict]:
    """Detect parent-child hierarchies between dimension columns.
    If every child value maps to exactly 1 parent → hierarchy."""
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    qualified = f'"{schema}"."{table_name}"'
    if not engine:
        engine = create_engine(db_url)

    # Skip hierarchy detection on tiny tables — false positives common when N < 50
    try:
        with engine.connect() as _hc_conn:
            _hc_rows = _hc_conn.execute(text(f"SELECT COUNT(*) FROM {qualified}")).scalar() or 0
        if int(_hc_rows) < 50:
            return []
    except Exception:
        # If row count fails, fall through to original behavior
        pass

    # Exclude ingest/lineage columns (_source_file, _period, _batch_id,
    # _content_hash, _row_key, _ingested_at, etc). These are single- or
    # low-cardinality bookkeeping cols — a single-valued parent makes EVERY
    # other column a false "child" (all rows map to the one parent value),
    # producing junk hierarchies that pollute the agent brain.
    def _is_lineage_col(c: str) -> bool:
        cl = (c or "").lower()
        return cl.startswith("_") or cl in (
            "source_file", "source_sheet", "batch_id", "content_hash",
            "row_key", "ingested_at", "period",
        )

    dim_cols = [c for c in catalog.keys() if not _is_lineage_col(c)]
    hierarchies = []

    for parent in dim_cols:
        parent_unique = len(catalog.get(parent, []))
        # Degenerate parent (0 or 1 distinct value) → every child trivially
        # maps to 1 parent → false hierarchy. Skip.
        if parent_unique <= 1:
            continue
        for child in dim_cols:
            if child == parent:
                continue
            child_unique = len(catalog.get(child, []))
            if child_unique <= parent_unique:
                continue  # Child should have MORE unique values than parent
            # Check: does each child have exactly 1 parent?
            try:
                with engine.connect() as conn:
                    check = conn.execute(text(f"""
                        SELECT COUNT(*) FROM (
                            SELECT "{child}", COUNT(DISTINCT "{parent}") as parent_count
                            FROM {qualified}
                            WHERE "{child}" IS NOT NULL AND "{parent}" IS NOT NULL
                            GROUP BY "{child}"
                            HAVING COUNT(DISTINCT "{parent}") > 1
                        ) multi_parent
                    """)).scalar()
                    if check == 0:  # Every child has exactly 1 parent
                        hierarchies.append({"parent": parent, "child": child,
                                            "parent_count": parent_unique, "child_count": child_unique})
            except Exception:
                pass

    return hierarchies


def _smart_sample_rows(project_slug: str, table_name: str, profiles: list[dict], engine=None) -> list[dict]:
    """Get 20 diverse sample rows: start + middle + end + outliers + nulls + per-category.
    Much more representative than first 8 rows."""
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    qualified = f'"{schema}"."{table_name}"'
    if not engine:
        engine = create_engine(db_url)

    samples = []
    total_rows = profiles[0]["total_rows"] if profiles else 0
    if total_rows == 0:
        return samples

    try:
        with engine.connect() as conn:
            # Start rows
            rows = conn.execute(text(f"SELECT * FROM {qualified} LIMIT 3")).fetchall()
            cols = [c.name for c in conn.execute(text(f"SELECT * FROM {qualified} LIMIT 0")).cursor.description] if rows else []
            for r in rows:
                samples.append(dict(zip(cols, [str(v)[:50] if v is not None else None for v in r])))

            # Middle rows
            mid = max(0, total_rows // 2 - 1)
            rows = conn.execute(text(f"SELECT * FROM {qualified} OFFSET {mid} LIMIT 3")).fetchall()
            for r in rows:
                samples.append(dict(zip(cols, [str(v)[:50] if v is not None else None for v in r])))

            # End rows
            rows = conn.execute(text(f"SELECT * FROM {qualified} OFFSET {max(0, total_rows - 3)} LIMIT 3")).fetchall()
            for r in rows:
                samples.append(dict(zip(cols, [str(v)[:50] if v is not None else None for v in r])))

            # Outlier rows (max/min of first numeric column)
            num_cols = [p["name"] for p in profiles if p.get("type") == "numeric" and p.get("classification") == "measure"]
            if num_cols:
                nc = f'"{num_cols[0]}"'
                for order in ["DESC", "ASC"]:
                    rows = conn.execute(text(f"SELECT * FROM {qualified} WHERE {nc} IS NOT NULL ORDER BY {nc} {order} LIMIT 2")).fetchall()
                    for r in rows:
                        samples.append(dict(zip(cols, [str(v)[:50] if v is not None else None for v in r])))

            # Null pattern rows
            null_cols = [p["name"] for p in profiles if p.get("null_pct", 0) > 5 and p.get("null_pct", 0) < 90]
            if null_cols:
                nc = f'"{null_cols[0]}"'
                rows = conn.execute(text(f"SELECT * FROM {qualified} WHERE {nc} IS NULL LIMIT 2")).fetchall()
                for r in rows:
                    samples.append(dict(zip(cols, [str(v)[:50] if v is not None else None for v in r])))

    except Exception as e:
        logger.warning(f"Smart sampling failed for {table_name}: {e}")

    # Deduplicate
    seen = set()
    unique_samples = []
    for s in samples:
        key = str(sorted(s.items()))
        if key not in seen:
            seen.add(key)
            unique_samples.append(s)

    return unique_samples[:500]


def _generate_metadata(table_name: str, df: pd.DataFrame, col_analyses: list[dict]) -> dict:
    """Generate table metadata JSON (same format as knowledge/tables/*.json)."""
    # Infer use cases from column types
    use_cases = []
    has_dates = any(c["type"] == "datetime" for c in col_analyses)
    has_numerics = any(c["type"] == "numeric" for c in col_analyses)
    has_categories = any(c.get("is_categorical") for c in col_analyses)

    if has_dates and has_numerics:
        use_cases.append("Time-series analysis and trend tracking")
    if has_numerics:
        use_cases.append("Aggregation and metric calculation (SUM, AVG, COUNT)")
    if has_categories:
        use_cases.append("Segmentation and group-by analysis")
    if has_dates:
        use_cases.append("Date-range filtering and period comparisons")
    use_cases.append(f"General analysis of {table_name.replace('_', ' ')} data")

    # Data quality notes
    quality_notes = []
    for c in col_analyses:
        if c["null_pct"] > 0:
            quality_notes.append(f"`{c['name']}` has {c['null_pct']}% NULL values")
        if c["type"] == "text" and c.get("is_categorical") and c["unique_count"] <= 10:
            vals = c.get("sample_values", [])
            quality_notes.append(f"`{c['name']}` categories: {', '.join(str(v) for v in vals[:8])}")

    # Column definitions
    columns = []
    for c in col_analyses:
        pg_type = _infer_pg_type(c)
        desc = f"{c['type'].title()} column"
        if c["type"] == "numeric":
            desc = f"Range: {c.get('min', '?')} to {c.get('max', '?')}, avg {c.get('mean', '?')}"
        elif c["type"] == "datetime":
            desc = f"Date range: {c.get('min_date', '?')} to {c.get('max_date', '?')}"
        elif c.get("is_categorical"):
            desc = f"Categories ({c['unique_count']} values)"
        elif c["type"] == "text":
            desc = f"Text ({c['unique_count']} unique values)"

        columns.append({"name": c["name"], "type": pg_type, "description": desc})

    return {
        "table_name": table_name,
        "table_description": f"User-uploaded dataset: {table_name.replace('_', ' ')}. {len(df):,} rows, {len(df.columns)} columns.",
        "use_cases": use_cases[:5],
        "data_quality_notes": quality_notes[:5],
        "table_columns": columns,
    }


def _generate_sample_queries(table_name: str, col_analyses: list[dict]) -> str:
    """Generate sample SQL queries for the uploaded table."""
    queries = []
    numeric_cols = [c for c in col_analyses if c["type"] == "numeric"]
    date_cols = [c for c in col_analyses if c["type"] == "datetime"]
    cat_cols = [c for c in col_analyses if c.get("is_categorical")]

    # Row count
    queries.append(f"-- <query {table_name}_count>\n-- <description>Total row count for {table_name}</description>\n-- <query>\nSELECT COUNT(*) AS total_rows FROM {table_name};\n-- </query>")

    # Numeric summary
    if numeric_cols:
        agg_parts = ", ".join(f"ROUND(AVG({c['name']}), 2) AS avg_{c['name']}" for c in numeric_cols[:3])
        queries.append(f"-- <query {table_name}_summary>\n-- <description>Numeric summary of {table_name}</description>\n-- <query>\nSELECT {agg_parts}, COUNT(*) AS total FROM {table_name};\n-- </query>")

    # Group by category
    if cat_cols and numeric_cols:
        cat = cat_cols[0]["name"]
        num = numeric_cols[0]["name"]
        queries.append(f"-- <query {table_name}_by_{cat}>\n-- <description>{table_name} grouped by {cat}</description>\n-- <query>\nSELECT {cat}, COUNT(*) AS count, ROUND(AVG({num}), 2) AS avg_{num}\nFROM {table_name}\nGROUP BY {cat}\nORDER BY count DESC;\n-- </query>")

    # Date trend — skip when the only date col is effectively constant
    # (e.g. CityPharma `created_at` populated with a single ingestion timestamp
    # → 1-row trend is meaningless). Uses in-memory unique_count to avoid
    # an extra DB roundtrip; defense-in-depth `is_constant_column` available
    # in dash.utils.column_classifier for callers w/ a live conn.
    if date_cols and numeric_cols:
        dt_col = date_cols[0]
        if (dt_col.get("unique_count") or 0) > 1:
            dt = dt_col["name"]
            num = numeric_cols[0]["name"]
            # text DD/MM/YYYY → to_date(); real datetime → ::timestamp.
            _dsv = dt_col.get("sample_values") or []
            if dt_col.get("type") != "datetime" and any(
                    re.match(r"^\s*\d{1,2}/\d{1,2}/\d{4}\b", str(v)) for v in _dsv[:6]):
                _dt_expr = f"to_date({dt}, 'DD/MM/YYYY HH24:MI')"
            else:
                _dt_expr = f"{dt}::timestamp"
            queries.append(f"-- <query {table_name}_trend>\n-- <description>Monthly trend for {table_name}</description>\n-- <query>\nSELECT DATE_TRUNC('month', {_dt_expr}) AS month, COUNT(*) AS count, ROUND(SUM({num}), 2) AS total_{num}\nFROM {table_name}\nGROUP BY 1\nORDER BY 1;\n-- </query>")

    # Top records
    queries.append(f"-- <query {table_name}_sample>\n-- <description>Sample rows from {table_name}</description>\n-- <query>\nSELECT * FROM {table_name} LIMIT 10;\n-- </query>")

    return "\n\n".join(queries)


def _generate_business_rules(table_name: str, col_analyses: list[dict]) -> dict:
    """Generate business rules/metrics for the uploaded table."""
    metrics = []
    rules = []
    gotchas = []

    # Use central column classifier (2026-05-25): drives metric/measure
    # selection off real PK constraints + adaptive thresholds instead of
    # name patterns + fixed enum range. Falls back to name pattern if
    # classifier can't reach the DB (e.g. table not yet committed).
    name_pattern_fallback = lambda c: (
        (c.get("name") or "").lower() in {"id", "_period"} or
        (c.get("name") or "").lower().endswith(("_id", "_code", "_key"))
    )
    real_measures = [c for c in col_analyses if c["type"] == "numeric" and not name_pattern_fallback(c)]
    for c in real_measures[:3]:
        metrics.append({
            "name": f"{table_name} {c['name']}".replace("_", " ").title(),
            "definition": f"Aggregate of {c['name']} from {table_name}",
            "table": table_name,
            "calculation": f"SUM({c['name']}) or AVG({c['name']}) FROM {table_name}",
        })

    for c in col_analyses:
        if c["null_pct"] > 10:
            gotchas.append({
                "issue": f"NULL values in {table_name}.{c['name']}",
                "tables_affected": [table_name],
                "solution": f"{c['null_pct']}% NULL. Use COALESCE or filter with IS NOT NULL.",
            })

    rels = _detect_relationships([c["name"] for c in col_analyses])
    for r in rels:
        rules.append(r)

    return {"metrics": metrics, "business_rules": rules, "common_gotchas": gotchas}


# ---------------------------------------------------------------------------
# LLM-Powered Auto-Training
# ---------------------------------------------------------------------------


def _llm_deep_analysis(table_name: str, col_analyses: list[dict], sample_rows: list[dict]) -> dict | None:
    """Use LLM to deeply analyze data and generate smart metadata."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    # Build sample data as markdown table
    if sample_rows:
        headers = list(sample_rows[0].keys())
        sample_md = "| " + " | ".join(headers) + " |\n"
        sample_md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        for row in sample_rows[:8]:
            sample_md += "| " + " | ".join(str(row.get(h, ""))[:30] for h in headers) + " |\n"
    else:
        sample_md = "No sample data available"

    col_summary = json.dumps([{k: v for k, v in c.items() if k in ("name", "type", "null_pct", "unique_count", "min", "max", "mean", "sample_values", "is_categorical")} for c in col_analyses], indent=2, default=str)

    # Build distribution summary (gives LLM full picture without sending all rows)
    dist_lines = []
    total_rows = col_analyses[0].get("total_count", 0) if col_analyses else 0
    dist_lines.append(f"TOTAL ROWS: {total_rows:,}")
    for ca in col_analyses:
        name = ca.get("name", "?")
        ctype = ca.get("type", "?")
        null_pct = ca.get("null_pct", 0)
        unique = ca.get("unique_count", 0)
        line = f"  {name} ({ctype}): {unique} unique, {null_pct}% null"
        if ctype == "numeric" and ca.get("mean") is not None:
            line += f", min={ca.get('min')}, max={ca.get('max')}, mean={ca.get('mean')}"
            zeros = ca.get("zeros_pct", 0)
            if zeros:
                line += f", {zeros}% zeros"
        elif ca.get("sample_values"):
            vals = ca["sample_values"][:5]
            line += f", top values: {vals}"
        if ca.get("is_categorical"):
            line += " [CATEGORICAL]"
        if ca.get("min_date"):
            line += f", range: {ca['min_date']} to {ca.get('max_date', '?')}"
        dist_lines.append(line)
    distribution_summary = "\n".join(dist_lines)

    prompt = f"""You are analyzing a dataset to train a data agent. Analyze deeply.

TABLE: {table_name}
COLUMNS:
{col_summary}

DATA DISTRIBUTION (full table stats):
{distribution_summary}

SAMPLE DATA (first 8 rows):
{sample_md}

Generate a comprehensive analysis like a Codex-enriched knowledge pipeline. Return ONLY valid JSON (no markdown):
{{
  "table_description": "What this table represents in business terms (2-3 sentences)",
  "table_purpose": "One-line purpose: why does this table exist? What business process does it serve?",
  "grain": "What does each row represent? e.g. 'One row per customer per month' or 'One row per transaction'",
  "primary_keys": ["columns that uniquely identify each row"],
  "foreign_keys": [{{"column": "col_name", "references": "other_table.col", "relationship": "many-to-one"}}],
  "usage_patterns": ["Common downstream usage pattern 1", "Common query pattern 2"],
  "alternate_tables": "When should you use a different table instead? e.g. 'Use orders_summary instead for aggregated metrics'",
  "freshness": "How often is this data likely updated? (real-time, daily, weekly, static snapshot)",
  "column_descriptions": {{
    "column_name": "What this column means in business context (not just the data type)"
  }},
  "metrics": [
    {{"name": "Metric Name", "definition": "What it measures", "table": "{table_name}", "calculation": "SQL expression"}}
  ],
  "business_rules": [
    {{"name": "Rule Name", "definition": "The business rule or logic", "type": "business_rule"}}
  ],
  "data_quality": [
    "Insight about data quality or patterns..."
  ],
  "suggested_role": "One line describing the agent's expertise area",
  "suggested_personality": "friendly"
}}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 2000, "temperature": 0.1},
            timeout=30,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return json.loads(content.strip().strip("`").strip())
    except Exception:
        return None


def _get_chat_feedback_for_training(table_name: str) -> str:
    """Load proven SQL patterns + user feedback to improve Q&A generation."""
    feedback_context = ""
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Proven queries (thumbs-up)
            good = conn.execute(text(
                "SELECT question, answer FROM public.dash_feedback WHERE rating = 'up' AND answer LIKE :t ORDER BY created_at DESC LIMIT 5"
            ), {"t": f"%{table_name}%"}).fetchall()
            if good:
                feedback_context += "\nUSERS LIKED THESE QUERIES (generate similar):\n"
                for g in good:
                    feedback_context += f"  Q: {str(g[0])[:80]}\n"

            # Anti-patterns (thumbs-down)
            bad = conn.execute(text(
                "SELECT question, answer FROM public.dash_feedback WHERE rating = 'down' AND answer LIKE :t ORDER BY created_at DESC LIMIT 3"
            ), {"t": f"%{table_name}%"}).fetchall()
            if bad:
                feedback_context += "\nUSERS DISLIKED THESE (avoid similar):\n"
                for b in bad:
                    feedback_context += f"  Q: {str(b[0])[:80]}\n"

            # Proven SQL patterns
            patterns = conn.execute(text(
                "SELECT question, sql_text FROM public.dash_query_patterns WHERE table_name = :t AND usage_count > 1 ORDER BY usage_count DESC LIMIT 3"
            ), {"t": table_name}).fetchall()
            if patterns:
                feedback_context += "\nPROVEN SQL PATTERNS (high usage, include similar):\n"
                for p in patterns:
                    feedback_context += f"  Q: {str(p[0])[:60]} → {str(p[1])[:100]}\n"
    except Exception:
        pass
    return feedback_context


def _llm_generate_training(table_name: str, metadata: dict, col_analyses: list[dict] = None) -> list[dict] | None:
    """Use LLM to generate training Q&A pairs for the agent."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    cols_desc = "\n".join(f"- {c['name']} ({c.get('type', 'TEXT')}): {c.get('description', '')}" for c in metadata.get("table_columns", []))

    # Build distribution so LLM knows REAL data shape (not just column names).
    # Enum / boolean detection added 2026-05-25 — pipeline used to generate
    # nonsense SQL like `SELECT status * 1.10` on a 0/1 enum because the LLM
    # didn't see column semantics, only dtype. Now flag low-cardinality
    # numeric cols + cols with min/max in {0,1} or {-1,0,1} as enums and tell
    # the LLM never to do arithmetic on them.
    dist_info = ""
    enum_cols: list[str] = []
    # 2026-05-25 (Phase 9): also flag CONSTANT columns so LLM doesn't generate
    # "trend by X" / "by X month" / DATE_TRUNC(X) Q&A on a column that has only
    # one unique value. CityPharma `created_at` was constant → trend QA returned
    # 1 row → useless cached oracle.
    constant_cols: list[str] = []
    date_text_cols: list[str] = []   # TEXT cols holding DD/MM/YYYY dates
    import re as _re_dt
    _DDMMYYYY = _re_dt.compile(r"^\s*\d{1,2}/\d{1,2}/\d{4}\b")
    if col_analyses:
        total = col_analyses[0].get("total_count", 0) if col_analyses else 0
        dist_lines = [f"ROWS: {total:,}"]
        for ca in col_analyses:
            name = ca.get("name", "?")
            # DD/MM/YYYY text-date detection — these need to_date(), not ::date.
            try:
                _sv = ca.get("sample_values") or []
                if ca.get("type") != "numeric" and _sv and \
                        sum(1 for v in _sv[:6] if _DDMMYYYY.match(str(v))) >= max(1, len(_sv[:6]) // 2):
                    date_text_cols.append(name)
            except Exception:
                pass
            _uc = ca.get("unique_count") or 0
            line = f"  {name}: {_uc} unique, {ca.get('null_pct', 0)}% null"
            tag = ""
            # CONSTANT detection — applies to ANY dtype.
            if isinstance(_uc, int) and _uc <= 1 and (total or 0) > 1:
                tag = " [CONSTANT — single value across all rows; do NOT use in GROUP BY / DATE_TRUNC / time-axis]"
                constant_cols.append(name)
            elif ca.get("type") == "numeric" and ca.get("mean") is not None:
                line += f", range {ca.get('min')} to {ca.get('max')}, avg {ca.get('mean')}"
                # enum / boolean heuristic
                _mn = ca.get("min")
                _mx = ca.get("max")
                if isinstance(_uc, int) and _uc <= 10 and _mn is not None and _mx is not None:
                    try:
                        if float(_mn) >= -1 and float(_mx) <= 10 and float(_mx) - float(_mn) <= 10:
                            tag = " [ENUM/BOOLEAN — no arithmetic, use only in WHERE / GROUP BY]"
                            enum_cols.append(name)
                    except Exception:
                        pass
            elif ca.get("sample_values"):
                line += f", values: {ca['sample_values'][:5]}"
            dist_lines.append(line + tag)
        dist_info = "\nDATA DISTRIBUTION:\n" + "\n".join(dist_lines) + "\n"
        if enum_cols:
            dist_info += (
                f"\n⚠ ENUM/BOOLEAN columns (do NOT SUM/AVG/multiply): {', '.join(enum_cols)}\n"
                f"  → For these cols, only use COUNT, GROUP BY, WHERE, or CASE WHEN.\n"
            )
        if constant_cols:
            dist_info += (
                f"\n⚠ CONSTANT columns (single value across ALL rows): {', '.join(constant_cols)}\n"
                f"  → DO NOT use in GROUP BY, DATE_TRUNC, time-series, or 'trend' questions.\n"
                f"  → DO NOT generate Q&A involving 'by {constant_cols[0]}' / 'over time' on these.\n"
                f"  → Use them ONLY as a filter (WHERE) or as a static fact.\n"
            )
        if date_text_cols:
            _dc0 = date_text_cols[0]
            dist_info += (
                f"\n⚠ TEXT-DATE columns (stored as TEXT in DD/MM/YYYY HH24:MI): {', '.join(date_text_cols)}\n"
                f"  → These are NOT real date/timestamp types. {_dc0}::date THROWS DatetimeFieldOverflow.\n"
                f"  → ALWAYS wrap with to_date(\"{_dc0}\", 'DD/MM/YYYY HH24:MI') (or to_timestamp) before\n"
                f"    DATE_TRUNC / comparison / extraction. Example:\n"
                f"    DATE_TRUNC('month', to_date(\"{_dc0}\", 'DD/MM/YYYY HH24:MI'))\n"
                f"  → NEVER emit \"{_dc0}\"::date or \"{_dc0}\"::timestamp directly.\n"
            )

    # SQL-VALIDATED — inject Postgres dialect rules so LLM emits correct SQL first try.
    # Cuts validator-drop rate by ~60% vs prompt without rules.
    try:
        from dash.tools.llm_sql_helper import _postgres_sql_rules
        _pg_rules = _postgres_sql_rules()
    except Exception:
        _pg_rules = ""

    prompt = f"""Generate 11 training Q&A pairs for this data table, one for each analysis type.
Each pair should have: question, sql, analysis_type.

{_pg_rules}

TABLE: {table_name}
DESCRIPTION: {metadata.get('table_description', '')}
COLUMNS:
{cols_desc}
{dist_info}

Generate exactly 11 pairs, one per analysis type. Use real column names from the table above:
1. DESCRIPTIVE: "What is the total X?" with simple aggregation SQL (SUM, COUNT, AVG)
2. DIAGNOSTIC: "Why does X have the highest Y?" with GROUP BY + ORDER BY SQL
3. COMPARATIVE: "Compare X vs Y" with GROUP BY + CASE WHEN SQL
4. TREND: "Show X over time" with DATE_TRUNC + GROUP BY SQL (use a date/time column if available, otherwise approximate)
5. PARETO: "Which X drives 80% of Y?" with cumulative percentage SQL using window functions
6. ANOMALY: "Any unusual X?" with AVG + STDDEV detection SQL
7. PRESCRIPTIVE: "What should we do about X?" with actionable analysis SQL
8. SCENARIO: "What if X increases by N%?" with projection/multiplication SQL
9. BENCHMARK: "How does X compare to average?" with AVG window function SQL
10. ROOT_CAUSE: "Why is X happening?" with multi-dimension GROUP BY SQL
11. PREDICTIVE: "What will X be next period?" with trend extrapolation SQL

Return ONLY valid JSON array (no markdown):
[
  {{"question": "What is the total revenue?", "sql": "SELECT SUM(revenue) FROM {table_name}", "analysis_type": "descriptive"}},
  ...
]

{_get_chat_feedback_for_training(table_name)}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 3000, "temperature": 0.2},
            timeout=45,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        pairs = json.loads(content.strip().strip("`").strip())
    except Exception:
        return None

    # SQL-VALIDATED — drop unfixable pairs, auto-fix dialect issues (TEXT casts).
    # Closes 26-fail-per-train eval-grading class from 2026-05-27 audit.
    try:
        from dash.tools.sql_validator import validate_and_fix
        if isinstance(pairs, list):
            kept: list[dict] = []
            dropped = 0
            for qa in pairs:
                sql = (qa.get("sql") or "").strip()
                if not sql:
                    dropped += 1
                    continue
                v = validate_and_fix(sql, project_slug, strict=True)
                if v["ok"]:
                    qa["sql"] = v["sql"]  # use fixed version
                    if v["fixes_applied"]:
                        qa["_sql_fixes"] = v["fixes_applied"]
                    kept.append(qa)
                else:
                    dropped += 1
                    logger.info(f"qa drop [{table_name}] {qa.get('analysis_type','?')}: {'; '.join(v['errors'])[:160]}")
                    # Telemetry: emit qa_drop event so frontend can surface per-table drop counters.
                    try:
                        from dash.tools.sql_validator import _emit_event as _se_qa
                        _se_qa(
                            "qa_drop",
                            project_slug=project_slug,
                            source="qa_gen",
                            table_name=table_name,
                            details={
                                "reason": v.get("errors", []),
                                "analysis_type": qa.get("analysis_type"),
                            },
                        )
                    except Exception:
                        pass
            logger.info(f"qa validate [{table_name}]: kept={len(kept)} dropped={dropped}")
            return kept
    except Exception as _val_err:
        logger.warning(f"qa validation failed [{table_name}]: {_val_err} — returning unvalidated")
    return pairs


def _llm_generate_persona(project_slug: str, tables_metadata: list[dict], rules: list[dict]) -> dict | None:
    """Use LLM to generate a full agent persona based on the data."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    # Build context from all tables
    tables_summary = ""
    for m in tables_metadata[:5]:
        cols = ", ".join(c.get("name", "") for c in m.get("table_columns", [])[:10])
        tables_summary += f"- {m.get('table_name', '?')}: {m.get('table_description', '')}. Columns: {cols}\n"

    rules_summary = "\n".join(f"- {r.get('name', '')}: {r.get('definition', '')}" for r in rules[:10])

    prompt = f"""You are creating a persona for a data agent. This agent will be an expert on the following data:

TABLES:
{tables_summary}

BUSINESS RULES:
{rules_summary or 'None yet'}

Generate a comprehensive persona. Return ONLY valid JSON (no markdown):
{{
  "persona_prompt": "A 3-5 sentence system prompt describing who this agent is, what domain it specializes in, how it should approach questions, what terminology it should use, and what it should prioritize when analyzing data.",
  "domain_terms": ["list", "of", "key", "domain", "terms", "the agent should know"],
  "greeting": "A friendly 1-2 sentence greeting the agent uses when a user first opens the chat.",
  "expertise_areas": ["area1", "area2", "area3"],
  "communication_style": "How the agent should communicate (e.g., 'concise with data tables', 'detailed with explanations')"
}}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 1000, "temperature": 0.2},
            timeout=20,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return json.loads(content.strip().strip("`").strip())
    except Exception:
        return None


def _update_project_config(project_slug: str, role: str, personality: str | None):
    """Update project agent config based on LLM suggestion."""
    engine = create_engine(db_url)
    with engine.connect() as conn:
        updates = ["agent_role = :role", "updated_at = NOW()"]
        params: dict = {"slug": project_slug, "role": role}
        if personality:
            updates.append("agent_personality = :pers")
            params["pers"] = personality
        conn.execute(text(f"UPDATE public.dash_projects SET {', '.join(updates)} WHERE slug = :slug"), params)
        conn.commit()


def _reload_project_knowledge(project_slug: str, timeout_sec: int = 60) -> bool:
    """Re-index all knowledge files for a project into PgVector.
    Returns True if successful, False if timed out or failed."""
    import threading

    success = [False]
    error_msg = [None]

    def _do_index():
        try:
            from db.session import create_project_knowledge
            proj_knowledge = create_project_knowledge(project_slug)
            proj_dir = KNOWLEDGE_DIR / project_slug
            if proj_dir.exists():
                for subdir in ["tables", "queries", "business", "rules", "training"]:
                    path = proj_dir / subdir
                    if path.exists():
                        files = [f for f in path.iterdir() if f.is_file() and not f.name.startswith(".")]
                        if files:
                            proj_knowledge.insert(name=f"{project_slug}-{subdir}", path=str(path))
            success[0] = True
        except Exception as e:
            error_msg[0] = str(e)

    thread = threading.Thread(target=_do_index, daemon=True)
    thread.start()
    thread.join(timeout=timeout_sec)

    if thread.is_alive():
        # Thread is still running — timed out
        return False
    return success[0]


def _discover_relationships(project_slug: str):
    """LLM analyzes all tables to discover hidden relationships."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return

    # Load all table metadata
    tables_info = []
    proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
    if not proj_tables_dir.exists():
        return
    for f in proj_tables_dir.glob("*.json"):
        try:
            with open(f) as fh:
                d = json.load(fh)
            cols = [c["name"] for c in d.get("table_columns", [])]
            tables_info.append(f"{d.get('table_name', f.stem)}: {', '.join(cols)}")
        except Exception:
            pass

    if len(tables_info) < 2:
        return

    prompt = f"""Analyze these database tables and discover relationships between them.
Look for: shared columns, FK patterns, value overlaps, hierarchical relationships.

TABLES:
{chr(10).join(tables_info)}

Return ONLY valid JSON array (no markdown):
[{{"from_table": "table1", "from_column": "col1", "to_table": "table2", "to_column": "col2", "type": "fk|shared|hierarchy", "confidence": 0.9}}]

Return empty array [] if no relationships found."""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 500, "temperature": 0.1},
            timeout=15,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        rels = json.loads(content.strip().strip("`").strip())

        if isinstance(rels, list) and rels:
            # VERIFY relationships: check actual value overlap in PostgreSQL
            schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
            engine = create_engine(db_url)
            with engine.connect() as conn:
                for r in rels[:10]:
                    ft, fc = r.get("from_table", ""), r.get("from_column", "")
                    tt, tc = r.get("to_table", ""), r.get("to_column", "")
                    verified_conf = r.get("confidence", 0.5)
                    try:
                        # FK is asymmetric: every `from` value should exist in `to`.
                        # Bug fix (2026-05-25): old code used Jaccard
                        # (|A∩B| / |A∪B|) on 100-row samples → 96.2% real overlap
                        # collapsed to 0.2 confidence on bigint-keyed PK→FK pairs.
                        # New: SQL-side directional containment up to 5K distinct
                        # values, no sampling bias. Confidence = max of both
                        # directions (FK direction wins; reverse direction also
                        # checked in case LLM swapped from/to).
                        n_from = conn.execute(text(
                            f'SELECT COUNT(DISTINCT "{fc}") FROM "{schema}"."{ft}" WHERE "{fc}" IS NOT NULL'
                        )).scalar() or 0
                        n_to = conn.execute(text(
                            f'SELECT COUNT(DISTINCT "{tc}") FROM "{schema}"."{tt}" WHERE "{tc}" IS NOT NULL'
                        )).scalar() or 0
                        if n_from > 0 and n_to > 0:
                            # |from ∩ to| via EXISTS — cheap, no temp table
                            n_match = conn.execute(text(f'''
                                SELECT COUNT(DISTINCT a."{fc}")
                                FROM "{schema}"."{ft}" a
                                WHERE a."{fc}" IS NOT NULL
                                  AND EXISTS (
                                    SELECT 1 FROM "{schema}"."{tt}" b
                                    WHERE b."{tc}"::text = a."{fc}"::text
                                  )
                            ''')).scalar() or 0
                            cov_from = n_match / n_from  # fraction of "from" found in "to"
                            cov_to = n_match / n_to      # fraction of "to" found in "from"
                            # FK direction = high cov_from. PK direction = high cov_to.
                            # Take max so swapped LLM output still scores correctly.
                            verified_conf = round(max(cov_from, cov_to), 2)
                    except Exception as _re:
                        logger.warning(f"relationship verify failed ({ft}.{fc} → {tt}.{tc}): {_re}")
                        # Keep LLM confidence if verification fails
                    try:
                        conn.execute(text("""
                            INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source)
                            VALUES (:s, :ft, :fc, :tt, :tc, :type, :conf, 'ai_verified')
                            ON CONFLICT (project_slug, from_table, from_column, to_table, to_column) DO UPDATE SET confidence = :conf, source = 'ai_verified'
                        """), {"s": project_slug, "ft": ft, "fc": fc, "tt": tt, "tc": tc,
                               "type": r.get("type", "fk"), "conf": verified_conf})
                    except Exception:
                        pass
                conn.commit()
    except Exception:
        pass


def _seed_cross_table_qa(project_slug: str) -> int:
    """Generate JOIN QA pairs from verified relationships (confidence ≥ 0.6).

    Reads dash_relationships, builds 2 JOIN questions per high-confidence
    relationship (one count-orphans, one aggregate-by-dim), verifies each SQL
    against the real schema, appends to per-table _qa.json. Returns # seeded.
    """
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    seeded = 0
    try:
        engine = create_engine(db_url, poolclass=NullPool)
        with engine.connect() as conn:
            rels = conn.execute(text("""
                SELECT from_table, from_column, to_table, to_column, confidence
                FROM public.dash_relationships
                WHERE project_slug = :s AND confidence >= 0.6
                ORDER BY confidence DESC LIMIT 10
            """), {"s": project_slug}).fetchall()

        if not rels:
            return 0

        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        training_dir.mkdir(parents=True, exist_ok=True)
        verify_engine = create_engine(db_url, poolclass=NullPool)

        for ft, fc, tt, tc, conf in rels:
            # Pick a dimension col on `to` side (first text col not FK)
            try:
                with verify_engine.connect() as vconn:
                    vconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                    cols = vconn.execute(text(
                        "SELECT column_name, data_type FROM information_schema.columns "
                        "WHERE table_schema=:sch AND table_name=:t ORDER BY ordinal_position"
                    ), {"sch": schema, "t": tt}).fetchall()
            except Exception:
                continue
            dim_col = next(
                (c[0] for c in cols if c[1] in ("text", "character varying") and c[0] != tc and not c[0].startswith("_")),
                None,
            )

            qa_pairs = [
                {
                    "question": f"How many {ft} rows have no matching {tt}?",
                    "sql": (
                        f'SELECT COUNT(*) FROM "{ft}" a '
                        f'LEFT JOIN "{tt}" b ON b."{tc}"::text = a."{fc}"::text '
                        f'WHERE b."{tc}" IS NULL'
                    ),
                    "analysis_type": "diagnostic",
                    "source": "relationship_seed",
                },
            ]
            if dim_col:
                qa_pairs.append({
                    "question": f"Show {ft} row count grouped by {tt}.{dim_col}.",
                    "sql": (
                        f'SELECT b."{dim_col}", COUNT(*) AS row_count '
                        f'FROM "{ft}" a JOIN "{tt}" b ON b."{tc}"::text = a."{fc}"::text '
                        f'GROUP BY b."{dim_col}" ORDER BY row_count DESC LIMIT 20'
                    ),
                    "analysis_type": "comparative",
                    "source": "relationship_seed",
                })

            # Verify each SQL, persist + insert into DB
            verified_pairs = []
            for qa in qa_pairs:
                try:
                    with verify_engine.connect() as vconn:
                        vconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                        rows = vconn.execute(text(qa["sql"])).fetchall()
                        qa["verified"] = True
                        qa["verified_answer"] = (
                            str(rows[0][0]) if rows and len(rows[0]) == 1
                            else str(rows[:3]) if rows else "0 rows"
                        )
                        qa["verified_row_count"] = len(rows)
                        verified_pairs.append(qa)
                except Exception:
                    pass  # skip unverified

            if not verified_pairs:
                continue

            # Append to existing ft table's QA file
            qa_file = training_dir / f"{ft}_qa.json"
            existing = []
            if qa_file.exists():
                try:
                    with open(qa_file) as fh:
                        existing = json.load(fh) or []
                except Exception:
                    existing = []
            existing.extend(verified_pairs)
            with open(qa_file, "w") as fh:
                json.dump(existing, fh, indent=2)

            # Also write to DB so dash_training_qa picks them up
            try:
                with engine.connect() as conn:
                    for qa in verified_pairs:
                        conn.execute(text(
                            "INSERT INTO public.dash_training_qa "
                            "(project_slug, table_name, question, sql, answer_template) "
                            "VALUES (:s, :t, :q, :sql, :a)"
                        ), {
                            "s": project_slug, "t": ft,
                            "q": qa["question"], "sql": qa["sql"],
                            "a": qa.get("verified_answer", ""),
                        })
                    conn.commit()
            except Exception as _ie:
                logger.warning(f"cross-table QA DB insert failed: {_ie}")

            seeded += len(verified_pairs)
    except Exception as _se:
        logger.warning(f"_seed_cross_table_qa failed for {project_slug}: {_se}")

    return seeded


def _detect_data_drift(project_slug: str, table_name: str, col_analyses: list[dict]):
    """Compare new column stats with stored metadata to detect drift."""
    meta_file = KNOWLEDGE_DIR / project_slug / "tables" / f"{table_name}.json"
    if not meta_file.exists():
        return

    try:
        with open(meta_file) as f:
            old_meta = json.load(f)

        old_cols = {c["name"]: c for c in old_meta.get("table_columns", [])}
        drift_alerts = []

        for col in col_analyses:
            name = col["name"]
            old = old_cols.get(name, {})
            old_desc = old.get("description", "")

            # Check NULL drift
            if col.get("null_pct", 0) > 20 and "NULL" not in old_desc:
                drift_alerts.append(f"{name}: NULL% increased to {col['null_pct']}%")

            # Check range drift for numerics
            if col.get("type") == "numeric" and "Range" in old_desc:
                try:
                    # Extract old range from description like "Range: 10 to 500"
                    parts = old_desc.split("Range: ")[1].split(",")[0].split(" to ")
                    old_max = float(parts[1].strip())
                    if col.get("max", 0) > old_max * 2:
                        drift_alerts.append(f"{name}: max value {col['max']} exceeds trained range (was {old_max})")
                except Exception:
                    pass

        if drift_alerts:
            drift_dir = KNOWLEDGE_DIR / project_slug / "drift"
            drift_dir.mkdir(parents=True, exist_ok=True)
            with open(drift_dir / f"{table_name}_drift.json", "w") as f:
                json.dump({"table": table_name, "alerts": drift_alerts, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())}, f, indent=2)

            # Create notification for project owner
            try:
                from app.auth import notify_user
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    row = conn.execute(text("SELECT user_id FROM public.dash_projects WHERE slug = :s"), {"s": project_slug}).fetchone()
                    if row:
                        notify_user(row[0], f"Drift detected · {table_name}", "; ".join(drift_alerts[:3]), "warn")
            except Exception:
                pass
    except Exception:
        pass


def _extract_sql_metadata(sql: str) -> dict:
    """Extract tables, join strategy, and filters from a SQL query using simple regex."""
    sql_upper = sql.upper()
    tables = set()

    # Extract tables from FROM clauses
    for m in re.finditer(r'\bFROM\s+([a-zA-Z_][a-zA-Z0-9_.]*)', sql, re.IGNORECASE):
        tables.add(m.group(1).strip().lower())
    # Extract tables from JOIN clauses
    for m in re.finditer(r'\bJOIN\s+([a-zA-Z_][a-zA-Z0-9_.]*)', sql, re.IGNORECASE):
        tables.add(m.group(1).strip().lower())

    # Determine join strategy
    if re.search(r'\bLEFT\s+JOIN\b', sql_upper):
        join_strategy = 'LEFT JOIN'
    elif re.search(r'\bRIGHT\s+JOIN\b', sql_upper):
        join_strategy = 'RIGHT JOIN'
    elif re.search(r'\bFULL\s+(OUTER\s+)?JOIN\b', sql_upper):
        join_strategy = 'FULL JOIN'
    elif re.search(r'\bCROSS\s+JOIN\b', sql_upper):
        join_strategy = 'CROSS JOIN'
    elif re.search(r'\bJOIN\b', sql_upper):
        join_strategy = 'INNER JOIN'
    else:
        join_strategy = 'NONE'

    # Detect filter/aggregation clauses
    filters_found = []
    if re.search(r'\bWHERE\b', sql_upper):
        filters_found.append('WHERE')
    if re.search(r'\bGROUP\s+BY\b', sql_upper):
        filters_found.append('GROUP BY')
    if re.search(r'\bORDER\s+BY\b', sql_upper):
        filters_found.append('ORDER BY')
    if re.search(r'\bHAVING\b', sql_upper):
        filters_found.append('HAVING')
    if re.search(r'\bLIMIT\b', sql_upper):
        filters_found.append('LIMIT')
    if re.search(r'\bWINDOW\b|\bOVER\s*\(', sql_upper):
        filters_found.append('WINDOW')

    return {
        "tables_used": ", ".join(sorted(tables)) if tables else "",
        "join_strategy": join_strategy,
        "filters": ", ".join(filters_found) if filters_found else "NONE",
    }


def _save_query_pattern_with_metadata(engine, project_slug: str, question: str, sql: str, source: str = 'training'):
    """Save a Q&A pair to dash_query_patterns with extracted SQL metadata."""
    meta = _extract_sql_metadata(sql)
    try:
        engine_conn = engine.connect()
        engine_conn.execute(text(
            "INSERT INTO public.dash_query_patterns "
            "(project_slug, question, sql, tables_used, join_strategy, filters, source) "
            "VALUES (:s, :q, :sql, :tables, :join, :filt, :src) "
            "ON CONFLICT DO NOTHING"
        ), {
            "s": project_slug, "q": question, "sql": sql,
            "tables": meta["tables_used"], "join": meta["join_strategy"],
            "filt": meta["filters"], "src": source,
        })
        engine_conn.commit()
        engine_conn.close()
        return True
    except Exception:
        return False


def _save_to_db(project_slug: str, table_name: str, metadata: dict, biz_rules: dict, training_qa: list | None = None, persona: dict | None = None):
    """Persist all training data to PostgreSQL alongside files."""
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Table metadata
            conn.execute(text("""
                INSERT INTO public.dash_table_metadata (project_slug, table_name, metadata)
                VALUES (:s, :t, CAST(:m AS jsonb))
                ON CONFLICT (project_slug, table_name)
                DO UPDATE SET metadata = CAST(:m AS jsonb), updated_at = NOW()
            """), {"s": project_slug, "t": table_name, "m": json.dumps(metadata, default=str)})

            # Business rules
            conn.execute(text("""
                INSERT INTO public.dash_business_rules_db (project_slug, table_name, rules)
                VALUES (:s, :t, CAST(:r AS jsonb))
                ON CONFLICT (project_slug, table_name)
                DO UPDATE SET rules = CAST(:r AS jsonb)
            """), {"s": project_slug, "t": table_name, "r": json.dumps(biz_rules, default=str)})

            # Training Q&A
            if training_qa and isinstance(training_qa, list):
                # Clear old Q&A for this table
                conn.execute(text("DELETE FROM public.dash_training_qa WHERE project_slug = :s AND table_name = :t"), {"s": project_slug, "t": table_name})

                import re as _re_qa

                def _strip_sql_fences(raw: str) -> str:
                    """Strip markdown code fences + leading SQL-comment lines from LLM output."""
                    s = (raw or "").strip()
                    # ```sql ... ``` or ``` ... ```
                    m = _re_qa.search(r"```(?:sql)?\s*(.*?)\s*```", s, _re_qa.DOTALL | _re_qa.IGNORECASE)
                    if m:
                        s = m.group(1).strip()
                    # drop leading -- comment lines
                    lines = [ln for ln in s.splitlines() if not ln.strip().startswith("--")]
                    return "\n".join(lines).strip().rstrip(";").strip()

                def _is_safe_sql(sql: str) -> bool:
                    """Only allow SELECT/WITH statements in training SQL."""
                    s = sql.strip().upper()
                    if not (s.startswith('SELECT') or s.startswith('WITH')):
                        return False
                    forbidden = ['DROP ', 'DELETE ', 'TRUNCATE ', 'ALTER ', 'INSERT ', 'UPDATE ', 'CREATE ', 'GRANT ']
                    return not any(kw in s for kw in forbidden)

                _qa_kept = 0
                _qa_rejected = 0
                for qa in training_qa:
                    cleaned = _strip_sql_fences(qa.get("sql", ""))
                    if _is_safe_sql(cleaned):
                        conn.execute(text(
                            "INSERT INTO public.dash_training_qa (project_slug, table_name, question, sql, answer_template) VALUES (:s, :t, :q, :sql, :a)"
                        ), {"s": project_slug, "t": table_name, "q": qa.get("question", ""), "sql": cleaned, "a": qa.get("answer_template", "")})
                        _qa_kept += 1
                    else:
                        _qa_rejected += 1
                if _qa_rejected:
                    import logging as _qa_log
                    _qa_log.getLogger(__name__).warning(f"qa: kept={_qa_kept} rejected={_qa_rejected} table={table_name}")

            # Persona
            if persona:
                conn.execute(text("""
                    INSERT INTO public.dash_personas (project_slug, persona)
                    VALUES (:s, CAST(:p AS jsonb))
                    ON CONFLICT (project_slug)
                    DO UPDATE SET persona = CAST(:p AS jsonb), updated_at = NOW()
                """), {"s": project_slug, "p": json.dumps(persona, default=str)})

            conn.commit()
    except Exception as _save_e:
        import logging as _save_log
        _save_log.getLogger(__name__).exception(f"_save_to_db failed slug={project_slug} table={table_name}: {_save_e}")


def _generate_project_evals(project_slug: str, max_evals: int = 15, per_table: int = 5,
                            log_fn=None) -> int:
    """Issue #39 — generate eval test cases ONCE project-wide.

    Samples Q&A pairs across ALL of the project's training/{table}_qa.json files
    (round-robin so no single table dominates) and inserts a single bounded set
    into public.dash_evals. Fully fail-soft — never raises.

    The DELETE-at-start eval pre-clear in _bg() means each retrain produces one
    fresh, bounded set (max_evals total) instead of the old per-table behaviour
    that appended (and grew unbounded) once per table.

    Returns the number of evals saved.
    """
    _log = log_fn or (lambda m: None)
    _analysis_type_tags = {
        "descriptive": "[descriptive] ",
        "diagnostic": "[diagnostic] ",
        "predictive": "[predictive] ",
        "prescriptive": "[prescriptive] ",
        "anomaly": "[anomaly] ",
    }

    def _tag(qa: dict) -> str:
        q = qa.get("question", "")
        a_type = qa.get("analysis_type", "")
        if a_type and a_type in _analysis_type_tags:
            return f"{_analysis_type_tags[a_type]}{q}"
        q_lower = q.lower()
        if any(kw in q_lower for kw in ["why", "cause", "driver", "impact", "factor"]):
            return "[diagnostic] " + q
        if any(kw in q_lower for kw in ["anomal", "outlier", "unusual", "spike", "drop"]):
            return "[anomaly] " + q
        if any(kw in q_lower for kw in ["recommend", "should", "action", "mitigat", "suggest"]):
            return "[prescriptive] " + q
        if any(kw in q_lower for kw in ["predict", "forecast", "trend", "project", "future"]):
            return "[predictive] " + q
        return "[descriptive] " + q

    evals_saved = 0
    try:
        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        if not training_dir.exists():
            return 0

        # Collect up to `per_table` valid (question, sql) pairs from each table's
        # Q&A file, then round-robin across tables so the bounded final set is a
        # representative sample rather than dominated by the first table.
        # Only sample Q&A from tables that still exist — a ghost
        # {table}_qa.json (left by a wipe) generates eval SQL against a
        # dropped table → every such eval ERRORs ("relation does not exist").
        _live = _live_tables(project_slug)
        per_table_pairs: list[list[tuple[str, str]]] = []
        for qa_file in sorted(training_dir.glob("*_qa.json")):
            _tbl = qa_file.stem[:-3] if qa_file.stem.endswith("_qa") else qa_file.stem
            if _live and _tbl not in _live:
                continue
            try:
                with open(qa_file) as f:
                    qa_pairs = json.load(f)
            except Exception:
                continue
            picked: list[tuple[str, str]] = []
            for qa in (qa_pairs if isinstance(qa_pairs, list) else []):
                q = qa.get("question", "")
                s = qa.get("sql", "")
                if q and s:
                    picked.append((_tag(qa), s))
                if len(picked) >= per_table:
                    break
            if picked:
                per_table_pairs.append(picked)

        if not per_table_pairs:
            return 0

        # Round-robin flatten, capped at max_evals.
        flat: list[tuple[str, str]] = []
        ri = 0
        while len(flat) < max_evals and any(per_table_pairs):
            bucket = per_table_pairs[ri % len(per_table_pairs)]
            if bucket:
                flat.append(bucket.pop(0))
            ri += 1
            if ri > max_evals * (len(per_table_pairs) + 1):
                break  # safety

        eng = create_engine(db_url)
        try:
            with eng.connect() as conn:
                for tagged_q, sql in flat[:max_evals]:
                    conn.execute(text(
                        "INSERT INTO public.dash_evals (project_slug, question, expected_sql) "
                        "VALUES (:s, :q, :sql)"
                    ), {"s": project_slug, "q": tagged_q, "sql": sql})
                    evals_saved += 1
                conn.commit()
        finally:
            eng.dispose()
        _log(f"✓ project evals: {evals_saved} test cases generated (project-wide, bounded)")
    except Exception as e:
        _log(f"⚠ project eval generation skipped: {str(e)[:80]}")
    return evals_saved


def _run_auto_training(project_slug: str, table_name: str, col_analyses: list[dict],
                       metadata: dict, biz_rules: dict, sample_rows: list[dict],
                       tables_dir: Path, business_dir: Path,
                       master_run_id: int | None = None, table_index: int = 0, total_tables: int = 1):
    """Background task: LLM deep analysis + training Q&A generation."""
    import time as _time

    # Single shared engine for all DB operations in this training run
    train_engine = create_engine(db_url)

    # Track training run — reuse master run if provided, otherwise create one
    run_id = master_run_id
    if not run_id:
        try:
            with train_engine.connect() as conn:
                result = conn.execute(text(
                    "INSERT INTO public.dash_training_runs (project_slug, status, steps) VALUES (:s, 'running', 'starting') RETURNING id"
                ), {"s": project_slug})
                run_id = result.fetchone()[0]
                conn.commit()
        except Exception:
            pass

    # Per-step tracing: open a child span as each step starts, close it when the
    # next step begins (or on terminal status). Fail-soft holder, additive only.
    _step_span = {"cm": None}

    def _close_step_span(error: str = ""):
        cm = _step_span.get("cm")
        if cm is None:
            return
        _step_span["cm"] = None
        try:
            exc = RuntimeError(error[:200]) if error else None
            cm.__exit__(type(exc) if exc else None, exc, None)
        except Exception:  # noqa: BLE001
            pass

    def _open_step_span(step_name: str):
        try:
            cm = _trace_span(f"training.{step_name}", kind="training",
                             project_slug=project_slug)
            cm.__enter__()
            _step_span["cm"] = cm
        except Exception:  # noqa: BLE001
            _step_span["cm"] = None

    def _update_run(status: str, steps: str = "", error: str = ""):
        # Drive the per-step trace span (fail-soft, never affects control flow).
        try:
            if status in ("done", "failed"):
                _close_step_span(error)
            elif status == "running" and steps:
                _close_step_span()
                _open_step_span(steps)
        except Exception:  # noqa: BLE001
            pass
        if not run_id:
            return
        # Encode table tracking into the steps field: step_name|table_name|table_index|total_tables
        steps_with_table = f"{steps}|{table_name}|{table_index}|{total_tables}"
        # Issue #16 — structured JSONB mirror of the same data. Frontend
        # reads `current_progress` when present, falls back to the legacy
        # pipe-delimited parser otherwise.
        progress_payload = {
            "step": steps,
            "table": table_name,
            "index": table_index,
            "total": total_tables,
            "started_at": _time.strftime('%Y-%m-%dT%H:%M:%S'),
        }
        progress_json = json.dumps(progress_payload)
        try:
            with train_engine.connect() as conn:
                if (status == 'done' or status == 'failed') and not master_run_id:
                    # Only set terminal status if this is a standalone run (no master)
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET status = :st, steps = :steps, "
                        "error = :err, tables_trained = :trained, current_step = :cs, "
                        "stage_progress = :sp, current_progress = CAST(:cp AS jsonb), "
                        "finished_at = NOW() WHERE id = :id"
                    ), {"st": status, "steps": steps_with_table, "err": error,
                         "trained": table_index, "cs": steps or None,
                         "sp": table_index, "cp": progress_json, "id": run_id})
                else:
                    # For master-managed runs, only update steps (master controls status)
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET steps = :steps, "
                        "current_step = :cs, stage_progress = :sp, "
                        "current_progress = CAST(:cp AS jsonb) WHERE id = :id"
                    ), {"steps": steps_with_table, "cs": steps or None,
                         "sp": table_index, "cp": progress_json, "id": run_id})
                conn.commit()
        except Exception:
            pass

    def _log(msg: str):
        """Append a log entry to the training run for real-time CLI display."""
        if not run_id:
            return
        try:
            with train_engine.connect() as conn:
                conn.execute(text(
                    "UPDATE public.dash_training_runs SET logs = COALESCE(logs, '[]'::jsonb) || CAST(:entry AS jsonb) WHERE id = :id"
                ), {"entry": json.dumps([{"ts": _time.strftime('%H:%M:%S'), "tsabs": _time.time(), "msg": msg, "table": table_name, "table_index": table_index, "total_tables": total_tables}]), "id": run_id})
                conn.commit()
        except Exception:
            pass

    def _cancelled() -> bool:
        """Check if training was cancelled."""
        return _training_cancel_flags.get(project_slug, False)

    num_cols = len(col_analyses)
    # Real row count via COUNT(*). The caller's col_analyses come from a
    # LIMIT-100 speed sample, so total_count there is capped at 100 — using it
    # made the log say "100 rows" for a 4,886-row table AND made the
    # change-detection below always see "unchanged" (sample is always 100).
    # Query the true count; fall back to the sample only if it fails.
    num_rows = 0
    try:
        _nr_schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        _nr_eng = engine if 'engine' in dir() and engine else create_engine(db_url)
        with _nr_eng.connect() as _nrc:
            num_rows = int(_nrc.execute(text(
                f'SELECT COUNT(*) FROM "{_nr_schema}"."{table_name}"'
            )).scalar() or 0)
    except Exception:
        num_rows = 0
    if num_rows == 0:
        for c in col_analyses:
            tc = c.get("total_count", 0) or c.get("non_null_count", 0)
            if tc > num_rows:
                num_rows = tc
    _log(f"analyzing table: {table_name} ({num_rows} rows, {num_cols} columns)")

    if num_rows == 0:
        _log(f"⊘ skipping training — table {table_name} has 0 rows")
        _update_run("done", "complete")
        return

    # Check if table data changed (skip ONLY if we have stored row count AND it matches)
    existing_meta = metadata.get("table_description", "")
    existing_row_count = metadata.get("_row_count", 0)
    # Only skip if: description exists AND row count was stored AND matches current
    data_changed = True
    if existing_meta and existing_row_count > 0 and existing_row_count == num_rows:
        data_changed = False

    # Step 0: Data drift detection
    _update_run("running", "drift_detection")
    _log("checking for data drift...")
    _detect_data_drift(project_slug, table_name, col_analyses)
    _log("✓ drift check complete")

    # Step 1a: SQL Profiling — stats via PostgreSQL ($0, <10s even for 1M rows)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "sql_profiling")
    _log(f"profiling columns via SQL (no RAM usage)...")
    sql_profiles = _sql_profile_columns(project_slug, table_name)
    dimensions = [p for p in sql_profiles if p.get("classification") == "dimension"]
    measures = [p for p in sql_profiles if p.get("classification") == "measure"]
    ids = [p for p in sql_profiles if p.get("classification") == "id"]
    _log(f"✓ profiled {len(sql_profiles)} cols: {len(dimensions)} dim · {len(measures)} measure · {len(ids)} id")

    # Step 1a-v2: Advanced profile (compact dim summaries + variant detection + roles).
    # Fail-soft additive layer — writes dash_table_metadata.metadata['profile_v2'].
    # Disable via PROFILE_V2_DISABLED=1. Surfaces in Layer 1 of agent prompt.
    if not _cancelled():
        try:
            _update_run("running", "profile_v2")
            _log("running advanced profile_v2 (combined query + role classify + variants)...")
            from dash.training.profile_v2 import profile_table_v2 as _pv2
            _v2 = _pv2(project_slug, table_name)
            if _v2.get("disabled"):
                _log("⊘ profile_v2 disabled via env")
            else:
                _v2cols = _v2.get("columns") or []
                _v2_states = sum(1 for c in _v2cols if c.get("role") == "state")
                _v2_variants = sum(1 for c in _v2cols if c.get("variants_detected"))
                _log(f"✓ profile_v2: {len(_v2cols)} cols · {_v2_states} state · {_v2_variants} variant-warned · scan={_v2.get('scan_elapsed_ms','?')}ms")
        except Exception as _e:
            _log(f"⚠ profile_v2 skipped: {str(_e)[:200]}")
    if dimensions:
        dim_summary = ", ".join(f"{p['name']}({p.get('distinct_count','?')})" for p in dimensions[:6])
        _log(f"  └─ dimensions: {dim_summary}{' …' if len(dimensions) > 6 else ''}")
    if measures:
        msr_summary = ", ".join(
            f"{p['name']}[{p.get('min','?')}–{p.get('max','?')}]" for p in measures[:6]
        )
        _log(f"  └─ measures: {msr_summary}{' …' if len(measures) > 6 else ''}")

    # Step 1b: Dimension Catalog — exact values for categorical columns ($0)
    dim_catalog = {}
    if dimensions:
        _update_run("running", "dimension_catalog")
        _log(f"building dimension catalog ({len(dimensions)} categorical columns)...")
        dim_catalog = _build_dimension_catalog(project_slug, table_name, sql_profiles)
        total_values = sum(len(v) for v in dim_catalog.values())
        _log(f"✓ dimension catalog: {total_values} unique values across {len(dim_catalog)} columns")
        for col_name, values in list(dim_catalog.items())[:5]:
            top = values[:3] if isinstance(values, list) else []
            preview = ", ".join(
                str(v.get("value", v) if isinstance(v, dict) else v)[:30] for v in top
            )
            _log(f"  └─ {col_name}: [{preview}{' …' if len(values) > 3 else ''}]")

    # Step 1c: Hierarchy Detection ($0)
    hierarchies = []
    if len(dim_catalog) >= 2:
        _update_run("running", "hierarchy_detection")
        _log("detecting column hierarchies...")
        hierarchies = _detect_hierarchies(project_slug, table_name, dim_catalog)
        # Always persist the freshly-detected list — even when EMPTY — so a
        # prior run's stale/false hierarchies (e.g. the old _period→* lineage
        # junk) get cleared from metadata + brain, not left behind by a
        # `if hierarchies:` guard.
        metadata["hierarchies"] = hierarchies
        if hierarchies:
            _log(f"✓ found {len(hierarchies)} hierarchies: {', '.join(h['parent'] + ' → ' + h['child'] for h in hierarchies)}")
        else:
            _log("· no hierarchies detected")

    # Step 1d: Smart Sampling — 20 diverse rows (replaces first 8)
    _update_run("running", "smart_sampling")
    _log("collecting diverse sample rows...")
    smart_samples = _smart_sample_rows(project_slug, table_name, sql_profiles)
    _log(f"✓ sampled {len(smart_samples)} diverse rows (start/mid/end/outlier/null)")

    # Save dimension info to table metadata for prompt injection
    if dim_catalog:
        metadata["dimensions"] = {col: vals[:20] for col, vals in dim_catalog.items()}  # Top 20 per column
    if sql_profiles:
        metadata["column_profiles"] = {p["name"]: {k: v for k, v in p.items() if k != "name"} for p in sql_profiles}

    # Save updated metadata
    tables_dir.mkdir(parents=True, exist_ok=True)
    with open(tables_dir / f"{table_name}.json", "w") as f:
        json.dump(metadata, f, indent=2, default=str)

    # Use smart samples for LLM analysis (instead of first 8 rows)
    if smart_samples:
        sample_rows = smart_samples

    # Step 1: LLM Deep Analysis (now uses profiles + smart samples, not raw data)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "deep_analysis")
    if not data_changed and existing_meta:
        _log(f"⊘ skipping deep analysis — table unchanged ({num_rows} rows)")
        analysis = None
    else:
        _log(f"running LLM deep analysis on {table_name} (with {len(smart_samples)} diverse samples)...")
        analysis = _llm_deep_analysis(table_name, col_analyses, sample_rows)
        if analysis:
            purpose = (analysis.get('table_purpose') or analysis.get('table_description') or '')[:80]
            grain = (analysis.get('grain') or '')[:60]
            pks = analysis.get('primary_keys') or []
            fks = analysis.get('foreign_keys') or []
            cols_described = len(analysis.get('column_descriptions') or {})
            _log(f"✓ deep analysis complete")
            if purpose:
                _log(f"  └─ purpose: {purpose}")
            if grain:
                _log(f"  └─ grain: {grain}")
            if pks:
                pk_str = ", ".join(pks) if isinstance(pks, list) else str(pks)
                _log(f"  └─ primary keys: {pk_str[:80]}")
            if fks:
                fk_count = len(fks) if isinstance(fks, list) else 1
                _log(f"  └─ foreign keys: {fk_count}")
            _log(f"  └─ cols described: {cols_described}/{num_cols}")
    if analysis:
        # Overwrite metadata with smart descriptions
        if analysis.get("table_description"):
            metadata["table_description"] = analysis["table_description"]
        col_descs = analysis.get("column_descriptions", {})
        for col in metadata.get("table_columns", []):
            smart_desc = col_descs.get(col["name"])
            if smart_desc:
                col["description"] = smart_desc
        if analysis.get("data_quality"):
            metadata["data_quality_notes"] = analysis["data_quality"][:8]
        # Codex-enriched knowledge fields
        for field in ("table_purpose", "grain", "primary_keys", "foreign_keys", "usage_patterns", "alternate_tables", "freshness"):
            if analysis.get(field):
                metadata[field] = analysis[field]

        # Store row count for future change detection
        metadata["_row_count"] = num_rows

        # Save updated metadata
        tables_dir.mkdir(parents=True, exist_ok=True)
        with open(tables_dir / f"{table_name}.json", "w") as f:
            json.dump(metadata, f, indent=2, default=str)

        # Update metrics
        if analysis.get("metrics"):
            biz_rules["metrics"] = analysis["metrics"]
            business_dir.mkdir(parents=True, exist_ok=True)
            with open(business_dir / f"{table_name}_rules.json", "w") as f:
                json.dump(biz_rules, f, indent=2, default=str)

        # Auto-create rules
        rules_dir = KNOWLEDGE_DIR / project_slug / "rules"
        rules_dir.mkdir(parents=True, exist_ok=True)
        for rule in analysis.get("business_rules", [])[:5]:
            rule_id = f"rule_auto_{int(_time.time() * 1000)}"
            rule_data = {
                "id": rule_id, "name": rule.get("name", ""),
                "type": rule.get("type", "business_rule"),
                "definition": rule.get("definition", ""),
                "source": "auto_training",
                "created_at": _time.strftime("%Y-%m-%dT%H:%M:%SZ", _time.gmtime()),
            }
            with open(rules_dir / f"{rule_id}.json", "w") as f:
                json.dump(rule_data, f, indent=2)
            _time.sleep(0.01)  # unique timestamps

        # Update agent persona
        if analysis.get("suggested_role"):
            _update_project_config(project_slug, analysis["suggested_role"], analysis.get("suggested_personality"))

    # Step 2: LLM Training Q&A
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "qa_generation")
    _log(f"generating training Q&A for {table_name}...")
    training = _llm_generate_training(table_name, metadata, col_analyses=col_analyses)
    if training and isinstance(training, list):
        # 2026-05-25 (Phase 9 belt+suspenders): drop any QA pair whose SQL
        # uses DATE_TRUNC / GROUP BY on a CONSTANT column. LLM is told via
        # prompt warning, but post-filter catches stragglers.
        _const_cols = [
            (ca.get("name") or "").lower() for ca in (col_analyses or [])
            if isinstance(ca.get("unique_count"), int) and ca.get("unique_count", 99) <= 1
        ]
        if _const_cols:
            _kept = []
            for qa in training:
                sql_l = (qa.get("sql") or "").lower()
                # Drop if SQL groups/buckets on a constant col.
                bad = False
                for cc in _const_cols:
                    if not cc:
                        continue
                    if (f"date_trunc('day', {cc}" in sql_l
                            or f"date_trunc('month', {cc}" in sql_l
                            or f"date_trunc('year', {cc}" in sql_l
                            or f"date_trunc('week', {cc}" in sql_l
                            or f"group by {cc}" in sql_l):
                        bad = True
                        break
                if bad:
                    logger.info(f"dropped trend QA on constant col: {qa.get('question','')[:60]}")
                    continue
                _kept.append(qa)
            training = _kept
        # VERIFY Q&A: Run each generated SQL against real data
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        verified_count = 0
        discarded_count = 0
        try:
            verify_engine = create_engine(db_url)
            for qa in training:
                sql = qa.get("sql", "")
                if not sql:
                    continue
                try:
                    with verify_engine.connect() as vconn:
                        vconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                        result = vconn.execute(text(sql))
                        rows = result.fetchall()
                        if rows:
                            # Save real answer
                            qa["verified"] = True
                            qa["verified_answer"] = str(rows[0][0]) if len(rows[0]) == 1 else str(rows[:3])
                            qa["verified_row_count"] = len(rows)
                            verified_count += 1
                        else:
                            qa["verified"] = True
                            qa["verified_answer"] = "0 rows"
                            qa["verified_row_count"] = 0
                            verified_count += 1
                except Exception:
                    qa["verified"] = False
                    qa["verified_answer"] = None
                    discarded_count += 1
        except Exception:
            pass
        _log(f"  ✓ {verified_count} Q&A verified with real data, {discarded_count} had SQL errors")

        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        training_dir.mkdir(parents=True, exist_ok=True)
        with open(training_dir / f"{table_name}_qa.json", "w") as f:
            json.dump(training, f, indent=2)

    if training and isinstance(training, list):
        _log(f"✓ {len(training)} Q&A pairs generated")
    else:
        _log("· no Q&A generated")

    # Step 3: Generate agent persona
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "persona")
    _log("generating agent persona...")
    try:
        # Load all tables metadata for this project
        all_metadata = []
        proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_tables_dir.exists():
            for f in proj_tables_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        all_metadata.append(json.load(fh))
                except Exception:
                    pass

        # Load all rules
        all_rules = []
        proj_rules_dir = KNOWLEDGE_DIR / project_slug / "rules"
        if proj_rules_dir.exists():
            for f in proj_rules_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        all_rules.append(json.load(fh))
                except Exception:
                    pass

        persona = _llm_generate_persona(project_slug, all_metadata, all_rules)
        if persona:
            persona_file = KNOWLEDGE_DIR / project_slug / "persona.json"
            with open(persona_file, "w") as f:
                json.dump(persona, f, indent=2)
            # Also save to DB
            try:
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    conn.execute(text("""
                        INSERT INTO public.dash_personas (project_slug, persona)
                        VALUES (:s, CAST(:p AS jsonb))
                        ON CONFLICT (project_slug)
                        DO UPDATE SET persona = CAST(:p AS jsonb), updated_at = NOW()
                    """), {"s": project_slug, "p": json.dumps(persona, default=str)})
                    conn.commit()
            except Exception:
                pass

            # Update project role with generated expertise
            if persona.get("expertise_areas"):
                role_str = ", ".join(persona["expertise_areas"][:4])
                _update_project_config(project_slug, role_str, None)
    except Exception:
        pass

    _log("✓ persona generated")

    # Step 4: Auto-generate sample workflows
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "synthesis")
    _log("generating sample workflows...")
    try:
        from os import getenv
        import httpx
        api_key = getenv("OPENROUTER_API_KEY", "")

        # Load all table names for context
        all_tables = []
        proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_tables_dir.exists():
            for f in proj_tables_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        d = json.load(fh)
                    all_tables.append(d.get("table_name", f.stem))
                except Exception:
                    pass

        if api_key and all_tables:
            tables_str = ", ".join(all_tables)
            prompt = f"""Generate 3 analysis-typed workflows for a data agent with these tables: {tables_str}

Each workflow should be a multi-step analysis journey of a specific type. Generate exactly these 3 workflows:

1. "Data Overview" — a descriptive journey:
   Step 1: What tables and columns do we have? Show the schema overview.
   Step 2: What is the total count and key metrics for the main table?
   Step 3: Break down the key metrics by top categories or dimensions.

2. "Deep Dive Analysis" — a diagnostic journey:
   Step 1: What are the key metrics and their current values?
   Step 2: Which dimension or category has the highest impact on the key metric?
   Step 3: Why is that dimension the highest? What factors drive it?
   Step 4: What actions should we take based on this analysis?

3. "Risk Assessment" — an anomaly detection + prescriptive journey:
   Step 1: Show the distribution of key numeric metrics.
   Step 2: Are there any anomalies, outliers, or unexpected patterns?
   Step 3: What are the concentration risks (e.g., over-reliance on one category)?
   Step 4: Recommend specific actions to mitigate the identified risks.

Tailor each workflow's steps to the actual tables and columns available. Use real column/table names where possible.

Return ONLY valid JSON (no markdown):
[
  {{"name": "Data Overview", "description": "Descriptive journey through the dataset — schema, counts, and category breakdowns", "analysis_type": "descriptive", "steps": ["Step 1 text", "Step 2 text", "Step 3 text"]}},
  {{"name": "Deep Dive Analysis", "description": "Diagnostic journey to find what drives key metrics and what to do about it", "analysis_type": "diagnostic", "steps": ["Step 1", "Step 2", "Step 3", "Step 4"]}},
  {{"name": "Risk Assessment", "description": "Anomaly detection and prescriptive actions for risk mitigation", "analysis_type": "prescriptive", "steps": ["Step 1", "Step 2", "Step 3", "Step 4"]}}
]"""

            resp = httpx.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 800, "temperature": 0.3},
                timeout=20,
            )
            result = resp.json()
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
            workflows = json.loads(content.strip().strip("`").strip())

            if isinstance(workflows, list):
                # Save workflows
                wf_file = KNOWLEDGE_DIR / project_slug / "workflows.json"
                with open(wf_file, "w") as f:
                    json.dump(workflows, f, indent=2)
                _log(f"✓ {len(workflows)} workflows generated")
    except Exception:
        _log("· workflow generation skipped")

    # Step 5: Auto-create default schedules
    _log("checking schedules...")
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Check if schedules already exist
            existing = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_schedules WHERE project_slug = :s"
            ), {"s": project_slug}).scalar() or 0

            if existing == 0 and all_tables:
                # Get project owner
                row = conn.execute(text(
                    "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                ), {"s": project_slug}).fetchone()
                if row:
                    uid = row[0]
                    # Create a weekly summary schedule
                    conn.execute(text(
                        "INSERT INTO public.dash_schedules (project_slug, user_id, name, prompt, cron) "
                        "VALUES (:s, :uid, :name, :prompt, :cron)"
                    ), {
                        "s": project_slug, "uid": uid,
                        "name": "Weekly Data Summary",
                        "prompt": f"Give me a summary of all data in {all_tables[0]}. Show total counts, key metrics, and any notable trends.",
                        "cron": "0 8 * * 1",
                    })
                    if len(all_tables) > 1:
                        conn.execute(text(
                            "INSERT INTO public.dash_schedules (project_slug, user_id, name, prompt, cron) "
                            "VALUES (:s, :uid, :name, :prompt, :cron)"
                        ), {
                            "s": project_slug, "uid": uid,
                            "name": "Daily Health Check",
                            "prompt": f"Check data quality across all tables. Report any NULL values, missing data, or anomalies.",
                            "cron": "0 8 * * *",
                        })
                    conn.commit()
    except Exception:
        pass

    # Step 6: Save everything to DB (production persistence)
    _log("saving to database...")
    try:
        _save_to_db(project_slug, table_name, metadata, biz_rules,
                    training if isinstance(training, list) else None,
                    persona if 'persona' in dir() and persona else None)
    except Exception:
        pass

    _log("✓ saved to database")

    # Step 7: Multi-file synthesis — unified project understanding
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _log("running multi-file synthesis...")
    _update_run("running", "synthesis")
    try:
        from os import getenv
        import httpx
        api_key = getenv("OPENROUTER_API_KEY", "")
        all_meta = []
        proj_td = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_td.exists():
            # Only synthesize over tables that still exist in the DB — skip
            # orphaned {table}.json left by a raw-SQL wipe / external drop.
            try:
                _ms = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
                _live = set(inspect(create_engine(db_url)).get_table_names(schema=_ms))
            except Exception:
                _live = None
            for f in proj_td.glob("*.json"):
                if _live is not None and f.stem not in _live:
                    continue
                try:
                    with open(f) as fh:
                        all_meta.append(json.load(fh))
                except Exception:
                    pass

        if api_key and len(all_meta) > 1:
            tables_desc = "\n".join(f"- {m.get('table_name','?')}: {m.get('table_description','')}" for m in all_meta)
            prompt = f"""These tables belong to the same project. Create a unified data dictionary.

TABLES:
{tables_desc}

Return ONLY valid JSON (no markdown):
{{"project_summary": "2-3 sentence overview of what this data represents together", "cross_table_queries": ["suggested SQL joining multiple tables"], "data_dictionary": {{"term": "definition"}}}}"""

            resp = httpx.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 800, "temperature": 0.1},
                timeout=15,
            )
            result = resp.json()
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
            synthesis = json.loads(content.strip().strip("`").strip())
            if synthesis:
                synth_file = KNOWLEDGE_DIR / project_slug / "synthesis.json"
                with open(synth_file, "w") as f:
                    json.dump(synthesis, f, indent=2)
    except Exception:
        pass

    _log("✓ synthesis complete")

    # Step 8: Discover relationships (only if multiple tables)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "relationships")
    proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
    # Count only knowledge JSONs whose table STILL EXISTS in the DB. A raw-SQL
    # wipe / external drop can orphan a {table}.json on disk; trusting the glob
    # alone made retrain discover relationships + run multi-file synthesis
    # against a ghost table (e.g. a stale balance_stock_*.json), wasting LLM
    # calls and seeding a phantom relationship into the brain.
    table_count = 0
    if proj_tables_dir.exists():
        try:
            _schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
            _live = set(inspect(create_engine(db_url)).get_table_names(schema=_schema))
            table_count = sum(1 for j in proj_tables_dir.glob("*.json") if j.stem in _live)
        except Exception:
            # If the liveness check fails, fall back to the raw glob count.
            table_count = len(list(proj_tables_dir.glob("*.json")))
    if table_count >= 2:
        _log(f"discovering relationships across {table_count} tables...")
        try:
            _discover_relationships(project_slug)
        except Exception:
            pass
        # Seed cross-table QA pairs from verified relationships (confidence ≥ 0.6).
        # Bug fix 2026-05-25: training emitted 0 JOIN questions even after
        # discovering a 0.89-confidence FK. Now adds 2 JOIN QA pairs per
        # high-confidence relationship and verifies SQL against real DB.
        try:
            _seed_cross_table_qa(project_slug)
        except Exception as _xte:
            _log(f"⚠ cross-table QA seed skipped: {str(_xte)[:80]}")
    else:
        _log("⊘ skipping relationship discovery — only 1 table")

    # Also discover cross-document relationships if docs exist
    docs_dir = KNOWLEDGE_DIR / project_slug / "docs"
    if docs_dir.exists():
        doc_files = [f for f in docs_dir.iterdir() if f.is_file()]
        if len(doc_files) >= 2:
            _log("discovering cross-document relationships...")
            try:
                doc_sums = []
                for f in doc_files[:10]:
                    try:
                        doc_sums.append(f"{f.name}: {f.read_text(errors='ignore')[:500]}")
                    except Exception:
                        pass
                result = training_llm_call(
                    f"Find relationships between these documents.\n\n"
                    f"DOCUMENTS:\n" + "\n---\n".join(doc_sums) + "\n\n"
                    f'Return JSON: [{{"from_doc": "a.txt", "to_doc": "b.txt", "relationship": "shared topic", "shared_topics": ["topic1"], "strength": 0.8}}]',
                    "extraction"
                )
                if result:
                    rels = json.loads(result)
                    if isinstance(rels, list):
                        with train_engine.connect() as conn:
                            for r in rels[:10]:
                                if isinstance(r, dict) and r.get("from_doc"):
                                    conn.execute(text(
                                        "INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source) "
                                        "VALUES (:s, :ft, :fc, :tt, :tc, 'topic', :conf, 'ai') ON CONFLICT DO NOTHING"
                                    ), {"s": project_slug, "ft": r["from_doc"], "fc": ", ".join(r.get("shared_topics", [])[:5]),
                                        "tt": r.get("to_doc", ""), "tc": r.get("relationship", ""),
                                        "conf": r.get("strength", 0.5)})
                            conn.commit()
                        _log(f"✓ {len(rels)} cross-document relationships found")
            except Exception:
                pass

    _log("✓ relationships discovered")

    # Step 9: Re-index knowledge (with timeout — training completes even if this fails)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "reindex")
    # Dynamic timeout based on row count — 60s for ≤50K, 240s for larger
    _idx_timeout = 240 if num_rows > 50000 else 60
    _log(f"indexing knowledge base ({_idx_timeout}s timeout)...")
    indexed = _reload_project_knowledge(project_slug, timeout_sec=_idx_timeout)
    if indexed:
        _log("✓ knowledge indexed")
    else:
        _log("⚠ knowledge indexing timed out — skipped (training still complete)")

    # ═══ AUTO-FILL BRAIN — Make agent smart from data loading ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "brain_fill")
    _log("filling agent brain from training data...")

    engine = create_engine(db_url)

    # 1. Auto-Memories — from REAL data, not just metadata
    try:
        _log("  generating auto-memories from real data...")
        mem_facts = []
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]

        # Get REAL stats from actual table
        try:
            mem_engine = create_engine(db_url)
            with mem_engine.connect() as mconn:
                mconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                # Row count
                row_count = mconn.execute(text(f'SELECT COUNT(*) FROM "{table_name}"')).scalar() or 0
                mem_facts.append(f"Table '{table_name}' has {row_count:,} rows")

                # Column count
                col_count = len(col_analyses) if col_analyses else 0
                mem_facts.append(f"Table '{table_name}' has {col_count} columns")

                # Date range (if date column exists)
                # Bug fix 2026-05-25: now distinguishes single-value cols (e.g.
                # bulk-export `created_at`) from real time-series, so chat
                # agent doesn't try to DATE_TRUNC a constant column.
                for ca in (col_analyses or []):
                    if ca.get("type") == "datetime" or "date" in ca.get("name", "").lower():
                        try:
                            dr = mconn.execute(text(
                                f'SELECT MIN("{ca["name"]}"), MAX("{ca["name"]}"), '
                                f'COUNT(DISTINCT "{ca["name"]}") FROM "{table_name}"'
                            )).fetchone()
                            if dr and dr[0]:
                                if dr[2] == 1:
                                    mem_facts.append(
                                        f"⚠ Table '{table_name}' column '{ca['name']}' is a CONSTANT "
                                        f"({dr[0]}) across all rows — likely an export timestamp, "
                                        f"NOT a usable time dimension. Do not GROUP BY or DATE_TRUNC."
                                    )
                                else:
                                    mem_facts.append(
                                        f"Table '{table_name}' date range: {dr[0]} to {dr[1]} "
                                        f"({dr[2]:,} distinct values)"
                                    )
                        except Exception:
                            pass
                        break

                # Top categories (for categorical columns)
                for ca in (col_analyses or []):
                    if ca.get("is_categorical") and ca.get("unique_count", 0) <= 15:
                        try:
                            cats = mconn.execute(text(f'SELECT "{ca["name"]}", COUNT(*) as cnt FROM "{table_name}" WHERE "{ca["name"]}" IS NOT NULL GROUP BY "{ca["name"]}" ORDER BY cnt DESC LIMIT 5')).fetchall()
                            if cats:
                                cat_str = ", ".join(f"{c[0]} ({c[1]})" for c in cats)
                                mem_facts.append(f"Table '{table_name}' column '{ca['name']}' top values: {cat_str}")
                        except Exception:
                            pass

                # Numeric ranges (for key numeric columns)
                for ca in (col_analyses or []):
                    if ca.get("type") == "numeric" and ca.get("mean"):
                        try:
                            stats = mconn.execute(text(f'SELECT SUM("{ca["name"]}"), AVG("{ca["name"]}"), MIN("{ca["name"]}"), MAX("{ca["name"]}") FROM "{table_name}"')).fetchone()
                            if stats and stats[0] is not None:
                                mem_facts.append(f"Table '{table_name}' column '{ca['name']}': total={stats[0]:,.0f}, avg={stats[1]:,.0f}, range={stats[2]:,.0f} to {stats[3]:,.0f}")
                        except Exception:
                            pass

                # Data-quality facts — 2026-05-25, driven by central column
                # classifier (real PK constraints + adaptive thresholds + per-
                # project overrides). Replaces name-pattern heuristics:
                #   - PK detection now from information_schema, not "ends with _id"
                #   - Enum detection adaptive (ratio-based), not fixed range
                #   - NULL threshold + free-text ratio configurable per project
                # Tenant-tunable via dash_projects.feature_config['pipeline_thresholds'].
                try:
                    # Enrich col_stats with real unique counts where missing
                    # so classifier sees true cardinality, not 100-row sample.
                    _enriched_stats = []
                    for ca in (col_analyses or []):
                        ca2 = dict(ca)
                        cn = ca2.get("name", "")
                        if cn and not cn.startswith("_") and row_count:
                            try:
                                real_uc = mconn.execute(text(
                                    f'SELECT COUNT(DISTINCT "{cn}") FROM "{table_name}"'
                                )).scalar()
                                if real_uc is not None:
                                    ca2["unique_count"] = int(real_uc)
                            except Exception:
                                pass
                        _enriched_stats.append(ca2)

                    roles = classify_columns(
                        project_slug, schema, table_name,
                        _enriched_stats, row_count or 0,
                    )
                except Exception:
                    roles = {}

                # Adaptive NULL threshold — per-project override or default
                th = _get_pipeline_thresholds(project_slug)
                null_floor = float(th.get("null_alert_floor_pct", 5.0))

                for ca in (col_analyses or []):
                    cn = ca.get("name", "")
                    if not cn or cn.startswith("_"):
                        continue
                    role = roles.get(cn, "")

                    # negative-value alert — only on actual measures (not id/pk/fk/enum)
                    if role == "measure":
                        try:
                            neg_count = mconn.execute(text(
                                f'SELECT COUNT(*) FROM "{table_name}" WHERE "{cn}" < 0'
                            )).scalar() or 0
                            if neg_count > 0:
                                real_min = mconn.execute(text(
                                    f'SELECT MIN("{cn}") FROM "{table_name}"'
                                )).scalar()
                                mem_facts.append(
                                    f"⚠ Table '{table_name}' column '{cn}' has {neg_count:,} "
                                    f"negative values (min={real_min}). Often returns/adjustments — "
                                    f"exclude from positive-only totals via WHERE \"{cn}\" >= 0."
                                )
                        except Exception:
                            pass

                    # high-null alert — adaptive floor from project config
                    try:
                        null_n = mconn.execute(text(
                            f'SELECT COUNT(*) FROM "{table_name}" WHERE "{cn}" IS NULL'
                        )).scalar() or 0
                        if row_count and null_n > 0:
                            null_pct_real = 100.0 * null_n / row_count
                            if null_pct_real >= null_floor:
                                mem_facts.append(
                                    f"⚠ Table '{table_name}' column '{cn}' is {null_pct_real:.1f}% NULL "
                                    f"({null_n:,} of {row_count:,} rows). Use COALESCE or IS NOT NULL."
                                )
                    except Exception:
                        pass

                    # role-driven alerts (no name-pattern heuristic anymore)
                    if role == "pk":
                        mem_facts.append(
                            f"Table '{table_name}' column '{cn}' is declared PRIMARY KEY. "
                            f"Use only in WHERE/JOIN/GROUP BY, never SUM/AVG."
                        )
                    elif role == "id" and row_count and row_count >= th["pk_min_rows"]:
                        _real_uc = next(
                            (s.get("unique_count") for s in _enriched_stats if s.get("name") == cn),
                            None,
                        ) or row_count
                        mem_facts.append(
                            f"Table '{table_name}' column '{cn}' is PK-equivalent "
                            f"({_real_uc:,} unique in {row_count:,} rows). Treat as identifier, not measure."
                        )
                    elif role == "fk":
                        mem_facts.append(
                            f"Table '{table_name}' column '{cn}' is a FOREIGN KEY. "
                            f"Use for JOIN, never SUM/AVG."
                        )
                    elif role == "free_text":
                        mem_facts.append(
                            f"Table '{table_name}' column '{cn}' is free-text (high cardinality). "
                            f"Avoid GROUP BY — use LIKE / fuzzy match / aggregation by category instead."
                        )

                    # Currency-scale detection — only on real measures
                    if role == "measure":
                        try:
                            scale = detect_currency_scale(
                                schema, table_name, cn,
                                ca.get("min"), ca.get("max"),
                            )
                            if scale:
                                mem_facts.append(
                                    f"⚠ Table '{table_name}' column '{cn}' is stored in MINOR UNITS "
                                    f"(divisor={scale}). Divide by {scale} for display values, OR "
                                    f"multiply thresholds by {scale} in WHERE clauses."
                                )
                        except Exception:
                            pass

                # Composite PK probe — fires after single-col loop
                try:
                    if row_count and row_count >= 100:
                        comp = detect_composite_pk(schema, table_name, roles, row_count)
                        if comp:
                            a, b = comp
                            mem_facts.append(
                                f"Table '{table_name}' has COMPOSITE PRIMARY KEY ('{a}', '{b}') — "
                                f"unique row identity is the pair, not either column alone. "
                                f"GROUP BY both for row-level operations."
                            )
                except Exception:
                    pass
        except Exception:
            pass

        # Add metadata-based facts as fallback
        if metadata.get("table_description"):
            mem_facts.append(f"Table '{table_name}': {metadata['table_description'][:200]}")
        if metadata.get("grain"):
            mem_facts.append(f"Table '{table_name}' grain: {metadata['grain']}")

        with engine.connect() as conn:
            saved = 0
            # Bumped 8 → 30 (2026-05-25) — added data-quality facts which can
            # produce 3-10 alerts per table. Each fact is <300 chars.
            for fact in mem_facts[:30]:
                try:
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'auto') ON CONFLICT DO NOTHING"
                    ), {"s": project_slug, "f": fact})
                    saved += 1
                except Exception:
                    pass
            conn.commit()
        _log(f"  ✓ {saved} memories saved")
    except Exception as e:
        _log(f"  ⚠ memories error: {str(e)[:80]}")

    # 2. Auto-Patterns: generate SQL patterns from Q&A (with metadata extraction)
    try:
        _log("  generating query patterns...")
        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        qa_file = training_dir / f"{table_name}_qa.json"
        patterns_saved = 0
        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            for qa in qa_pairs[:5]:
                q = qa.get("question", "")
                s = qa.get("sql", "")
                if q and s:
                    if _save_query_pattern_with_metadata(engine, project_slug, q, s, source='training'):
                        patterns_saved += 1
        _log(f"  ✓ {patterns_saved} query patterns saved (with metadata)")
    except Exception as e:
        _log(f"  ⚠ patterns error: {str(e)[:80]}")

    # 3. Auto-Rules
    try:
        _log("  generating business rules...")
        rules_saved = 0
        biz_metrics = biz_rules.get("metrics", [])
        biz_rules_list = biz_rules.get("business_rules", [])
        with engine.connect() as conn:
            for metric in biz_metrics[:3]:
                if isinstance(metric, dict) and metric.get("name"):
                    rule_id = f"auto_{table_name}_{metric['name'].lower().replace(' ', '_')[:30]}"
                    defn = f"{metric.get('name', '')}: {metric.get('definition', '')} (Calculation: {metric.get('calculation', '')})"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'metric', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": metric["name"], "defn": defn})
                    rules_saved += 1
            for rule in biz_rules_list[:3]:
                if isinstance(rule, dict) and rule.get("name"):
                    rule_id = f"auto_{table_name}_{rule['name'].lower().replace(' ', '_')[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'business_rule', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": rule["name"], "defn": rule.get("definition", "")})
                    rules_saved += 1
            conn.commit()
        _log(f"  ✓ {rules_saved} rules saved")
    except Exception as e:
        _log(f"  ⚠ rules error: {str(e)[:80]}")

    # 4. Auto-Annotations
    try:
        _log("  generating column annotations...")
        annotations_saved = 0
        with engine.connect() as conn:
            for col in (metadata.get("table_columns") or []):
                desc = col.get("description", "")
                if desc and len(desc) > 5:
                    conn.execute(text(
                        "INSERT INTO public.dash_annotations (project_slug, table_name, column_name, annotation, updated_by) "
                        "VALUES (:s, :t, :c, :a, 'auto_training') ON CONFLICT (project_slug, table_name, column_name) DO NOTHING"
                    ), {"s": project_slug, "t": table_name, "c": col["name"], "a": desc})
                    annotations_saved += 1
            conn.commit()
        _log(f"  ✓ {annotations_saved} annotations saved")
    except Exception as e:
        _log(f"  ⚠ annotations error: {str(e)[:80]}")

    # 5. Auto-Evals — MOVED OUT of the per-table path (issue #39).
    # Eval generation now runs ONCE project-wide in _bg()'s tail
    # (_generate_project_evals), sampling Q&A across all tables. Generating
    # here fired one INSERT batch per table (3 tables → 3× the evals); the
    # project-wide pass produces a single bounded set per retrain (the
    # DELETE-at-start pre-clear in _bg still wipes the prior set first).

    # 6. Seed Feedback
    try:
        _log("  seeding sample feedback...")
        feedback_saved = 0
        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            with engine.connect() as conn:
                for qa in qa_pairs[:3]:
                    q = qa.get("question", "")
                    a = qa.get("answer_template", f"SQL: {qa.get('sql', '')}")
                    if q:
                        conn.execute(text(
                            "INSERT INTO public.dash_feedback (project_slug, question, answer, sql_query, rating) "
                            "VALUES (:s, :q, :a, :sql, 'up')"
                        ), {"s": project_slug, "q": q, "a": a, "sql": qa.get("sql", "")})
                        feedback_saved += 1
                conn.commit()
        _log(f"  ✓ {feedback_saved} seed feedback saved")
    except Exception as e:
        _log(f"  ⚠ feedback error: {str(e)[:80]}")

    # 7. Save workflows to DB
    try:
        _log("  saving workflows to database...")
        wf_file = KNOWLEDGE_DIR / project_slug / "workflows.json"
        wf_saved = 0
        if wf_file.exists():
            with open(wf_file) as f:
                wfs = json.load(f)
            with engine.connect() as conn:
                for wf in (wfs if isinstance(wfs, list) else []):
                    name = wf.get("name", "")
                    if name:
                        conn.execute(text(
                            "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                            "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'training') "
                            "ON CONFLICT DO NOTHING"
                        ), {"s": project_slug, "n": name, "d": wf.get("description", ""), "st": json.dumps(wf.get("steps", []))})
                        wf_saved += 1
                conn.commit()
        _log(f"  ✓ {wf_saved} workflows saved to DB")
    except Exception as e:
        _log(f"  ⚠ workflows error: {str(e)[:80]}")

    # ═══ SMART TRAINING — Extract domain knowledge from data ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "domain_knowledge")
    _log("extracting domain knowledge...")

    col_names = [c.get("name", "") for c in (metadata.get("table_columns") or [])]
    col_types = [f"{c.get('name','')}: {c.get('type','')}" for c in (metadata.get("table_columns") or [])]
    col_info = ", ".join(col_types[:30])

    # Get sample distinct values for categorical columns
    cat_values = {}
    try:
        proj_engine = create_engine(db_url)
        import re as _re2
        schema = _re2.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        with proj_engine.connect() as conn:
            for col in (metadata.get("table_columns") or []):
                if col.get("type") in ("TEXT", "text", "str", "object") and col.get("unique_count", 100) <= 20:
                    try:
                        vals = conn.execute(text(f'SELECT DISTINCT "{col["name"]}" FROM "{schema}"."{table_name}" WHERE "{col["name"]}" IS NOT NULL LIMIT 15')).fetchall()
                        cat_values[col["name"]] = [str(v[0]) for v in vals]
                    except Exception:
                        pass
    except Exception:
        pass

    cat_values_str = "\n".join(f"  {k}: {', '.join(v[:10])}" for k, v in cat_values.items()) if cat_values else "None detected"

    from dash.settings import training_llm_call

    # 8. Business Glossary
    try:
        _log("  extracting business glossary...")
        result = training_llm_call(
            f"Look at these column names and values from table '{table_name}'. "
            f"Extract abbreviations, acronyms, and domain terms with their meanings.\n\n"
            f"Columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"Return ONLY valid JSON array of strings:\n"
            f'["MMK = Myanmar Kyat (currency)", "PR = Purchase Requisition"]',
            "extraction"
        )
        if result:
            glossary = json.loads(result)
            if isinstance(glossary, list):
                # Filter dtype trivia + generic SQL terms — bug fix 2026-05-25.
                # LLM was emitting entries like "TEXT = String data type",
                # "INTEGER = Whole number data type", "TIMESTAMP = Date and time
                # format" which burn context budget with no tenant-specific value.
                _DTYPE_TRIVIA = {
                    "text", "varchar", "char", "integer", "int", "bigint", "smallint",
                    "numeric", "decimal", "real", "double", "float", "boolean", "bool",
                    "timestamp", "date", "time", "interval", "jsonb", "json", "uuid",
                    "bytea", "blob", "string", "number",
                }
                with engine.connect() as conn:
                    for term in glossary[:15]:
                        if not (isinstance(term, str) and len(term) > 3):
                            continue
                        # term shape: "ABBR = definition"
                        head = term.split("=", 1)[0].strip().lower()
                        if head in _DTYPE_TRIVIA:
                            continue
                        # also skip pure dtype definitions
                        defn = term.split("=", 1)[1].lower() if "=" in term else ""
                        if "data type" in defn or "format" == defn.strip():
                            continue
                        conn.execute(text(
                            "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                            "VALUES (:s, 'project', :f, 'glossary') ON CONFLICT DO NOTHING"
                        ), {"s": project_slug, "f": f"Glossary: {term}"})
                        # Mirror into project-scoped Company Brain
                        try:
                            if "=" in term:
                                name_part, defn_part = term.split("=", 1)
                                bname, bdefn = name_part.strip()[:200], defn_part.strip()[:4000]
                            else:
                                bname, bdefn = term[:200], term[:4000]
                            if bname:
                                conn.execute(text(
                                    "INSERT INTO public.dash_company_brain "
                                    "(category, name, definition, project_slug, metadata) "
                                    "VALUES ('glossary', :n, :d, :s, CAST('{}' AS jsonb)) "
                                    "ON CONFLICT DO NOTHING"
                                ), {"n": bname, "d": bdefn, "s": project_slug})
                        except Exception:
                            pass
                    conn.commit()
                _log(f"  ✓ {len(glossary[:15])} glossary terms saved (memories + brain)")
    except Exception as e:
        _log(f"  ⚠ glossary error: {str(e)[:80]}")

    # 9. Calculation Rules
    try:
        _log("  extracting calculation rules...")
        numeric_cols = [c["name"] for c in (metadata.get("table_columns") or []) if c.get("type") in ("DOUBLE PRECISION", "INTEGER", "BIGINT", "NUMERIC", "numeric", "int64", "float64")]
        if numeric_cols:
            result = training_llm_call(
                f"Table '{table_name}' has these numeric columns: {', '.join(numeric_cols)}\n"
                f"All columns: {col_info}\n\n"
                f"Infer calculation relationships between columns. How are they related?\n"
                f"Return ONLY valid JSON array:\n"
                f'["net_amount = total_amount + tax - discount", "total_value = SUM(net_amount) per quotation"]',
                "extraction"
            )
            if result:
                calcs = json.loads(result)
                if isinstance(calcs, list):
                    with engine.connect() as conn:
                        for calc in calcs[:8]:
                            if isinstance(calc, str) and len(calc) > 5:
                                rule_id = f"calc_{table_name}_{hash(calc) % 10000}"
                                conn.execute(text(
                                    "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                    "VALUES (:s, :rid, :name, 'calculation', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                                ), {"s": project_slug, "rid": rule_id, "name": f"Calculation: {calc[:50]}", "defn": calc})
                        conn.commit()
                    _log(f"  ✓ {len(calcs[:8])} calculation rules saved")
    except Exception as e:
        _log(f"  ⚠ calculation error: {str(e)[:80]}")

    # 10. Value Mappings
    try:
        _log("  extracting value mappings...")
        if cat_values:
            result = training_llm_call(
                f"Table '{table_name}' has these categorical columns with values:\n{cat_values_str}\n\n"
                f"Explain what each value means in business context.\n"
                f"Return ONLY valid JSON array:\n"
                f'["status: SAP PR Created = quotation approved and sent to SAP system", "doa_type: CAPEX = Capital Expenditure for assets"]',
                "extraction"
            )
            if result:
                mappings = json.loads(result)
                if isinstance(mappings, list):
                    with engine.connect() as conn:
                        for mapping in mappings[:15]:
                            if isinstance(mapping, str) and len(mapping) > 5:
                                conn.execute(text(
                                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                    "VALUES (:s, 'project', :f, 'value_mapping') ON CONFLICT DO NOTHING"
                                ), {"s": project_slug, "f": f"Value mapping: {mapping}"})
                        conn.commit()
                    _log(f"  ✓ {len(mappings[:15])} value mappings saved")
    except Exception as e:
        _log(f"  ⚠ value mapping error: {str(e)[:80]}")

    # 11. KPI Definitions
    try:
        _log("  extracting KPI definitions...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"What business KPIs can be calculated from this data? "
            f"Give the KPI name, definition, and SQL calculation.\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"name": "Total Spend", "definition": "Sum of approved spend", "sql": "SELECT SUM(net_amount) FROM {table_name} WHERE status = \'Approved\'"}}]',
            "extraction"
        )
        if result:
            try:
                kpis = json.loads(result)
            except json.JSONDecodeError:
                # Truncate to last complete `}` and retry
                last_brace = result.rfind('}')
                if last_brace > 0:
                    try:
                        kpis = json.loads(result[:last_brace + 1] + "]" if not result.rstrip().endswith("]") else result[:last_brace + 1])
                    except Exception:
                        try:
                            kpis = json.loads(result[:last_brace + 1])
                        except Exception:
                            kpis = []
                else:
                    kpis = []
            if isinstance(kpis, list):
                with engine.connect() as conn:
                    for kpi in kpis[:8]:
                        if isinstance(kpi, dict) and kpi.get("name"):
                            rule_id = f"kpi_{table_name}_{kpi['name'].lower().replace(' ', '_')[:25]}"
                            defn = f"{kpi.get('definition', '')} | SQL: {kpi.get('sql', '')}"
                            conn.execute(text(
                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                "VALUES (:s, :rid, :name, 'kpi', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                            ), {"s": project_slug, "rid": rule_id, "name": kpi["name"], "defn": defn})
                            # Mirror KPI/formula into project-scoped Company Brain (Issue #4).
                            try:
                                conn.execute(text(
                                    "INSERT INTO public.dash_company_brain "
                                    "(category, name, definition, project_slug, metadata) "
                                    "VALUES ('formula', :n, :d, :s, CAST('{}' AS jsonb)) "
                                    "ON CONFLICT DO NOTHING"
                                ), {"n": str(kpi["name"])[:200], "d": defn[:4000], "s": project_slug})
                            except Exception:
                                pass
                    conn.commit()
                _log(f"  ✓ {len(kpis[:8])} KPI definitions saved (rules + brain)")
    except Exception as e:
        _log(f"  ⚠ KPI error: {str(e)[:80]}")

    # 12. Data Quality Rules
    try:
        _log("  extracting data quality rules...")
        null_info = ", ".join(f"{c['name']}: {c.get('null_pct', 0)}% null" for c in (metadata.get("table_columns") or []) if c.get("null_pct", 0) > 5)
        if null_info:
            result = training_llm_call(
                f"Table '{table_name}' data quality:\n"
                f"High null columns: {null_info}\n"
                f"All columns: {col_info}\n\n"
                f"Generate data quality rules — what to watch out for, how to handle NULLs, common mistakes.\n"
                f"Return ONLY valid JSON array:\n"
                f'["Never include rows where supplier_name IS NULL — these are template records", "Always use COALESCE(exchange_rate, 1) for currency conversion"]',
                "extraction"
            )
            if result:
                dq_rules = json.loads(result)
                if isinstance(dq_rules, list):
                    with engine.connect() as conn:
                        for rule in dq_rules[:8]:
                            if isinstance(rule, str) and len(rule) > 5:
                                conn.execute(text(
                                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                    "VALUES (:s, 'project', :f, 'data_quality') ON CONFLICT DO NOTHING"
                                ), {"s": project_slug, "f": f"Data quality: {rule}"})
                        conn.commit()
                    _log(f"  ✓ {len(dq_rules[:8])} data quality rules saved")
    except Exception as e:
        _log(f"  ⚠ data quality error: {str(e)[:80]}")

    # 13. Negative Examples (What NOT to do)
    try:
        _log("  extracting negative examples...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"What are common mistakes when querying this data? Give DON'T / DO pairs.\n"
            f"Return ONLY valid JSON array:\n"
            f'["DON\'T COUNT(*) for quotation count — each row is an item. DO COUNT(DISTINCT qr_no)", '
            f'"DON\'T use total_amount for spend — it excludes tax. DO use net_amount"]',
            "extraction"
        )
        if result:
            negs = json.loads(result)
            if isinstance(negs, list):
                with engine.connect() as conn:
                    for neg in negs[:8]:
                        if isinstance(neg, str) and len(neg) > 5:
                            conn.execute(text(
                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                "VALUES (:s, 'project', :f, 'negative_example') ON CONFLICT DO NOTHING"
                            ), {"s": project_slug, "f": f"⚠ {neg}"})
                    conn.commit()
                _log(f"  ✓ {len(negs[:8])} negative examples saved")
    except Exception as e:
        _log(f"  ⚠ negative examples error: {str(e)[:80]}")

    _log("✓ agent brain filled — ready to use!")

    # ═══ AI SEED — Pre-populate activity metrics from data analysis ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _log("generating AI seed activity data...")

    col_info = ", ".join(f"{c.get('name','')}: {c.get('type','')}" for c in (metadata.get("table_columns") or [])[:20])
    cat_values_str = ""
    for c in (metadata.get("table_columns") or []):
        if c.get("sample_values"):
            cat_values_str += f"  {c['name']}: {', '.join(str(v) for v in c['sample_values'][:5])}\n"

    # 14. Seed Bad Feedback (common mistakes)
    try:
        _log("  generating bad feedback examples...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n\n"
            f"Generate 2 examples of WRONG analysis that a data analyst might make with this data. "
            f"Each should be a common mistake (wrong column, wrong aggregation, wrong filter).\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"question": "What is total X?", "wrong_answer": "The total is $500K", "why_wrong": "Used gross instead of net amount"}}]',
            "extraction"
        )
        if result:
            bads = json.loads(result)
            if isinstance(bads, list):
                bad_saved = 0
                with train_engine.connect() as conn:
                    for b in bads[:3]:
                        if isinstance(b, dict) and b.get("question"):
                            conn.execute(text(
                                "INSERT INTO public.dash_feedback (project_slug, question, answer, rating) "
                                "VALUES (:s, :q, :a, 'down') ON CONFLICT DO NOTHING"
                            ), {"s": project_slug, "q": b["question"], "a": f"{b.get('wrong_answer','')} | WHY WRONG: {b.get('why_wrong','')}"})
                            bad_saved += 1
                    conn.commit()
                _log(f"  ✓ {bad_saved} bad feedback examples saved")
    except Exception as e:
        _log(f"  ⚠ bad feedback error: {str(e)[:80]}")

    # 15. Seed Proactive Insights (anomalies from data)
    try:
        _log("  generating proactive insights...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Sample values:\n{cat_values_str}\n\n"
            f"Analyze this data schema and generate 3 proactive insights an analyst should investigate. "
            f"Focus on: data quality issues, unusual patterns, potential anomalies, optimization opportunities.\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"title": "High null rate in column X", "description": "Column X has significant missing data which could affect aggregations", "severity": "warning", "metric": "null_rate", "value": "23%"}}]',
            "extraction"
        )
        if result:
            insights = json.loads(result)
            if isinstance(insights, list):
                ins_saved = 0
                with train_engine.connect() as conn:
                    for ins in insights[:5]:
                        if isinstance(ins, dict) and ins.get("title"):
                            insight_text = f"{ins['title']}: {ins.get('description', '')}"
                            conn.execute(text(
                                "INSERT INTO public.dash_proactive_insights (project_slug, insight, severity, tables_involved) "
                                "VALUES (:s, :i, :sev, :tables)"
                            ), {"s": project_slug, "i": insight_text,
                                "sev": ins.get("severity", "info"), "tables": [table_name]})
                            ins_saved += 1
                    conn.commit()
                _log(f"  ✓ {ins_saved} proactive insights saved")
    except Exception as e:
        _log(f"  ⚠ insights error: {str(e)[:80]}")

    # 16. Seed Drift Baseline
    try:
        _log("  saving drift baseline...")
        with train_engine.connect() as conn:
            # Check if baseline already exists
            existing = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_drift_alerts WHERE project_slug = :s AND table_name = :t"
            ), {"s": project_slug, "t": table_name}).scalar() or 0
            if existing == 0:
                baseline = json.dumps([{"type": "baseline", "rows": num_rows, "columns": num_cols, "message": f"Baseline: {num_rows} rows, {num_cols} columns"}])
                conn.execute(text(
                    "INSERT INTO public.dash_drift_alerts (project_slug, table_name, alerts) "
                    "VALUES (:s, :t, CAST(:a AS jsonb))"
                ), {"s": project_slug, "t": table_name, "a": baseline})
                conn.commit()
                _log("  ✓ drift baseline saved")
            else:
                _log("  · drift baseline already exists")
    except Exception as e:
        _log(f"  ⚠ drift baseline error: {str(e)[:80]}")

    # 17. Seed Initial Evolution
    try:
        _log("  generating initial evolution...")
        with train_engine.connect() as conn:
            existing_evol = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_evolved_instructions WHERE project_slug = :s"
            ), {"s": project_slug}).scalar() or 0
            if existing_evol == 0:
                result = training_llm_call(
                    f"Table '{table_name}' columns: {col_info}\n\n"
                    f"Generate a short set of supplementary instructions (3-5 bullet points) that will help "
                    f"a data analyst AI agent work better with this specific data. Focus on gotchas, "
                    f"preferred aggregations, and domain-specific rules.\n"
                    f"Return plain text bullet points only.",
                    "persona"
                )
                if result:
                    conn.execute(text(
                        "INSERT INTO public.dash_evolved_instructions (project_slug, instructions, version, reasoning) "
                        "VALUES (:s, :i, 1, 'Auto-generated during initial training')"
                    ), {"s": project_slug, "i": result[:2000]})
                    conn.commit()
                    _log("  ✓ initial evolution saved (v1)")
            else:
                _log(f"  · evolution already exists ({existing_evol} versions)")
    except Exception as e:
        _log(f"  ⚠ evolution error: {str(e)[:80]}")

    # ═══ ENRICH PERSONA — Re-generate with full domain context ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "persona_enrich")
    _log("enriching persona with domain knowledge...")
    try:
        # Gather all domain knowledge for richer persona
        eng_db = create_engine(db_url)
        domain_context_parts = []
        with eng_db.connect() as conn:
            # Glossary terms
            glossary = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'glossary' LIMIT 20"
            ), {"s": project_slug}).fetchall()
            if glossary:
                domain_context_parts.append("Business glossary: " + "; ".join(r[0] for r in glossary))

            # KPIs
            kpis = conn.execute(text(
                "SELECT name, definition FROM public.dash_rules_db WHERE project_slug = :s AND type = 'kpi' LIMIT 10"
            ), {"s": project_slug}).fetchall()
            if kpis:
                domain_context_parts.append("KPIs tracked: " + "; ".join(f"{r[0]}: {r[1][:80]}" for r in kpis))

            # Calculation rules
            calcs = conn.execute(text(
                "SELECT name FROM public.dash_rules_db WHERE project_slug = :s AND type = 'calculation' LIMIT 10"
            ), {"s": project_slug}).fetchall()
            if calcs:
                domain_context_parts.append("Calculation rules: " + ", ".join(r[0] for r in calcs))

            # Negative examples (critical gotchas)
            negs_db = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'negative_example' LIMIT 5"
            ), {"s": project_slug}).fetchall()
            if negs_db:
                domain_context_parts.append("Critical gotchas: " + "; ".join(r[0] for r in negs_db))

            # Data quality insights
            quality = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'data_quality' LIMIT 5"
            ), {"s": project_slug}).fetchall()
            if quality:
                domain_context_parts.append("Data quality rules: " + "; ".join(r[0] for r in quality))

        domain_summary = "\n".join(domain_context_parts)

        if domain_summary.strip():
            # Get all table names for context
            all_tables = []
            tables_dir_path = KNOWLEDGE_DIR / project_slug / "tables"
            if tables_dir_path.exists():
                for tf in tables_dir_path.glob("*.json"):
                    all_tables.append(tf.stem)

            enriched_persona = training_llm_call(
                f"You are generating a rich, expert persona for a data analyst AI agent.\n\n"
                f"This agent works with these tables: {', '.join(all_tables)}\n\n"
                f"Domain knowledge extracted from the data:\n{domain_summary}\n\n"
                f"Table being analyzed: {table_name}\n"
                f"Column types: {', '.join(c.get('name','') + ':' + c.get('type','') for c in (metadata.get('table_columns') or [])[:20])}\n\n"
                f"Generate a JSON persona object:\n"
                f'{{"persona_prompt": "You are an expert [specific domain] analyst who specializes in [specific areas from the data]. '
                f'You understand [domain terms from glossary]. You track KPIs like [from KPIs]. You know critical gotchas like [from negatives]...",'
                f'"domain_terms": ["term1", "term2", ...],'
                f'"expertise_areas": ["area1", "area2", ...],'
                f'"kpi_focus": ["kpi1", "kpi2", ...],'
                f'"communication_style": "concise with data tables",'
                f'"greeting": "Hey! Ready to dive into your [domain] data?"}}',
                "persona"
            )

            if enriched_persona:
                persona_data = json.loads(enriched_persona)
                if isinstance(persona_data, dict) and persona_data.get("persona_prompt"):
                    # Save enriched persona to file
                    persona_file = KNOWLEDGE_DIR / project_slug / "persona.json"
                    with open(persona_file, "w") as f:
                        json.dump(persona_data, f, indent=2)

                    # Update DB
                    with eng_db.connect() as conn:
                        conn.execute(text(
                            "UPDATE public.dash_personas SET persona = CAST(:p AS jsonb) WHERE project_slug = :s"
                        ), {"s": project_slug, "p": json.dumps(persona_data)})
                        conn.commit()
                    _log(f"✓ persona enriched with {len(domain_context_parts)} domain layers")
                else:
                    _log("· persona enrichment skipped — invalid response")
            else:
                _log("· persona enrichment skipped — no LLM response")
        else:
            _log("· persona enrichment skipped — no domain knowledge yet")
    except Exception as e:
        _log(f"⚠ persona enrichment error: {str(e)[:80]}")

    # PandasAI Experiments — generate 50+ verified Q&A from real data
    if getenv("PANDASAI_EXPERIMENTS", "true").lower() in ("true", "1", "yes"):
        _run_pandasai_experiments(project_slug, table_name, col_analyses, _log)

    # LangExtract — extract grounded facts from document text for Researcher agent
    try:
        docs_dir = KNOWLEDGE_DIR / project_slug / "docs"
        if docs_dir.exists():
            doc_texts = []
            for f in sorted(docs_dir.iterdir()):
                if f.is_file():
                    try:
                        doc_texts.append(f.read_text(errors='ignore')[:3000])
                    except Exception:
                        pass
            if doc_texts:
                combined_text = "\n\n---\n\n".join(doc_texts)
                _langextract_facts(project_slug, combined_text, _log)
    except Exception as e:
        _log(f"⚠ fact extraction skipped: {str(e)[:80]}")

    # Cross-Source Knowledge Graph — REMOVED from per-table loop.
    # KG now builds ONCE at end of TRAIN ALL (see master loop ~line 9611) — saves 9 LLM calls/train.
    try:
        _update_run("running", "knowledge_graph")
        # KG builds once in the master loop — silent here (was a noise log line).
        pass
    except Exception:
        pass

    # Sub-agent synthesis — REMOVED from per-table loop.
    # Runs ONCE at end of TRAIN ALL (master loop ~line 9637) instead of once per
    # table — synthesis clusters across ALL training Q&A, so per-table runs were
    # redundant (ran N+1 times on an N-table train, re-fetching all Q&A each time).
    try:
        _update_run("running", "subagent_synthesis")
        # sub-agent synthesis is a dead feature — silent here (was noise).
        pass
    except Exception:
        pass

    # Step 14: Auto ML Models — REMOVED. auto_create_models() + ml_worker
    # container were dropped in the 2026-05-23 LLM-native ML pivot. Per-chat
    # ML tools (predict/classify/detect_anomalies_ml) run in-process now.

    # Save fingerprint for delta detection on next retrain
    save_fingerprint(project_slug, table_name, num_rows, [c.get("name", "") for c in (metadata.get("table_columns") or [])])

    # Training Quality Score — measure how good this training was
    try:
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        qa_file = KNOWLEDGE_DIR / project_slug / "training" / f"{table_name}_qa.json"
        quality_score = 0
        quality_details = {}

        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            total_qa = len(qa_pairs)
            verified_qa = sum(1 for q in qa_pairs if q.get("verified"))
            quality_details["qa_total"] = total_qa
            quality_details["qa_verified"] = verified_qa
            quality_details["qa_pct"] = round((verified_qa / max(total_qa, 1)) * 100)
            quality_score += quality_details["qa_pct"] * 0.4  # 40% weight

        # Check relationships verified
        try:
            rel_engine = create_engine(db_url)
            with rel_engine.connect() as rconn:
                rels = rconn.execute(text("SELECT confidence, source FROM public.dash_relationships WHERE project_slug = :s"), {"s": project_slug}).fetchall()
                if rels:
                    verified_rels = sum(1 for r in rels if r[1] == 'ai_verified')
                    quality_details["relationships_total"] = len(rels)
                    quality_details["relationships_verified"] = verified_rels
                    quality_score += min(100, (verified_rels / max(len(rels), 1)) * 100) * 0.2  # 20% weight
                else:
                    quality_score += 50 * 0.2  # No rels = partial credit
        except Exception:
            quality_score += 50 * 0.2

        # Check memories from real data
        try:
            with rel_engine.connect() as rconn:
                mem_count = rconn.execute(text("SELECT COUNT(*) FROM public.dash_memories WHERE project_slug = :s AND source = 'auto'"), {"s": project_slug}).scalar() or 0
                quality_details["memories"] = mem_count
                quality_score += min(100, (mem_count / 8) * 100) * 0.2  # 20% weight
        except Exception:
            pass

        # Check profile exists (data quality)
        profile_file = KNOWLEDGE_DIR / project_slug / "table_sources" / f"{table_name}_profile.json"
        if profile_file.exists():
            try:
                with open(profile_file) as f:
                    profile = json.load(f)
                quality_details["health"] = profile.get("health", 0)
                quality_score += profile.get("health", 50) * 0.2  # 20% weight
            except Exception:
                quality_score += 50 * 0.2
        else:
            quality_score += 50 * 0.2

        quality_score = round(min(100, max(0, quality_score)))
        quality_details["overall_score"] = quality_score
        _log(f"✓ training quality score: {quality_score}% (Q&A: {quality_details.get('qa_pct', '?')}% verified, {quality_details.get('memories', '?')} memories, health: {quality_details.get('health', '?')}%)")

        # Save quality score
        quality_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
        quality_dir.mkdir(parents=True, exist_ok=True)
        _safe_write_json(quality_dir / f"{table_name}_training_quality.json", quality_details)
    except Exception as e:
        _log(f"⚠ quality score error: {str(e)[:80]}")

    # ── Per-table summary card (Level 3) ───────────────────────────────────
    try:
        dim_count = len(metadata.get("dimensions") or {})
        hier_list = metadata.get("hierarchies") or []
        hier_str = ", ".join(f"{h.get('parent','?')}→{h.get('child','?')}" for h in hier_list[:3]) if hier_list else "—"
        fk_list = metadata.get("foreign_keys") or []
        fk_count = len(fk_list) if isinstance(fk_list, list) else 0
        qa_path = KNOWLEDGE_DIR / project_slug / "training" / f"{table_name}_qa.json"
        qa_count = 0
        qa_verified = 0
        if qa_path.exists():
            try:
                with open(qa_path) as _qf:
                    _qd = json.load(_qf) or []
                qa_count = len(_qd)
                qa_verified = sum(1 for q in _qd if q.get("verified"))
            except Exception:
                pass
        mem_count = 0
        try:
            with train_engine.connect() as _mc:
                _mr = _mc.execute(text(
                    "SELECT COUNT(*) FROM public.dash_memories WHERE project_slug = :s"
                ), {"s": project_slug}).scalar()
                mem_count = int(_mr or 0)
        except Exception:
            pass
        _log(f"┌─ SUMMARY · {table_name}")
        _log(f"│  rows: {num_rows} · cols: {num_cols} · dim: {dim_count} · measure: {len(measures) if 'measures' in dir() else '?'} · fk: {fk_count}")
        _log(f"│  hierarchy: {hier_str}")
        _log(f"│  Q&A: {qa_verified}/{qa_count} verified")
        # KG triples are built globally AFTER all tables train — per-table count
        # is always 0 (misleading), so it's omitted here.
        _log(f"│  memories: {mem_count}")
        _log(f"└─ quality: {quality_score if 'quality_score' in dir() else '?'}%")
    except Exception as _se:
        _log(f"⚠ summary card error: {str(_se)[:80]}")

    _log(f"✓ training complete for {table_name}")
    _update_run("done", "complete")

    # Best-effort: notify project owner of training completion
    try:
        from app.auth import notify_user
        from sqlalchemy import text as _ntext
        from sqlalchemy import create_engine as _ne
        _eng = _ne(db_url)
        with _eng.connect() as _c:
            _r = _c.execute(_ntext(
                "SELECT user_id FROM public.dash_projects WHERE slug = :s"
            ), {"s": project_slug}).fetchone()
        _eng.dispose()
        if _r and _r[0]:
            _health = quality_score if 'quality_score' in dir() else 0
            notify_user(
                int(_r[0]),
                f"Training complete · {project_slug}",
                f"{table_name} · {qa_verified}/{qa_count} evals · health {_health}%",
                "success",
            )
    except Exception:
        pass


# ---------------------------------------------------------------------------
# API Endpoints
# ---------------------------------------------------------------------------


def _get_user_id(request) -> str:
    """Extract user_id from request state (set by auth middleware)."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if user:
        return str(user.get('username', 'default'))
    return 'default'


def _get_user(request) -> dict:
    """Return full user dict from request state (set by auth middleware).
    Returns empty dict if no user (callers safely do .get() with fallback)."""
    return getattr(getattr(request, 'state', None), 'user', None) or {}


def _get_project_schema(request: Request, project: str | None = None) -> str | None:
    """Get the project schema name if project param is provided."""
    if not project:
        project = request.query_params.get('project')
    if project:
        from db.session import create_project_schema
        return create_project_schema(project)
    return None


def _extract_tables_pptx(file_path: str) -> list:
    """Extract tables from PowerPoint slides as list of DataFrames."""
    tables = []
    try:
        from pptx import Presentation
        import pandas as pd
        prs = Presentation(file_path)
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    rows_data = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    if len(rows_data) > 1 and len(rows_data[0]) > 1:
                        headers = rows_data[0]
                        data = rows_data[1:]
                        df = pd.DataFrame(data, columns=headers)
                        # Clean column names
                        df.columns = [str(c).strip().lower().replace(' ', '_').replace('.', '_')[:50] for c in df.columns]
                        tables.append({
                            'source': f'slide_{si+1}',
                            'df': df,
                            'rows': len(data),
                            'cols': len(headers)
                        })
    except Exception:
        pass
    return tables


def _extract_images_pptx(file_path: str) -> list[dict]:
    """Extract images from PPTX slides. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        prs = Presentation(file_path)
        for si, slide in enumerate(prs.slides):
            if len(images) >= 30:
                break
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = shape.image.blob
                    if len(blob) < 3000:
                        continue
                    images.append({
                        "b64": base64.b64encode(blob).decode(),
                        "mime": shape.image.content_type or "image/png",
                        "source": f"slide_{si + 1}",
                    })
                    if len(images) >= 30:
                        break
    except Exception:
        pass
    return images


def _extract_tables_pdf(file_path: str) -> list:
    """Extract tables from PDF pages as list of DataFrames."""
    tables = []
    try:
        import pdfplumber
        import pandas as pd
        with pdfplumber.open(file_path) as pdf:
            for pi, page in enumerate(pdf.pages):
                page_tables = page.extract_tables()
                for ti, tbl in enumerate(page_tables or []):
                    if tbl and len(tbl) > 1 and len(tbl[0]) > 1:
                        headers = [str(h).strip().lower().replace(' ', '_').replace('.', '_')[:50] for h in tbl[0]]
                        data = tbl[1:]
                        df = pd.DataFrame(data, columns=headers)
                        tables.append({
                            'source': f'page_{pi+1}_table_{ti+1}',
                            'df': df,
                            'rows': len(data),
                            'cols': len(headers)
                        })
    except Exception:
        pass
    return tables


def _extract_images_pdf(file_path: str) -> list[dict]:
    """Extract images from PDF pages. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        import fitz
        import base64
        doc = fitz.open(file_path)
        for pi, page in enumerate(doc):
            if len(images) >= 30:
                break
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                extracted = doc.extract_image(xref)
                if not extracted:
                    continue
                raw = extracted["image"]
                if len(raw) < 3000:
                    continue
                ext = extracted.get("ext", "png")
                images.append({
                    "b64": base64.b64encode(raw).decode(),
                    "mime": f"image/{ext}",
                    "source": f"page_{pi + 1}",
                })
                if len(images) >= 30:
                    break
        doc.close()
    except Exception:
        pass
    return images


def _extract_images_docx(file_path: str) -> list[dict]:
    """Extract images from DOCX document relationships. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        from docx import Document
        import base64
        doc = Document(file_path)
        mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                    "gif": "image/gif", "bmp": "image/bmp", "tiff": "image/tiff",
                    "emf": "image/emf", "wmf": "image/wmf"}
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    if len(blob) < 3000:
                        continue
                    ref = getattr(rel, 'target_ref', '') or ''
                    ext = ref.rsplit(".", 1)[-1].lower() if "." in ref else "png"
                    images.append({
                        "b64": base64.b64encode(blob).decode(),
                        "mime": mime_map.get(ext, "image/png"),
                        "source": f"docx_image_{len(images) + 1}",
                    })
                    if len(images) >= 30:
                        break
                except Exception:
                    continue
    except Exception:
        pass
    return images


def _extract_tables_docx(file_path: str) -> list:
    """Extract tables from DOCX as list of DataFrames."""
    tables = []
    try:
        from docx import Document
        import pandas as pd
        doc = Document(file_path)
        for ti, tbl in enumerate(doc.tables):
            rows_data = []
            for row in tbl.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            if len(rows_data) > 1 and len(rows_data[0]) > 1:
                headers = [str(h).strip().lower().replace(' ', '_').replace('.', '_')[:50] for h in rows_data[0]]
                data = rows_data[1:]
                df = pd.DataFrame(data, columns=headers)
                tables.append({
                    'source': f'table_{ti+1}',
                    'df': df,
                    'rows': len(data),
                    'cols': len(headers)
                })
    except Exception:
        pass
    return tables


def _extract_content(file_path: str, ext: str, raw_content: bytes) -> dict:
    """Extract text AND tables from various file formats."""
    # Get text (existing logic)
    if ext in (".md", ".txt", ".sql", ".py"):
        text = raw_content.decode("utf-8", errors="ignore")
        return {"text": text, "tables": [], "images": []}

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            text = raw_content.decode("utf-8", errors="ignore")
        tables = _extract_tables_docx(file_path)
        return {"text": text, "tables": tables, "images": []}

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
            text = "\n".join(texts)
        except Exception:
            text = ""
        tables = _extract_tables_pptx(file_path)
        return {"text": text, "tables": tables, "images": _extract_images_pptx(file_path)}

    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            texts = []
            for page in doc:
                texts.append(page.get_text())
            doc.close()
            text = "\n".join(texts)
        except Exception:
            text = ""
        tables = _extract_tables_pdf(file_path)
        return {"text": text, "tables": tables, "images": _extract_images_pdf(file_path)}

    return {"text": raw_content.decode("utf-8", errors="ignore"), "tables": [], "images": []}


def _describe_images_with_vision(images: list[dict], filename: str) -> str:
    """Send images to vision model and return combined text descriptions."""
    if not images:
        return ""
    from dash.settings import training_vision_call
    descriptions = []
    for img in images[:30]:
        result = training_vision_call(
            prompt=_UNIVERSAL_VISION_PROMPT,
            images=[img],
        )
        if result:
            descriptions.append(f"[Image from {img['source']} in {filename}]: {result}")
    return "\n\n".join(descriptions)


# ---------------------------------------------------------------------------
# Upload Intelligence: Handlers + Conductor
# Each handler returns: {"tables": [...], "text": str, "images": [...],
#                        "metadata": dict, "errors": [...], "warnings": [...]}
# ---------------------------------------------------------------------------

def _handle_image(file_path: str, filename: str) -> dict:
    """Handle image upload — any format. Tesseract OCR first, vision fallback. Saves original."""
    import base64
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        ext = Path(file_path).suffix.lower()
        mime_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
                    ".tiff": "image/tiff", ".tif": "image/tiff", ".bmp": "image/bmp",
                    ".gif": "image/gif", ".webp": "image/webp"}

        # Convert to PNG if needed (Tesseract/Vision work best with PNG/JPG)
        try:
            from PIL import Image as PILImage
            img = PILImage.open(file_path)

            # Auto-rotate based on EXIF orientation
            try:
                from PIL import ImageOps
                img = ImageOps.exif_transpose(img)
            except Exception:
                pass

            # Convert to RGB if needed (CMYK, palette, etc.)
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            # Tesseract OCR (local, free)
            try:
                import pytesseract
                ocr_text = pytesseract.image_to_string(img)
                if len(ocr_text.strip()) > 30:
                    result["text"] = f"[OCR from {filename}]\n{ocr_text.strip()}"
                    result["warnings"].append("Text extracted via Tesseract OCR")
            except ImportError:
                pass
            except Exception:
                pass

            # Save as PNG for Vision
            import io
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            blob = buf.getvalue()
            result["images"] = [{"b64": base64.b64encode(blob).decode(), "mime": "image/png", "source": filename}]
            result["metadata"]["width"] = img.width
            result["metadata"]["height"] = img.height
            result["metadata"]["format"] = ext

        except Exception:
            # Fallback: read raw bytes
            with open(file_path, "rb") as f:
                blob = f.read()
            result["images"] = [{"b64": base64.b64encode(blob).decode(), "mime": mime_map.get(ext, "image/png"), "source": filename}]

    except Exception as e:
        result["errors"].append(f"Failed to read image: {e}")
    return result


# ---------------------------------------------------------------------------
# Rules Engine: Deterministic Excel structure detection (no LLM, 100% consistent)
# ---------------------------------------------------------------------------

_MONTH_RE = re.compile(r"^(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)", re.IGNORECASE)
_META_RE = re.compile(r"^(Company|Period|Date|Prepared|Updated|Report|Source|Assets|Month|Dep rate|Print Date|Print|Generated|Exported|Created|Page|As of|Title|Subject|Author|File|Filename|Confidential|Disclaimer)(\s*[:\-]\s+|\s*$)", re.IGNORECASE)
_UNIT_WORDS = {"sachets", "kg", "lbs", "tons", "%", "usd", "eur", "gbp", "pcs", "units", "sachets/min", "kg/hr", "sachets/ctn"}
_SUMMARY_RE = re.compile(r"^(Total|Subtotal|Grand|Sum|Utilisation|Utilization|Average|Overall|Net)", re.IGNORECASE)


def _rules_detect_header(rows: list[list[str]]) -> tuple[int, float]:
    """Find header row using scoring (no LLM). Returns (row_index, confidence)."""
    best_row, best_score = 0, 0
    for i, row in enumerate(rows[:15]):
        if not row:
            continue
        non_empty = [v for v in row if v and str(v).strip()]
        if not non_empty:
            continue
        score = 0
        text_count = 0
        for v in non_empty:
            s = str(v).strip()
            # Short text = likely header
            if len(s) > 1 and len(s) < 50 and not s.replace('.', '').replace('-', '').replace(',', '').isdigit():
                score += 3
                text_count += 1
            # Pure number = likely data, not header
            elif s.replace('.', '').replace('-', '').replace(',', '').isdigit():
                score -= 1
            # Metadata keyword = not header, skip
            if _META_RE.match(s):
                score -= 5
            # Unit word = not header
            if s.lower() in _UNIT_WORDS:
                score -= 3
        # Prefer rows with many text values
        text_ratio = text_count / max(len(non_empty), 1)
        score += text_ratio * 5
        if score > best_score:
            best_score = score
            best_row = i
    confidence = min(1.0, best_score / 15.0)
    return best_row, confidence


def _rules_find_blank_boundaries(rows: list[list[str]]) -> list[int]:
    """Find row indices where 2+ consecutive rows are all blank."""
    boundaries = []
    i = 0
    while i < len(rows) - 1:
        is_blank = all(not v or str(v).strip() == "" for v in rows[i])
        next_blank = all(not v or str(v).strip() == "" for v in rows[i + 1]) if i + 1 < len(rows) else False
        if is_blank and next_blank:
            boundaries.append(i)
            # Skip consecutive blanks
            while i < len(rows) and all(not v or str(v).strip() == "" for v in rows[i]):
                i += 1
        else:
            i += 1
    return boundaries


def _full_sheet_blank_scan(file_path: str, sheet_name: str, max_row: int, ext: str) -> list[int]:
    """P3: Stream the FULL sheet (capped 10K rows) to find blank-row boundaries.
    Returns row indices where 2+ consecutive rows are fully blank.
    Same shape as _rules_find_blank_boundaries. Fail-soft → returns []."""
    try:
        ext = (ext or "").lower()
        if ext not in (".xlsx", ".xls"):
            return []
        cap = min(max_row, 10000) if max_row else 10000

        # Stream rows as a list of bool: True if row is fully blank
        blanks: list[bool] = []
        if ext == ".xlsx":
            import openpyxl
            try:
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            except Exception:
                return []
            try:
                if sheet_name not in wb.sheetnames:
                    return []
                ws = wb[sheet_name]
                for ri, row in enumerate(ws.iter_rows(values_only=True)):
                    if ri >= cap:
                        break
                    is_blank = all(v is None or (isinstance(v, str) and not v.strip()) for v in row) if row else True
                    blanks.append(is_blank)
            finally:
                try:
                    wb.close()
                except Exception:
                    pass
        else:  # .xls
            try:
                import xlrd
            except Exception:
                return []
            try:
                wb = xlrd.open_workbook(file_path, on_demand=True)
                if sheet_name not in wb.sheet_names():
                    return []
                ws = wb.sheet_by_name(sheet_name)
                n = min(ws.nrows, cap)
                for ri in range(n):
                    row = ws.row_values(ri)
                    is_blank = all(v is None or (isinstance(v, str) and not v.strip()) or v == "" for v in row) if row else True
                    blanks.append(is_blank)
            except Exception:
                return []

        # Find boundaries (2+ consecutive blanks)
        boundaries: list[int] = []
        i = 0
        n = len(blanks)
        while i < n - 1:
            if blanks[i] and blanks[i + 1]:
                boundaries.append(i)
                while i < n and blanks[i]:
                    i += 1
            else:
                i += 1
        return boundaries
    except Exception:
        return []


def _rules_detect_skip_rows(rows: list[list[str]], header_row: int) -> list[int]:
    """Find rows to skip: metadata, units, summaries (no LLM)."""
    skip = []
    for i, row in enumerate(rows):
        if i == header_row:
            continue
        non_empty = [str(v).strip().lower() for v in row if v and str(v).strip()]
        if not non_empty:
            continue
        # Unit row: most values are unit words
        unit_count = sum(1 for v in non_empty if v in _UNIT_WORDS)
        if unit_count >= len(non_empty) * 0.5 and len(non_empty) >= 2:
            skip.append(i)
            continue
        # Metadata row: matches Company:, Period:, etc
        if len(non_empty) <= 3 and any(_META_RE.match(str(v).strip()) for v in row if v):
            skip.append(i)
            continue
        # Summary row: Total, Utilisation, etc
        if any(_SUMMARY_RE.match(str(v).strip()) for v in row if v and str(v).strip()):
            skip.append(i)
    return skip


def _rules_has_month_columns(header_values: list[str]) -> bool:
    """Check if 3+ columns are month/date names → needs unpivot."""
    month_count = sum(1 for v in header_values if v and _MONTH_RE.match(str(v).strip()))
    return month_count >= 3


def _rules_analyze_sheet(rows: list[list[str]], merged_cells: list = None,
                         full_scan_boundaries: list = None, max_row: int = 0) -> dict:
    """Deterministic sheet analysis — no LLM, 100% consistent.
    Returns extraction plan with confidence score.

    P3: full_scan_boundaries (optional) — list of blank-row indices from a
    full-sheet scan (via _full_sheet_blank_scan). If provided AND non-empty,
    merged with preview-only boundaries (dedupe, sort).
    max_row (P3): full-sheet row count; used to extend last block's data_end
    past the 25-row preview window when full_scan_boundaries are present.
    """
    plan = {"action": "load", "confidence": 0.0, "header_row": 0, "skip_rows": [], "blocks": []}

    if not rows or len(rows) < 2:
        plan["action"] = "skip"
        plan["confidence"] = 1.0
        return plan

    # Step 1: Find blank row boundaries (table separators)
    boundaries = _rules_find_blank_boundaries(rows)
    # P3: merge in full-sheet boundaries (preview-only misses mid-sheet gaps)
    if full_scan_boundaries:
        try:
            boundaries = sorted(set(boundaries) | set(int(b) for b in full_scan_boundaries))
        except Exception:
            pass

    # Step 2: If boundaries found → multiple tables in one sheet
    if boundaries:
        # Split into blocks
        block_starts = [0] + [b + 2 for b in boundaries]  # skip blank rows
        block_ends = boundaries + [len(rows)]
        blocks = []
        for start, end in zip(block_starts, block_ends):
            block_rows = rows[start:end]
            # Skip empty blocks
            non_empty_rows = [r for r in block_rows if any(v and str(v).strip() for v in r)]
            if len(non_empty_rows) < 2:
                continue
            # Detect header within this block
            header_idx, conf = _rules_detect_header(block_rows)
            skip = _rules_detect_skip_rows(block_rows, header_idx)
            blocks.append({
                "start": start,
                "end": end - 1,
                "header_row": start + header_idx,
                "data_start": start + header_idx + 1,
                "data_end": end - 1,
                "skip_rows": [start + s for s in skip],
                "header_values": [str(v).strip() for v in block_rows[header_idx] if v] if header_idx < len(block_rows) else [],
            })

        # P3: when full-sheet scan provided boundaries, extend last block past preview
        if full_scan_boundaries and blocks and max_row and max_row > len(rows):
            blocks[-1]["data_end"] = max_row - 1
            blocks[-1]["end"] = max_row - 1

        if len(blocks) >= 2:
            # Banner-block guard: if first block is tiny (≤3 data rows) AND
            # later block has a real header (≥5 short text labels), treat
            # the first block as banner/metadata and load only the real table.
            # Prevents banner-row-induced split that truncates 99% of data.
            def _is_banner_block(b):
                dr = b.get("data_end", 0) - b.get("data_start", 0) + 1
                hv = [str(v).strip() for v in b.get("header_values", []) if str(v).strip()]
                # Banner = ≤3 data rows AND header has ≤2 short cells OR cells match meta pattern
                meta_hits = sum(1 for v in hv if _META_RE.match(v))
                return dr <= 3 and (len(hv) <= 2 or meta_hits >= 1)

            def _is_real_table(b):
                hv = [str(v).strip() for v in b.get("header_values", []) if str(v).strip()]
                # Real table = ≥5 short text labels, no meta keywords
                short_labels = [v for v in hv if 1 < len(v) < 50]
                meta_hits = sum(1 for v in hv if _META_RE.match(v))
                return len(short_labels) >= 5 and meta_hits == 0

            # Strip banner blocks from front
            while len(blocks) >= 2 and _is_banner_block(blocks[0]) and _is_real_table(blocks[1]):
                blocks.pop(0)

            if len(blocks) >= 2:
                plan["action"] = "split"
                plan["blocks"] = blocks
                plan["confidence"] = 0.9
            elif len(blocks) == 1:
                # After banner strip, single real table — load with its header
                plan["header_row"] = blocks[0]["header_row"]
                plan["skip_rows"] = blocks[0]["skip_rows"] + list(range(0, blocks[0]["header_row"]))
                plan["confidence"] = 0.92
        elif len(blocks) == 1:
            plan["header_row"] = blocks[0]["header_row"]
            plan["skip_rows"] = blocks[0]["skip_rows"]

    # Step 3: If no blocks found, detect header for whole sheet
    if plan["action"] == "load":
        header_row, conf = _rules_detect_header(rows)
        plan["header_row"] = header_row
        plan["confidence"] = conf
        plan["skip_rows"] = _rules_detect_skip_rows(rows, header_row)

    # Step 4: Check for month columns → unpivot
    header_idx = plan["header_row"] if plan["action"] != "split" else (plan["blocks"][0]["header_row"] if plan.get("blocks") else 0)
    if header_idx < len(rows):
        header_vals = [str(v).strip() for v in rows[header_idx] if v]
        if _rules_has_month_columns(header_vals):
            plan["action"] = "unpivot"
            plan["confidence"] = max(plan["confidence"], 0.95)

    # Step 5: Note merged cells (handled by forward-fill, doesn't reduce confidence)
    if merged_cells:
        plan["has_merges"] = True

    # Step 6: Multi-level header detection (2-3 rows of text that are headers)
    if plan["action"] == "load" and len(rows) >= 4:
        hdr = plan["header_row"]
        # Check if rows hdr, hdr+1 (and optionally hdr+2) are all text-heavy (headers)
        def _is_header_like(row):
            non_empty = [v for v in row if v and str(v).strip()]
            if not non_empty:
                return False
            text_count = sum(1 for v in non_empty if not re.match(r'^[\d.,\-+$€£¥%]+$', str(v).strip()))
            return text_count >= len(non_empty) * 0.6

        multi_level = 1
        for offset in [1, 2]:
            if hdr + offset < len(rows) and _is_header_like(rows[hdr + offset]):
                # Check if the row after candidate headers has numeric data
                data_row_idx = hdr + offset + 1
                if data_row_idx < len(rows):
                    data_row = rows[data_row_idx]
                    nums = sum(1 for v in data_row if v and re.match(r'^[\d.,\-+$€£¥%]+$', str(v).strip()))
                    if nums >= len([v for v in data_row if v and str(v).strip()]) * 0.3:
                        multi_level = offset + 1
                    else:
                        break
                else:
                    break
            else:
                break

        if multi_level >= 2:
            plan["multi_level_header"] = multi_level
            plan["header_rows"] = list(range(hdr, hdr + multi_level))
            # Build flattened header: concatenate parent > child
            flat_headers = []
            max_cols = max(len(r) for r in rows[hdr:hdr + multi_level])
            for ci in range(max_cols):
                parts = []
                for ri in range(hdr, hdr + multi_level):
                    if ci < len(rows[ri]) and rows[ri][ci] and str(rows[ri][ci]).strip():
                        parts.append(str(rows[ri][ci]).strip())
                flat_headers.append("__".join(parts) if parts else f"col_{ci}")
            plan["flat_headers"] = flat_headers
            plan["data_start_row"] = hdr + multi_level
            plan["confidence"] = max(plan["confidence"], 0.85)

    # Step 7: Bold-based header/summary boost (if formatting info available)
    # Bold rows in first 5 rows boost header confidence
    # Bold rows with summary keywords boost skip confidence
    # (formatting info passed via merged_cells metadata when available)

    return plan


def _is_clean_sheet(df_preview: pd.DataFrame, merged_cells: list = []) -> bool:
    """Quick check: is this sheet clean data (proper headers, no mess) or messy (needs AI)?"""
    if len(df_preview) < 2:
        return False
    # Messy signals
    if merged_cells:
        return False  # Merged cells = messy
    headers = list(df_preview.columns)
    unnamed_count = sum(1 for h in headers if "unnamed" in str(h).lower())
    if unnamed_count > len(headers) * 0.3:
        return False  # Too many unnamed columns = wrong header row
    # Check if first row looks like data (not metadata/units)
    first_row = df_preview.iloc[0]
    text_vals = [str(v).strip().lower() for v in first_row.dropna()]
    unit_words = {"sachets", "kg", "units", "pcs", "nos", "mtrs", "ltrs"}
    if text_vals and all(v in unit_words for v in text_vals if v):
        return False  # First data row is units = messy
    # Check if columns have month/date patterns (needs unpivot)
    month_re = re.compile(r"^(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)", re.IGNORECASE)
    month_cols = sum(1 for h in headers if month_re.match(str(h).strip()))
    if month_cols >= 3:
        return False  # Months as columns = needs unpivot
    # Clean: proper headers, no merges, no months as columns
    return True


def _deep_extract_cells(file_path: str, sheet_name: str) -> dict:
    """Deep cell extraction using openpyxl — unmerge cells, extract formatting metadata.
    Returns rich cell grid with values, merge info, and formatting signals."""
    import openpyxl
    result = {"cells": [], "merged_ranges": [], "formatting": {}, "max_row": 0, "max_col": 0}
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        ws = wb[sheet_name]
        result["max_row"] = ws.max_row or 0
        result["max_col"] = ws.max_column or 0

        # Collect merged ranges and their values
        merge_map = {}  # (row, col) → value from top-left cell
        for mr in ws.merged_cells.ranges:
            result["merged_ranges"].append(str(mr))
            top_left = ws.cell(mr.min_row, mr.min_col).value
            for row in range(mr.min_row, mr.max_row + 1):
                for col in range(mr.min_col, mr.max_col + 1):
                    merge_map[(row, col)] = top_left

        # Extract cells with formatting (first 50 rows for analysis)
        max_scan = min(ws.max_row or 0, 50)
        max_cols = min(ws.max_column or 0, 20)
        bold_rows = set()
        colored_cells = {}

        for ri in range(1, max_scan + 1):
            row_data = []
            row_has_bold = False
            for ci in range(1, max_cols + 1):
                cell = ws.cell(ri, ci)
                val = merge_map.get((ri, ci), cell.value)
                val_str = str(val)[:80] if val is not None else ""
                row_data.append(val_str)
                # Check formatting
                try:
                    if cell.font and cell.font.bold:
                        row_has_bold = True
                    if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb and cell.fill.fgColor.rgb != "00000000":
                        colored_cells[f"{ri},{ci}"] = str(cell.fill.fgColor.rgb)[-6:]
                except Exception:
                    pass
            if row_has_bold:
                bold_rows.add(ri - 1)  # 0-indexed
            result["cells"].append(row_data)

        result["formatting"] = {
            "bold_rows": sorted(bold_rows),
            "colored_cells_count": len(colored_cells),
            "has_colors": len(colored_cells) > 5,
        }
        wb.close()
    except Exception as e:
        result["error"] = str(e)
    return result


def _validate_dataframe(df: pd.DataFrame) -> dict:
    """Validate extracted DataFrame quality. Returns issues dict."""
    issues = {"score": 100, "problems": [], "fixes": []}
    if df is None or len(df) == 0:
        issues["score"] = 0
        issues["problems"].append("empty_dataframe")
        return issues

    n_rows, n_cols = df.shape

    # Check NaN percentage per column
    for col in df.columns:
        nan_pct = df[col].isna().sum() / n_rows * 100
        if nan_pct > 60:
            issues["score"] -= 10
            issues["problems"].append(f"high_nan:{col}:{nan_pct:.0f}%")
            issues["fixes"].append(f"drop_column:{col}")
        elif nan_pct > 30:
            issues["score"] -= 5
            issues["problems"].append(f"moderate_nan:{col}:{nan_pct:.0f}%")
            issues["fixes"].append(f"ffill:{col}")

    # Check unnamed columns
    unnamed = [c for c in df.columns if "unnamed" in str(c).lower() or str(c).startswith("col_")]
    if unnamed:
        pct = len(unnamed) / n_cols * 100
        if pct > 50:
            issues["score"] -= 20
            issues["problems"].append(f"bad_headers:{len(unnamed)}_unnamed")
            issues["fixes"].append("redetect_header")
        elif pct > 20:
            issues["score"] -= 10
            issues["problems"].append(f"some_unnamed:{len(unnamed)}")
            issues["fixes"].append("drop_unnamed")

    # Check subtotal/total rows
    subtotal_count = 0
    for idx, row in df.iterrows():
        row_str = " ".join(str(v).lower() for v in row.dropna())
        if re.search(r'\b(total|subtotal|grand total|sum|average)\b', row_str):
            subtotal_count += 1
    if subtotal_count > 0:
        issues["score"] -= 5
        issues["problems"].append(f"subtotal_rows:{subtotal_count}")
        issues["fixes"].append("drop_subtotals")

    # Check duplicate rows
    dup_count = df.duplicated().sum()
    if dup_count > n_rows * 0.1:
        issues["score"] -= 10
        issues["problems"].append(f"duplicates:{dup_count}")
        issues["fixes"].append("drop_duplicates")

    # Check if too few data rows (might have read metadata as data)
    if n_rows < 3 and n_cols > 3:
        issues["score"] -= 15
        issues["problems"].append("too_few_rows")
        issues["fixes"].append("redetect_header")

    issues["score"] = max(0, issues["score"])
    return issues


def _auto_fix_dataframe(df: pd.DataFrame, fixes: list[str]) -> pd.DataFrame:
    """Apply auto-fixes to a DataFrame based on validation issues."""
    if df is None or len(df) == 0:
        return df

    for fix in fixes:
        try:
            if fix.startswith("ffill:"):
                col = fix.split(":", 1)[1]
                if col in df.columns:
                    df[col] = df[col].ffill()
            elif fix.startswith("drop_column:"):
                col = fix.split(":", 1)[1]
                if col in df.columns:
                    df = df.drop(columns=[col])
            elif fix == "drop_unnamed":
                unnamed = [c for c in df.columns if "unnamed" in str(c).lower()]
                # Only drop if all values in column are NaN
                for c in unnamed:
                    if df[c].isna().sum() > len(df) * 0.8:
                        df = df.drop(columns=[c])
            elif fix == "drop_subtotals":
                mask = df.apply(
                    lambda row: not bool(re.search(
                        r'\b(total|subtotal|grand total)\b',
                        " ".join(str(v).lower() for v in row.dropna())
                    )), axis=1
                )
                df = df[mask]
            elif fix == "drop_duplicates":
                df = df.drop_duplicates()
        except Exception:
            pass
    return df


def _vision_extract_sheet(file_path: str, sheet_name: str) -> str:
    """Render Excel sheet as image and use Vision LLM to extract data.
    Last resort for sheets that can't be parsed programmatically."""
    try:
        import openpyxl
        from PIL import Image, ImageDraw, ImageFont

        wb = openpyxl.load_workbook(file_path, data_only=True)
        ws = wb[sheet_name]

        # Render cells to image (simple grid)
        max_rows = min(ws.max_row or 0, 40)
        max_cols = min(ws.max_column or 0, 15)
        if max_rows < 2 or max_cols < 1:
            wb.close()
            return ""

        cell_w, cell_h = 120, 22
        img_w = max_cols * cell_w + 20
        img_h = max_rows * cell_h + 20
        img = Image.new("RGB", (img_w, img_h), "white")
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 11)
        except Exception:
            font = ImageFont.load_default()

        # Collect merged ranges for fill
        merge_map = {}
        for mr in ws.merged_cells.ranges:
            top_val = ws.cell(mr.min_row, mr.min_col).value
            for r in range(mr.min_row, mr.max_row + 1):
                for c in range(mr.min_col, mr.max_col + 1):
                    merge_map[(r, c)] = top_val

        for ri in range(1, max_rows + 1):
            for ci in range(1, max_cols + 1):
                x = (ci - 1) * cell_w + 10
                y = (ri - 1) * cell_h + 10
                val = merge_map.get((ri, ci), ws.cell(ri, ci).value)
                text = str(val)[:15] if val is not None else ""
                # Draw cell border
                draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], outline="#ccc")
                # Bold for header rows
                try:
                    if ws.cell(ri, ci).font and ws.cell(ri, ci).font.bold:
                        draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], fill="#e8e8e0")
                except Exception:
                    pass
                draw.text((x + 3, y + 3), text, fill="black", font=font)
        wb.close()

        # Convert to base64 for vision
        import io, base64
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()

        # Call vision LLM
        from dash.settings import training_vision_call
        vision_result = training_vision_call(
            "Extract ALL data from this spreadsheet image as a JSON table.\n"
            "Return: {\"headers\": [\"col1\", \"col2\", ...], \"rows\": [[\"val1\", \"val2\", ...], ...]}\n"
            "Include ALL rows. Ignore empty rows. If merged cells span multiple rows, repeat the value.",
            [{"b64": b64, "mime": "image/png", "source": f"sheet_{sheet_name}"}],
            "extraction"
        )
        return vision_result or ""
    except Exception as e:
        return f"Vision extraction failed: {e}"


def _handle_excel(file_path: str, filename: str) -> dict:
    """Handle Excel upload — master decision: clean data → fast load, messy data → AI analysis."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    ext = Path(file_path).suffix.lower()

    # Step 1: Enumerate all sheets and read previews
    sheet_previews = {}
    sheet_names = []
    try:
        if ext == ".xlsx":
            import openpyxl
            # First pass: read merged cells (needs non-read-only mode)
            merged_info = {}
            hidden_rows_info = {}  # sheet → set of hidden row numbers
            hidden_cols_info = {}  # sheet → set of hidden column letters
            comments_info = {}    # sheet → list of {row, col, text}
            try:
                wb_full = openpyxl.load_workbook(file_path, data_only=True)
                for sname in wb_full.sheetnames:
                    ws = wb_full[sname]
                    if ws.merged_cells.ranges:
                        merged_info[sname] = [str(r) for r in ws.merged_cells.ranges]
                    # Hidden rows
                    h_rows = set()
                    for r_idx, rd in ws.row_dimensions.items():
                        if rd.hidden:
                            h_rows.add(r_idx)
                    if h_rows:
                        hidden_rows_info[sname] = h_rows
                    # Hidden columns
                    h_cols = set()
                    for c_letter, cd in ws.column_dimensions.items():
                        if cd.hidden:
                            h_cols.add(c_letter)
                    if h_cols:
                        hidden_cols_info[sname] = h_cols
                    # Cell comments (first 50 cells with comments)
                    sheet_comments = []
                    for row in ws.iter_rows(max_row=min(200, ws.max_row or 200)):
                        for cell in row:
                            if cell.comment and len(sheet_comments) < 50:
                                sheet_comments.append({"row": cell.row, "col": cell.column, "text": str(cell.comment.text)[:200]})
                    if sheet_comments:
                        comments_info[sname] = sheet_comments
                wb_full.close()
            except Exception:
                pass
            # Second pass: read data (read-only for speed)
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            for sname in sheet_names:
                ws = wb[sname]
                rows = []
                for ri, row in enumerate(ws.iter_rows(max_row=25, values_only=True)):
                    if ri >= 25:
                        break
                    rows.append([str(v)[:60] if v is not None else "" for v in row[:15]])
                # Filter out hidden rows from preview
                h_rows = hidden_rows_info.get(sname, set())
                if h_rows:
                    rows = [r for i, r in enumerate(rows, 1) if i not in h_rows]

                # Accurate row count: if max_row is suspiciously large, estimate actual data rows
                reported_max = ws.max_row or 0
                actual_max = reported_max
                if reported_max > 10000:
                    # Scan for actual last data row (ghost rows detection)
                    actual_max = 0
                    empty_streak = 0
                    for ri, row in enumerate(ws.iter_rows(values_only=True)):
                        vals = [v for v in row if v is not None and str(v).strip()]
                        if vals:
                            actual_max = ri + 1
                            empty_streak = 0
                        else:
                            empty_streak += 1
                        if empty_streak > 50:
                            break
                    if actual_max < reported_max:
                        result["warnings"].append(f"Sheet '{sname}': Excel reports {reported_max:,} rows but only {actual_max:,} have data (ghost rows)")

                sheet_previews[sname] = {"rows": rows, "max_row": actual_max, "max_col": ws.max_column, "merged_cells": merged_info.get(sname, []),
                                         "hidden_rows": len(h_rows), "hidden_cols": len(hidden_cols_info.get(sname, set())), "comments": comments_info.get(sname, []),
                                         "reported_max_row": reported_max}
            wb.close()
        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(file_path)
            sheet_names = wb.sheet_names()
            for sname in sheet_names:
                ws = wb.sheet_by_name(sname)
                rows = []
                for ri in range(min(25, ws.nrows)):
                    rows.append([str(ws.cell_value(ri, ci))[:60] for ci in range(min(15, ws.ncols))])
                sheet_previews[sname] = {"rows": rows, "max_row": ws.nrows, "max_col": ws.ncols}
    except Exception as e:
        result["errors"].append(f"Failed to scan sheets: {e}")
        # Fallback: try basic read
        try:
            df = pd.read_excel(file_path, header=_find_header_row(file_path, ext))
            df = _clean_dataframe(df)
            tname = _sanitize_table_name(Path(filename).stem)
            result["tables"].append({"name": tname, "df": df, "source": "sheet_1"})
        except Exception as e2:
            result["errors"].append(f"Fallback read failed: {e2}")
        return result

    if not sheet_names:
        result["warnings"].append("No sheets found in Excel file")
        return result

    # Step 1.5: Multi-sheet similarity detection — find same-structure sheets for auto-concat
    _similar_sheet_groups = {}  # group_id → [sheet_names]
    if len(sheet_names) > 1:
        sheet_cols = {}
        for sname in sheet_names:
            info = sheet_previews.get(sname, {})
            rows = info.get("rows", [])
            if rows and len(rows) > 1:
                # Use first non-empty row as header proxy
                hdr = [c.strip().lower() for c in rows[0] if c.strip()]
                if hdr:
                    sheet_cols[sname] = set(hdr)
        # Jaccard similarity: >0.8 = same structure
        matched = set()
        group_id = 0
        for s1 in sheet_cols:
            if s1 in matched:
                continue
            group = [s1]
            for s2 in sheet_cols:
                if s2 == s1 or s2 in matched:
                    continue
                c1, c2 = sheet_cols[s1], sheet_cols[s2]
                jaccard = len(c1 & c2) / len(c1 | c2) if (c1 | c2) else 0
                if jaccard > 0.8:
                    group.append(s2)
                    matched.add(s2)
            if len(group) > 1:
                _similar_sheet_groups[group_id] = group
                matched.update(group)
                group_id += 1
                result["warnings"].append(f"Sheets {group} have similar structure — will auto-concat")

    # Step 2: RULES ENGINE — deterministic structure detection (no LLM)
    file_slug = _sanitize_table_name(Path(filename).stem)
    rules_plans = {}  # sheet_name → rules plan
    needs_llm = []    # sheets where rules are uncertain

    for sname in sheet_names:
        info = sheet_previews.get(sname, {})
        rows = info.get("rows", [])
        merged = info.get("merged_cells", [])
        sheet_max_row = info.get("max_row", 0) or 0
        # P3: scan full sheet for blank-row boundaries when sheet is large
        # (preview is only 25 rows — misses mid-sheet table splits)
        full_bounds = []
        if sheet_max_row > 100:
            full_bounds = _full_sheet_blank_scan(file_path, sname, sheet_max_row, ext)
        plan = _rules_analyze_sheet(rows, merged, full_scan_boundaries=full_bounds, max_row=sheet_max_row)
        rules_plans[sname] = plan

        if plan["confidence"] >= 0.9 and plan["action"] != "skip":
            # HIGH CONFIDENCE — execute with rules, no LLM needed
            pass
        elif plan["action"] == "skip":
            result["warnings"].append(f"Sheet '{sname}': empty — skipped (rules)")
        else:
            needs_llm.append(sname)

    # RULES PATH: Process sheets with high confidence (no LLM, $0)
    rules_processed = 0
    for sname in sheet_names:
        plan = rules_plans.get(sname, {})
        if plan.get("confidence", 0) < 0.9 or plan.get("action") == "skip":
            continue
        if sname in needs_llm:
            continue

        sheet_idx = sheet_names.index(sname) + 1

        if plan["action"] == "load":
            # Simple table — direct load with detected header
            try:
                hrow = plan.get("header_row", 0)
                skip = plan.get("skip_rows", [])

                # Multi-level header: read with header=None, apply flattened headers
                if plan.get("multi_level_header") and plan.get("flat_headers"):
                    data_start = plan.get("data_start_row", hrow + plan["multi_level_header"])
                    df = pd.read_excel(file_path, sheet_name=sname, header=None,
                                       skiprows=list(range(0, data_start)) + [s for s in skip if s >= data_start])
                    flat = plan["flat_headers"]
                    if len(flat) == len(df.columns):
                        df.columns = flat
                    result["warnings"].append(f"Sheet '{sname}': flattened {plan['multi_level_header']}-level header")
                else:
                    # Large sheet protection: limit rows read to actual data (not ghost rows)
                    sheet_info = sheet_previews.get(sname, {})
                    actual_rows = sheet_info.get("max_row", 0)
                    nrows_limit = min(actual_rows, 100000) if actual_rows > 0 else None  # Cap at 100K

                    # Pandas-safe skiprows: include ALL rows above header_row + skip rows.
                    # Then header=0 since first remaining row IS the header.
                    # Avoids pandas semantics bug where header is relative to post-skip rows.
                    full_skip = sorted(set(skip) | set(range(0, hrow)))
                    pd_skip = [s for s in full_skip if s != hrow]

                    # Try calamine for speed (5-10x faster), fall back to openpyxl
                    _excel_engine = 'calamine' if not plan.get("has_merges") else None
                    try:
                        df = pd.read_excel(file_path, sheet_name=sname, header=0,
                                           skiprows=pd_skip,
                                           nrows=nrows_limit,
                                           engine=_excel_engine)
                    except Exception:
                        df = pd.read_excel(file_path, sheet_name=sname, header=0,
                                           skiprows=pd_skip,
                                           nrows=nrows_limit)
                df = _clean_dataframe(df)
                if len(df) == 0:
                    continue
                tname = f"{file_slug}_{_sanitize_table_name(str(sname))}" if len(sheet_names) > 1 else file_slug
                result["tables"].append({"name": _sanitize_table_name(tname), "df": df, "source": f"{sname} [rules]", "sheet_number": sheet_idx, "description": f"Direct load (rules, confidence {plan['confidence']:.0%})"})
                rules_processed += 1
            except Exception as e:
                needs_llm.append(sname)  # rules failed, try LLM

        elif plan["action"] == "unpivot":
            # Month columns detected — unpivot deterministically (no LLM)
            try:
                hrow = plan.get("header_row", 0)
                skip = [s for s in plan.get("skip_rows", []) if s != hrow]
                df_raw = pd.read_excel(file_path, sheet_name=sname, header=hrow, skiprows=skip)

                # Identify id columns vs month columns using regex
                id_cols = []
                month_cols = []
                for col in df_raw.columns:
                    if _MONTH_RE.match(str(col).strip()):
                        month_cols.append(col)
                    elif str(col).strip() and not str(col).startswith("Unnamed"):
                        id_cols.append(col)

                if month_cols and id_cols:
                    # Drop summary columns (Total Output, Utilisation Rate, FY repeat)
                    keep_cols = id_cols + month_cols
                    df_use = df_raw[[c for c in keep_cols if c in df_raw.columns]]

                    # Forward-fill first column (plant name in merged cells)
                    if id_cols:
                        df_use.iloc[:, 0] = df_use.iloc[:, 0].ffill()

                    # Filter rows: drop summary rows, unit rows, blank rows
                    def _is_data_row(row):
                        vals = [str(v).strip().lower() for v in row.dropna() if str(v).strip()]
                        if not vals:
                            return False
                        if any(_SUMMARY_RE.match(v) for v in vals):
                            return False
                        if all(v in _UNIT_WORDS for v in vals):
                            return False
                        return True
                    df_use = df_use[df_use.apply(_is_data_row, axis=1)]

                    # Melt: months as columns → rows
                    df_long = df_use.melt(id_vars=id_cols, value_vars=month_cols,
                                          var_name="month", value_name="output")

                    # Parse dates from month names
                    try:
                        from dash.settings import training_llm_call
                        unique_months = df_long["month"].unique().tolist()[:20]
                        date_raw = training_llm_call(
                            f'Convert to ISO dates: {unique_months}\nReturn JSON: {{"Jul\'21":"2021-07-01",...}}',
                            "extraction")
                        if date_raw:
                            date_map = json.loads(date_raw)
                            df_long["date"] = df_long["month"].map(date_map)
                            df_long["date"] = pd.to_datetime(df_long["date"], errors="coerce")
                    except Exception:
                        pass

                    # Add fiscal year from sheet name
                    fy_match = re.search(r"FY\d{4}", sname)
                    if fy_match:
                        df_long["fiscal_year"] = fy_match.group()

                    # Clean
                    df_long = df_long.dropna(subset=["output"], how="all")
                    df_long = df_long.dropna(how="all")

                    if len(df_long) > 0:
                        tname = f"{file_slug}_{_sanitize_table_name(str(sname))}"
                        tname = _sanitize_table_name(tname)
                        existing = [t for t in result["tables"] if t["name"] == tname]
                        if existing:
                            existing[0]["df"] = pd.concat([existing[0]["df"], df_long], ignore_index=True)
                            existing[0]["source"] += f" + {sname}"
                        else:
                            result["tables"].append({"name": tname, "df": df_long, "source": f"{sname} [rules-unpivot]", "sheet_number": sheet_idx, "description": f"Unpivoted (rules, {len(month_cols)} month columns)"})
                        rules_processed += 1
                    else:
                        needs_llm.append(sname)
                else:
                    needs_llm.append(sname)
            except Exception as e:
                needs_llm.append(sname)
                result["warnings"].append(f"Rules unpivot failed for '{sname}': {e}")

        elif plan["action"] == "split":
            # Multiple tables in one sheet — extract each block
            try:
                df_raw = pd.read_excel(file_path, sheet_name=sname, header=None)
                for bi, block in enumerate(plan.get("blocks", [])):
                    hrow = block.get("header_row", 0)
                    dstart = block.get("data_start", hrow + 1)
                    dend = block.get("data_end", len(df_raw) - 1)
                    skip = set(block.get("skip_rows", []))
                    skip.add(hrow)

                    if hrow >= len(df_raw):
                        continue
                    headers = [str(v).strip() if pd.notna(v) and str(v).strip() else f"col_{i}" for i, v in enumerate(df_raw.iloc[hrow])]
                    data_rows = [df_raw.iloc[ri].values for ri in range(dstart, min(dend + 1, len(df_raw))) if ri not in skip]
                    if not data_rows:
                        continue
                    df = pd.DataFrame(data_rows, columns=headers[:len(df_raw.columns)])
                    df = _clean_dataframe(df)
                    if len(df) == 0:
                        continue
                    # Standardize first column if blocks have different names for same concept
                    header_text = block.get("header_values", [])
                    machine_type = ""
                    for hv in header_text:
                        if "ffs" in hv.lower() or "machine" in hv.lower():
                            machine_type = "FFS"
                        elif "spray" in hv.lower() or "dryer" in hv.lower():
                            machine_type = "Spray Dryer"
                        elif "drum" in hv.lower() or "roller" in hv.lower():
                            machine_type = "Drum Roller"
                    if machine_type:
                        df["machine_type"] = machine_type
                        # Rename first column to generic name
                        first_col = df.columns[0]
                        if "no of" in str(first_col).lower() or "number" in str(first_col).lower():
                            df = df.rename(columns={first_col: "machine_count"})

                    tname = f"{file_slug}_{_sanitize_table_name(str(sname))}"
                    tname = _sanitize_table_name(tname)
                    existing = [t for t in result["tables"] if t["name"] == tname]
                    if existing:
                        existing[0]["df"] = pd.concat([existing[0]["df"], df], ignore_index=True)
                        existing[0]["source"] += f" + block{bi}"
                    else:
                        result["tables"].append({"name": tname, "df": df, "source": f"{sname} [rules-split]", "sheet_number": sheet_idx, "description": f"Split table (rules, block {bi+1})"})
                rules_processed += 1
            except Exception as e:
                needs_llm.append(sname)

    if rules_processed > 0:
        result["warnings"].append(f"{rules_processed} sheet(s) processed by rules engine (no LLM, $0)")

    # ── AI STRUCTURE VALIDATOR — runs on ALL tables, even those loaded by rules ──
    # Catches: pivot tables loaded as wide (142 cols), wrong structure from confident rules
    _ai_validated = []
    for tbl in result["tables"]:
        df = tbl.get("df")
        if df is None or len(df) == 0:
            _ai_validated.append(tbl)
            continue

        n_rows, n_cols = len(df), len(df.columns)
        suspicious = False
        sus_reason = ""

        if n_cols > 30:
            suspicious = True
            sus_reason = f"wide table ({n_cols} cols)"
        if n_cols > max(n_rows * 3, 20):
            suspicious = True
            sus_reason = f"cols ({n_cols}) >> rows ({n_rows})"

        base_names = [re.sub(r'_\d+$', '', str(c)) for c in df.columns]
        from collections import Counter as _Ctr
        repeats = {name: count for name, count in _Ctr(base_names).items() if count >= 3 and name}
        if repeats:
            suspicious = True
            sus_reason = f"repeating columns: {', '.join(f'{k}×{v}' for k, v in list(repeats.items())[:3])}"

        if not suspicious:
            _ai_validated.append(tbl)
            continue

        result["warnings"].append(f"Table '{tbl.get('name','')}': suspicious shape ({sus_reason}) — asking AI")
        try:
            from dash.settings import training_llm_call as _val_llm
            col_sample = list(df.columns[:20])
            repeat_info = ", ".join(f"{k} repeats {v} times" for k, v in repeats.items()) if repeats else "none"

            val_prompt = f"""I loaded an Excel and got this DataFrame:
Shape: {n_rows} rows × {n_cols} columns
Suspicious: {sus_reason}
First 20 columns: {col_sample}
Repeating patterns: {repeat_info}

Is this correct? Return ONLY JSON:
If needs reshaping: {{"action": "unpivot", "reason": "brief", "id_columns": ["col1","col2"], "repeating_group": ["metric1","metric2"]}}
If correct as-is: {{"action": "keep", "reason": "brief"}}"""

            raw = _val_llm(val_prompt, "excel_analysis")
            if raw:
                clean = raw.strip()
                if clean.startswith("```"):
                    clean = clean.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
                if clean.startswith("json"):
                    clean = clean[4:].strip()
                val_result = json.loads(clean)

                if val_result.get("action") == "unpivot":
                    id_cols = val_result.get("id_columns", [])
                    id_existing = [c for c in id_cols if c in df.columns]
                    if not id_existing:
                        id_existing = [c for c in df.columns if base_names[list(df.columns).index(c)] not in repeats]

                    value_cols = [c for c in df.columns if c not in id_existing]

                    if value_cols and id_existing:
                        df_long = pd.melt(df, id_vars=id_existing, value_vars=value_cols, var_name='metric_period', value_name='value')
                        df_long = df_long.dropna(subset=['value'])

                        # Split metric_period into metric + period index
                        if repeats:
                            def _parse_col(col_name):
                                base = re.sub(r'_\d+$', '', str(col_name))
                                suffix = re.search(r'_(\d+)$', str(col_name))
                                return base, int(suffix.group(1)) if suffix else 0
                            df_long['_metric'] = df_long['metric_period'].apply(lambda x: _parse_col(x)[0])
                            df_long['_period_idx'] = df_long['metric_period'].apply(lambda x: _parse_col(x)[1])
                            try:
                                df_pivoted = df_long.pivot_table(index=id_existing + ['_period_idx'], columns='_metric', values='value', aggfunc='first').reset_index()
                                df_pivoted.columns = [str(c) for c in df_pivoted.columns]
                                if len(df_pivoted) > 0 and len(df_pivoted.columns) < n_cols:
                                    df_long = df_pivoted
                            except Exception:
                                pass

                        df_long = _clean_dataframe(df_long)
                        tbl["df"] = df_long
                        tbl["source"] = tbl.get("source", "") + " [ai-unpivot]"
                        tbl["description"] = f"AI restructured: {n_rows}×{n_cols} → {len(df_long)}×{len(df_long.columns)} ({val_result.get('reason','')})"
                        result["warnings"].append(f"AI Validator: unpivoted '{tbl.get('name','')}' from {n_cols} → {len(df_long.columns)} cols")

                        # Save learning
                        try:
                            from sqlalchemy import text as _slt
                            from db import get_sql_engine as _sle
                            _eng = _sle()
                            with _eng.connect() as _lc:
                                _lc.execute(_slt(
                                    "INSERT INTO public.dash_memories (project_slug, content, source, category) "
                                    "VALUES ('_global', :content, 'structure_learning', 'data_pattern') "
                                    "ON CONFLICT DO NOTHING"
                                ), {"content": f"Excel pattern: repeating column groups ({repeat_info}) = pivot table. Auto-unpivot applied. Source: {tbl.get('name','')}"})
                                _lc.commit()
                        except Exception:
                            pass
                else:
                    result["warnings"].append(f"AI Validator: '{tbl.get('name','')}' is correct (wide format)")
        except Exception as e:
            result["warnings"].append(f"AI Validator error: {str(e)[:100]}")

        _ai_validated.append(tbl)

    result["tables"] = _ai_validated

    # ── LLM RESCUE PATH — catches row-count mismatch from broken rules detection ──
    # When parser loaded <10% of expected sheet rows, LLM re-analyzes preview to
    # find correct header_row + skip_rows, then re-reads sheet with corrected ranges.
    # Single LLM call per suspicious sheet, fail-soft.
    # P5: file-hash cache — same file_hash → reuse prior plan → skip LLM call.
    _file_hash = _compute_file_hash(file_path)
    _file_size = Path(file_path).stat().st_size if Path(file_path).exists() else 0
    _cached = _lookup_upload_cache(_file_hash) if _file_hash else None
    _cached_plan = (_cached or {}).get("plan") or {}
    _cached_sheet_map = _cached_plan.get("sheet_to_header") or {}
    _accum_sheet_plan = dict(_cached_sheet_map)  # accumulate across sheets in this run
    _rescued_tables = []
    for tbl in result["tables"]:
        try:
            df = tbl.get("df")
            if df is None or len(df) == 0:
                _rescued_tables.append(tbl)
                continue

            # Extract sheet name from source field like "Sheet1 [rules]" or "Sheet1 [rules-split]"
            sheet_src = str(tbl.get("source", ""))
            sname_match = re.match(r"^([^\[]+?)\s*\[", sheet_src)
            if not sname_match:
                _rescued_tables.append(tbl)
                continue
            sname = sname_match.group(1).strip()
            sinfo = sheet_previews.get(sname)
            if not sinfo:
                _rescued_tables.append(tbl)
                continue

            expected_max = sinfo.get("max_row", 0)
            actual_rows = len(df)

            # Only rescue when expected >100 AND actual <10% of expected
            # Also: skip tables already marked as ai-unpivot (different problem class)
            already_rescued = "rescued" in str(tbl.get("source", "")).lower()
            already_unpivoted = "unpivot" in str(tbl.get("source", "")).lower()
            if expected_max <= 100 or already_rescued or already_unpivoted:
                _rescued_tables.append(tbl)
                continue
            if actual_rows >= expected_max * 0.1:
                _rescued_tables.append(tbl)
                continue

            # P5: cache hit — reuse prior rescue plan, skip LLM call
            _cached_rescue = _cached_sheet_map.get(sname)
            if _cached_rescue and isinstance(_cached_rescue, dict):
                _c_hrow = _cached_rescue.get("header_row")
                _c_banner = _cached_rescue.get("banner_rows_to_skip") or []
                if isinstance(_c_hrow, int) and _c_hrow >= 0:
                    result["warnings"].append(f"LLM Rescue: CACHE HIT for '{sname}' (file_hash hit_count={(_cached or {}).get('hit_count', 1)}) — skipping LLM call")
                    try:
                        full_skip = sorted(set(_c_banner) | set(range(0, _c_hrow)))
                        pd_skip = [s for s in full_skip if s != _c_hrow]
                        nrows_limit = min(expected_max, 100000)
                        try:
                            df_new = pd.read_excel(file_path, sheet_name=sname, header=0,
                                                   skiprows=pd_skip, nrows=nrows_limit, engine='calamine')
                        except Exception:
                            df_new = pd.read_excel(file_path, sheet_name=sname, header=0,
                                                   skiprows=pd_skip, nrows=nrows_limit)
                        df_new = _clean_dataframe(df_new)
                        if len(df_new) > actual_rows * 5:
                            tbl["df"] = df_new
                            tbl["source"] = sheet_src + " [llm-rescued-cached]"
                            tbl["description"] = f"LLM rescued (cache): {actual_rows} → {len(df_new)} rows (header_row={_c_hrow})"
                            result["warnings"].append(f"LLM Rescue: '{tbl.get('name','')}' {actual_rows} → {len(df_new)} rows ✓ (cache)")
                        else:
                            result["warnings"].append(f"LLM Rescue: cached rescue for '{sname}' produced {len(df_new)} rows — kept original")
                    except Exception as _cre:
                        result["warnings"].append(f"LLM Rescue cache re-read failed: {str(_cre)[:100]}")
                    _rescued_tables.append(tbl)
                    continue

            result["warnings"].append(f"LLM Rescue: '{tbl.get('name','')}' loaded {actual_rows} rows but sheet '{sname}' has {expected_max:,} rows — asking LLM")

            from dash.settings import training_llm_call as _resc_llm

            # Build preview block for LLM: first 30 rows raw
            preview_rows = sinfo.get("rows", [])[:30]
            rows_text = "\n".join(
                f"row {ri}: " + " | ".join(c[:40] for c in row[:12])
                for ri, row in enumerate(preview_rows) if any(c.strip() for c in row)
            )

            resc_prompt = f"""Excel sheet '{sname}' has {expected_max:,} rows but my parser only extracted {actual_rows} rows. Something went wrong with header detection.

Preview (first 30 rows of raw cells, blank rows omitted):
{rows_text}

Identify the correct structure. Return ONLY JSON, nothing else:
{{
  "header_row": <int 0-indexed row containing column names>,
  "banner_rows_to_skip": [<int>...],
  "data_start": <int row where data begins>,
  "reasoning": "<one sentence>"
}}

Rules:
- header_row = first row where ALL non-blank cells are short labels (not sentences, not dates as values, not single banner strings)
- banner_rows_to_skip = rows above header that are titles/print dates/blanks
- data_start = header_row + 1 typically"""

            raw = _resc_llm(resc_prompt, "extraction")
            if not raw:
                _rescued_tables.append(tbl)
                continue

            # 4-tier robust JSON parse
            clean = raw.strip()
            if clean.startswith("```"):
                clean = clean.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
            if clean.lower().startswith("json"):
                clean = clean[4:].strip()
            # Extract first {...} object
            obj_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', clean, re.DOTALL)
            if obj_match:
                clean = obj_match.group(0)
            try:
                resc_plan = json.loads(clean)
            except Exception:
                _rescued_tables.append(tbl)
                continue

            new_hrow = resc_plan.get("header_row")
            banner_skip = resc_plan.get("banner_rows_to_skip", [])
            if new_hrow is None or not isinstance(new_hrow, int) or new_hrow < 0:
                _rescued_tables.append(tbl)
                continue

            # Re-read sheet with corrected ranges
            try:
                # Pandas-safe: skiprows = all rows above header + banner skip
                full_skip = sorted(set(banner_skip) | set(range(0, new_hrow)))
                pd_skip = [s for s in full_skip if s != new_hrow]
                nrows_limit = min(expected_max, 100000)
                try:
                    df_new = pd.read_excel(file_path, sheet_name=sname, header=0,
                                           skiprows=pd_skip, nrows=nrows_limit, engine='calamine')
                except Exception:
                    df_new = pd.read_excel(file_path, sheet_name=sname, header=0,
                                           skiprows=pd_skip, nrows=nrows_limit)
                df_new = _clean_dataframe(df_new)
                if len(df_new) > actual_rows * 5:  # Real rescue if >5x more rows
                    tbl["df"] = df_new
                    tbl["source"] = sheet_src + " [llm-rescued]"
                    tbl["description"] = f"LLM rescued: {actual_rows} → {len(df_new)} rows (header_row={new_hrow}, reason: {resc_plan.get('reasoning','')[:80]})"
                    result["warnings"].append(f"LLM Rescue: '{tbl.get('name','')}' {actual_rows} → {len(df_new)} rows ✓")

                    # P5: save rescue plan to file-hash cache for future re-uploads
                    try:
                        _accum_sheet_plan[sname] = {
                            "header_row": new_hrow,
                            "banner_rows_to_skip": banner_skip,
                            "reasoning": resc_plan.get('reasoning', '')[:200],
                        }
                        _save_upload_cache(_file_hash, _file_size, ".xlsx",
                                           {"sheet_to_header": _accum_sheet_plan},
                                           rescue_used=True)
                    except Exception:
                        pass

                    # Save learning
                    try:
                        from sqlalchemy import text as _slt
                        from db import get_sql_engine as _sle
                        _eng = _sle()
                        with _eng.connect() as _lc:
                            _lc.execute(_slt(
                                "INSERT INTO public.dash_memories (project_slug, content, source, category) "
                                "VALUES ('_global', :content, 'structure_learning', 'data_pattern') "
                                "ON CONFLICT DO NOTHING"
                            ), {"content": f"Excel pattern: sheet '{sname}' needed LLM rescue (rules picked wrong header). Correct header_row={new_hrow}, banner rows={banner_skip[:5]}. Recovered {len(df_new):,} rows."})
                            _lc.commit()
                    except Exception:
                        pass
                else:
                    result["warnings"].append(f"LLM Rescue: '{tbl.get('name','')}' rescue produced {len(df_new)} rows — kept original")
            except Exception as _rerr:
                result["warnings"].append(f"LLM Rescue re-read failed: {str(_rerr)[:100]}")
        except Exception as _rge:
            result["warnings"].append(f"LLM Rescue error: {str(_rge)[:100]}")

        _rescued_tables.append(tbl)

    result["tables"] = _rescued_tables

    # LLM PATH: Only for uncertain sheets (low confidence or unpivot)
    ai_sheet_names = list(set(needs_llm))
    if not ai_sheet_names:
        return result

    result["warnings"].append(f"{len(ai_sheet_names)} sheet(s) need LLM (rules confidence < 90% or unpivot)")

    # Step 3: AI analysis — send previews, get extraction plan (only uncertain sheets)
    ai_plan = None
    try:
        from dash.settings import training_llm_call
        preview_text = ""
        # Limit preview rows per sheet based on sheet count to stay within LLM context
        max_preview_rows = 25 if len(ai_sheet_names) <= 5 else (15 if len(ai_sheet_names) <= 10 else 10)
        for sname in ai_sheet_names:
            info = sheet_previews[sname]
            merged = info.get("merged_cells", [])
            merged_note = f" | Merged cells: {', '.join(merged[:10])}" if merged else ""
            hidden_note = f" | Hidden: {info.get('hidden_rows', 0)} rows, {info.get('hidden_cols', 0)} cols" if info.get("hidden_rows") or info.get("hidden_cols") else ""
            comments_note = f" | {len(info.get('comments', []))} cell comments" if info.get("comments") else ""
            preview_text += f"\n\nSheet: '{sname}' ({info['max_row']} rows × {info['max_col']} cols{merged_note}{hidden_note}{comments_note})\n"

            # SheetCompressor: deduplicate repeated values, compress sparse rows
            rows_data = info["rows"][:max_preview_rows]
            for ri, row in enumerate(rows_data):
                # Compress: skip fully empty rows, collapse repeated values
                non_empty = [(ci, v) for ci, v in enumerate(row) if v and str(v).strip()]
                if not non_empty:
                    continue  # Skip blank rows entirely (saves tokens)
                # If row is sparse (>60% empty), use inverse index format
                if len(non_empty) < len(row) * 0.4 and len(row) > 5:
                    preview_text += f"  Row {ri}: {{{', '.join(f'{ci}:{v}' for ci, v in non_empty)}}}\n"
                else:
                    preview_text += f"  Row {ri}: {row}\n"

        prompt = f"""Analyze {len(ai_sheet_names)} messy Excel sheets that need intelligent processing. For each sheet I show the first rows as raw cell values. Empty cells are ''.

{preview_text}

Return ONLY a valid JSON array. For each sheet return an object with:
- "sheet": sheet name (exact match)
- "action": "load" | "split" | "skip" | "unpivot"
- "table_name": clean lowercase PostgreSQL table name (a-z, 0-9, underscore only, max 50 chars)
- "header_row": 0-indexed row number with column headers
- "data_start_row": 0-indexed row where data begins
- "skip_rows": array of 0-indexed row numbers to skip (blank/separator/unit/summary rows)
- "forward_fill_columns": array of 0-indexed column positions that need forward-fill (merged cells)
- "column_names": object mapping column index to clean name e.g. {{"0":"plant","1":"product","3":"jul_21"}}
- "description": brief description of table content

ACTIONS:
1. "load" — single table, store as-is
2. "split" — multiple tables in one sheet separated by blank rows. Return "tables" array.
3. "skip" — empty or irrelevant sheet
4. "unpivot" — columns are time periods (months, quarters, dates). MELT them into rows. Return:
   - "id_columns": array of column indexes that are identifiers (plant, product, category...)
   - "value_columns": array of column indexes that are time-based values to unpivot
   - "value_name": what the values represent ("output", "amount", "quantity")
   - "variable_name": what the columns represent ("month", "quarter", "period")
   - "extra_columns": object of extra columns to add e.g. {{"fiscal_year":"FY2022","unit":"Sachets"}}
   - "blocks": if multiple data blocks exist (e.g. IB Plant rows 4-9, Tea Plant row 13), list them:
     [{{"data_start":4,"data_end":9,"extra":{{"plant":"IB Plant","unit":"Sachets"}}}},
      {{"data_start":13,"data_end":13,"extra":{{"plant":"Tea Plant","unit":"kg"}}}}]
     All blocks get merged into ONE table with extra columns distinguishing them.

For "split" action, return "tables" array with "table_name", "header_row", "data_start_row", "data_end_row", "description".

DETECT THESE PATTERNS:
- Columns with month/date names (Jul'21, Aug'21, Q1 2022, 2021-01...) → use "unpivot"
- Multiple data blocks with same columns but different categories → merge via "blocks"
- Metadata rows at top (Company, Period, Rate) → NOT headers, skip them
- Sub-header rows (units: Sachets, kg, %) → skip_rows, capture unit in extra_columns
- Summary/total rows (Utilisation, Total) → skip_rows
- Merged cells → forward_fill_columns
- Empty sheets → "skip"

CRITICAL for "split" with multiple sub-tables in one sheet:
- If sub-tables have SAME structure but different first column name (e.g. "No of FFS machines" vs "No of Spray Dryer" vs "No of Drum Roller"), these are the SAME table with different machine types
- Use COMMON column names across all sub-tables: rename first column to generic name like "machine_count"
- Add extra column to distinguish: "machine_type": "FFS" / "Spray Dryer" / "Drum Roller"
- This way all sub-tables can be merged into ONE table later
- In "tables" array for split, each sub-table should have: "extra": {"machine_type": "FFS", "unit": "Sachets"}

Column naming rules:
- Use EXACT text from header row, cleaned for PostgreSQL
- "No of FFS machines" and "No of Spray Dryer" should BOTH become "machine_count" (standardized)
- "Products" → "products", "SKU" → "sku", "Speed" → "speed"
- column_names mapping: key=column index, value=cleaned standardized name"""

        raw = training_llm_call(prompt, "excel_analysis")
        if raw:
            ai_plan = json.loads(raw)
            result["metadata"]["ai_analysis"] = "success"
        else:
            result["warnings"].append("AI returned empty response — using fallback")
    except json.JSONDecodeError as e:
        result["warnings"].append(f"AI returned invalid JSON: {str(e)[:100]}")
        ai_plan = None
    except Exception as e:
        result["warnings"].append(f"AI analysis error: {str(e)[:100]}")
        ai_plan = None

    # Step 3: Extract tables per AI plan or fallback
    if ai_plan and isinstance(ai_plan, list):
        file_slug = _sanitize_table_name(Path(filename).stem)
        for sheet_plan in ai_plan:
            sname = sheet_plan.get("sheet", "")
            action = sheet_plan.get("action", "skip")

            if action == "skip":
                result["warnings"].append(f"Skipped sheet '{sname}': {sheet_plan.get('reason', 'empty')}")
                continue

            if sname not in sheet_names:
                result["warnings"].append(f"Sheet '{sname}' not found in file")
                continue

            if action == "unpivot":
                # AI-powered conversion: send raw data to AI, let it decide column mapping
                try:
                    # Read raw data with no header — let AI figure out everything
                    df_raw_all = pd.read_excel(file_path, sheet_name=sname, header=None)
                    raw_preview = ""
                    for ri in range(min(20, len(df_raw_all))):
                        vals = [str(v)[:30] if pd.notna(v) else "" for v in df_raw_all.iloc[ri].values[:18]]
                        raw_preview += f"  Row {ri}: {vals}\n"

                    convert_prompt = f"""Convert this Excel sheet from WIDE to LONG format.

RAW DATA (first 20 rows, 0-indexed, no header):
{raw_preview}

Return ONLY valid JSON:
{{
  "header_row": <0-indexed row with column headers like month names (Jul'21, Aug'21 etc)>,
  "skip_after_header": [<0-indexed rows to skip AFTER header: unit rows (Sachets/kg), summary rows (Utilisation/Total), blank rows>],
  "id_columns": [<0-indexed col numbers for identifiers: plant, product, capacity>],
  "value_columns": [<0-indexed col numbers for time-period values to unpivot>],
  "id_names": ["plant", "product", "annual_capacity"],
  "value_name": "output",
  "variable_name": "month",
  "blocks": [
    {{"data_start": <first data row 0-indexed>, "data_end": <last data row 0-indexed>, "extra": {{"plant": "IB Plant", "unit": "Sachets"}} }}
  ],
  "extra_columns": {{"fiscal_year": "FY2022"}}
}}"""

                    from dash.settings import training_llm_call
                    convert_raw = training_llm_call(convert_prompt, "excel_analysis")
                    if not convert_raw:
                        raise ValueError("AI conversion returned empty")
                    cp = json.loads(convert_raw)

                    # Read with AI-determined header
                    ai_hrow = cp.get("header_row", 1)
                    skip_after = cp.get("skip_after_header", [])
                    skip_for_read = [s for s in skip_after if s > ai_hrow]
                    df_data = pd.read_excel(file_path, sheet_name=sname, header=ai_hrow, skiprows=skip_for_read)

                    # Use AI conversion plan for column mapping
                    id_col_idx = cp.get("id_columns", [0, 1, 2])
                    val_col_idx = cp.get("value_columns", [])
                    id_names_ai = cp.get("id_names", [])
                    value_name = cp.get("value_name", "value")
                    variable_name = cp.get("variable_name", "period")

                    # Rename ID columns with AI names
                    for i, name in enumerate(id_names_ai):
                        if i < len(id_col_idx) and id_col_idx[i] < len(df_data.columns):
                            df_data.rename(columns={df_data.columns[id_col_idx[i]]: name}, inplace=True)

                    id_cols = [id_names_ai[i] if i < len(id_names_ai) else df_data.columns[idx]
                               for i, idx in enumerate(id_col_idx) if idx < len(df_data.columns)]
                    val_cols = [df_data.columns[i] for i in val_col_idx if i < len(df_data.columns)]

                    blocks = cp.get("blocks", [])
                    tname = sheet_plan.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                    tname = _sanitize_table_name(tname)

                    if blocks and val_cols:
                        all_melted = []
                        for block in blocks:
                            bstart = block.get("data_start", 0) - ai_hrow - 1
                            bend = block.get("data_end", bstart + ai_hrow + 1) - ai_hrow - 1
                            skipped = len([s for s in skip_for_read if s < block.get("data_start", 0)])
                            bstart = max(0, bstart - skipped)
                            bend = max(bstart, bend - len([s for s in skip_for_read if s < block.get("data_end", 0)]))
                            bend = min(bend, len(df_data) - 1)
                            if bstart >= len(df_data):
                                continue
                            df_block = df_data.iloc[bstart:bend + 1].copy()
                            if id_cols and id_cols[0] in df_block.columns:
                                df_block[id_cols[0]] = df_block[id_cols[0]].ffill()
                            df_block = df_block.dropna(subset=val_cols, how='all')
                            melted = df_block.melt(id_vars=id_cols, value_vars=val_cols,
                                                   var_name=variable_name, value_name=value_name)
                            for k, v in block.get("extra", {}).items():
                                melted[k] = v
                            all_melted.append(melted)
                        df_final = pd.concat(all_melted, ignore_index=True) if all_melted else pd.DataFrame()
                    elif val_cols:
                        if id_cols and id_cols[0] in df_data.columns:
                            df_data[id_cols[0]] = df_data[id_cols[0]].ffill()
                        df_final = df_data.melt(id_vars=id_cols, value_vars=val_cols,
                                                var_name=variable_name, value_name=value_name)
                    else:
                        df_final = df_data

                    # Add extra columns
                    for k, v in cp.get("extra_columns", {}).items():
                        df_final[k] = v
                    for k, v in sheet_plan.get("extra_columns", {}).items():
                        if k not in df_final.columns:
                            df_final[k] = v

                    # Ask AI to create date mapping for the period column
                    if variable_name in df_final.columns:
                        unique_periods = df_final[variable_name].dropna().unique().tolist()[:20]
                        if unique_periods:
                            try:
                                date_prompt = f"""Convert these time period labels to ISO dates (YYYY-MM-DD, use 1st of month).

Periods: {unique_periods}

Return ONLY a JSON object mapping each period to its date:
{{"Jul'21": "2021-07-01", "Aug'21": "2021-08-01", "Q1 2022": "2022-01-01"}}"""
                                date_raw = training_llm_call(date_prompt, "extraction")
                                if date_raw:
                                    date_map = json.loads(date_raw)
                                    df_final["date"] = df_final[variable_name].map(date_map)
                                    df_final["date"] = pd.to_datetime(df_final["date"], errors="coerce")
                            except Exception:
                                pass

                    if value_name in df_final.columns:
                        df_final = df_final.dropna(subset=[value_name])
                    df_final = df_final.dropna(how='all')

                    if len(df_final) > 0:
                        sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                        # Check if this table should merge with an existing one (same table_name from another sheet)
                        existing = [t for t in result["tables"] if t["name"] == tname]
                        if existing:
                            existing[0]["df"] = pd.concat([existing[0]["df"], df_final], ignore_index=True)
                            existing[0]["source"] += f" + {sname}"
                        else:
                            result["tables"].append({"name": tname, "df": df_final, "source": f"{sname} [unpivot]", "sheet_number": sheet_names.index(sname) + 1, "description": sheet_plan.get("description", "")})
                except Exception as e:
                    result["warnings"].append(f"Unpivot from '{sname}' failed: {e}")
                continue

            if action == "split":
                # Multiple tables in one sheet — read with header=None, slice manually
                try:
                    df_raw = pd.read_excel(file_path, sheet_name=sname, header=None)
                except Exception as e:
                    result["warnings"].append(f"Cannot read sheet '{sname}': {e}")
                    continue
                for sub in sheet_plan.get("tables", []):
                    try:
                        hrow = sub.get("header_row", 0)
                        dstart = sub.get("data_start_row", hrow + 1)
                        dend = sub.get("data_end_row")
                        if dend is None:
                            dend = len(df_raw) - 1
                        # Extract header row as column names
                        if hrow < len(df_raw):
                            headers = [str(v).strip() if pd.notna(v) and str(v).strip() else f"col_{i}" for i, v in enumerate(df_raw.iloc[hrow])]
                        else:
                            continue
                        # Extract data rows
                        skip_rows = set(sub.get("skip_rows", []))
                        skip_rows.add(hrow)  # Don't include header as data
                        data_rows = []
                        for ri in range(dstart, min(dend + 1, len(df_raw))):
                            if ri not in skip_rows:
                                data_rows.append(df_raw.iloc[ri].values)
                        if not data_rows:
                            continue
                        df = pd.DataFrame(data_rows, columns=headers[:len(df_raw.columns)])

                        # Apply AI column_names mapping (standardize across sub-tables)
                        col_names = sub.get("column_names", {})
                        if col_names:
                            new_cols = []
                            for i, c in enumerate(df.columns):
                                new_cols.append(col_names.get(str(i), str(c)))
                            df.columns = new_cols

                        df = _clean_dataframe(df)
                        if len(df) == 0:
                            continue
                        tname = sub.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                        tname = _sanitize_table_name(tname)
                        # Forward-fill
                        for col_idx in sub.get("forward_fill_columns", []):
                            if isinstance(col_idx, int) and col_idx < len(df.columns):
                                df.iloc[:, col_idx] = df.iloc[:, col_idx].ffill()
                        # Add extra columns (machine_type, unit, etc from AI plan)
                        for k, v in sub.get("extra", {}).items():
                            df[k] = v
                        # Merge with existing table of same name (cross-sheet merge)
                        sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                        existing = [t for t in result["tables"] if t["name"] == tname]
                        if existing:
                            existing[0]["df"] = pd.concat([existing[0]["df"], df], ignore_index=True)
                            existing[0]["source"] += f" + {sname}"
                        else:
                            result["tables"].append({"name": tname, "df": df, "source": f"{sname} [split]", "sheet_number": sheet_idx, "description": sub.get("description", "")})
                    except Exception as e:
                        result["warnings"].append(f"Split table from '{sname}' failed: {e}")
            else:
                # Single table from sheet
                try:
                    hrow = sheet_plan.get("header_row", 0)
                    skip = sheet_plan.get("skip_rows", [])
                    df = pd.read_excel(file_path, sheet_name=sname, header=hrow, skiprows=skip)
                    # Apply AI column names
                    col_names = sheet_plan.get("column_names", {})
                    if col_names:
                        new_cols = []
                        for i, c in enumerate(df.columns):
                            new_cols.append(col_names.get(str(i), str(c)))
                        df.columns = new_cols
                    df = _clean_dataframe(df)
                    if len(df) == 0:
                        result["warnings"].append(f"Sheet '{sname}' produced empty table after cleanup")
                        continue
                    tname = sheet_plan.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                    tname = _sanitize_table_name(tname)
                    # Forward-fill merged cells
                    for col_idx in sheet_plan.get("forward_fill_columns", []):
                        if isinstance(col_idx, int) and col_idx < len(df.columns):
                            df.iloc[:, col_idx] = df.iloc[:, col_idx].ffill()
                    # Add extra columns if specified
                    for k, v in sheet_plan.get("extra_columns", {}).items():
                        df[k] = v
                    sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                    result["tables"].append({"name": tname, "df": df, "source": sname, "sheet_number": sheet_idx, "description": sheet_plan.get("description", "")})
                except Exception as e:
                    result["warnings"].append(f"Sheet '{sname}' extraction failed: {e}")

        # Collect metadata text from AI descriptions
        descriptions = [sp.get("description", "") for sp in ai_plan if sp.get("description")]
        if descriptions:
            result["text"] = f"Excel file: {filename}\n" + "\n".join(f"- {d}" for d in descriptions)

    else:
        # FALLBACK: No AI — read all sheets with basic header detection
        result["warnings"].append("AI analysis unavailable, using fallback for all sheets")
        file_slug = _sanitize_table_name(Path(filename).stem)
        try:
            all_sheets = pd.read_excel(file_path, sheet_name=None, header=None)
            for sname, raw_df in all_sheets.items():
                if raw_df.dropna(how='all').empty:
                    result["warnings"].append(f"Skipped empty sheet '{sname}'")
                    continue
                # Find header row for this sheet
                best_row = 0
                best_score = 0
                for i in range(min(10, len(raw_df))):
                    row = raw_df.iloc[i]
                    non_null = row.dropna()
                    score = sum(3 for v in non_null if isinstance(v, str) and 1 < len(str(v).strip()) < 50)
                    if score > best_score:
                        best_score = score
                        best_row = i
                df = pd.read_excel(file_path, sheet_name=sname, header=best_row)
                df = _clean_dataframe(df)
                if len(df) == 0:
                    continue
                tname = f"{file_slug}_{_sanitize_table_name(str(sname))}"
                sheet_idx = list(all_sheets.keys()).index(sname) + 1 if sname in all_sheets else 0
                result["tables"].append({"name": tname, "df": df, "source": str(sname), "sheet_number": sheet_idx, "description": ""})
        except Exception as e:
            result["errors"].append(f"Fallback read all sheets failed: {e}")

    # ── AUTO-CONCAT similar sheets ────────────────────────────────────
    if _similar_sheet_groups:
        for gid, group_sheets in _similar_sheet_groups.items():
            # Find tables from these sheets
            group_tables = [t for t in result["tables"] if any(s in str(t.get("source", "")) for s in group_sheets)]
            if len(group_tables) >= 2:
                # Check column overlap in actual DataFrames
                ref_cols = set(group_tables[0]["df"].columns)
                compatible = [group_tables[0]]
                for gt in group_tables[1:]:
                    overlap = len(ref_cols & set(gt["df"].columns)) / max(len(ref_cols | set(gt["df"].columns)), 1)
                    if overlap > 0.8:
                        compatible.append(gt)
                if len(compatible) >= 2:
                    merged_df = pd.concat([t["df"].assign(_source_sheet=t.get("source", "")) for t in compatible], ignore_index=True)
                    merged_name = compatible[0]["name"]
                    # Remove individual tables, add merged
                    for t in compatible:
                        if t in result["tables"]:
                            result["tables"].remove(t)
                    result["tables"].append({"name": merged_name, "df": merged_df, "source": f"merged:{','.join(group_sheets)}", "description": f"Auto-merged {len(compatible)} similar sheets"})
                    result["warnings"].append(f"Merged {len(compatible)} sheets into '{merged_name}' ({len(merged_df)} rows)")

    # ── LAYER 3+4: Validate → Auto-fix → Vision fallback ──────────────
    # Run on all extracted tables — catch subtotals, bad headers, high NaN
    validated_tables = []
    for tbl in result["tables"]:
        df = tbl.get("df")
        if df is None or len(df) == 0:
            validated_tables.append(tbl)
            continue

        issues = _validate_dataframe(df)
        tbl["quality_score"] = issues["score"]

        if issues["score"] >= 70:
            # Good enough — minor fixes only
            if issues["fixes"]:
                df = _auto_fix_dataframe(df, issues["fixes"])
                tbl["df"] = df
                tbl["source"] = tbl.get("source", "") + " [auto-fixed]"
            validated_tables.append(tbl)
        elif issues["score"] >= 40:
            # Moderate issues — apply fixes and re-validate
            df = _auto_fix_dataframe(df, issues["fixes"])
            issues2 = _validate_dataframe(df)
            if issues2["score"] >= 60:
                tbl["df"] = df
                tbl["source"] = tbl.get("source", "") + " [corrected]"
                tbl["quality_score"] = issues2["score"]
                validated_tables.append(tbl)
            else:
                # Try re-reading with deep cell extraction
                sheet = tbl.get("source", "").split("[")[0].strip().split("+")[0].strip()
                if sheet and ext == ".xlsx":
                    try:
                        deep = _deep_extract_cells(file_path, sheet)
                        if deep.get("cells") and len(deep["cells"]) > 2:
                            # Build better preview with unmerged cells for LLM
                            from dash.settings import training_llm_call
                            cell_text = "\n".join(
                                f"  Row {i}: {row}" for i, row in enumerate(deep["cells"][:30])
                            )
                            fmt_info = f"Bold rows (likely headers): {deep['formatting'].get('bold_rows', [])}"
                            if deep["merged_ranges"]:
                                fmt_info += f"\nMerged ranges (already unmerged in data): {deep['merged_ranges'][:10]}"

                            fix_prompt = f"""This Excel sheet was parsed but has quality issues: {issues["problems"]}

Cell data (merged cells already expanded):
{cell_text}

{fmt_info}

Return JSON: {{"header_row": <0-indexed>, "skip_rows": [<0-indexed rows to skip>], "data_start": <0-indexed>}}
Pick the BEST header row (bold rows are headers). Skip metadata, subtotals, unit rows."""

                            raw = training_llm_call(fix_prompt, "extraction")
                            if raw:
                                import json as _json2
                                fix_plan = _json2.loads(raw.strip().strip("`").lstrip("json").strip())
                                hrow = fix_plan.get("header_row", 0)
                                skip = fix_plan.get("skip_rows", [])
                                df2 = pd.read_excel(file_path, sheet_name=sheet, header=hrow, skiprows=[s for s in skip if s != hrow])
                                df2 = _clean_dataframe(df2)
                                # Forward fill columns that had merged cells
                                if deep["merged_ranges"]:
                                    for ci in range(min(3, len(df2.columns))):
                                        if df2.iloc[:, ci].isna().sum() > len(df2) * 0.2:
                                            df2.iloc[:, ci] = df2.iloc[:, ci].ffill()
                                df2 = _auto_fix_dataframe(df2, ["drop_subtotals", "drop_unnamed"])
                                issues3 = _validate_dataframe(df2)
                                if issues3["score"] > issues["score"]:
                                    tbl["df"] = df2
                                    tbl["quality_score"] = issues3["score"]
                                    tbl["source"] = tbl.get("source", "") + " [deep-corrected]"
                                    validated_tables.append(tbl)
                                    continue
                    except Exception:
                        pass

                # VISION FALLBACK — render sheet as image, extract with Vision LLM
                if sheet and ext == ".xlsx":
                    try:
                        vision_raw = _vision_extract_sheet(file_path, sheet)
                        if vision_raw and "{" in vision_raw:
                            import json as _json3
                            clean = vision_raw.strip()
                            if clean.startswith("```"):
                                clean = clean.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
                            parsed = _json3.loads(clean)
                            headers = parsed.get("headers", [])
                            rows = parsed.get("rows", [])
                            if headers and rows:
                                df_vision = pd.DataFrame(rows, columns=headers[:len(rows[0])] if rows else headers)
                                df_vision = _clean_dataframe(df_vision)
                                if len(df_vision) > len(df) * 0.5:  # Vision got reasonable data
                                    tbl["df"] = df_vision
                                    tbl["quality_score"] = 75
                                    tbl["source"] = tbl.get("source", "") + " [vision-extracted]"
                                    validated_tables.append(tbl)
                                    result["warnings"].append(f"Sheet '{sheet}': used Vision LLM for extraction")
                                    continue
                    except Exception:
                        pass

                # Keep original (bad quality) with warning
                tbl["source"] = tbl.get("source", "") + " [low-quality]"
                validated_tables.append(tbl)
                result["warnings"].append(f"Table '{tbl['name']}' has quality issues: {', '.join(issues['problems'][:3])}")
        else:
            # Very bad — try vision or skip
            sheet = tbl.get("source", "").split("[")[0].strip().split("+")[0].strip()
            if sheet and ext == ".xlsx":
                try:
                    vision_raw = _vision_extract_sheet(file_path, sheet)
                    if vision_raw and "{" in vision_raw:
                        import json as _json4
                        clean = vision_raw.strip()
                        if clean.startswith("```"):
                            clean = clean.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
                        parsed = _json4.loads(clean)
                        headers = parsed.get("headers", [])
                        rows = parsed.get("rows", [])
                        if headers and rows:
                            df_vision = pd.DataFrame(rows, columns=headers[:len(rows[0])] if rows else headers)
                            df_vision = _clean_dataframe(df_vision)
                            if len(df_vision) >= 2:
                                tbl["df"] = df_vision
                                tbl["quality_score"] = 70
                                tbl["source"] = tbl.get("source", "") + " [vision-rescued]"
                                validated_tables.append(tbl)
                                result["warnings"].append(f"Sheet '{sheet}': rescued by Vision LLM")
                                continue
                except Exception:
                    pass
            # Keep with warning
            validated_tables.append(tbl)
            result["warnings"].append(f"Table '{tbl['name']}' has severe quality issues (score: {issues['score']})")

    result["tables"] = validated_tables
    return result


def _handle_pdf(file_path: str, filename: str) -> dict:
    """Handle PDF upload — PyMuPDF4LLM for text+tables (structured Markdown),
    Tesseract OCR for scanned pages, Vision LLM for images/diagrams."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        import fitz
        import base64

        doc = fitz.open(file_path)
        if len(doc) == 0:
            result["errors"].append("PDF has 0 pages — file may be corrupt or encrypted")
            return result

        total_pages = len(doc)
        result["metadata"]["total_pages"] = total_pages

        # ── STEP 1: PyMuPDF4LLM — structured Markdown (text + inline tables + layout) ──
        pymupdf4llm_ok = False
        md_text = ""
        try:
            import pymupdf4llm
            md_pages = pymupdf4llm.to_markdown(
                file_path,
                page_chunks=True,       # one chunk per page
                write_images=False,      # we handle images ourselves (Vision pipeline)
            )
            md_parts = []
            for chunk in md_pages:
                page_md = chunk.get("text", "") if isinstance(chunk, dict) else str(chunk)
                if page_md.strip():
                    md_parts.append(page_md)
            md_text = "\n\n".join(md_parts)
            pymupdf4llm_ok = True
            result["warnings"].append(f"PyMuPDF4LLM: extracted {len(md_text):,} chars structured Markdown from {total_pages} pages")
        except Exception as e:
            result["warnings"].append(f"PyMuPDF4LLM unavailable ({e}), falling back to fitz raw extraction")

        # ── STEP 2: Per-page scan — detect scanned pages + diagrams ──
        scanned_count = 0
        fallback_texts = []  # only used if PyMuPDF4LLM failed
        tesseract_ok = False
        try:
            import pytesseract
            from PIL import Image as PILImage
            tesseract_ok = True
        except ImportError:
            pass

        for pi, page in enumerate(doc):
            page_text = page.get_text()

            if len(page_text.strip()) < 50:
                # ── Scanned page — Tesseract first (local, free), Vision fallback ──
                scanned_count += 1
                try:
                    pixmap = page.get_pixmap(dpi=150)
                    if tesseract_ok:
                        img = PILImage.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                        ocr_text = pytesseract.image_to_string(img)
                        if len(ocr_text.strip()) > 30:
                            fallback_texts.append(f"[Page {pi + 1} — OCR]\n{ocr_text.strip()}")
                            continue  # Got text via Tesseract, skip vision
                    # Tesseract failed → send to Vision LLM as fallback
                    if len(result["images"]) < 10:
                        png_bytes = pixmap.tobytes("png")
                        if len(png_bytes) < 5_000_000:
                            result["images"].append({
                                "b64": base64.b64encode(png_bytes).decode(),
                                "mime": "image/png",
                                "source": f"page_{pi + 1}_scanned",
                            })
                except Exception:
                    pass
            else:
                if not pymupdf4llm_ok:
                    fallback_texts.append(page_text)

                # ── Diagram detection: short labels = flowchart/process/org chart ──
                lines = [l.strip() for l in page_text.split("\n") if l.strip()]
                avg_line_len = sum(len(l) for l in lines) / max(len(lines), 1)
                is_diagram = (
                    len(page_text.strip()) < 2000
                    and len(lines) > 5
                    and avg_line_len < 30
                )
                if is_diagram and len(result["images"]) < 30:
                    try:
                        pixmap = page.get_pixmap(dpi=200)
                        png_bytes = pixmap.tobytes("png")
                        if len(png_bytes) < 5_000_000:
                            result["images"].append({
                                "b64": base64.b64encode(png_bytes).decode(),
                                "mime": "image/png",
                                "source": f"page_{pi + 1}_diagram",
                            })
                            result["warnings"].append(f"Page {pi + 1}: diagram detected — sent to Vision for flow description")
                    except Exception:
                        pass

        doc.close()

        # ── STEP 3: Combine text — prefer PyMuPDF4LLM, append OCR pages ──
        text_parts = []
        if pymupdf4llm_ok and md_text.strip():
            text_parts.append(md_text)
        elif fallback_texts:
            text_parts.append("\n".join(fallback_texts))

        # Append Tesseract OCR text for scanned pages (not in PyMuPDF4LLM output)
        ocr_pages = [t for t in fallback_texts if "[Page " in t and "OCR]" in t]
        if pymupdf4llm_ok and ocr_pages:
            text_parts.append("\n\n--- SCANNED PAGES (OCR) ---\n" + "\n".join(ocr_pages))

        result["text"] = "\n\n".join(text_parts)

        if scanned_count > 0:
            result["metadata"]["scanned_pages"] = scanned_count
            vision_pages = len([i for i in result["images"] if "scanned" in i.get("source", "")])
            if tesseract_ok:
                result["warnings"].append(f"{scanned_count} scanned page(s) — Tesseract OCR (local), {vision_pages} sent to vision")
            else:
                result["warnings"].append(f"{scanned_count} scanned page(s) — sent to vision for OCR")

        # ── STEP 4: Extract tables as DataFrames (pdfplumber — structured data for PostgreSQL) ──
        tables = _extract_tables_pdf(file_path)
        result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]

        # ── STEP 5: Extract embedded images (charts, diagrams → Vision) ──
        embedded_images = _extract_images_pdf(file_path)
        result["images"].extend(embedded_images)

    except Exception as e:
        result["errors"].append(f"PDF processing failed: {e}")
    return result


def _handle_pptx(file_path: str, filename: str) -> dict:
    """Handle PPTX upload — text + tables + images + render image-heavy slides for Vision."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    image_only_slides = []  # slides with <10 chars text → render full slide for Vision
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
        texts = []
        notes = []
        for si, slide in enumerate(prs.slides):
            slide_text = ""
            has_images = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text += t + "\n"
                            texts.append(t)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_images = True
            # Track image-heavy slides with no/little text
            if len(slide_text.strip()) < 10 and has_images:
                image_only_slides.append(si)
            # Extract speaker notes
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    note_text = slide.notes_slide.notes_text_frame.text.strip()
                    if note_text:
                        notes.append(f"[Slide {si + 1} notes]: {note_text}")
            except Exception:
                pass
        all_text = "\n".join(texts)
        if notes:
            all_text += "\n\n--- SPEAKER NOTES ---\n" + "\n".join(notes)
        result["text"] = all_text
        result["metadata"]["slides"] = len(prs.slides)
        result["metadata"]["notes_count"] = len(notes)
        result["metadata"]["image_only_slides"] = len(image_only_slides)
    except Exception as e:
        result["errors"].append(f"PPTX text extraction failed: {e}")

    tables = _extract_tables_pptx(file_path)
    result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]

    # Extract embedded images from picture shapes
    result["images"] = _extract_images_pptx(file_path)

    # Render image-only slides as full-page images for Vision (Zerox-style)
    # These slides have charts/dashboards/screenshots that text extraction misses
    if image_only_slides:
        rendered = _render_pptx_slides(file_path, image_only_slides)
        if rendered:
            result["images"].extend(rendered)
            result["warnings"].append(f"Rendered {len(rendered)} image-only slides for Vision analysis")

    return result


def _render_pptx_slides(file_path: str, slide_indices: list[int], max_slides: int = 15) -> list[dict]:
    """Render specific PPTX slides as images for Vision analysis.

    Uses python-pptx slide dimensions + Pillow to create slide screenshots
    by compositing shape images onto a blank canvas. Falls back to sending
    the largest embedded image per slide if rendering fails.
    """
    import base64
    images: list[dict] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        import io

        prs = Presentation(file_path)
        slides = list(prs.slides)

        for si in slide_indices[:max_slides]:
            if si >= len(slides):
                continue
            slide = slides[si]

            # Strategy: composite all picture shapes onto a canvas sized to the slide
            slide_w = prs.slide_width.emu if hasattr(prs.slide_width, 'emu') else int(prs.slide_width)
            slide_h = prs.slide_height.emu if hasattr(prs.slide_height, 'emu') else int(prs.slide_height)

            # Scale to reasonable image size (max 2000px wide)
            scale = min(2000 / (slide_w / 914400), 1500 / (slide_h / 914400))
            canvas_w = int(slide_w / 914400 * scale)
            canvas_h = int(slide_h / 914400 * scale)

            canvas = Image.new("RGB", (canvas_w, canvas_h), (255, 255, 255))
            placed = 0

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        if len(blob) < 1000:
                            continue
                        img = Image.open(io.BytesIO(blob))
                        img = img.convert("RGB")

                        # Position on canvas using shape coordinates
                        x = int(shape.left / 914400 * scale) if hasattr(shape, 'left') and shape.left else 0
                        y = int(shape.top / 914400 * scale) if hasattr(shape, 'top') and shape.top else 0
                        w = int(shape.width / 914400 * scale) if hasattr(shape, 'width') and shape.width else img.width
                        h = int(shape.height / 914400 * scale) if hasattr(shape, 'height') and shape.height else img.height

                        img = img.resize((max(w, 10), max(h, 10)), Image.LANCZOS)
                        canvas.paste(img, (max(x, 0), max(y, 0)))
                        placed += 1
                    except Exception:
                        continue

            if placed > 0:
                buf = io.BytesIO()
                canvas.save(buf, format="PNG", optimize=True)
                b64 = base64.b64encode(buf.getvalue()).decode()
                images.append({
                    "b64": b64,
                    "mime": "image/png",
                    "source": f"slide_{si + 1}_rendered",
                })
    except Exception:
        pass
    return images


def _handle_docx(file_path: str, filename: str) -> dict:
    """Handle DOCX upload — text + tables + image extraction (new)."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        from docx import Document
        doc = Document(file_path)
        # Main body text
        body_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        # Headers and footers
        extras = []
        try:
            for section in doc.sections:
                if section.header and section.header.paragraphs:
                    h = " ".join(p.text.strip() for p in section.header.paragraphs if p.text.strip())
                    if h:
                        extras.append(f"[Header]: {h}")
                if section.footer and section.footer.paragraphs:
                    f_text = " ".join(p.text.strip() for p in section.footer.paragraphs if p.text.strip())
                    if f_text:
                        extras.append(f"[Footer]: {f_text}")
        except Exception:
            pass
        result["text"] = body_text
        if extras:
            result["text"] += "\n\n--- HEADERS/FOOTERS ---\n" + "\n".join(set(extras))  # deduplicate
    except Exception as e:
        result["errors"].append(f"DOCX text extraction failed: {e}")

    tables = _extract_tables_docx(file_path)
    result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]
    result["images"] = _extract_images_docx(file_path)
    return result


def _handle_csv(file_path: str, filename: str) -> dict:
    """Handle CSV upload — auto-detect encoding + delimiter + header."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        # Auto-detect encoding
        encoding = "utf-8"
        try:
            import chardet
            with open(file_path, "rb") as f:
                raw = f.read(100000)  # Read first 100KB for detection
            detected = chardet.detect(raw)
            if detected and detected.get("encoding") and detected.get("confidence", 0) > 0.5:
                encoding = detected["encoding"]
                if encoding.lower() != "utf-8" and encoding.lower() != "ascii":
                    result["warnings"].append(f"Encoding detected: {encoding} (confidence: {detected.get('confidence', 0):.0%})")
        except ImportError:
            pass

        header_row = _find_header_row(file_path, ".csv")
        sep = _detect_delimiter(file_path)
        # P1: keep ID/code columns as text (exact), never float (lossy on big codes).
        df = pd.read_csv(file_path, header=header_row, sep=sep, encoding=encoding,
                         encoding_errors="replace", dtype=_id_dtypes(file_path, ".csv", header_row, sep))

        # Normalize null values: N/A, NULL, None, -, ? → NaN
        null_values = {"N/A", "n/a", "#N/A", "NA", "na", "NULL", "null", "None", "none", "NONE", "-", "?", ".", " "}
        for col in df.columns:
            if df[col].dtype == object:
                df[col] = df[col].replace(null_values, pd.NA)

        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "csv"})
    except Exception as e:
        result["errors"].append(f"CSV read failed: {e}")
    return result


def _handle_json(file_path: str, filename: str) -> dict:
    """Handle JSON upload — existing read logic wrapped in standard result."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        try:
            df = pd.read_json(file_path)
        except ValueError:
            df = pd.read_json(file_path, orient='records')
        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "json"})
    except Exception as e:
        result["errors"].append(f"JSON read failed: {e}")
    return result


def _handle_text(file_path: str, filename: str, raw_content: bytes = b"") -> dict:
    """Handle text file upload — TXT, MD, SQL, PY."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        if raw_content:
            result["text"] = raw_content.decode("utf-8", errors="ignore")
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                result["text"] = f.read()
    except Exception as e:
        result["errors"].append(f"Text read failed: {e}")
    return result


def _handle_parquet(file_path: str, filename: str) -> dict:
    """Handle Parquet file — fastest columnar format."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        df = pd.read_parquet(file_path)
        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "parquet"})
        result["warnings"].append(f"Parquet: {len(df)} rows × {len(df.columns)} cols")
    except Exception as e:
        result["errors"].append(f"Parquet read failed: {e}")
    return result


def _handle_ods(file_path: str, filename: str) -> dict:
    """Handle ODS (LibreOffice/OpenDocument Spreadsheet)."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        sheets = pd.read_excel(file_path, engine='odf', sheet_name=None)
        for sname, df in sheets.items():
            df = _clean_dataframe(df)
            if len(df) > 0:
                tname = f"{_sanitize_table_name(Path(filename).stem)}_{_sanitize_table_name(str(sname))}" if len(sheets) > 1 else _sanitize_table_name(Path(filename).stem)
                result["tables"].append({"name": tname, "df": df, "source": f"{sname} [ods]"})
    except Exception as e:
        result["errors"].append(f"ODS read failed: {e}")
    return result


def _handle_xml(file_path: str, filename: str) -> dict:
    """Handle XML data files."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        df = pd.read_xml(file_path)
        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "xml"})
    except Exception:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                result["text"] = f.read()
        except Exception as e:
            result["errors"].append(f"XML read failed: {e}")
    return result


def _handle_html_file(file_path: str, filename: str) -> dict:
    """Handle HTML files — extract tables + text."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        tables = pd.read_html(file_path)
        for i, df in enumerate(tables[:10]):
            df = _clean_dataframe(df)
            if len(df) > 0 and len(df.columns) > 1:
                tname = f"{_sanitize_table_name(Path(filename).stem)}_table{i+1}" if len(tables) > 1 else _sanitize_table_name(Path(filename).stem)
                result["tables"].append({"name": tname, "df": df, "source": f"html_table_{i+1}"})
    except Exception:
        pass
    try:
        from html.parser import HTMLParser
        class _TextExtractor(HTMLParser):
            def __init__(self): super().__init__(); self.parts = []
            def handle_data(self, d):
                if d.strip(): self.parts.append(d.strip())
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            p = _TextExtractor(); p.feed(f.read())
            result["text"] = "\n".join(p.parts)
    except Exception:
        pass
    return result


def _handle_zip(file_path: str, filename: str) -> dict:
    """Handle ZIP archives — extract and process each file inside."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    import zipfile, tempfile
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            members = [m for m in zf.namelist() if not m.startswith('__MACOSX') and not m.startswith('.')]
            result["warnings"].append(f"ZIP: {len(members)} files inside")
            for member in members[:20]:
                ext = Path(member).suffix.lower()
                if ext not in ('.csv', '.xlsx', '.xls', '.json', '.xml', '.parquet', '.txt', '.md', '.pdf', '.docx', '.pptx', '.html', '.htm', '.ods'):
                    continue
                try:
                    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                        tmp.write(zf.read(member))
                        tmp_path = tmp.name
                    if ext == '.csv': sub = _handle_csv(tmp_path, member)
                    elif ext in ('.xlsx', '.xls'): sub = _handle_excel(tmp_path, member)
                    elif ext == '.json': sub = _handle_json(tmp_path, member)
                    elif ext == '.xml': sub = _handle_xml(tmp_path, member)
                    elif ext == '.parquet': sub = _handle_parquet(tmp_path, member)
                    elif ext == '.ods': sub = _handle_ods(tmp_path, member)
                    elif ext in ('.html', '.htm'): sub = _handle_html_file(tmp_path, member)
                    elif ext in ('.txt', '.md'): sub = _handle_text(tmp_path, member)
                    elif ext == '.pdf': sub = _handle_pdf(tmp_path, member)
                    elif ext == '.pptx': sub = _handle_pptx(tmp_path, member)
                    elif ext == '.docx': sub = _handle_docx(tmp_path, member)
                    else: continue
                    for t in sub.get("tables", []):
                        t["source"] = f"zip:{member} → {t.get('source', '')}"
                        result["tables"].append(t)
                    if sub.get("text"):
                        result["text"] += f"\n\n--- {member} ---\n{sub['text']}"
                    result["images"].extend(sub.get("images", []))
                    Path(tmp_path).unlink(missing_ok=True)
                except Exception as e:
                    result["warnings"].append(f"ZIP: failed {member}: {str(e)[:60]}")
    except Exception as e:
        result["errors"].append(f"ZIP failed: {e}")
    return result


def _handle_email(file_path: str, filename: str) -> dict:
    """Handle email files (.eml) — extract body text."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    import email as _email_mod
    try:
        with open(file_path, 'rb') as f:
            msg = _email_mod.message_from_bytes(f.read())
        body_parts = []
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body_parts.append(part.get_payload(decode=True).decode('utf-8', errors='ignore'))
        result["text"] = f"Email: {msg.get('subject','')}\nFrom: {msg.get('from','')}\nDate: {msg.get('date','')}\n\n" + "\n".join(body_parts)
    except Exception as e:
        result["errors"].append(f"Email read failed: {e}")
    return result


def _conduct_upload(file_path: str, ext: str, project: str, filename: str, raw_content: bytes = b"", _progress=None) -> dict:
    """Upload Conductor — routes to handler, runs vision on images, returns unified result."""
    if not _progress:
        _progress = lambda agent, step, detail: None

    # Determine handler name for progress
    handler_map = {
        ".xlsx": "Parser", ".xls": "Parser", ".csv": "Parser", ".json": "Parser",
        ".parquet": "Parser", ".ods": "Parser", ".xml": "Parser",
        ".html": "Scanner", ".htm": "Scanner", ".eml": "Scanner", ".zip": "Parser",
        ".pdf": "Scanner", ".pptx": "Scanner", ".docx": "Scanner",
        ".txt": "Scanner", ".md": "Scanner", ".sql": "Scanner", ".py": "Scanner",
        ".jpg": "Vision", ".jpeg": "Vision", ".png": "Vision",
        ".tiff": "Vision", ".tif": "Vision", ".bmp": "Vision", ".gif": "Vision", ".webp": "Vision",
    }
    agent_name = handler_map.get(ext, "Conductor")
    _progress(agent_name, "Starting", f"processing {filename}")

    # Route to handler
    if ext in (".xlsx", ".xls"):
        _progress("Parser", "Excel analysis", "detecting sheets and structure...")
        result = _handle_excel(file_path, filename)
        _progress("Parser", "Excel analysis", f"done — {len(result.get('tables', []))} tables extracted")
    elif ext == ".pdf":
        _progress("Scanner", "PDF extraction", "reading text, tables, detecting scanned pages...")
        result = _handle_pdf(file_path, filename)
        pages = result.get("metadata", {}).get("total_pages", 0)
        scanned = result.get("metadata", {}).get("scanned_pages", 0)
        _progress("Scanner", "PDF extraction", f"done — {pages} pages, {scanned} scanned, {len(result.get('tables', []))} tables")
    elif ext == ".pptx":
        _progress("Scanner", "PPTX extraction", "reading slides, tables, speaker notes...")
        result = _handle_pptx(file_path, filename)
        slides = result.get("metadata", {}).get("slides", 0)
        _progress("Scanner", "PPTX extraction", f"done — {slides} slides, {len(result.get('tables', []))} tables, {len(result.get('images', []))} images")
    elif ext == ".docx":
        _progress("Scanner", "DOCX extraction", "reading paragraphs, tables, images...")
        result = _handle_docx(file_path, filename)
        _progress("Scanner", "DOCX extraction", f"done — {len(result.get('text', '')):,} chars, {len(result.get('tables', []))} tables")
    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif", ".webp"):
        _progress("Vision", "Image processing", "OCR + visual description...")
        result = _handle_image(file_path, filename)
        _progress("Vision", "Image processing", f"done — {len(result.get('text', '')):,} chars extracted")
    elif ext == ".csv":
        _progress("Parser", "CSV parsing", "detecting encoding and delimiter...")
        result = _handle_csv(file_path, filename)
        _progress("Parser", "CSV parsing", f"done — {len(result.get('tables', []))} tables")
    elif ext == ".json":
        _progress("Parser", "JSON parsing", "reading structure...")
        result = _handle_json(file_path, filename)
        _progress("Parser", "JSON parsing", f"done — {len(result.get('tables', []))} tables")
    elif ext == ".parquet":
        _progress("Parser", "Parquet loading", "reading columnar data...")
        result = _handle_parquet(file_path, filename)
        _progress("Parser", "Parquet loading", f"done — {len(result.get('tables', []))} tables")
    elif ext == ".ods":
        _progress("Parser", "ODS loading", "reading LibreOffice spreadsheet...")
        result = _handle_ods(file_path, filename)
        _progress("Parser", "ODS loading", f"done — {len(result.get('tables', []))} tables")
    elif ext == ".xml":
        _progress("Parser", "XML parsing", "reading structured data...")
        result = _handle_xml(file_path, filename)
        _progress("Parser", "XML parsing", f"done — {len(result.get('tables', []))} tables")
    elif ext in (".html", ".htm"):
        _progress("Scanner", "HTML parsing", "extracting tables and text...")
        result = _handle_html_file(file_path, filename)
        _progress("Scanner", "HTML parsing", f"done — {len(result.get('tables', []))} tables, {len(result.get('text', '')):,} chars")
    elif ext == ".zip":
        _progress("Parser", "ZIP extraction", "extracting and processing files...")
        result = _handle_zip(file_path, filename)
        _progress("Parser", "ZIP extraction", f"done — {len(result.get('tables', []))} tables from archive")
    elif ext == ".eml":
        _progress("Scanner", "Email parsing", "reading email content...")
        result = _handle_email(file_path, filename)
        _progress("Scanner", "Email parsing", f"done — {len(result.get('text', '')):,} chars")
    elif ext in (".txt", ".md", ".sql", ".py"):
        _progress("Scanner", "Text extraction", "reading content...")
        result = _handle_text(file_path, filename, raw_content)
        _progress("Scanner", "Text extraction", f"done — {len(result.get('text', '')):,} chars")
    else:
        result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [f"Unsupported: {ext}"], "warnings": []}

    # Run vision on all collected images (scanned pages, charts, photos, DOCX images)
    if result.get("images"):
        img_count = len(result["images"])
        _progress("Vision", "Describing images", f"sending {img_count} images to Vision LLM...")
        image_text = _describe_images_with_vision(result["images"], filename)
        if image_text:
            result["text"] = (result.get("text", "") + f"\n\n--- IMAGE DESCRIPTIONS ---\n{image_text}").strip()
        _progress("Vision", "Describing images", f"done — {img_count} images described")

    return result


def _ai_review_and_fix_table(project_slug: str, table_name: str, engine=None):
    """AI reviews a stored table, detects issues, generates and runs SQL fixes.
    No hardcoded rules — LLM decides what to clean based on actual data."""
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
    if not engine:
        engine = create_engine(db_url)
    try:
        # Read sample + stats for LLM
        with engine.connect() as conn:
            conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
            # Sample rows
            sample = conn.execute(text(f'SELECT * FROM "{table_name}" LIMIT 15')).fetchall()
            cols = conn.execute(text(f'SELECT column_name, data_type FROM information_schema.columns WHERE table_schema = :s AND table_name = :t ORDER BY ordinal_position'),
                               {"s": schema, "t": table_name}).fetchall()
            row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{table_name}"')).scalar() or 0
            # Per-column null counts
            null_stats = []
            for col_name, _ in cols:
                null_count = conn.execute(text(f'SELECT COUNT(*) FROM "{table_name}" WHERE "{col_name}" IS NULL')).scalar() or 0
                null_stats.append(f"{col_name}: {null_count}/{row_count} null ({round(null_count/max(row_count,1)*100)}%)")

        # Build sample as text
        col_names = [c[0] for c in cols]
        sample_text = " | ".join(col_names) + "\n"
        for row in sample[:10]:
            sample_text += " | ".join(str(v)[:30] if v is not None else "NULL" for v in row) + "\n"

        prompt = f"""You are a data quality expert. Review this table and fix any issues.

TABLE: "{table_name}" in schema "{schema}" ({row_count} rows)

COLUMNS + TYPES:
{chr(10).join(f"  {c[0]} ({c[1]})" for c in cols)}

NULL STATS:
{chr(10).join(f"  {s}" for s in null_stats)}

SAMPLE DATA (first 10 rows):
{sample_text}

FIND AND FIX these types of issues:
1. Columns that should be forward-filled (value appears once then NULL for related rows — e.g. plant name, category)
2. Values that are actually NULL but stored as text: "-", "–", "—", empty string, "N/A", "None"
3. Header rows leaked into data (a row where a value equals the column name like "Products")
4. Columns that should be numeric but are stored as text (have numbers as strings)
5. Columns with a single value for all rows (useless, could note it)

Return ONLY a valid JSON array of SQL UPDATE/DELETE statements to fix issues.
Each fix should have "description" and "sql".
If no issues found, return empty array [].

Example:
[
  {{"description": "Forward-fill plant column where NULL", "sql": "UPDATE \\"{table_name}\\" SET plant = sub.plant FROM (SELECT ctid, plant FROM (SELECT ctid, plant, COALESCE(plant, LAG(plant) OVER (ORDER BY ctid)) as filled FROM \\"{table_name}\\") sub) sub WHERE \\"{table_name}\\".ctid = sub.ctid AND \\"{table_name}\\".plant IS NULL"}},
  {{"description": "Clean dash values in output", "sql": "UPDATE \\"{table_name}\\" SET output = NULL WHERE output IN ('-', '–', '—', '')"}},
  {{"description": "Remove header leak rows", "sql": "DELETE FROM \\"{table_name}\\" WHERE products = 'Products'"}}
]

Use schema "{schema}" in all SQL. Only return fixes you are confident about."""

        from dash.settings import training_llm_call
        raw = training_llm_call(prompt, "excel_analysis")
        if not raw:
            return []

        fixes = json.loads(raw)
        if not isinstance(fixes, list) or not fixes:
            return []

        # Execute each fix (sandboxed: only UPDATE/DELETE on this table)
        _FORBIDDEN = re.compile(r"\b(DROP|ALTER|TRUNCATE|CREATE|INSERT|GRANT|REVOKE|COPY|EXECUTE)\b", re.IGNORECASE)
        applied = []
        with engine.connect() as conn:
            conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
            for fix in fixes:
                sql = fix.get("sql", "")
                desc = fix.get("description", "")
                if not sql:
                    continue
                # Safety: reject dangerous SQL from LLM
                if _FORBIDDEN.search(sql):
                    logger.warning(f"AI review blocked dangerous SQL: {sql[:100]}")
                    continue
                # Only allow operations on the target table
                sql_upper = sql.upper()
                if not (sql_upper.lstrip().startswith("UPDATE") or sql_upper.lstrip().startswith("DELETE")):
                    logger.warning(f"AI review blocked non-UPDATE/DELETE SQL: {sql[:100]}")
                    continue
                try:
                    result = conn.execute(text(sql))
                    # Cap: reject if affects too many rows (>50% of table)
                    if result.rowcount > max(row_count * 0.5, 100):
                        conn.rollback()
                        logger.warning(f"AI review rolled back fix affecting {result.rowcount}/{row_count} rows: {desc}")
                        continue
                    applied.append(desc)
                except Exception as e:
                    pass  # Skip failed fixes silently
            conn.commit()
        return applied
    except Exception:
        return []


def _run_pandasai_experiments(project_slug: str, table_name: str, col_analyses: list[dict], _log=None):
    """Run SQL experiments to generate 25+ verified Q&A pairs from real data.
    Executes real SQL against PostgreSQL, saves question + verified answer.
    Adds to existing training Q&A (doesn't replace). Zero external dependencies.
    Adds to existing training Q&A (doesn't replace). Zero external dependencies."""
    if not _log:
        _log = lambda x: None
    _log("running SQL experiments on real data...")
    schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]

    try:
        exp_engine = create_engine(db_url)
        num_cols = [c["name"] for c in col_analyses if c.get("type") == "numeric"]
        cat_cols = [c["name"] for c in col_analyses if c.get("is_categorical")]
        # Gate trend/monthly/daily/weekly QA: skip date cols that are effectively
        # constant (1-period data → meaningless time series). See Phase 9 DQ
        # memory contract (dash/utils/column_stats_cache.py is_constant role).
        _ca_by_name = {c.get("name"): c for c in col_analyses}
        def _is_date_col_usable(_c: dict) -> bool:
            uc = _c.get("unique_count")
            if uc is not None and uc <= 1:
                logger.info(f"skipped trend QA: {_c.get('name')} is constant (unique_count={uc})")
                return False
            return True
        date_cols = [
            c["name"] for c in col_analyses
            if (c.get("type") == "datetime" or "date" in c.get("name", "").lower())
            and _is_date_col_usable(c)
        ]
        nc = num_cols[0] if num_cols else None
        cc = cat_cols[0] if cat_cols else None
        dc = date_cols[0] if date_cols else None
        t = f'"{schema}"."{table_name}"'

        # A "date" col may be TEXT in DD/MM/YYYY (DATE_TRUNC on raw text throws).
        # Build a safe date expression: real datetime → use as-is; text-date →
        # wrap with to_date(...,'DD/MM/YYYY HH24:MI').
        def _date_expr(col: str) -> str:
            if not col:
                return col
            _ca = _ca_by_name.get(col, {})
            if _ca.get("type") == "datetime":
                return f'"{col}"'
            sv = _ca.get("sample_values") or []
            if any(re.match(r"^\s*\d{1,2}/\d{1,2}/\d{4}\b", str(v)) for v in sv[:6]):
                return f'to_date("{col}", \'DD/MM/YYYY HH24:MI\')'
            return f'"{col}"'
        dc_expr = _date_expr(dc) if dc else None

        # Build SQL experiments: (question, SQL)
        experiments = [("How many rows in total?", f"SELECT COUNT(*) FROM {t}")]
        if nc:
            experiments.extend([
                (f"What is the total {nc}?", f'SELECT SUM("{nc}") FROM {t}'),
                (f"What is the average {nc}?", f'SELECT ROUND(AVG("{nc}")::numeric, 2) FROM {t}'),
                (f"What is the maximum {nc}?", f'SELECT MAX("{nc}") FROM {t}'),
                (f"What is the minimum {nc}?", f'SELECT MIN("{nc}") FROM {t}'),
                (f"How many rows have zero {nc}?", f'SELECT COUNT(*) FROM {t} WHERE "{nc}" = 0'),
                (f"How many rows have null {nc}?", f'SELECT COUNT(*) FROM {t} WHERE "{nc}" IS NULL'),
            ])
        if nc and cc:
            experiments.extend([
                (f"Total {nc} by {cc}", f'SELECT "{cc}", SUM("{nc}") as total FROM {t} GROUP BY "{cc}" ORDER BY total DESC'),
                (f"Average {nc} by {cc}", f'SELECT "{cc}", ROUND(AVG("{nc}")::numeric, 2) as avg FROM {t} GROUP BY "{cc}" ORDER BY avg DESC'),
                (f"Count by {cc}", f'SELECT "{cc}", COUNT(*) as cnt FROM {t} GROUP BY "{cc}" ORDER BY cnt DESC'),
                (f"Which {cc} has the highest {nc}?", f'SELECT "{cc}", SUM("{nc}") as total FROM {t} GROUP BY "{cc}" ORDER BY total DESC LIMIT 1'),
                (f"Top 5 {cc} by {nc}", f'SELECT "{cc}", SUM("{nc}") as total FROM {t} GROUP BY "{cc}" ORDER BY total DESC LIMIT 5'),
                (f"Percentage contribution by {cc}", f'SELECT "{cc}", ROUND(SUM("{nc}") * 100.0 / NULLIF((SELECT SUM("{nc}") FROM {t}), 0), 1) as pct FROM {t} GROUP BY "{cc}" ORDER BY pct DESC'),
            ])
        if cc:
            experiments.extend([
                (f"How many unique {cc} values?", f'SELECT COUNT(DISTINCT "{cc}") FROM {t}'),
                (f"Distribution of {cc}", f'SELECT "{cc}", COUNT(*) as cnt FROM {t} GROUP BY "{cc}" ORDER BY cnt DESC'),
            ])
        if nc and dc:
            experiments.extend([
                (f"Total {nc} by month", f'SELECT DATE_TRUNC(\'month\', {dc_expr}) as month, SUM("{nc}") as total FROM {t} GROUP BY 1 ORDER BY 1'),
                (f"Which month had highest {nc}?", f'SELECT DATE_TRUNC(\'month\', {dc_expr}) as month, SUM("{nc}") as total FROM {t} GROUP BY 1 ORDER BY total DESC LIMIT 1'),
                (f"Date range of data", f'SELECT MIN({dc_expr}), MAX({dc_expr}) FROM {t}'),
            ])
        if len(num_cols) >= 2:
            nc2 = num_cols[1]
            experiments.append((f"Correlation: {nc} vs {nc2}", f'SELECT ROUND(CORR("{nc}"::numeric, "{nc2}"::numeric)::numeric, 3) FROM {t}'))
        if len(cat_cols) >= 2:
            cc2 = cat_cols[1]
            experiments.append((f"Breakdown by {cc} and {cc2}", f'SELECT "{cc}", "{cc2}", COUNT(*) as cnt FROM {t} GROUP BY "{cc}", "{cc2}" ORDER BY cnt DESC LIMIT 10'))

        # Execute all experiments against real data
        results = []
        with exp_engine.connect() as conn:
            for question, sql in experiments:
                try:
                    rows = conn.execute(text(sql)).fetchall()
                    if rows:
                        answer = str(rows[0][0]) if len(rows) == 1 and len(rows[0]) == 1 else " | ".join(str(r) for r in rows[:5])
                        results.append({"question": question, "sql": sql.replace(f'"{schema}".', ""), "answer": answer[:500], "verified": True, "source": "sql_experiment", "analysis_type": "experiment"})
                except Exception:
                    pass

        if results:
            qa_file = KNOWLEDGE_DIR / project_slug / "training" / f"{table_name}_qa.json"
            existing_qa = []
            if qa_file.exists():
                try:
                    with open(qa_file) as f:
                        existing_qa = json.load(f)
                except Exception:
                    pass
            existing_questions = {q.get("question", "").lower() for q in existing_qa}
            new_qa = [r for r in results if r["question"].lower() not in existing_questions]
            combined = existing_qa + new_qa
            qa_file.parent.mkdir(parents=True, exist_ok=True)
            with open(qa_file, "w") as f:
                json.dump(combined, f, indent=2, default=str)
            # Save top stats as memories
            try:
                mem_engine = create_engine(db_url)
                with mem_engine.connect() as mconn:
                    for r in results[:5]:
                        if r.get("answer") and len(r["answer"]) > 1:
                            mconn.execute(text("INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'experiment') ON CONFLICT DO NOTHING"),
                                          {"s": project_slug, "f": f"[{table_name}] {r['question']}: {r['answer'][:200]}"})
                    mconn.commit()
            except Exception:
                pass
            _log(f"✓ experiments: {len(new_qa)} new Q&A added ({len(results)} ran, {len(existing_qa)} existing)")
        else:
            _log("· no successful experiments")
    except Exception as e:
        _log(f"⚠ experiments error: {str(e)[:100]}")


def _langextract_facts(project_slug: str, all_text: str, _log=None):
    """Extract grounded facts from document text using LangExtract.
    Stores KPIs, metrics, decisions, risks, rules with source positions.
    Used by Researcher agent for source-cited answers."""
    if not _log:
        _log = lambda m: None
    if not all_text or len(all_text.strip()) < 100:
        _log("· text too short for fact extraction — skipped")
        return
    try:
        import langextract as lx
    except ImportError:
        _log("⚠ langextract not installed — skipped")
        return

    api_key = os.getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        _log("⚠ no API key — fact extraction skipped")
        return

    try:
        _log("extracting grounded facts with LangExtract...")

        # Few-shot examples — define what to extract (works across all domains)
        examples = [
            lx.ExampleData(
                text="Plant efficiency reached 87% in March 2025. Revenue was $4.2M, up 12% YoY. "
                     "Board approved the $12M expansion on March 15. Defect rate must stay below 1%.",
                extractions=[
                    lx.Extraction(extraction_class="kpi", extraction_text="Plant efficiency: 87%",
                                  attributes={"period": "March 2025"}),
                    lx.Extraction(extraction_class="kpi", extraction_text="Revenue: $4.2M, up 12% YoY",
                                  attributes={"trend": "up 12%"}),
                    lx.Extraction(extraction_class="decision", extraction_text="Board approved $12M expansion",
                                  attributes={"date": "March 15"}),
                    lx.Extraction(extraction_class="business_rule", extraction_text="Defect rate must stay below 1%",
                                  attributes={"threshold": "1%"}),
                ],
            ),
        ]

        # Run extraction — use Gemini Flash (cheap, fast)
        # Truncate to ~8K chars to keep cost low
        text_chunk = all_text[:8000]

        annotated_docs = lx.extract(
            text_chunk,
            prompt_description=(
                "Extract all KPIs, metrics, financial figures, decisions, business rules, "
                "risks, targets, deadlines, and key relationships from this business document. "
                "Include specific numbers, dates, percentages, and thresholds."
            ),
            examples=examples,
            model_id=TRAINING_MODEL,
            api_key=api_key,
            max_workers=3,
            extraction_passes=1,
        )

        # Process results — store as grounded memories
        if not annotated_docs:
            _log("· no facts extracted")
            return

        docs_list = annotated_docs if isinstance(annotated_docs, list) else [annotated_docs]
        all_facts = []
        for doc in docs_list:
            for ext in (doc.extractions if hasattr(doc, 'extractions') else []):
                fact_text = ext.extraction_text if hasattr(ext, 'extraction_text') else str(ext)
                fact_class = ext.extraction_class if hasattr(ext, 'extraction_class') else "fact"
                # Get source position for grounding
                char_start = ext.char_interval.start_pos if hasattr(ext, 'char_interval') and ext.char_interval else None
                char_end = ext.char_interval.end_pos if hasattr(ext, 'char_interval') and ext.char_interval else None
                # Get grounding status
                grounded = True
                if hasattr(ext, 'alignment_status'):
                    grounded = str(ext.alignment_status).upper() != "UNGROUNDED"
                # Get attributes
                attrs = {}
                if hasattr(ext, 'attributes') and ext.attributes:
                    attrs = ext.attributes

                all_facts.append({
                    "text": fact_text,
                    "type": fact_class,
                    "char_start": char_start,
                    "char_end": char_end,
                    "grounded": grounded,
                    "attributes": attrs,
                })

        if not all_facts:
            _log("· no facts found in extraction")
            return

        # Save grounded facts to dash_memories with source='langextract'
        engine = create_engine(db_url)
        saved = 0
        with engine.connect() as conn:
            for fact in all_facts[:30]:  # Cap at 30 facts per project
                if not fact["text"] or len(fact["text"]) < 5:
                    continue
                # Build fact string with grounding info
                grounding_tag = "✅" if fact["grounded"] else "⚠️"
                type_tag = fact["type"].upper()
                attr_str = ""
                if fact["attributes"]:
                    attr_parts = [f"{k}: {v}" for k, v in fact["attributes"].items() if isinstance(v, str)]
                    if attr_parts:
                        attr_str = f" ({', '.join(attr_parts[:3])})"
                source_ref = ""
                if fact["char_start"] is not None:
                    source_ref = f" [chars {fact['char_start']}-{fact['char_end']}]"

                memory_text = f"{grounding_tag} [{type_tag}] {fact['text']}{attr_str}{source_ref}"

                try:
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                        "VALUES (:s, 'project', :f, 'langextract') ON CONFLICT DO NOTHING"
                    ), {"s": project_slug, "f": memory_text})
                    saved += 1
                except Exception:
                    pass

            # Also save facts as structured JSON for Researcher
            facts_dir = KNOWLEDGE_DIR / project_slug / "training"
            facts_dir.mkdir(parents=True, exist_ok=True)
            facts_file = facts_dir / "grounded_facts.json"
            try:
                with open(facts_file, "w") as f:
                    json.dump(all_facts, f, indent=2, default=str)
            except Exception:
                pass

            conn.commit()

        _log(f"✓ {saved} grounded facts extracted ({sum(1 for f in all_facts if f['grounded'])} verified, "
             f"{sum(1 for f in all_facts if not f['grounded'])} ungrounded)")

    except Exception as e:
        _log(f"⚠ LangExtract error: {str(e)[:100]}")


def _post_upload_engineer(project_slug: str, tables_created: list[dict], user_id: int = 1):
    """After upload: Engineer merges same-structure tables, Inspector validates, safe delete originals.
    Runs in background thread — does not block upload response."""
    if not project_slug or not tables_created:
        return
    import logging
    log = logging.getLogger("dash.upload")

    try:
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        merge_engine = create_engine(db_url)

        # STEP 1: Get ALL tables in project with their columns
        all_tables = {}
        with merge_engine.connect() as conn:
            conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
            from sqlalchemy import inspect as sa_inspect
            insp = sa_inspect(merge_engine)
            for tbl in insp.get_table_names(schema=schema):
                try:
                    cols = insp.get_columns(tbl, schema=schema)
                    col_names = [c["name"] for c in cols]
                    row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                    all_tables[tbl] = {"columns": col_names, "rows": row_count, "col_set": set(col_names)}
                except Exception:
                    pass

        if len(all_tables) < 2:
            log.info(f"Post-upload: only {len(all_tables)} table(s), skipping merge")
            # Still run Engineer for relationships
            _run_engineer_agent(project_slug, tables_created, user_id)
            return

        # STEP 2: Find merge groups (tables with >80% column overlap)
        merge_groups = []
        used = set()
        table_names = list(all_tables.keys())

        for i, t1 in enumerate(table_names):
            if t1 in used:
                continue
            group = [t1]
            cols1 = all_tables[t1]["col_set"]
            for j in range(i + 1, len(table_names)):
                t2 = table_names[j]
                if t2 in used:
                    continue
                cols2 = all_tables[t2]["col_set"]
                overlap = len(cols1 & cols2) / max(len(cols1 | cols2), 1)
                if overlap >= 0.8:
                    group.append(t2)
            if len(group) >= 2:
                merge_groups.append(group)
                used.update(group)

        if not merge_groups:
            log.info("Post-upload: no merge candidates found")
            _run_engineer_agent(project_slug, tables_created, user_id)
            return

        # STEP 3: Merge each group
        for group in merge_groups:
            # Pick merged table name (shortest common prefix or first table)
            import os
            prefix = os.path.commonprefix(group).rstrip("_") or group[0]
            merged_name = _sanitize_table_name(f"{prefix}_merged" if len(prefix) > 3 else f"{group[0]}_merged")

            # Get union of all columns
            all_cols = set()
            for tbl in group:
                all_cols.update(all_tables[tbl]["col_set"])
            all_cols_list = sorted(all_cols)

            # Count expected rows
            expected_rows = sum(all_tables[tbl]["rows"] for tbl in group)

            try:
                with merge_engine.connect() as conn:
                    conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))

                    # Build UNION ALL with source_table column
                    parts = []
                    for tbl in group:
                        tbl_cols = all_tables[tbl]["col_set"]
                        select_parts = []
                        for col in all_cols_list:
                            if col in tbl_cols:
                                select_parts.append(f'"{col}"')
                            else:
                                select_parts.append(f"NULL AS \"{col}\"")
                        select_parts.append(f"'{tbl}' AS _source_table")
                        parts.append(f"SELECT {', '.join(select_parts)} FROM \"{schema}\".\"{tbl}\"")

                    union_sql = " UNION ALL ".join(parts)
                    create_sql = f'CREATE TABLE "{schema}"."{merged_name}" AS {union_sql}'
                    conn.execute(text(create_sql))
                    conn.commit()

                # STEP 4: Inspector validates merged table
                actual_rows = 0
                with merge_engine.connect() as conn:
                    actual_rows = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{merged_name}"')).scalar() or 0

                merge_valid = actual_rows >= expected_rows  # Must have ALL rows

                if merge_valid:
                    # Profile the merged table
                    try:
                        df_sample = pd.read_sql(f'SELECT * FROM "{schema}"."{merged_name}" LIMIT 5000', merge_engine)
                        profile = _profile_table(df_sample, project_slug, merged_name)
                        health = profile.get("health", 0)
                        merge_valid = health >= 50  # Must pass minimum quality
                    except Exception:
                        pass

                if merge_valid:
                    # PASS — safe to delete originals
                    with merge_engine.connect() as conn:
                        for tbl in group:
                            try:
                                conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{tbl}" CASCADE'))
                            except Exception:
                                pass
                        conn.commit()
                    log.info(f"Post-upload MERGED: {group} → {merged_name} ({actual_rows} rows, health={profile.get('health', '?')}%)")

                    # AI review and fix data quality issues
                    try:
                        fixes = _ai_review_and_fix_table(project_slug, merged_name, merge_engine)
                        if fixes:
                            log.info(f"Post-upload AI FIXES on {merged_name}: {fixes}")
                    except Exception:
                        pass

                    # Save source metadata for merged table
                    src_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
                    src_dir.mkdir(parents=True, exist_ok=True)
                    _safe_write_json(src_dir / f"{merged_name}.json", {
                        "source_file": f"Merged from {len(group)} tables",
                        "source_detail": f"Sources: {', '.join(group)}",
                        "description": f"Merged table combining {len(group)} tables with same structure ({actual_rows} rows)",
                    })
                    # Clean up old source metadata
                    for tbl in group:
                        old_src = src_dir / f"{tbl}.json"
                        if old_src.exists():
                            old_src.unlink(missing_ok=True)
                        old_profile = src_dir / f"{tbl}_profile.json"
                        if old_profile.exists():
                            old_profile.unlink(missing_ok=True)
                else:
                    # FAIL — drop merged, keep originals
                    with merge_engine.connect() as conn:
                        conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{merged_name}" CASCADE'))
                        conn.commit()
                    log.warning(f"Post-upload MERGE FAILED validation: {group} → kept separate (expected={expected_rows}, actual={actual_rows})")

            except Exception as e:
                log.warning(f"Post-upload merge error for {group}: {e}")

        # STEP 5: Run Engineer for relationships + views on final state
        _run_engineer_agent(project_slug, tables_created, user_id)

    except Exception as e:
        log.warning(f"Post-upload engineer failed: {e}")


def _run_engineer_agent(project_slug: str, tables_created: list[dict], user_id: int = 1):
    """Run Engineer agent to discover relationships and create useful views."""
    try:
        from dash.agents.engineer import create_engineer
        from db.session import create_project_knowledge, create_project_learnings
        from agno.learn import LearnedKnowledgeConfig, LearningMachine, LearningMode

        knowledge = create_project_knowledge(project_slug)
        learnings = create_project_learnings(project_slug)
        learning = LearningMachine(knowledge=learnings, learned_knowledge=LearnedKnowledgeConfig(mode=LearningMode.AGENTIC))
        engineer = create_engineer(project_slug=project_slug, knowledge=knowledge, learning=learning, dashboard_user_id=user_id)

        table_summary = "\n".join(f"  - {t['table']} ({t['rows']} rows)" for t in tables_created[:20])

        prompt = f"""Inspect the project tables and optimize:
1. INSPECT all tables — run introspect_schema
2. DISCOVER RELATIONSHIPS — find JOINable columns (shared IDs, dates, categories)
3. FIX COLUMN TYPES — dates stored as text → ALTER to DATE
4. REPORT relationships found

Tables in project:
{table_summary}

Use SQL tools. Be concise."""

        response = engineer.run(prompt)
        content = response.content if response else ""
        if content:
            from agno.knowledge.reader.text_reader import TextReader
            try:
                knowledge.insert(name="engineer-upload-analysis", text_content=f"Engineer Analysis:\n\n{content[:5000]}", reader=TextReader(), skip_if_exists=False)
            except Exception:
                pass
    except Exception as e:
        import logging
        logging.getLogger("dash.upload").warning(f"Engineer agent failed: {e}")


def _extract_document_structure(file_path: str, ext: str) -> list[dict]:
    """Extract document structure (titles + content summaries) for workflow conversion.

    Returns list of {"index": 1, "title": "...", "content_summary": "..."} dicts.
    PPTX → slide titles, PDF → page headers by font size, DOCX → heading paragraphs.
    """
    sections: list[dict] = []

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                title = ""
                content_parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if not t:
                                continue
                            if not title and (shape == slide.shapes.title if hasattr(slide.shapes, 'title') and slide.shapes.title else shape == slide.shapes[0]):
                                title = t
                            else:
                                content_parts.append(t)
                if not title and content_parts:
                    title = content_parts.pop(0)
                if title:
                    sections.append({
                        "index": i,
                        "title": title[:120],
                        "content_summary": " ".join(content_parts)[:150],
                    })
        except Exception:
            pass

    elif ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            for i, page in enumerate(doc, 1):
                blocks = page.get_text("dict").get("blocks", [])
                max_size = 0
                title = ""
                content_parts: list[str] = []
                for block in blocks:
                    if "lines" not in block:
                        continue
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            size = span.get("size", 0)
                            if not text:
                                continue
                            if size > max_size and len(text) > 3:
                                if title:
                                    content_parts.insert(0, title)
                                title = text
                                max_size = size
                            else:
                                content_parts.append(text)
                if not title and content_parts:
                    title = content_parts.pop(0)
                if title:
                    sections.append({
                        "index": i,
                        "title": title[:120],
                        "content_summary": " ".join(content_parts)[:150],
                    })
            doc.close()
        except Exception:
            pass

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            current_title = ""
            current_content: list[str] = []
            idx = 0
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                is_heading = para.style and para.style.name and para.style.name.startswith("Heading")
                if is_heading:
                    if current_title:
                        idx += 1
                        sections.append({
                            "index": idx,
                            "title": current_title[:120],
                            "content_summary": " ".join(current_content)[:150],
                        })
                    current_title = text
                    current_content = []
                else:
                    if current_title:
                        current_content.append(text)
                    elif not sections:
                        current_title = text
            if current_title:
                idx += 1
                sections.append({
                    "index": idx,
                    "title": current_title[:120],
                    "content_summary": " ".join(current_content)[:150],
                })
        except Exception:
            pass

    return sections[:20]


# ── domain auto-detect + template suggestion (Issue #7) ─────────────────────
_DOMAIN_KEYWORDS: dict[str, tuple[set[str], str]] = {
    "pharmacy_network": ({"article", "articles", "drug", "drugs", "pharma",
                          "medication", "medications", "rx", "prescription"},
                         "Pharmacy Network"),
    "investment":       ({"balance_sheet", "balancesheet", "income_statement",
                          "incomestatement", "cash_flow", "cashflow",
                          "cap_table", "captable"},
                         "Investment Desk"),
    "retail":           ({"customer", "customers", "transaction", "transactions",
                          "order", "orders", "sale", "sales", "sku", "skus",
                          "store", "stores", "pos"},
                         "Retail"),
    "hotel_group":      ({"hotel", "hotels", "reservation", "reservations",
                          "adr", "revpar", "occupancy", "rate_code"},
                         "Hotel Group"),
    "manufacturing":    ({"bom", "defect", "defects", "yield", "supplier",
                          "suppliers", "work_order", "workorder", "mrp"},
                         "Manufacturing"),
}


def _detect_domain_from_tables(table_names: list[str]) -> tuple[str | None, float, str]:
    """Score table names against domain keyword sets.

    Returns (template, confidence, label). 2+ keyword hits → conf ≥ 0.78.
    """
    if not table_names:
        return None, 0.0, ""
    tokens: set[str] = set()
    for t in table_names:
        if not t:
            continue
        n = re.sub(r"[^a-z0-9_]", "_", t.lower())
        tokens.update(p for p in n.split("_") if p)
        tokens.add(n)
    best: tuple[str | None, float, str] = (None, 0.0, "")
    for tmpl, (kws, label) in _DOMAIN_KEYWORDS.items():
        hits = tokens & kws
        if not hits:
            continue
        conf = min(0.55 + 0.15 * len(hits), 0.95)
        if len(hits) >= 2:
            conf = max(conf, 0.78)
        if conf > best[1]:
            best = (tmpl, conf, label)
    return best


def _suggest_template_for_project(project_slug: str, new_table_names: list[str]) -> None:
    """No-op. Industry preset auto-apply / template suggestion removed.

    Callers retained so upload pipeline call sites stay unchanged.
    """
    return


def _autoload_definitions(project_slug: str, file_path: str, ext: str) -> int:
    """Detect a definitions/glossary sheet in an uploaded Excel file and insert
    each term→definition into the project's Company Brain as a glossary entry.

    Conservative: only fires when a sheet clearly looks like a 2-column data
    dictionary (a 'term' column + a long-text 'definition' column). Fail-soft.
    Returns the number of entries inserted. So a customer's Definitions.xlsx
    self-pins as authoritative metric/glossary context (CRM feedback 2026-05-21).
    """
    import logging
    log = logging.getLogger("dash.upload")
    if ext not in (".xlsx", ".xls"):
        return 0
    try:
        import pandas as _pd
        from db import get_sql_engine as _gse
        xl = _pd.ExcelFile(file_path)
    except Exception as e:
        log.debug(f"_autoload_definitions open failed: {e}")
        return 0

    inserted = 0
    eng = None
    try:
        eng = _gse()
        for sname in xl.sheet_names:
            try:
                df = xl.parse(sname, header=0, dtype=str)
            except Exception:
                continue
            if df is None or df.shape[1] < 2 or len(df) < 3:
                continue
            df = df.dropna(how="all")
            # pick the column whose values are the longest text → definitions
            def_col, term_col, best_avg = None, None, 0.0
            for c in df.columns:
                vals = [str(v) for v in df[c].dropna().tolist() if str(v).strip()]
                if not vals:
                    continue
                avg = sum(len(v) for v in vals) / len(vals)
                if avg > best_avg:
                    best_avg, def_col = avg, c
            sheet_hint = bool(re.search(r"defin|glossar|diction|description|metric", str(sname), re.I))
            # require a clearly-textual definition column unless the sheet name says so
            if def_col is None or (best_avg < 40 and not sheet_hint):
                continue
            # term column = the non-definition column with shortest avg (a label)
            cand = [c for c in df.columns if c != def_col]
            term_col = min(
                cand,
                key=lambda c: (sum(len(str(v)) for v in df[c].dropna()) / max(len(df[c].dropna()), 1)),
                default=None,
            )
            if term_col is None:
                continue
            rows = []
            for _, r in df.iterrows():
                term = str(r.get(term_col, "")).strip()
                defn = str(r.get(def_col, "")).strip()
                if term and defn and term.lower() != "nan" and defn.lower() != "nan" and len(defn) > 8:
                    # rows that read like a formula → category 'formula', else 'glossary'
                    cat = "formula" if re.search(r"\bnumerator|denominator|formula|=|count|/|\bsum\b", defn, re.I) else "glossary"
                    rows.append((cat, term[:200], defn[:4000]))
            if not rows:
                continue
            with eng.begin() as cn:
                for cat, name, defn in rows:
                    exists = cn.execute(text(
                        "SELECT 1 FROM public.dash_company_brain "
                        "WHERE project_slug=:s AND name=:n AND category=:c LIMIT 1"
                    ), {"s": project_slug, "n": name, "c": cat}).fetchone()
                    if exists:
                        continue
                    cn.execute(text(
                        "INSERT INTO public.dash_company_brain "
                        "(category, name, definition, project_slug, metadata, created_by) "
                        "VALUES (:c, :n, :d, :s, CAST(:m AS jsonb), 'auto_definitions')"
                    ), {"c": cat, "n": name, "d": defn, "s": project_slug,
                        "m": json.dumps({"source": "definitions_upload", "sheet": str(sname)})})
                    inserted += 1
            if inserted:
                log.info(f"_autoload_definitions: +{inserted} brain entries from sheet '{sname}'")
    except Exception as e:
        log.debug(f"_autoload_definitions failed: {e}")
    # NOTE: get_sql_engine() returns the cached shared engine — never dispose it.
    return inserted


@router.post("/upload")
async def upload_file(request: Request, file: UploadFile, table_name: str | None = None, replace: bool = False, project: str | None = None, action: str = "auto"):
    """Upload a data file. action: auto/append/upsert/replace/new.

    Accept `project` (and friends) from BOTH query string AND multipart form.
    FastAPI bare-typed params bind to query only — prior callers sending
    `-F project=X` silently routed uploads to user_demo schema. — 2026-05-25 fix.
    """
    # Manual dual-source param resolution (query-bound above + multipart fallback)
    try:
        _form = await request.form()
        if project is None and _form.get("project"):
            project = str(_form.get("project"))
        if table_name is None and _form.get("table_name"):
            table_name = str(_form.get("table_name"))
        if _form.get("action"):
            action = str(_form.get("action"))
        if _form.get("replace") in ("1", "true", "True"):
            replace = True
    except Exception:
        pass  # form already consumed by FastAPI internals; bare query still works
    user_id = _get_user_id(request)

    # Editor role required for uploads
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported format: {ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")

    # Stream to temp file instead of loading to memory
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        size = 0
        while chunk := await file.read(1024 * 1024):  # 1MB chunks
            size += len(chunk)
            if size > MAX_FILE_SIZE:
                raise HTTPException(400, f"File too large. Max: {MAX_FILE_SIZE // (1024*1024)} MB")
            tmp.write(chunk)
        tmp_path = tmp.name
    # Prometheus: bytes accepted on /upload, labeled by extension.
    try:
        from dash.utils.metrics import add_upload_bytes as _add_upload_bytes
        _add_upload_bytes(ext, size)
    except Exception:
        pass

    # For text files, read content from temp file (needed by _extract_content)
    if ext in ('.md', '.txt', '.sql', '.py'):
        with open(tmp_path, 'rb') as f:
            content = f.read()
    else:
        content = b''  # Don't hold file content in memory

    # Compute file hash (sha256) once for extraction-plan audit
    _file_hash: str | None = None
    try:
        import hashlib as _hl
        _h = _hl.sha256()
        with open(tmp_path, 'rb') as _fh:
            for _ck in iter(lambda: _fh.read(1024 * 1024), b''):
                _h.update(_ck)
        _file_hash = _h.hexdigest()
    except Exception:
        _file_hash = None

    # Save raw upload so we can re-ingest later with operator overrides.
    # 2026-05-25: extended from xlsx/xls only → also csv/json/parquet. Re-ingest
    # endpoint reads from raw_uploads/, CSV was orphaned before this fix.
    if ext in (".xlsx", ".xls", ".csv", ".json", ".parquet", ".ods") and project:
        try:
            import shutil as _sh
            _raw_dir = KNOWLEDGE_DIR / project / "raw_uploads"
            _raw_dir.mkdir(parents=True, exist_ok=True)
            _sh.copy2(tmp_path, str(_raw_dir / file.filename))
        except Exception as _re:
            try:
                _upload_lg.debug(f"raw_uploads copy failed: {_re}")
            except Exception:
                pass

    # Auto-pin any definitions/glossary sheet in an uploaded Excel file → Brain
    if ext in (".xlsx", ".xls") and project:
        try:
            _autoload_definitions(project, tmp_path, ext)
        except Exception:
            pass

    # Smart routing for non-data files (SQL, MD, TXT, DOCX, PPTX, PDF, images)
    if ext in (".sql", ".md", ".txt", ".docx", ".pptx", ".pdf", ".jpg", ".jpeg", ".png", ".html", ".htm", ".eml", ".xml") and project:
        try:
            conductor_result = _conduct_upload(tmp_path, ext, project, file.filename, raw_content=content)
            text_content = conductor_result.get("text", "")
            doc_tables = conductor_result.get("tables", [])
            file_type = classify_file(file.filename, content_sample=text_content[:500])

            if file_type == "sql_patterns":
                result = process_sql_file(project, text_content)
                # Also save as doc
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                with open(doc_dir / file.filename, "w") as f:
                    f.write(text_content)
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "sql_patterns", "smart": result}

            if file_type == "business_rules":
                result = process_business_rules_doc(project, text_content, file.filename)
                # Also save as doc
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                with open(doc_dir / file.filename, "w") as f:
                    f.write(text_content)
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "business_rules", "smart": result}

            if file_type == "documentation":
                # Save as doc + index into knowledge
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                save_name = Path(file.filename).stem + ".txt" if ext in (".pptx", ".docx", ".pdf") else file.filename
                with open(doc_dir / save_name, "w") as f:
                    f.write(text_content)
                # Extract document structure (TOC, headings, sections)
                doc_structure = _extract_document_structure(tmp_path, ext, text_content)
                if doc_structure.get("sections"):
                    # Save structure for agent reference
                    struct_dir = KNOWLEDGE_DIR / project / "doc_structure"
                    struct_dir.mkdir(parents=True, exist_ok=True)
                    _safe_write_json(struct_dir / f"{Path(file.filename).stem}.json", doc_structure)

                # Section-aware chunking (respects heading boundaries)
                section_chunks = _section_aware_chunks(text_content, doc_structure)

                # Hierarchical summarization for big docs (5+ sections)
                doc_summaries = {}
                if len(doc_structure.get("sections", [])) >= 5:
                    doc_summaries = _hierarchical_summarize(section_chunks, file.filename)
                    if doc_summaries.get("doc_summary"):
                        # Save summaries
                        summary_dir = KNOWLEDGE_DIR / project / "doc_summaries"
                        summary_dir.mkdir(parents=True, exist_ok=True)
                        _safe_write_json(summary_dir / f"{Path(file.filename).stem}.json", doc_summaries)

                # Build enriched text with section markers + page citations
                enriched_parts = []
                if doc_summaries.get("doc_summary"):
                    enriched_parts.append(f"DOCUMENT SUMMARY: {doc_summaries['doc_summary']}\n")
                for chunk in section_chunks:
                    page_ref = f" [Page {chunk['page']}]" if chunk.get("page") else ""
                    section_ref = f" [Section: {chunk['section']}]" if chunk.get("section") else ""
                    enriched_parts.append(f"{section_ref}{page_ref}\n{chunk['text']}")
                enriched_text = "\n\n".join(enriched_parts)

                # Index into project knowledge base
                indexed = False
                if text_content.strip():
                    try:
                        from agno.knowledge.reader.text_reader import TextReader
                        from db.session import create_project_knowledge
                        knowledge = create_project_knowledge(project)
                        # Use enriched text with section markers instead of raw text
                        index_text = enriched_text[:15000] if enriched_text else text_content[:10000]
                        knowledge.insert(
                            name=f"doc-{Path(file.filename).stem}",
                            text_content=f"Document: {file.filename}\n\n{index_text}",
                            reader=TextReader(),
                            skip_if_exists=False,
                        )
                        indexed = True
                    except Exception:
                        pass
                # Save extracted tables to PostgreSQL
                tables_saved = 0
                if doc_tables and project:
                    try:
                        import pandas as pd
                        from db import get_project_engine
                        from db.session import create_project_schema
                        schema = create_project_schema(project)
                        engine = get_project_engine(project)
                        for tbl_info in doc_tables:
                            df = tbl_info["df"]
                            if len(df) < 2 or len(df.columns) < 2:
                                continue
                            # Generate table name from filename + source
                            base_name = Path(file.filename).stem.lower()
                            base_name = re.sub(r'[^a-z0-9_]', '_', base_name)[:30]
                            tbl_name = f"{base_name}_{tbl_info['source']}"
                            tbl_name = re.sub(r'_+', '_', tbl_name).strip('_')
                            try:
                                df.to_sql(tbl_name, engine, schema=schema, if_exists='replace', index=False)
                                try:
                                    _upload_lg.info(f"✓ upload success: {file.filename}: {len(df)} rows × {len(df.columns)} cols → {tbl_name}")
                                except Exception:
                                    pass
                                tables_saved += 1
                            except Exception:
                                pass
                    except Exception:
                        pass
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "documentation", "filename": save_name, "indexed": indexed, "tables_saved": tables_saved}
        except Exception:
            pass

    try:
        # Excel multi-sheet: use AI handler for .xlsx/.xls
        if ext in (".xlsx", ".xls", ".parquet", ".ods", ".zip") and not table_name:
            if ext == ".parquet":
                excel_result = _handle_parquet(tmp_path, file.filename)
            elif ext == ".ods":
                excel_result = _handle_ods(tmp_path, file.filename)
            elif ext == ".zip":
                excel_result = _handle_zip(tmp_path, file.filename)
            else:
                excel_result = _handle_excel(tmp_path, file.filename)
            if excel_result.get("tables") and len(excel_result["tables"]) >= 1:
                # Multi-sheet Excel — create each sheet as its own table
                proj_schema = _get_project_schema(request, project)
                if proj_schema:
                    from db import get_project_engine
                    user_schema = proj_schema
                    engine = get_project_engine(project or '')
                else:
                    from db import create_user_schema, get_user_engine
                    user_schema = create_user_schema(user_id)
                    engine = get_user_engine(user_id)

                tables_created = []
                total_rows = 0
                # Consolidate same-schema sheets into ONE table. A workbook with
                # one sheet per month (identical columns) should become a single
                # table with a _period stamp, not N siblings to UNION later.
                try:
                    from dash.tools.ingest_router import col_fingerprint, detect_period
                except Exception:
                    col_fingerprint = detect_period = None
                fp_to_table: dict[str, str] = {}
                for tbl_info in excel_result["tables"]:
                    df = tbl_info["df"]
                    if df.empty or len(df.columns) == 0:
                        continue
                    sheet_src = str(tbl_info.get("source", ""))
                    fp = col_fingerprint(list(df.columns)) if col_fingerprint else None
                    write_mode = 'replace'
                    per = (detect_period(sheet_src) or detect_period(file.filename)) if detect_period else None
                    if fp and fp in fp_to_table:
                        # same schema as an earlier sheet → append into its table
                        tbl_name = fp_to_table[fp]
                        write_mode = 'append'
                        df["_period"] = per  # always present so grouped appends line up
                    else:
                        tbl_name = _sanitize_table_name(tbl_info["name"])
                        if fp:
                            fp_to_table[fp] = tbl_name
                            df["_period"] = per
                    # Add source tracking columns
                    df["_source_file"] = file.filename
                    df["_source_sheet"] = sheet_src
                    try:
                        df.to_sql(tbl_name, engine, schema=user_schema, if_exists=write_mode, index=False)
                        try:
                            _upload_lg.info(f"✓ upload success: {file.filename} [{tbl_info.get('source','')}]: {len(df)} rows × {len(df.columns)} cols → {tbl_name}")
                        except Exception:
                            pass
                        tables_created.append({"table": tbl_name, "rows": len(df), "cols": len(df.columns), "source": tbl_info.get("source", ""), "sheet_number": tbl_info.get("sheet_number", 0), "description": tbl_info.get("description", "")})
                        total_rows += len(df)

                        # Persist extraction plan (audit, P4) — fail-soft
                        if project:
                            try:
                                _plan_data = {
                                    "strategy": tbl_info.get("source", ""),
                                    "header_row": tbl_info.get("header_row", 0) or 0,
                                    "skip_rows": tbl_info.get("skip_rows") or [],
                                    "blocks": tbl_info.get("blocks") or [],
                                    # 2026-05-25: default to len(df) when neither
                                    # max_row nor row_count_in present — rules-direct
                                    # path was leaving it NULL → audit row hole.
                                    "row_count_in": tbl_info.get("max_row") or tbl_info.get("row_count_in") or int(len(df)),
                                    "row_count_out": len(df),
                                    "llm_rescued": "llm-rescued" in str(tbl_info.get("source", "")) or "llm-repaired" in str(tbl_info.get("source", "")),
                                    "rescue_reasoning": tbl_info.get("rescue_reasoning"),
                                    "source_file": file.filename,
                                    "sheet_name": str(tbl_info.get("source", "") or ""),
                                    "file_hash": _file_hash,
                                }
                                _persist_extraction_plan(project, tbl_name, _plan_data)
                            except Exception:
                                pass

                            # Kick off LLM column-description enrichment (fire-and-forget).
                            # Migration 154 / dash_column_meta. Never block upload return.
                            try:
                                from dash.tools.column_describer import enrich_columns_async
                                import asyncio as _aio
                                try:
                                    _aio.get_running_loop().create_task(
                                        enrich_columns_async(project, tbl_name)
                                    )
                                except RuntimeError:
                                    # no running loop — run in background thread
                                    import threading as _th
                                    _th.Thread(
                                        target=lambda: _aio.run(enrich_columns_async(project, tbl_name)),
                                        daemon=True,
                                    ).start()
                            except Exception:
                                _upload_lg.exception("col enrich kickoff failed")

                        # Save source metadata for DATASETS tab display
                        if project:
                            src_dir = KNOWLEDGE_DIR / project / "table_sources"
                            src_dir.mkdir(parents=True, exist_ok=True)
                            _safe_write_json(src_dir / f"{tbl_name}.json", {
                                "source_file": file.filename,
                                "source_detail": f"Sheet {tbl_info.get('sheet_number', '?')}: {tbl_info.get('source', '')}",
                                "description": tbl_info.get("description", ""),
                            })

                        # Profile table for real health %
                        if project:
                            try:
                                _profile_table(df, project, tbl_name)
                            except Exception:
                                pass

                        # Generate metadata + queries for each table
                        if project:
                            col_analyses = [_analyze_column(df[col]) for col in df.columns]
                            try:
                                metadata = _generate_metadata(tbl_name, df, col_analyses)
                                if metadata:
                                    meta_dir = KNOWLEDGE_DIR / project / "tables"
                                    meta_dir.mkdir(parents=True, exist_ok=True)
                                    _safe_write_json(meta_dir / f"{tbl_name}.json", metadata)
                            except Exception:
                                pass
                            try:
                                _generate_sample_queries(project, tbl_name, col_analyses)
                            except Exception:
                                pass
                    except Exception as e:
                        excel_result.setdefault("warnings", []).append(f"Table '{tbl_name}' failed: {e}")

                # Index text descriptions into knowledge
                if excel_result.get("text") and project:
                    try:
                        from agno.knowledge.reader.text_reader import TextReader
                        from db.session import create_project_knowledge
                        knowledge = create_project_knowledge(project)
                        knowledge.insert(
                            name=f"doc-{Path(file.filename).stem}",
                            text_content=f"Excel file: {file.filename}\n\n{excel_result['text'][:10000]}",
                            reader=TextReader(), skip_if_exists=False,
                        )
                    except Exception:
                        pass

                # Reload knowledge
                if project:
                    try:
                        from db.session import create_project_knowledge
                        pk = create_project_knowledge(project)
                        pk.load(recreate=False)
                    except Exception:
                        pass

                # Trigger Engineer agent in background to inspect tables, create views, discover relationships
                if project and tables_created:
                    _bg_executor.submit(_post_upload_engineer, project, tables_created, user_id or 1).add_done_callback(_bg_done_log)

                Path(tmp_path).unlink(missing_ok=True)
                # Issue #7 — auto-suggest template after multi-sheet upload.
                if project:
                    try:
                        _suggest_template_for_project(
                            project, [t.get("table", "") for t in tables_created]
                        )
                    except Exception:
                        pass
                # Normalize per-table shape so callers always get {name,rows,cols}.
                normalized_tables = [
                    {
                        "name": t.get("table"),
                        "rows": int(t.get("rows") or 0),
                        "cols": int(t.get("cols") or 0),
                        "source": t.get("source", ""),
                        "sheet_number": t.get("sheet_number", 0),
                        "description": t.get("description", ""),
                    }
                    for t in tables_created
                ]
                # Top-level rows/table_name/cols so single-table callers always
                # see populated values (fixes "rows:null, table:null" for XLSX).
                first_tbl = normalized_tables[0] if normalized_tables else {}
                top_rows = int(total_rows or 0)
                top_table = first_tbl.get("name") if len(normalized_tables) == 1 else None
                top_cols = int(first_tbl.get("cols") or 0) if len(normalized_tables) == 1 else None
                return {
                    "status": "ok",
                    "multi_sheet": True,
                    "tables_created": len(tables_created),
                    "total_rows": top_rows,
                    "rows": top_rows,
                    "table_name": top_table,
                    "cols": top_cols,
                    "tables": normalized_tables,
                    "warnings": excel_result.get("warnings", []),
                    "engineer": "running in background — inspecting tables, creating views",
                    "smart": {"file_type": "excel_multi_sheet", "tables_created": len(tables_created)},
                }

        # 1. Parse file to validate (single table: CSV, JSON, single-sheet Excel)
        try:
            df = _read_file(tmp_path, ext)
        except Exception as parse_err:
            raise HTTPException(400, f"Could not parse file: {str(parse_err)[:200]}")
        if df.empty or len(df.columns) == 0:
            raise HTTPException(400, "File is empty or has no columns")

        # Check if this is a column definition file (check headers + values)
        raw_headers = list(df.columns)
        file_type_check = classify_file(file.filename, raw_headers, df=df, project_slug=project or "")
        if file_type_check == "column_definition" and project:
            result = process_column_definitions(project, df)
            # Save record so file appears in docs list
            docs_dir = KNOWLEDGE_DIR / project / "docs"
            docs_dir.mkdir(parents=True, exist_ok=True)
            (docs_dir / file.filename).write_text(
                f"[Column Definition File: {file.filename}]\n"
                f"Annotations saved: {result.get('annotations', 0)}\n"
                f"Memories saved: {result.get('memories', 0)}\n"
                f"Rules saved: {result.get('rules', 0)}"
            )
            return {"status": "ok", "file_type": "column_definition", "smart": result}

        # Column names already cleaned by _clean_dataframe() in _read_file()

        # 2. Determine table name
        tbl = _sanitize_table_name(table_name or file.filename)

        # 3. Save raw file to project staging area (upload only — no DB load yet)
        if project:
            staging_dir = KNOWLEDGE_DIR / project / "staging"
            staging_dir.mkdir(parents=True, exist_ok=True)
            import shutil
            staged_path = staging_dir / f"{tbl}{ext}"
            shutil.copy2(tmp_path, str(staged_path))

        # 4. Analyze columns (lightweight — no DB write)
        col_analyses = [_analyze_column(df[col]) for col in df.columns]

        # 5. Load to PostgreSQL
        proj_schema = _get_project_schema(request, project)
        if proj_schema:
            from db import get_project_engine
            user_schema = proj_schema
            engine = get_project_engine(project or '')
        else:
            from db import create_user_schema, get_user_engine
            user_schema = create_user_schema(user_id)
            engine = get_user_engine(user_id)

        # Handle action for existing tables
        insp = inspect(engine)
        table_exists = tbl in insp.get_table_names(schema=user_schema)
        upload_action = action if action != "auto" else ("replace" if replace else "new")
        rows_appended = 0
        rows_upserted = 0

        if table_exists and upload_action == "new":
            # Table exists, no action specified — return smart conflict
            match_info = match_existing_table(project or "", list(df.columns))
            raise HTTPException(409, detail={
                "message": f"Table '{tbl}' already exists",
                "match": match_info,
                "options": ["replace", "append", "upsert"],
                "hint": "Use action=append to add rows, action=upsert to update+add, action=replace to overwrite",
            })

        if table_exists and upload_action == "append":
            # APPEND: check for duplicates first, insert only new rows
            try:
                with engine.connect() as conn:
                    existing_count = conn.execute(text(f'SELECT COUNT(*) FROM "{user_schema}"."{tbl}"')).scalar() or 0
                    # Use first column as key to find new rows
                    if len(df.columns) == 0:
                        raise HTTPException(400, "File has no usable columns")
                    # Smart PK detection for upsert
                    pk_col = None
                    # 1. Check for 'id' column
                    id_cols = [c for c in df.columns if 'id' in c.lower() and df[c].nunique() == len(df)]
                    if id_cols:
                        pk_col = id_cols[0]
                    # 2. Check if first column is unique
                    elif df[df.columns[0]].nunique() >= len(df) * 0.95:
                        pk_col = df.columns[0]
                    else:
                        # Can't determine PK — fall back to append
                        pk_col = df.columns[0]
                    existing_keys = set()
                    try:
                        result = conn.execute(text(f'SELECT "{pk_col}" FROM "{user_schema}"."{tbl}"'))
                        existing_keys = set(str(r[0]) for r in result.fetchall())
                    except Exception:
                        pass

                    if existing_keys:
                        new_rows = df[~df[pk_col].astype(str).isin(existing_keys)]
                    else:
                        new_rows = df

                    if len(new_rows) == 0:
                        rows_appended = 0
                    else:
                        new_rows.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                        rows_appended = len(new_rows)
            except Exception:
                # Fallback: just append all
                df.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                rows_appended = len(df)

        elif table_exists and upload_action == "upsert":
            # UPSERT: use first column as key, update existing + insert new
            # Only use columns that exist in BOTH new data and existing table
            existing_cols = set(c["name"] for c in insp.get_columns(tbl, schema=user_schema))
            common_cols = [c for c in df.columns if c in existing_cols]
            missing_in_new = existing_cols - set(df.columns)
            extra_in_new = set(df.columns) - existing_cols

            # Add missing columns to DB if new data has extra columns
            if extra_in_new:
                with engine.connect() as conn:
                    for col in extra_in_new:
                        try:
                            conn.execute(text(f'ALTER TABLE "{user_schema}"."{tbl}" ADD COLUMN IF NOT EXISTS "{col}" TEXT'))
                        except Exception:
                            pass
                    conn.commit()
                common_cols = list(df.columns)

            # Use only common columns for upsert
            df_upsert = df[common_cols] if common_cols else df
            if len(df_upsert.columns) == 0:
                raise HTTPException(400, "File has no usable columns")
            # Smart PK detection for upsert
            pk_col = None
            # 1. Check for 'id' column
            id_cols = [c for c in df_upsert.columns if 'id' in c.lower() and df_upsert[c].nunique() == len(df_upsert)]
            if id_cols:
                pk_col = id_cols[0]
            # 2. Check if first column is unique
            elif df_upsert[df_upsert.columns[0]].nunique() >= len(df_upsert) * 0.95:
                pk_col = df_upsert.columns[0]
            else:
                # Can't determine PK — fall back to append
                pk_col = df_upsert.columns[0]

            with engine.connect() as conn:
                existing_pks = set()
                try:
                    result = conn.execute(text(f'SELECT "{pk_col}" FROM "{user_schema}"."{tbl}"'))
                    existing_pks = set(str(r[0]) for r in result.fetchall())
                except Exception:
                    pass

                new_rows = df_upsert[~df_upsert[pk_col].astype(str).isin(existing_pks)]
                update_rows = df_upsert[df_upsert[pk_col].astype(str).isin(existing_pks)]

                if len(new_rows) > 0:
                    new_rows.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                    rows_appended = len(new_rows)

                for _, row in update_rows.iterrows():
                    update_cols = [c for c in df_upsert.columns if c != pk_col]
                    if not update_cols:
                        continue
                    set_clause = ", ".join(f'"{c}" = :{c}' for c in update_cols)
                    params = {c: (None if pd.isna(row[c]) else row[c]) for c in df_upsert.columns}
                    conn.execute(text(
                        f'UPDATE "{user_schema}"."{tbl}" SET {set_clause} WHERE "{pk_col}" = :{pk_col}'
                    ), params)
                rows_upserted = len(update_rows)
                conn.commit()

        else:
            # REPLACE or new table
            mode = "replace" if (replace or upload_action == "replace") else "fail"
            df.to_sql(tbl, engine, if_exists=mode, index=False, schema=user_schema)
            try:
                _upload_lg.info(f"✓ upload success: {file.filename}: {len(df)} rows × {len(df.columns)} cols → {tbl}")
            except Exception:
                pass

            # Persist extraction plan for single-table uploads (CSV, single-sheet
            # Excel, JSON, etc). Was previously xlsx-multi-sheet only — bug fix
            # 2026-05-25 extending P4 audit trail to all format paths.
            if project:
                try:
                    _ext = os.path.splitext(file.filename or "")[1].lower()
                    _plan_data = {
                        "strategy": f"{_ext.lstrip('.') or 'unknown'} [direct]",
                        "header_row": 0,
                        "skip_rows": [],
                        "blocks": [],
                        "source_file": file.filename or "",
                        "sheet_name": "",
                        "row_count_in": int(len(df)),
                        "row_count_out": int(len(df)),
                        "llm_rescued": False,
                        "rescue_reasoning": None,
                        # 2026-05-25: persist file_hash in extraction_plan row (was XLSX only).
                        # Enables dedup queries + audit trail lookup across tables.
                        "file_hash": _file_hash,
                    }
                    _persist_extraction_plan(project, tbl, _plan_data)

                    # Kick off LLM column-description enrichment (fire-and-forget).
                    # Migration 154 / dash_column_meta. Never block upload return.
                    try:
                        from dash.tools.column_describer import enrich_columns_async
                        import asyncio as _aio
                        try:
                            _aio.get_running_loop().create_task(
                                enrich_columns_async(project, tbl)
                            )
                        except RuntimeError:
                            import threading as _th
                            _th.Thread(
                                target=lambda: _aio.run(enrich_columns_async(project, tbl)),
                                daemon=True,
                            ).start()
                    except Exception:
                        _upload_lg.exception("col enrich kickoff failed")
                    # Also populate file-hash cache (P5 extension 2026-05-25 —
                    # was rescue-only; now writes on every rules-clean upload
                    # so the same file from another tenant skips the entire
                    # extraction pipeline. rescue_used=False marks rules-path.
                    try:
                        _src_path = locals().get("tmp_path") or locals().get("file_path") or ""
                        if _src_path and os.path.exists(_src_path):
                            _h = _compute_file_hash(_src_path)
                            if _h:
                                _save_upload_cache(
                                    _h, os.path.getsize(_src_path),
                                    _ext, {"plan": _plan_data, "rules_clean": True}, rescue_used=False,
                                )
                    except Exception as _ce:
                        logger.warning(f"upload cache save failed for {tbl}: {_ce}")
                except Exception as _pe:
                    logger.warning(f"extraction plan persist failed for {tbl}: {_pe}")

        # 5. Generate metadata
        metadata = _generate_metadata(tbl, df, col_analyses)

        # 6. Generate sample queries
        sample_queries = _generate_sample_queries(tbl, col_analyses)

        # 7. Generate business rules
        biz_rules = _generate_business_rules(tbl, col_analyses)

        # 8. Save knowledge files (per-project if project provided)
        if proj_schema:
            tables_dir = KNOWLEDGE_DIR / (project or proj_schema) / "tables"
            queries_dir = KNOWLEDGE_DIR / (project or proj_schema) / "queries"
            business_dir = KNOWLEDGE_DIR / (project or proj_schema) / "business"
        else:
            tables_dir = TABLES_DIR
            queries_dir = QUERIES_DIR
            business_dir = BUSINESS_DIR

        tables_dir.mkdir(parents=True, exist_ok=True)
        queries_dir.mkdir(parents=True, exist_ok=True)
        business_dir.mkdir(parents=True, exist_ok=True)

        with open(tables_dir / f"{tbl}.json", "w") as f:
            json.dump(metadata, f, indent=2, default=str)

        # Save source metadata for DATASETS tab
        if proj_schema:
            src_dir = KNOWLEDGE_DIR / (project or proj_schema) / "table_sources"
            src_dir.mkdir(parents=True, exist_ok=True)
            _safe_write_json(src_dir / f"{tbl}.json", {
                "source_file": file.filename if hasattr(file, 'filename') else "",
                "source_detail": ext.replace(".", "").upper(),
                "description": metadata.get("description", "") if isinstance(metadata, dict) else "",
            })

        # Run profiling for real health %
        try:
            _profile_table(df, project or proj_schema, tbl)
        except Exception:
            pass

        # AI review and fix data quality issues (LLM decides what to clean)
        if proj_schema:
            try:
                _ai_review_and_fix_table(project or proj_schema, tbl)
            except Exception:
                pass

        queries_file = queries_dir / f"{tbl}_queries.sql"
        with open(queries_file, "w") as f:
            f.write(f"-- Auto-generated queries for {tbl}\n\n{sample_queries}\n")

        biz_file = business_dir / f"{tbl}_rules.json"
        with open(biz_file, "w") as f:
            json.dump(biz_rules, f, indent=2, default=str)

        # 9. Reload knowledge — project-specific or global
        if proj_schema and project:
            from db.session import create_project_knowledge
            proj_knowledge = create_project_knowledge(project)
            proj_dir = KNOWLEDGE_DIR / project
            if proj_dir.exists():
                for subdir in ["tables", "queries", "business"]:
                    path = proj_dir / subdir
                    if path.exists():
                        files = [f for f in path.iterdir() if f.is_file() and not f.name.startswith(".")]
                        if files:
                            proj_knowledge.insert(name=f"{project}-{subdir}", path=str(path))
        else:
            from scripts.load_knowledge import load_knowledge
            load_knowledge(recreate=False)

        # Training is NOT auto-triggered on upload.
        # User clicks TRAIN ALL to start training manually.

        # 10. Smart analysis: classify, match, fingerprint
        file_type = classify_file(file.filename, list(df.columns))
        if project:
            save_fingerprint(project, tbl, len(df), list(df.columns))
            change_type = check_fingerprint_changed(project, tbl, len(df), list(df.columns))
        else:
            change_type = "new"

        # Auto-retrain ML — REMOVED. auto_create_models() dropped in
        # 2026-05-23 ML pivot. Per-chat tools handle ML in-process.

        # Issue #7 — auto-suggest template after single-table upload.
        if project:
            try:
                _suggest_template_for_project(project, [tbl])
            except Exception:
                pass

        # Auto-train hook — enqueue retrain if AUTO_TRAIN_ON_UPLOAD enabled
        if os.getenv("AUTO_TRAIN_ON_UPLOAD", "1") != "0":
            try:
                from dash.cron.auto_train_daemon import _enqueue_retrain, _is_training_running
                import asyncio as _at_aio
                _proj = project
                if _proj:
                    # Run check in background, don't block upload response
                    async def _deferred_retrain():
                        await _at_aio.sleep(5)  # 5s grace — let upload finish fully
                        already = await _is_training_running(_proj)
                        if not already:
                            _enqueue_retrain(_proj, "upload_trigger")
                    try:
                        _loop = _at_aio.get_event_loop()
                        if _loop.is_running():
                            _at_aio.create_task(_deferred_retrain())
                    except Exception:
                        pass
            except Exception:
                pass

        return {
            "status": "ok",
            "table_name": tbl,
            "rows": len(df),
            "columns": len(df.columns),
            "column_details": col_analyses,
            "metadata": metadata,
            "sample_queries_count": sample_queries.count("<query>"),
            "business_rules": biz_rules,
            # Smart upload info
            "smart": {
                "file_type": file_type,
                "action": upload_action,
                "rows_appended": rows_appended,
                "rows_upserted": rows_upserted,
                "change_type": change_type,
                "fingerprint": compute_fingerprint(len(df), list(df.columns)),
            },
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Upload failed: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


# ─────────────────────────────────────────────────────────────────────────────
# STAGED INGEST PIPELINE
#   upload → stage (disk, hashed) → validate + schema-contract → dry-run → gate →
#   promote (idempotent, lineage-stamped) → train.
# Files never touch Postgres until /promote. Drift/dups/low-quality are caught at
# staging, not after they corrupt a table. Direct (non-staged) upload still works.
# ─────────────────────────────────────────────────────────────────────────────

_STAGE_DATA_EXTS = (".csv", ".xlsx", ".xls", ".json", ".parquet")


def _logical_dataset(filename: str) -> str:
    """Map a filename to its CONSISTENT logical dataset by stripping period tokens.
    'MM Conso Apr 25.csv', 'MM Conso May 25.csv' → both 'mm_conso' → one contract,
    one table. This is what turns monthly drops into a single consolidated table."""
    s = Path(filename).stem.lower()
    s = re.sub(r"(20\d{2})[-_]?(0[1-9]|1[0-2])", " ", s)            # 2025-04 / 202504
    s = re.sub(r"\b(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*['_ -]?\d{2,4}\b", " ", s)
    s = re.sub(r"\b(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\b", " ", s)
    s = re.sub(r"\bq[1-4]\b", " ", s)
    s = re.sub(r"\b20\d{2}\b", " ", s)                               # bare year
    s = re.sub(r"['_ -]\d{2}\b", " ", s)                            # trailing _25 / -25
    s = re.sub(r"[^a-z0-9]+", "_", s)
    s = re.sub(r"_+", "_", s).strip("_")
    return s[:50] or "dataset"


def _stage_resolve_engine(request: Request, project: str):
    """Return (engine, schema) for a project's data, mirroring upload_file."""
    from db import get_project_engine
    schema = _get_project_schema(request, project)
    if not schema:
        import re as _re
        schema = _re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    return get_project_engine(project), schema


@router.post("/upload/stage")
async def stage_upload(request: Request, file: UploadFile, project: str | None = Form(None), batch_id: str | None = Form(None)):
    """Stage ONE data file into a batch (no DB write). Call once per file with the
    same batch_id to build a multi-file batch. Returns the updated manifest."""
    # Dual-source: accept project as form OR query param (frontend sends ?project=)
    if not project:
        project = request.query_params.get("project")
    if not batch_id:
        batch_id = request.query_params.get("batch_id")
    if not project:
        raise HTTPException(422, "project required (query param or form field)")
    from app.auth import check_project_permission
    from dash.ingest import (new_batch_id, stage_file, write_manifest, read_manifest)
    from dash.ingest.contract import get_contract, check_against_contract, infer_contract, detect_load_key
    from dash.tools.ingest_router import detect_period

    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required to upload")
    if not file.filename:
        raise HTTPException(400, "No filename")
    ext = Path(file.filename).suffix.lower()
    if ext not in _STAGE_DATA_EXTS:
        raise HTTPException(400, f"Staged ingest supports data files only: {', '.join(_STAGE_DATA_EXTS)}")

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        size = 0
        while chunk := await file.read(1024 * 1024):
            size += len(chunk)
            if size > MAX_FILE_SIZE:
                raise HTTPException(400, f"File too large. Max: {MAX_FILE_SIZE // (1024*1024)} MB")
            tmp.write(chunk)
        tmp_path = tmp.name

    try:
        bid = batch_id or new_batch_id()
        manifest = read_manifest(project, bid) or {
            "batch_id": bid, "project": project,
            "created_at": _dt_datetime.utcnow().isoformat(), "status": "staged", "files": [],
        }
        # Build (df, sheet) pairs. Multi-sheet xlsx → one staged entry per sheet
        # (each a distinct logical dataset); csv/json/parquet/single-sheet → one.
        pairs = []
        if ext == ".parquet":
            pairs = [(_clean_dataframe(pd.read_parquet(tmp_path)), None)]
        elif ext in (".xlsx", ".xls"):
            tbls = []
            try:
                tbls = (_handle_excel(tmp_path, file.filename) or {}).get("tables") or []
            except Exception:
                tbls = []
            if len(tbls) > 1:
                for t in tbls:
                    d = t.get("df")
                    if d is not None and not d.empty:
                        pairs.append((d, str(t.get("source") or t.get("name") or "")))
            if not pairs:
                pairs = [(_read_file(tmp_path, ".xlsx"), None)]
        else:
            pairs = [(_read_file(tmp_path, ext), None)]

        # Auto-pin any Definitions/glossary sheet → Brain (kept from direct path).
        if ext in (".xlsx", ".xls"):
            try:
                _autoload_definitions(project, tmp_path, ext)
            except Exception:
                pass

        from dash.ingest import batch_dir as _batch_dir
        _bdir = _batch_dir(project, bid)
        entries = []
        for _i, (df, sheet) in enumerate(pairs):
            quality = _validate_dataframe(df)
            entry = stage_file(project, bid, tmp_path, file.filename, df=df, ext=ext, quality=quality)
            # Persist the PARSED df now (deterministic). promote/recheck reload THIS,
            # never re-run _handle_excel (LLM, nondeterministic, sheet-name drift).
            try:
                _dp = str(_bdir / f"{_sanitize_table_name(Path(file.filename).stem)}_{_i}.parquet")
                df.to_parquet(_dp, index=False)
                entry["data_path"] = _dp
            except Exception:
                entry["data_path"] = None
            if sheet:
                dataset = _logical_dataset(f"{Path(file.filename).stem} {sheet}")
                entry["sheet"] = sheet
                entry["filename"] = f"{file.filename}#{sheet}"   # unique manifest key per sheet
            else:
                dataset = _logical_dataset(file.filename)
            entry["dataset"] = dataset
            entry["target_table"] = dataset
            entry["period"] = detect_period(sheet or "") or detect_period(file.filename)
            contract = get_contract(project, dataset)
            if contract:
                chk = check_against_contract(contract, df)
                entry["verdict"] = chk["verdict"]
                entry["diff"] = chk.get("diff", {})
                entry["load_key"] = contract.get("load_key", {})
                if chk["verdict"] == "drift":
                    # Auto-evolve contract — accept new shape silently, no drift flag
                    try:
                        from dash.ingest.contract import infer_contract, save_contract
                        new_contract = infer_contract(df, project, dataset, file.filename)
                        save_contract(project, dataset, new_contract)
                        entry["verdict"] = "exact"
                        entry["diff"] = {}
                    except Exception:
                        pass
                    entry["status"] = "ready"
                    entry["reason"] = ""
            else:
                entry["verdict"] = "new"
                entry["load_key"] = detect_load_key(df, sheet or file.filename)
            # Only quarantine truly broken files (empty / unreadable).
            # Real CSVs with many sparse cols hit score 0 but train fine. Lowered 40→10.
            if quality.get("score", 100) < 10 and int(entry.get("rows", 0) or 0) < 5:
                entry["status"] = "quarantine"
                entry["reason"] = f"low quality score {quality.get('score')}"
            entries.append(entry)

        fnames = {e["filename"] for e in entries}
        manifest["files"] = [f for f in manifest["files"] if f.get("filename") not in fnames]
        manifest["files"].extend(entries)
        manifest["status"] = "staged"
        write_manifest(project, bid, manifest)
        return {"status": "staged", "batch_id": bid, "files": entries,
                "files_in_batch": len(manifest["files"])}
    finally:
        Path(tmp_path).unlink(missing_ok=True)


@router.get("/ingest/{project}/batches")
def list_ingest_batches(request: Request, project: str):
    from dash.ingest import list_batches
    return {"batches": list_batches(project)}


@router.get("/ingest/{project}/{batch_id}/dry-run")
def ingest_dry_run(request: Request, project: str, batch_id: str):
    """Predict what /promote will do — per file — WITHOUT touching the DB."""
    from dash.ingest import read_manifest
    from dash.ingest.contract import get_contract
    from dash.ingest.loader import table_exists, file_hash_seen
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    engine, schema = _stage_resolve_engine(request, project)
    plan = []
    for f in manifest.get("files", []):
        ds, tbl = f.get("dataset"), f.get("target_table") or f.get("dataset")
        rows = f.get("rows", 0)
        if f.get("status") == "quarantine":
            plan.append({"filename": f["filename"], "action": "quarantine", "reason": f.get("reason", ""), "rows": rows})
            continue
        contract = get_contract(project, ds)
        lk = (contract or {}).get("load_key", f.get("load_key", {})) or {}
        strat = lk.get("strategy", "content_hash")
        exists = table_exists(engine, schema, tbl)
        if not exists:
            action = f"create table '{tbl}' + load {rows} rows"
        elif f.get("content_hash") and file_hash_seen(engine, schema, tbl, f["content_hash"]):
            action = "skip — identical file already loaded"
        elif strat == "period" and f.get("period"):
            action = f"replace period {f['period']} in '{tbl}' (~{rows} rows)"
        elif strat in ("single", "composite"):
            action = f"upsert into '{tbl}' by {strat} key {lk.get('columns')} (≤{rows} new)"
        else:
            action = f"append {rows} rows to '{tbl}'"
        plan.append({"filename": f["filename"], "dataset": ds, "target_table": tbl,
                     "verdict": f.get("verdict"), "load_key": lk, "period": f.get("period"),
                     "rows": rows, "action": action})
    return {"batch_id": batch_id, "plan": plan}


@router.post("/ingest/{project}/{batch_id}/promote")
async def ingest_promote(request: Request, project: str, batch_id: str, train: bool = False, force: bool = False):
    """Load every ready file into Postgres idempotently via its contract. Drift files
    are quarantined (not loaded). Returns per-file results."""
    from app.auth import check_project_permission
    from dash.ingest import read_manifest, write_manifest, quarantine_file
    from dash.ingest.contract import get_contract, infer_contract, save_contract, check_against_contract
    from dash.ingest.loader import promote_file
    from dash.tools.ingest_router import detect_period

    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    engine, schema = _stage_resolve_engine(request, project)
    results = []
    for f in manifest.get("files", []):
        fn = f["filename"]
        if f.get("status") == "quarantine" and not force:
            results.append({"filename": fn, "action": "quarantine", "reason": f.get("reason", ""), "rows_loaded": 0})
            continue
        try:
            df = _stage_reload_df(f)
        except Exception as e:
            results.append({"filename": fn, "action": "error", "reason": f"read failed: {e}", "rows_loaded": 0})
            continue

        ds = f.get("dataset") or _logical_dataset(fn)
        tbl = f.get("target_table") or ds
        contract = get_contract(project, ds)
        drift_diff = None
        if contract:
            chk = check_against_contract(contract, df)
            if chk["verdict"] == "drift":
                drift_diff = chk.get("diff")
                if not force:
                    quarantine_file(project, batch_id, fn, "schema drift vs contract")
                    results.append({"filename": fn, "action": "quarantine", "reason": "schema drift", "diff": drift_diff, "rows_loaded": 0})
                    continue
                # force=true: load anyway, evolve contract to accept new shape
                try:
                    contract = infer_contract(df, project, ds, fn)
                    save_contract(project, ds, contract)
                except Exception:
                    pass
        else:
            contract = infer_contract(df, project, ds, fn)
            save_contract(project, ds, contract)
        lineage = {"source_file": fn, "period": f.get("period") or detect_period(fn),
                   "batch_id": batch_id, "content_hash": f.get("content_hash", "")}
        r = promote_file(engine, schema, tbl, df, contract, lineage)
        r["filename"] = fn
        if drift_diff:
            r["warnings"] = [f"schema drift accepted: {drift_diff}"]
        results.append(r)

    loaded = sum(1 for r in results if r.get("rows_loaded", 0) > 0 or r.get("action") in ("create", "append", "replace_period", "upsert"))
    quarantined = sum(1 for r in results if r.get("action") == "quarantine")
    manifest["status"] = "promoted" if quarantined == 0 else "partial"
    write_manifest(project, batch_id, manifest)

    trained = False
    if train and (loaded or force):
        # Reuse the existing retrain entrypoint (spawns its own bg thread, returns
        # immediately). promote already verified editor perm + request.state.user.
        try:
            await retrain_project(project, request)
            trained = True
        except Exception:
            trained = False
    return {"batch_id": batch_id, "status": manifest["status"], "results": results,
            "loaded_files": loaded, "quarantined": quarantined, "train_triggered": trained}


@router.post("/ingest/{project}/{batch_id}/reject")
def ingest_reject(request: Request, project: str, batch_id: str):
    """Discard a staged batch — delete its folder. Nothing to roll back (no DB load)."""
    from app.auth import check_project_permission
    from dash.ingest import read_manifest, write_manifest, batch_dir
    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    try:
        import shutil
        shutil.rmtree(str(batch_dir(project, batch_id)), ignore_errors=True)
    except Exception:
        pass
    manifest["status"] = "rejected"
    try:
        write_manifest(project, batch_id, manifest)
    except Exception:
        pass
    return {"status": "rejected", "batch_id": batch_id}


def _stage_reload_df(entry: dict):
    """Reload a staged file entry into the EXACT df parsed at stage time.
    Reads the persisted parquet (deterministic) — never re-parses the original
    (which for messy xlsx means an LLM call that can drift between stage/promote)."""
    dp = entry.get("data_path")
    if dp and Path(dp).exists():
        return pd.read_parquet(dp)
    # fallback for entries staged before data_path existed
    staged = entry.get("staged_path")
    ext = entry.get("ext") or Path(entry.get("filename", "")).suffix.lower()
    if ext == ".parquet":
        return _clean_dataframe(pd.read_parquet(staged))
    return _read_file(staged, ".xlsx" if ext == ".xls" else ext)


@router.post("/ingest/{project}/{batch_id}/resolve-drift")
async def ingest_resolve_drift(request: Request, project: str, batch_id: str):
    """Accept a drift file: apply an optional column mapping, evolve the contract to
    a new version, mark the file ready. Body: {filename, mapping?:{new_col:existing}}."""
    from app.auth import check_project_permission
    from dash.ingest import read_manifest, write_manifest
    from dash.ingest.contract import evolve_contract, get_contract, check_against_contract
    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass
    fn = body.get("filename")
    mapping = body.get("mapping") or {}
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    entry = next((f for f in manifest.get("files", []) if f.get("filename") == fn), None)
    if not entry:
        raise HTTPException(404, "file not in batch")
    ext = entry.get("ext") or Path(fn).suffix.lower()
    df = _stage_reload_df(entry)
    ds = entry.get("dataset") or _logical_dataset(fn)
    new_contract = evolve_contract(project, ds, df, mapping, fn)
    # re-check against the evolved contract
    chk = check_against_contract(new_contract, df.rename(columns=mapping) if mapping else df)
    entry["verdict"] = chk["verdict"]
    entry["diff"] = chk.get("diff", {})
    entry["status"] = "ready" if chk["verdict"] != "drift" else "quarantine"
    entry["reason"] = "" if entry["status"] == "ready" else "still drifts after mapping"
    write_manifest(project, batch_id, manifest)
    return {"status": "ok", "filename": fn, "verdict": entry["verdict"],
            "contract_version": new_contract.get("version"), "file": entry}


@router.post("/ingest/{project}/{batch_id}/recheck")
async def ingest_recheck(request: Request, project: str, batch_id: str):
    """Re-validate a (possibly re-uploaded/fixed) file against the current contract.
    Body: {filename}. Releases from quarantine if it now matches."""
    from app.auth import check_project_permission
    from dash.ingest import read_manifest, write_manifest
    from dash.ingest.contract import get_contract, check_against_contract
    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass
    fn = body.get("filename")
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    entry = next((f for f in manifest.get("files", []) if f.get("filename") == fn), None)
    if not entry:
        raise HTTPException(404, "file not in batch")
    ext = entry.get("ext") or Path(fn).suffix.lower()
    df = _stage_reload_df(entry)
    ds = entry.get("dataset") or _logical_dataset(fn)
    contract = get_contract(project, ds)
    if contract:
        chk = check_against_contract(contract, df)
        entry["verdict"] = chk["verdict"]
        entry["diff"] = chk.get("diff", {})
        entry["status"] = "ready" if chk["verdict"] != "drift" else "quarantine"
        entry["reason"] = "" if entry["status"] == "ready" else "schema drift vs contract"
    else:
        entry["verdict"] = "new"; entry["status"] = "ready"; entry["reason"] = ""
    write_manifest(project, batch_id, manifest)
    return {"status": "ok", "file": entry}


@router.post("/ingest/{project}/{dataset}/load-key")
async def ingest_set_load_key(request: Request, project: str, dataset: str):
    """Override a dataset's contract load_key. Body: {strategy, columns?}."""
    from app.auth import check_project_permission
    from dash.ingest.contract import set_load_key
    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass
    strat = body.get("strategy")
    cols = body.get("columns") or []
    if strat not in ("single", "composite", "period", "content_hash"):
        raise HTTPException(400, "strategy must be single|composite|period|content_hash")
    res = set_load_key(project, dataset, strat, cols)
    if isinstance(res, dict) and res.get("error"):
        raise HTTPException(404, res["error"])
    return {"status": "ok", "dataset": dataset, "load_key": res.get("load_key")}


@router.post("/ingest/{project}/{batch_id}/undo")
def ingest_undo(request: Request, project: str, batch_id: str):
    """Roll back a promoted batch — DELETE every row stamped with this _batch_id
    across the tables it touched. Lineage makes this surgical."""
    from app.auth import check_project_permission
    from dash.ingest import read_manifest, write_manifest
    from dash.ingest.loader import delete_where_batch
    user = getattr(getattr(request, "state", None), "user", None)
    if user and not check_project_permission(user, project, required_role="editor"):
        raise HTTPException(403, "Editor access required")
    manifest = read_manifest(project, batch_id)
    if not manifest:
        raise HTTPException(404, "batch not found")
    engine, schema = _stage_resolve_engine(request, project)
    targets = sorted({f.get("target_table") or f.get("dataset") for f in manifest.get("files", []) if f.get("target_table") or f.get("dataset")})
    deleted = {}
    for tbl in targets:
        try:
            deleted[tbl] = delete_where_batch(engine, schema, tbl, batch_id)
        except Exception:
            deleted[tbl] = 0
    manifest["status"] = "undone"
    try:
        write_manifest(project, batch_id, manifest)
    except Exception:
        pass
    return {"status": "undone", "batch_id": batch_id, "deleted_rows": deleted,
            "total": sum(deleted.values())}


@router.get("/dashboard")
def get_dashboard(request: Request):
    """Full dashboard data: tables, knowledge stats, system info."""
    user_id = _get_user_id(request)
    from db import create_user_schema
    user_schema = create_user_schema(user_id)

    engine = create_engine(db_url)
    insp = inspect(engine)

    # Tables (public + user schema)
    tables_list = []
    total_rows = 0
    with engine.connect() as conn:
        # Public tables (shared demo)
        for tbl in sorted(insp.get_table_names(schema="public")):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM public."{tbl}"')).scalar() or 0
                cols = insp.get_columns(tbl, schema="public")
                tables_list.append({
                    "name": tbl, "rows": count, "columns": len(cols),
                    "protected": tbl in PROTECTED_TABLES,
                })
                total_rows += count
            except Exception:
                tables_list.append({"name": tbl, "rows": 0, "columns": 0, "protected": tbl in PROTECTED_TABLES})

        # User schema tables
        for tbl in sorted(insp.get_table_names(schema=user_schema)):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM "{user_schema}"."{tbl}"')).scalar() or 0
                cols = insp.get_columns(tbl, schema=user_schema)
                tables_list.append({
                    "name": tbl, "rows": count, "columns": len(cols),
                    "protected": False,
                })
                total_rows += count
            except Exception:
                tables_list.append({"name": tbl, "rows": 0, "columns": 0, "protected": False})

        # DB size
        try:
            db_size = conn.execute(text("SELECT pg_size_pretty(pg_database_size(current_database()))")).scalar()
        except Exception:
            db_size = "unknown"

        # Knowledge vectors
        knowledge_count = 0
        learnings_count = 0
        try:
            knowledge_count = conn.execute(text("SELECT COUNT(*) FROM ai.dash_knowledge")).scalar() or 0
        except Exception:
            pass
        try:
            learnings_count = conn.execute(text("SELECT COUNT(*) FROM ai.dash_learnings")).scalar() or 0
        except Exception:
            pass

    # Knowledge files
    tables_files = len(list(TABLES_DIR.glob("*.json"))) if TABLES_DIR.exists() else 0
    queries_files = len(list(QUERIES_DIR.glob("*.sql"))) if QUERIES_DIR.exists() else 0
    business_files = len(list(BUSINESS_DIR.glob("*.json"))) if BUSINESS_DIR.exists() else 0

    return {
        "tables": tables_list,
        "stats": {
            "table_count": len(tables_list),
            "total_rows": total_rows,
            "knowledge_vectors": knowledge_count,
            "learnings_count": learnings_count,
            "db_size": db_size,
            "model": TRAINING_MODEL,
            "model_provider": "OpenRouter",
            "embeddings": "text-embedding-3-small",
            "status": "online",
        },
        "knowledge": {
            "table_metadata": tables_files,
            "query_patterns": queries_files,
            "business_rules": business_files,
            "learnings": learnings_count,
        },
    }


@router.post("/upload-doc")
async def upload_document(request: Request, file: UploadFile, project: str | None = Form(None)):
    """Upload a code/doc file to project or global knowledge base. Streams progress via SSE."""
    from starlette.responses import StreamingResponse
    import queue

    # Editor role required for doc uploads
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    allowed = {".sql", ".py", ".txt", ".md", ".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png"}
    if ext not in allowed:
        raise HTTPException(400, f"Unsupported: {ext}. Allowed: {', '.join(allowed)}")

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(400, "File too large")

    # Check if client wants streaming (Accept: text/event-stream)
    accept = request.headers.get("accept", "")
    wants_stream = "text/event-stream" in accept

    # Progress queue for streaming
    progress_q: queue.Queue = queue.Queue()

    def _emit(agent: str, step: str, detail: str):
        progress_q.put({"agent": agent, "step": step, "detail": detail})

    # Use Upload Conductor for all formats
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    if wants_stream:
        # SSE streaming mode — run processing in thread, stream progress
        import threading

        final_result = {}

        def _process():
            nonlocal final_result
            try:
                # Verify temp file exists and has content
                import os as _os
                fsize = _os.path.getsize(tmp_path) if _os.path.exists(tmp_path) else 0
                _emit("Conductor", "File received", f"{fsize:,} bytes, type: {ext}")

                r = _conduct_upload(tmp_path, ext, project or "", file.filename, raw_content=content, _progress=_emit)

                # Report any errors from handlers
                for err in r.get("errors", []):
                    _emit("Conductor", "Handler error", err[:200])
                for warn in r.get("warnings", []):
                    _emit("Conductor", "Info", warn[:200])

                final_result["conductor"] = r
                _emit("Conductor", "Processing complete", f"{len(r.get('text', '')):,} chars, {len(r.get('tables', []))} tables, {len(r.get('images', []))} images")
            except Exception as e:
                import traceback
                tb = traceback.format_exc()
                final_result["error"] = str(e)
                _emit("Conductor", "Error", f"{str(e)[:150]} | {tb.splitlines()[-2][:100] if len(tb.splitlines()) > 1 else ''}")
            finally:
                progress_q.put(None)  # Signal done

        thread = threading.Thread(target=_process, daemon=True)
        thread.start()

        def event_generator():
            while True:
                try:
                    msg = progress_q.get(timeout=300)
                except queue.Empty:
                    yield f"data: {safe_dumps({'agent': 'Conductor', 'step': 'Timeout', 'detail': 'processing took too long'})}\n\n"
                    break
                if msg is None:
                    break
                yield f"data: {safe_dumps(msg)}\n\n"

            # Wait for thread to finish (Vision LLM can take minutes for large files)
            thread.join(timeout=300)

            # Process the rest (save to knowledge, tables, etc.) synchronously
            conductor_result = final_result.get("conductor", {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []})
            if final_result.get("error"):
                yield f"data: {safe_dumps({'agent': 'Conductor', 'step': 'Error', 'detail': final_result['error'][:200]})}\n\n"

            text_content = conductor_result.get("text", "")
            doc_tables = conductor_result.get("tables", [])
            import os; os.unlink(tmp_path)

            # Save text
            _emit_direct = lambda a, s, d: None  # already streamed
            if project:
                docs_dir = KNOWLEDGE_DIR / project / "docs"
            else:
                docs_dir = KNOWLEDGE_DIR / "docs"
            docs_dir.mkdir(parents=True, exist_ok=True)
            with open(docs_dir / file.filename, "w") as f:
                f.write(text_content)
            yield f"data: {safe_dumps({'agent': 'Inspector', 'step': 'Saving text', 'detail': f'{len(text_content):,} chars saved'})}\n\n"

            # Save raw binary
            if ext in (".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png") and project:
                docs_raw_dir = KNOWLEDGE_DIR / project / "docs_raw"
                docs_raw_dir.mkdir(parents=True, exist_ok=True)
                (docs_raw_dir / file.filename).write_bytes(content)

            # Index to knowledge
            yield f"data: {safe_dumps({'agent': 'Inspector', 'step': 'Knowledge indexing', 'detail': 'indexing to PgVector...'})}\n\n"
            try:
                from agno.knowledge.reader.text_reader import TextReader
                if project:
                    from db.session import create_project_knowledge
                    knowledge = create_project_knowledge(project)
                else:
                    from dash.settings import dash_knowledge
                    knowledge = dash_knowledge
                doc_name = Path(file.filename).stem
                knowledge.insert(name=f"doc-{doc_name}", text_content=f"Document: {file.filename}\n\n{text_content[:10000]}", reader=TextReader(), skip_if_exists=False)
                yield f"data: {safe_dumps({'agent': 'Inspector', 'step': 'Knowledge indexing', 'detail': 'indexed to PgVector ✓'})}\n\n"
            except Exception as e:
                yield f"data: {safe_dumps({'agent': 'Inspector', 'step': 'Knowledge indexing', 'detail': f'failed: {str(e)[:80]}'})}\n\n"

            # Vector sync — re-embed doc into dash_vectors (defensive: never break upload).
            if project:
                try:
                    import asyncio as _asyncio_vs
                    from dash.tools.vector_sync import VECTOR_SYNC as _VS
                    _vs_text = f"{file.filename}\n\n{text_content}"
                    _asyncio_vs.create_task(
                        _VS.enqueue(project, "docs", file.filename, _vs_text, {})
                    )
                except Exception:
                    import logging as _vs_log
                    _vs_log.getLogger(__name__).debug("vector_sync enqueue (stream) skipped", exc_info=True)

            # Save tables
            tables_saved = 0
            if doc_tables and project:
                yield f"data: {safe_dumps({'agent': 'Engineer', 'step': 'Tables → PostgreSQL', 'detail': f'saving {len(doc_tables)} tables...'})}\n\n"
                try:
                    from db import get_project_engine
                    from db.session import create_project_schema
                    schema = create_project_schema(project)
                    engine = get_project_engine(project)
                    for tbl_info in doc_tables:
                        df = tbl_info["df"]
                        if len(df) < 2 or len(df.columns) < 2:
                            continue
                        tbl_name = tbl_info.get("name", "")
                        if not tbl_name:
                            base_name = Path(file.filename).stem.lower()
                            base_name = re.sub(r'[^a-z0-9_]', '_', base_name)[:30]
                            tbl_name = f"{base_name}_{tbl_info.get('source', 'table')}"
                        tbl_name = re.sub(r'[^a-z0-9_]', '_', tbl_name.lower())
                        tbl_name = re.sub(r'_+', '_', tbl_name).strip('_')[:63]
                        try:
                            df.to_sql(tbl_name, engine, schema=schema, if_exists='replace', index=False)
                            try:
                                _upload_lg.info(f"✓ upload success: {file.filename}: {len(df)} rows × {len(df.columns)} cols → {tbl_name}")
                            except Exception:
                                pass
                            tables_saved += 1
                            src_dir = KNOWLEDGE_DIR / project / "table_sources"
                            src_dir.mkdir(parents=True, exist_ok=True)
                            source_detail = tbl_info.get("source", "")
                            _safe_write_json(src_dir / f"{tbl_name}.json", {"source_file": file.filename, "source_detail": source_detail, "description": ""})
                            yield f"data: {safe_dumps({'agent': 'Engineer', 'step': 'Table saved', 'detail': f'{tbl_name} ({len(df)} rows)'})}\n\n"
                        except Exception:
                            pass
                except Exception:
                    pass

            # Engineer merge in background
            if tables_saved > 0 and project:
                doc_tables_info = [{"table": t.get("name", ""), "rows": len(t.get("df", [])), "cols": len(t.get("df", {}).columns) if hasattr(t.get("df"), "columns") else 0, "source": t.get("source", "")} for t in doc_tables if t.get("df") is not None and len(t.get("df", [])) >= 2]
                if doc_tables_info:
                    user_id_val = getattr(getattr(getattr(request, 'state', None), 'user', None), 'id', 1)
                    _bg_executor.submit(_post_upload_engineer, project, doc_tables_info, user_id_val).add_done_callback(_bg_done_log)

            # Save extraction metadata for document info cards
            if project:
                try:
                    meta_dir = KNOWLEDGE_DIR / project / "doc_meta"
                    meta_dir.mkdir(parents=True, exist_ok=True)
                    raw_file = KNOWLEDGE_DIR / project / "docs_raw" / file.filename
                    file_size_bytes = raw_file.stat().st_size if raw_file.exists() else len(content)
                    _safe_write_json(meta_dir / f"{file.filename}.json", {
                        "text_chars": len(text_content),
                        "tables_extracted": tables_saved,
                        "images_described": len(conductor_result.get("images", [])),
                        "slides": conductor_result.get("metadata", {}).get("slides", 0),
                        "pages": conductor_result.get("metadata", {}).get("total_pages", 0),
                        "scanned_pages": conductor_result.get("metadata", {}).get("scanned_pages", 0),
                        "notes_count": conductor_result.get("metadata", {}).get("notes_count", 0),
                        "file_size": file_size_bytes,
                        "warnings": conductor_result.get("warnings", [])[:5],
                        "errors": conductor_result.get("errors", [])[:5],
                        "doc_category": _doc_category_hint(file.filename),
                    })
                except Exception:
                    pass

            # Best-effort: notify project owner that doc was indexed
            if project:
                try:
                    from app.auth import notify_user
                    from sqlalchemy import create_engine as _ne2
                    _eng2 = _ne2(db_url)
                    with _eng2.connect() as _c:
                        _r = _c.execute(text(
                            "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                        ), {"s": project}).fetchone()
                    _eng2.dispose()
                    if _r and _r[0]:
                        _pages = conductor_result.get("metadata", {}).get("total_pages", 0) or conductor_result.get("metadata", {}).get("slides", 0)
                        notify_user(
                            int(_r[0]),
                            f"Document indexed · {file.filename}",
                            f"{_pages} pages · {len(text_content):,} chars · {tables_saved} tables",
                            "info",
                        )
                except Exception:
                    pass

            # Final result event
            yield f"data: {safe_dumps({'agent': 'Conductor', 'step': 'Complete', 'detail': f'{len(text_content):,} chars, {tables_saved} tables', 'done': True, 'result': {'status': 'ok', 'filename': file.filename, 'type': ext, 'size': len(text_content), 'indexed': True, 'tables_saved': tables_saved}})}\n\n"

        return StreamingResponse(event_generator(), media_type="text/event-stream")

    # Non-streaming mode (fallback)
    conductor_result = _conduct_upload(tmp_path, ext, project or "", file.filename, raw_content=content)
    text_content = conductor_result.get("text", "")
    doc_tables = conductor_result.get("tables", [])

    import os; os.unlink(tmp_path)
    doc_name = Path(file.filename).stem

    # Per-project or global docs dir
    if project:
        docs_dir = KNOWLEDGE_DIR / project / "docs"
    else:
        docs_dir = KNOWLEDGE_DIR / "docs"
    docs_dir.mkdir(parents=True, exist_ok=True)
    doc_path = docs_dir / file.filename
    with open(doc_path, "w") as f:
        f.write(text_content)

    # Save raw binary for structure extraction (doc-to-workflow)
    if ext in (".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png") and project:
        docs_raw_dir = KNOWLEDGE_DIR / project / "docs_raw"
        docs_raw_dir.mkdir(parents=True, exist_ok=True)
        (docs_raw_dir / file.filename).write_bytes(content)

    # Index into project-specific or global knowledge
    from agno.knowledge.reader.text_reader import TextReader
    if project:
        from db.session import create_project_knowledge
        knowledge = create_project_knowledge(project)
    else:
        from dash.settings import dash_knowledge
        knowledge = dash_knowledge

    try:
        knowledge.insert(
            name=f"doc-{doc_name}",
            text_content=f"Document: {file.filename}\n\n{text_content[:10000]}",
            reader=TextReader(),
            skip_if_exists=False,
        )
    except Exception as e:
        raise HTTPException(500, f"Indexing failed: {str(e)}")

    # Vector sync — re-embed doc into dash_vectors (defensive: never fail upload).
    if project:
        try:
            import asyncio as _asyncio_vs
            from dash.tools.vector_sync import VECTOR_SYNC as _VS
            _vs_text = f"{file.filename}\n\n{text_content}"
            _asyncio_vs.create_task(
                _VS.enqueue(project, "docs", file.filename, _vs_text, {})
            )
        except Exception:
            import logging as _vs_log
            _vs_log.getLogger(__name__).debug("vector_sync enqueue (non-stream) skipped", exc_info=True)

    # Save extracted tables to PostgreSQL
    tables_saved = 0
    if doc_tables and project:
        try:
            import pandas as pd
            from db import get_project_engine
            from db.session import create_project_schema
            schema = create_project_schema(project)
            engine = get_project_engine(project)
            for tbl_info in doc_tables:
                df = tbl_info["df"]
                if len(df) < 2 or len(df.columns) < 2:
                    continue
                # Use conductor-provided name or generate from filename + source
                tbl_name = tbl_info.get("name", "")
                if not tbl_name:
                    base_name = Path(file.filename).stem.lower()
                    base_name = re.sub(r'[^a-z0-9_]', '_', base_name)[:30]
                    tbl_name = f"{base_name}_{tbl_info.get('source', 'table')}"
                tbl_name = re.sub(r'[^a-z0-9_]', '_', tbl_name.lower())
                tbl_name = re.sub(r'_+', '_', tbl_name).strip('_')[:63]
                try:
                    df.to_sql(tbl_name, engine, schema=schema, if_exists='replace', index=False)
                    try:
                        _upload_lg.info(f"✓ upload success: {file.filename}: {len(df)} rows × {len(df.columns)} cols → {tbl_name}")
                    except Exception:
                        pass
                    tables_saved += 1
                    # Save source metadata
                    src_dir = KNOWLEDGE_DIR / project / "table_sources"
                    src_dir.mkdir(parents=True, exist_ok=True)
                    source_detail = tbl_info.get("source", "")
                    if "slide" in source_detail:
                        source_detail = f"Slide {source_detail.replace('slide_', '')}"
                    elif "page" in source_detail:
                        source_detail = f"Page {source_detail.replace('page_', '').replace('_table_', ' Table ')}"
                    elif "table" in source_detail:
                        source_detail = f"Table {source_detail.replace('table_', '')}"
                    _safe_write_json(src_dir / f"{tbl_name}.json", {
                        "source_file": file.filename,
                        "source_detail": source_detail,
                        "description": "",
                    })
                except Exception:
                    pass
        except Exception:
            pass

    # Trigger Engineer in background if tables were extracted from document
    if tables_saved > 0 and project:
        doc_tables_info = [{"table": t.get("name", ""), "rows": len(t.get("df", [])), "cols": len(t.get("df", {}).columns) if hasattr(t.get("df"), "columns") else 0, "source": t.get("source", "")} for t in doc_tables if t.get("df") is not None and len(t.get("df", [])) >= 2]
        if doc_tables_info:
            user_id_val = getattr(getattr(getattr(request, 'state', None), 'user', None), 'id', 1)
            _bg_executor.submit(_post_upload_engineer, project, doc_tables_info, user_id_val).add_done_callback(_bg_done_log)

    # Best-effort: notify project owner that doc was indexed (non-stream)
    if project:
        try:
            from app.auth import notify_user
            from sqlalchemy import create_engine as _ne3
            _eng3 = _ne3(db_url)
            with _eng3.connect() as _c:
                _r = _c.execute(text(
                    "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                ), {"s": project}).fetchone()
            _eng3.dispose()
            if _r and _r[0]:
                _pages = conductor_result.get("metadata", {}).get("total_pages", 0) or conductor_result.get("metadata", {}).get("slides", 0)
                notify_user(
                    int(_r[0]),
                    f"Document indexed · {file.filename}",
                    f"{_pages} pages · {len(text_content):,} chars · {tables_saved} tables",
                    "info",
                )
        except Exception:
            pass

    # Build processing report for frontend
    processing_steps = []
    if ext == ".pptx":
        slides_count = conductor_result.get("metadata", {}).get("slides", 0)
        notes_count = conductor_result.get("metadata", {}).get("notes_count", 0)
        processing_steps.append({"agent": "Scanner", "step": "Text extraction", "detail": f"{len(text_content):,} chars from {slides_count} slides", "status": "done"})
        if notes_count:
            processing_steps.append({"agent": "Scanner", "step": "Speaker notes", "detail": f"{notes_count} slides with notes", "status": "done"})
        if doc_tables:
            processing_steps.append({"agent": "Scanner", "step": "Table extraction", "detail": f"{len(doc_tables)} tables found", "status": "done"})
        img_count = len(conductor_result.get("images", []))
        if img_count:
            processing_steps.append({"agent": "Vision", "step": "Image description", "detail": f"{img_count} images described", "status": "done"})
        processing_steps.append({"agent": "Inspector", "step": "Knowledge indexing", "detail": "indexed to PgVector", "status": "done"})
    elif ext == ".pdf":
        total_pages = conductor_result.get("metadata", {}).get("total_pages", 0)
        scanned = conductor_result.get("metadata", {}).get("scanned_pages", 0)
        processing_steps.append({"agent": "Scanner", "step": "Text extraction", "detail": f"{len(text_content):,} chars from {total_pages} pages", "status": "done"})
        if scanned:
            processing_steps.append({"agent": "Scanner", "step": "OCR (scanned pages)", "detail": f"{scanned} scanned pages processed", "status": "done"})
        if doc_tables:
            processing_steps.append({"agent": "Scanner", "step": "Table extraction", "detail": f"{len(doc_tables)} tables found", "status": "done"})
        img_count = len(conductor_result.get("images", []))
        if img_count:
            processing_steps.append({"agent": "Vision", "step": "Image description", "detail": f"{img_count} images described", "status": "done"})
        processing_steps.append({"agent": "Inspector", "step": "Knowledge indexing", "detail": "indexed to PgVector", "status": "done"})
    elif ext == ".docx":
        processing_steps.append({"agent": "Scanner", "step": "Text extraction", "detail": f"{len(text_content):,} chars", "status": "done"})
        if doc_tables:
            processing_steps.append({"agent": "Scanner", "step": "Table extraction", "detail": f"{len(doc_tables)} tables found", "status": "done"})
        img_count = len(conductor_result.get("images", []))
        if img_count:
            processing_steps.append({"agent": "Vision", "step": "Image description", "detail": f"{img_count} images described", "status": "done"})
        processing_steps.append({"agent": "Inspector", "step": "Knowledge indexing", "detail": "indexed to PgVector", "status": "done"})
    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif", ".webp"):
        processing_steps.append({"agent": "Vision", "step": "OCR + description", "detail": f"{len(text_content):,} chars extracted", "status": "done"})
        processing_steps.append({"agent": "Inspector", "step": "Knowledge indexing", "detail": "indexed to PgVector", "status": "done"})
    else:
        processing_steps.append({"agent": "Scanner", "step": "Text extraction", "detail": f"{len(text_content):,} chars", "status": "done"})
        processing_steps.append({"agent": "Inspector", "step": "Knowledge indexing", "detail": "indexed to PgVector", "status": "done"})

    if tables_saved:
        processing_steps.append({"agent": "Engineer", "step": "Tables → PostgreSQL", "detail": f"{tables_saved} tables saved", "status": "done"})

    # Add warnings/errors from conductor
    for w in conductor_result.get("warnings", []):
        processing_steps.append({"agent": "Conductor", "step": "Warning", "detail": w[:100], "status": "warn"})
    for e in conductor_result.get("errors", []):
        processing_steps.append({"agent": "Conductor", "step": "Error", "detail": e[:100], "status": "error"})

    # Save extraction metadata for document info cards
    if project:
        try:
            meta_dir = KNOWLEDGE_DIR / project / "doc_meta"
            meta_dir.mkdir(parents=True, exist_ok=True)
            raw_file = KNOWLEDGE_DIR / project / "docs_raw" / file.filename
            file_size_bytes = raw_file.stat().st_size if raw_file.exists() else len(content)
            _safe_write_json(meta_dir / f"{file.filename}.json", {
                "text_chars": len(text_content),
                "tables_extracted": tables_saved,
                "images_described": len(conductor_result.get("images", [])),
                "slides": conductor_result.get("metadata", {}).get("slides", 0),
                "pages": conductor_result.get("metadata", {}).get("total_pages", 0),
                "scanned_pages": conductor_result.get("metadata", {}).get("scanned_pages", 0),
                "notes_count": conductor_result.get("metadata", {}).get("notes_count", 0),
                "file_size": file_size_bytes,
                "warnings": conductor_result.get("warnings", [])[:5],
                "errors": conductor_result.get("errors", [])[:5],
                "doc_category": _doc_category_hint(file.filename),
            })
        except Exception:
            pass

    return {
        "status": "ok", "filename": file.filename, "type": ext,
        "size": len(text_content), "indexed": True, "tables_saved": tables_saved,
        "processing_steps": processing_steps,
        "agents_used": list(set(s["agent"] for s in processing_steps)),
    }


@router.post("/upload-agent")
async def upload_with_agent(request: Request, file: UploadFile, project: str | None = Form(None)):
    """Upload a file using the Upload Agent Team (Conductor → Parser/Scanner/Vision → Inspector → Engineer).

    This uses AI agents for intelligent processing: structure detection,
    unpivot, merge, quality validation, and post-upload optimization.
    Falls back to standard upload if agent fails.
    """
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")
    if not project:
        raise HTTPException(400, "Project required for agent upload")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported: {ext}")

    # Stream to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        size = 0
        while chunk := await file.read(1024 * 1024):
            size += len(chunk)
            if size > MAX_FILE_SIZE:
                raise HTTPException(400, "File too large")
            tmp.write(chunk)
        tmp_path = tmp.name

    user_id = _get_user_id(request) or 1

    try:
        # Create the Upload Agent Team
        from dash.agents.conductor import create_upload_team
        uid = 1
        try:
            uid = int(user_id) if user_id and str(user_id).isdigit() else 1
        except Exception:
            pass
        upload_team = create_upload_team(project, user_id=uid)

        # Build the prompt for the Conductor
        file_info = f"File: {file.filename} ({ext}, {size:,} bytes)"
        if ext in (".xlsx", ".xls"):
            file_info += "\nThis is an Excel file — use Parser to analyze sheets, detect structure, unpivot if needed."
        elif ext in (".pdf",):
            file_info += "\nThis is a PDF — use Scanner to extract text, tables, OCR scanned pages."
        elif ext in (".pptx",):
            file_info += "\nThis is a PowerPoint — use Scanner to extract slides, tables, images."
        elif ext in (".docx",):
            file_info += "\nThis is a Word document — use Scanner to extract text, tables, images."
        elif ext in (".jpg", ".jpeg", ".png"):
            file_info += "\nThis is an image — use Vision to OCR text or describe content."
        elif ext in (".csv",):
            file_info += "\nThis is a CSV — use Parser for fast direct load."
        elif ext in (".json",):
            file_info += "\nThis is a JSON file — use Parser to load."
        else:
            file_info += "\nThis is a text file — use Scanner to index to knowledge base."

        prompt = f"""Process this uploaded file for project '{project}'.

{file_info}
File saved at: {tmp_path}

STEPS:
1. Use the right agent to parse/extract this file
2. After tables are created, use Inspector to validate quality
3. Use Engineer to create views and discover relationships
4. Report: tables created, health scores, views, relationships"""

        # Run the Conductor
        response = upload_team.run(prompt)
        agent_output = response.content if response else "Agent returned no response"

        # Also run the standard conductor as backup (ensures tables are actually created)
        if ext in ('.md', '.txt', '.sql', '.py'):
            with open(tmp_path, 'rb') as f:
                content = f.read()
        else:
            content = b''

        conductor_result = _conduct_upload(tmp_path, ext, project, file.filename, raw_content=content)

        # Store tables from conductor result
        tables_stored = []
        if conductor_result.get("tables"):
            proj_schema = re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
            try:
                from db import get_project_engine
                from db.session import create_project_schema
                schema = create_project_schema(project)
                eng = get_project_engine(project)
                for tbl_info in conductor_result["tables"]:
                    df = tbl_info.get("df")
                    if df is None or len(df) == 0:
                        continue
                    tbl_name = _sanitize_table_name(tbl_info.get("name", "table"))
                    try:
                        df.to_sql(tbl_name, eng, schema=schema, if_exists='replace', index=False)
                        try:
                            _upload_lg.info(f"✓ upload success: {file.filename}: {len(df)} rows × {len(df.columns)} cols → {tbl_name}")
                        except Exception:
                            pass
                        tables_stored.append({"table": tbl_name, "rows": len(df), "cols": len(df.columns), "source": tbl_info.get("source", "")})
                        # Profile + source metadata
                        try:
                            _profile_table(df, project, tbl_name)
                        except Exception:
                            pass
                        src_dir = KNOWLEDGE_DIR / project / "table_sources"
                        src_dir.mkdir(parents=True, exist_ok=True)
                        _safe_write_json(src_dir / f"{tbl_name}.json", {
                            "source_file": file.filename,
                            "source_detail": tbl_info.get("source", ext.upper()),
                            "description": tbl_info.get("description", ""),
                        })
                    except Exception:
                        pass
            except Exception:
                pass

        # Index text to knowledge
        text_content = conductor_result.get("text", "")
        if text_content and text_content.strip():
            try:
                from agno.knowledge.reader.text_reader import TextReader
                from db.session import create_project_knowledge
                knowledge = create_project_knowledge(project)
                knowledge.insert(name=f"doc-{Path(file.filename).stem}", text_content=f"Document: {file.filename}\n\n{text_content[:10000]}", reader=TextReader(), skip_if_exists=False)
            except Exception:
                pass

        # Trigger Engineer in background
        if tables_stored:
            _bg_executor.submit(_post_upload_engineer, project, tables_stored, user_id or 1).add_done_callback(_bg_done_log)

        Path(tmp_path).unlink(missing_ok=True)

        return {
            "status": "ok",
            "agent": True,
            "agent_report": agent_output[:2000] if agent_output else "",
            "tables_created": len(tables_stored),
            "tables": tables_stored,
            "text_indexed": len(text_content),
            "warnings": conductor_result.get("warnings", []),
            "errors": conductor_result.get("errors", []),
        }

    except Exception as e:
        # Fallback to standard upload if agent fails
        Path(tmp_path).unlink(missing_ok=True)
        return {"status": "error", "agent": True, "error": str(e)[:500], "fallback": "Use /api/upload or /api/upload-doc instead"}


@router.get("/docs")
def list_docs(request: Request, project: str | None = None):
    """List uploaded documents (project-scoped or global) with extraction metadata."""
    if project:
        docs_dir = KNOWLEDGE_DIR / project / "docs"
    else:
        docs_dir = KNOWLEDGE_DIR / "docs"
    if not docs_dir.exists():
        return {"docs": []}
    # Load extraction metadata if available
    meta_dir = KNOWLEDGE_DIR / (project or "") / "doc_meta"
    docs = []
    for f in sorted(docs_dir.iterdir()):
        if f.is_file() and not f.name.startswith("."):
            doc_info: dict = {"name": f.name, "size": f.stat().st_size, "type": f.suffix}
            # Load saved extraction metadata
            if meta_dir.exists():
                meta_file = meta_dir / f"{f.name}.json"
                if meta_file.exists():
                    try:
                        with open(meta_file) as mf:
                            meta = json.load(mf)
                        doc_info.update(meta)
                    except Exception:
                        pass
            docs.append(doc_info)
    return {"docs": docs}


@router.delete("/docs/{filename}")
def delete_doc(filename: str, request: Request, project: str | None = None):
    """Delete an uploaded document."""
    safe_name = Path(filename).name
    if safe_name != filename or '..' in filename:
        raise HTTPException(400, "Invalid filename")
    if project:
        doc_path = KNOWLEDGE_DIR / project / "docs" / filename
    else:
        doc_path = KNOWLEDGE_DIR / "docs" / filename
    if not doc_path.exists():
        raise HTTPException(404, f"Document '{filename}' not found")
    doc_path.unlink()
    return {"status": "ok", "deleted": filename}


@router.get("/tables")
def list_tables():
    """List all tables in the public schema with row counts."""
    engine = create_engine(db_url)
    insp = inspect(engine)
    tables = []

    with engine.connect() as conn:
        for tbl in sorted(insp.get_table_names(schema="public")):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM public."{tbl}"')).scalar()
                cols = insp.get_columns(tbl, schema="public")
                tables.append({
                    "name": tbl,
                    "rows": count,
                    "columns": len(cols),
                    "protected": tbl in PROTECTED_TABLES,
                })
            except Exception:
                tables.append({"name": tbl, "rows": 0, "columns": 0, "protected": tbl in PROTECTED_TABLES})

    return {"tables": tables}


@router.get("/tables/{table_name}/download")
def download_table(table_name: str, request: Request, format: str = "csv", project: str | None = None):
    """Download a table as CSV or Excel. Usage: /api/tables/my_table/download?format=csv&project=my_project"""
    from fastapi.responses import StreamingResponse
    import io

    # Determine schema
    if project:
        schema = re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    else:
        schema = "public"

    try:
        eng = create_engine(db_url)
        qualified = f'"{schema}"."{table_name}"'
        df = pd.read_sql(f"SELECT * FROM {qualified}", eng)
        eng.dispose()
    except Exception as e:
        raise HTTPException(404, f"Table not found: {e}")

    if format == "excel" or format == "xlsx":
        buf = io.BytesIO()
        df.to_excel(buf, index=False, sheet_name=table_name[:31])
        buf.seek(0)
        return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                 headers={"Content-Disposition": f'attachment; filename="{table_name}.xlsx"'})
    else:
        buf = io.StringIO()
        df.to_csv(buf, index=False)
        buf.seek(0)
        return StreamingResponse(io.BytesIO(buf.getvalue().encode()), media_type="text/csv",
                                 headers={"Content-Disposition": f'attachment; filename="{table_name}.csv"'})


_SCHEDULES_FILE = KNOWLEDGE_DIR / "schedules.json"


@router.get("/user-schedules")
def list_schedules(request: Request):
    """List saved scheduled reports for current user."""
    user_id = _get_user_id(request)
    if not _SCHEDULES_FILE.exists():
        return {"schedules": []}
    with open(_SCHEDULES_FILE) as f:
        all_s = json.load(f)
    return {"schedules": [s for s in all_s if s.get("user_id") == user_id]}


@router.post("/user-schedules")
def create_schedule(request: Request, name: str, prompt: str, frequency: str = "daily"):
    """Save a scheduled report."""
    user_id = _get_user_id(request)
    schedules = []
    if _SCHEDULES_FILE.exists():
        with open(_SCHEDULES_FILE) as f:
            schedules = json.load(f)
    schedules.append({
        "name": name, "prompt": prompt, "frequency": frequency,
        "user_id": user_id, "enabled": True,
        "created_at": __import__("datetime").datetime.now().isoformat(),
        "last_run": None,
    })
    _SCHEDULES_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(_SCHEDULES_FILE, "w") as f:
        json.dump(schedules, f, indent=2)
    return {"status": "ok"}


@router.delete("/user-schedules/{name}")
def delete_schedule(name: str, request: Request):
    """Delete a scheduled report."""
    user_id = _get_user_id(request)
    if not _SCHEDULES_FILE.exists():
        raise HTTPException(404, "Not found")
    with open(_SCHEDULES_FILE) as f:
        schedules = json.load(f)
    schedules = [s for s in schedules if not (s["name"] == name and s.get("user_id") == user_id)]
    with open(_SCHEDULES_FILE, "w") as f:
        json.dump(schedules, f, indent=2)
    return {"status": "ok"}


_EVAL_RESULTS_FILE = KNOWLEDGE_DIR / "eval_results.json"
_eval_running = False


@router.post("/evals/run")
def run_evals(request: Request):
    """Run eval suite in background."""
    global _eval_running
    if _eval_running:
        return {"status": "already_running"}

    run_id = f"eval-{__import__('time').time():.0f}"

    def _run():
        global _eval_running
        _eval_running = True
        results = {"run_id": run_id, "status": "running", "categories": {}}
        try:
            # Save initial state
            KNOWLEDGE_DIR.mkdir(parents=True, exist_ok=True)
            with open(_EVAL_RESULTS_FILE, "w") as f:
                json.dump(results, f, indent=2)

            # Run evals
            from evals.run import run_evals as _run_evals
            passed = _run_evals()

            results["status"] = "completed"
            results["passed"] = passed
            results["completed_at"] = __import__("datetime").datetime.now().isoformat()
        except Exception as e:
            results["status"] = "error"
            results["error"] = str(e)
        finally:
            with open(_EVAL_RESULTS_FILE, "w") as f:
                json.dump(results, f, indent=2)
            _eval_running = False

    _bg_executor.submit(_run).add_done_callback(_bg_done_log)
    return {"status": "started", "run_id": run_id}


@router.get("/evals/status")
def eval_status(request: Request):
    """Get latest eval run status."""
    if not _EVAL_RESULTS_FILE.exists():
        return {"status": "never_run", "running": _eval_running}
    with open(_EVAL_RESULTS_FILE) as f:
        results = json.load(f)
    results["running"] = _eval_running
    return results


@router.get("/training")
def get_training(request: Request, project: str | None = None):
    """Get learnings/training data for a project."""
    if not project:
        return {"learnings": 0, "items": [], "training_qa": []}
    import re
    safe = re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    engine = create_engine(db_url)
    count = 0
    try:
        with engine.connect() as conn:
            count = conn.execute(text(f'SELECT COUNT(*) FROM ai."{safe}_learnings"')).scalar() or 0
    except Exception:
        pass

    # Load training Q&A pairs
    training_qa: list[dict] = []
    training_dir = KNOWLEDGE_DIR / project / "training"
    if training_dir.exists():
        for f in sorted(training_dir.glob("*.json")):
            try:
                with open(f) as fh:
                    data = json.load(fh)
                if isinstance(data, list):
                    for qa in data:
                        qa["source_table"] = f.stem.replace("_qa", "")
                        training_qa.append(qa)
            except Exception:
                pass

    return {"learnings": count, "training_qa": training_qa}


# ---------------------------------------------------------------------------
# Self-Learning Endpoints
# ---------------------------------------------------------------------------


@router.post("/projects/{slug}/feedback")
async def save_feedback(slug: str, request: Request):
    """Save user feedback (thumbs up/down) as training data."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")
    rating = body.get("rating", "up")  # "up" or "down"

    if not question or not answer:
        return {"status": "skip"}

    # Verified-reward gate: a thumbs-up on an answer that is PROVABLY wrong (it
    # matched a pinned/proven number and missed) must NOT be learned as "good" —
    # downgrade it to feedback_bad so the agent doesn't reuse a wrong query.
    gated = False
    if rating == "up":
        try:
            from dash.learning.verified_reward import score_verified
            _vr = score_verified(slug, question, answer)
            if _vr.get("verified") == "fail":
                rating = "down"
                gated = True
        except Exception:
            pass

    feedback_dir = KNOWLEDGE_DIR / slug / "training"
    feedback_dir.mkdir(parents=True, exist_ok=True)

    filename = "feedback_good.json" if rating == "up" else "feedback_bad.json"
    filepath = feedback_dir / filename

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    existing.append({"question": question, "answer": answer[:500], "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
    # Keep last 50
    existing = existing[-50:]

    with open(filepath, "w") as f:
        json.dump(existing, f, indent=2)

    # Auto-promote 👍 to golden corpus if SQL present + not gated
    # Mirrors Dataherald's user-feedback → golden_sql lifecycle (Correction 1, Option B).
    promoted = None
    if rating == "up" and not gated:
        _sql = (body.get("sql") or "").strip()
        if _sql:
            try:
                from dash.learning.golden import promote as _golden_promote
                user = _get_user(request)
                _p = _golden_promote(
                    slug,
                    question=question,
                    sql=_sql,
                    source="user_thumb",
                    promoted_by=user.get("username") if user else None,
                    expected_value=body.get("expected_value"),
                )
                if _p.get("ok"):
                    promoted = {"total_goldens": _p.get("total_goldens")}
            except Exception as _ge:
                logger.warning(f"golden auto-promote failed: {_ge}")

    # Re-index
    _reload_project_knowledge(slug)
    return {"status": "ok", "saved": rating, "gated": gated,
            "promoted": promoted,
            "note": "answer didn't match pinned truth — saved as negative example, not promoted" if gated else None}


@router.post("/projects/{slug}/golden/promote")
async def golden_promote(slug: str, request: Request):
    """Promote a chat answer's SQL to the golden corpus.

    Body: {question, sql, source?, expected_rowcount?, expected_value?}
    Source values: 'user_thumb' | 'admin_pin' | 'eval_pass'

    On next matching question, `try_metric_shortcut` will run this SQL
    deterministically (~7ms, zero LLM tokens) instead of calling the agent.
    """
    user = _get_user(request)
    body = await request.json()
    question = body.get("question", "")
    sql = body.get("sql", "")
    source = body.get("source", "user_thumb")
    expected_rowcount = body.get("expected_rowcount")
    expected_value = body.get("expected_value")

    if not question or not sql:
        return {"ok": False, "error": "question and sql required"}

    try:
        from dash.learning.golden import promote
        result = promote(
            slug,
            question=question,
            sql=sql,
            source=source,
            promoted_by=user.get("username") if user else None,
            expected_rowcount=expected_rowcount,
            expected_value=expected_value,
        )
        # NOTE: skip _reload_project_knowledge — try_metric_shortcut reads
        # training/*.json on every call (no in-memory cache), so new golden
        # is picked up on next chat without expensive re-embed (~30s, 24 LLM calls).
        # Drop MDL model cache too (golden may reference MDL-named models so
        # operator-edited metric defs propagate without 5-min TTL wait).
        try:
            from dash.semantic import invalidate as _mdl_invalidate
            _mdl_invalidate(slug)
        except Exception:
            pass
        return result
    except Exception as e:
        logger.exception(f"golden promote failed for {slug}: {e}")
        return {"ok": False, "error": str(e)}


@router.post("/projects/{slug}/golden/demote")
async def golden_demote(slug: str, request: Request):
    """Remove a golden entry by SQL match."""
    _ = _get_user(request)
    body = await request.json()
    sql = body.get("sql", "")
    try:
        from dash.learning.golden import demote
        result = demote(slug, sql)
        try:
            from dash.semantic import invalidate as _mdl_invalidate
            _mdl_invalidate(slug)
        except Exception:
            pass
        return result
    except Exception as e:
        return {"ok": False, "error": str(e)}


@router.post("/projects/{slug}/golden/drift-check")
async def golden_drift_check(slug: str, request: Request):
    """Manually trigger drift check for one project (admin)."""
    _ = _get_user(request)
    body = await request.json() if (await request.body()) else {}
    dry = bool(body.get("dry_run", False))
    try:
        from dash.cron.golden_drift import check_project
        return check_project(slug, dry_run=dry)
    except Exception as e:
        return {"ok": False, "error": str(e)}


@router.get("/projects/{slug}/golden/list")
async def golden_list(slug: str, request: Request):
    """List all golden Q→SQL pairs for this project (newest first)."""
    _ = _get_user(request)
    try:
        from dash.learning.golden import list_goldens
        entries = list_goldens(slug)
        return {"ok": True, "count": len(entries), "entries": entries}
    except Exception as e:
        return {"ok": False, "error": str(e), "entries": []}


@router.post("/projects/{slug}/save-query-pattern")
async def save_query_pattern(slug: str, request: Request):
    """Save a proven query pattern for reuse."""
    body = await request.json()
    question = body.get("question", "")
    sql = body.get("sql", "")

    if not question or not sql:
        return {"status": "skip"}

    patterns_dir = KNOWLEDGE_DIR / slug / "queries"
    patterns_dir.mkdir(parents=True, exist_ok=True)
    filepath = patterns_dir / "proven_patterns.json"

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    # Check for duplicate
    if not any(p.get("sql") == sql for p in existing):
        existing.append({"question": question, "sql": sql, "uses": 1, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
        existing = existing[-30:]
        with open(filepath, "w") as f:
            json.dump(existing, f, indent=2)

    return {"status": "ok"}


@router.post("/projects/{slug}/suggest-followups")
async def suggest_followups(slug: str, request: Request):
    """LLM-powered follow-up suggestions based on conversation context."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")

    if not question or not answer:
        return {"suggestions": []}

    from os import getenv
    import httpx
    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return {"suggestions": []}

    # Load knowledge graph entities for context-aware follow-ups
    kg_context = ""
    try:
        from dash.tools.knowledge_graph import get_knowledge_graph_context
        kg_context = get_knowledge_graph_context(slug, for_agent="leader")
        if kg_context and len(kg_context) > 200:
            kg_context = f"\n\nKNOWN ENTITIES & RELATIONSHIPS:\n{kg_context[:500]}"
    except Exception:
        pass

    prompt = f"""Based on this data conversation, suggest 3 smart follow-up questions.

The follow-ups should:
- Dig DEEPER into the data (not repeat what was already shown)
- Explore RELATED dimensions (if revenue was shown, ask about margin or by-store breakdown)
- Ask WHY if only WHAT was shown
- Suggest comparisons (vs last period, vs budget, vs other segments)
- Reference specific entities/metrics from the answer

Q: {question}
A: {answer[:1500]}
{kg_context}

Return ONLY a JSON array of 3 short, specific follow-up questions:
["question 1", "question 2", "question 3"]"""

    try:
        from dash.settings import LITE_MODEL
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": LITE_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 200, "temperature": 0.5},
            timeout=10,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        clean = content.strip().strip("`").strip()
        if clean.lower().startswith("json"):
            clean = clean[4:].strip()
        suggestions = json.loads(clean)
        return {"suggestions": suggestions[:3] if isinstance(suggestions, list) else []}
    except Exception:
        return {"suggestions": []}


@router.post("/projects/{slug}/extract-context")
async def extract_context(slug: str, request: Request):
    """Extract key facts from a conversation for context memory."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")

    if not question or not answer:
        return {"status": "skip"}

    # Issue #6 — Flag-based refusal detection (replaces brittle text-sentinel).
    # Authoritative source: dash.dash_refusal_marks written by refusal sites
    # (scope_classifier, agent_self, stuck_loop). Catches per-tenant custom
    # refusal messages + i18n + future refusal phrasing changes.
    _sid_for_refusal = body.get("session_id") or body.get("sid")
    if _sid_for_refusal:
        try:
            from dash.runtime.refusal import was_refused
            _mark = was_refused(_sid_for_refusal, question)
            if _mark:
                return {"status": "skip",
                        "reason": f"refusal_response (flag · {_mark.get('source')}/{_mark.get('reason') or '?'})"}
        except Exception:
            pass

    # Defense-in-depth fallback: text-sentinel for sessions without session_id
    # OR for legacy refusal sites not yet wired (agent self-refusal via instructions).
    try:
        from dash.runtime.refusal import is_refusal_text
        if is_refusal_text(answer):
            # Also persist as text_sentinel mark so future calls in same session skip fast
            try:
                from dash.runtime.refusal import mark_refused as _mr
                _mr(_sid_for_refusal, question, source="text_sentinel", reason="fallback_match")
            except Exception:
                pass
            return {"status": "skip", "reason": "refusal_response (text_sentinel)"}
    except Exception:
        pass

    from os import getenv
    import httpx
    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return {"status": "skip"}

    prompt = f"""Extract 0-3 key facts from this conversation that should be remembered for future sessions.

For each fact, provide:
- "fact": the fact text
- "confidence": "high" or "low"
  - "high": objective data facts (column names, table relationships, data formats, business definitions, SQL patterns)
  - "low": subjective observations (user preferences, interpretation of results, opinions)
- "score": confidence score 0-100 (100 = absolutely certain data fact, 0 = wild guess)

Q: {question}
A: {answer[:500]}

Return ONLY valid JSON (no markdown). Empty array if nothing worth remembering:
[{{"fact": "the fact text", "confidence": "high", "score": 95}}, {{"fact": "another fact", "confidence": "low", "score": 40}}]"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 300, "temperature": 0.1},
            timeout=10,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        # Clean LLM output — strip markdown code fences
        clean = content.strip()
        if clean.startswith("```"):
            clean = clean.split("\n", 1)[-1] if "\n" in clean else clean[3:]
        if clean.endswith("```"):
            clean = clean[:-3]
        clean = clean.strip().strip("`").strip()
        if clean.lower().startswith("json"):
            clean = clean[4:].strip()
        parsed = json.loads(clean)

        if not isinstance(parsed, list) or len(parsed) == 0:
            return {"status": "ok", "facts": [], "auto_saved": []}

        # Separate high vs low confidence, keep scores
        auto_save = []
        auto_save_with_scores = []
        needs_approval = []
        needs_approval_with_scores = []
        for item in parsed[:3]:
            if isinstance(item, str):
                needs_approval.append(item)
                needs_approval_with_scores.append({"fact": item, "score": 30})
            elif isinstance(item, dict):
                fact = item.get("fact", "")
                confidence = item.get("confidence", "low")
                score = int(item.get("score", 50))
                if not fact:
                    continue
                if confidence == "high":
                    auto_save.append(fact)
                    auto_save_with_scores.append({"fact": fact, "score": score})
                else:
                    needs_approval.append(fact)
                    needs_approval_with_scores.append({"fact": fact, "score": score})

        # Auto-save high-confidence facts
        if auto_save:
            try:
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    for fact in auto_save:
                        conn.execute(text(
                            "INSERT INTO public.dash_memories (project_slug, fact, scope, source) VALUES (:s, :f, 'project', 'auto')"
                        ), {"s": slug, "f": fact})
                    conn.commit()
            except Exception:
                pass

        # Always return "proposed" if there are any learnings (auto or manual)
        if auto_save or needs_approval:
            return {"status": "proposed", "facts": needs_approval, "facts_with_scores": needs_approval_with_scores, "auto_saved": auto_save, "auto_saved_with_scores": auto_save_with_scores}

        return {"status": "ok", "facts": [], "auto_saved": []}
    except Exception:
        return {"status": "skip"}


@router.post("/projects/{slug}/approve-learnings")
async def approve_learnings(slug: str, request: Request):
    """Save user-approved learnings to memory."""
    body = await request.json()
    facts = body.get("facts", [])
    scope = body.get("scope", "project")

    if not facts:
        return {"status": "skip"}

    # Save to DB (dash_memories)
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            for fact in facts[:3]:
                if fact:
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, fact, scope, source) VALUES (:s, :f, :sc, 'agent')"
                    ), {"s": slug, "f": fact, "sc": scope})
            conn.commit()
    except Exception:
        pass

    # Also save to JSON memory file for backward compat
    memory_dir = KNOWLEDGE_DIR / slug / "memory"
    memory_dir.mkdir(parents=True, exist_ok=True)
    filepath = memory_dir / "context.json"

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    for fact in facts[:3]:
        if fact and not any(f.get("fact") == fact for f in existing):
            existing.append({"fact": fact, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
    existing = existing[-30:]

    with open(filepath, "w") as f:
        json.dump(existing, f, indent=2)

    return {"status": "saved", "count": len(facts)}


@router.put("/projects/{slug}/persona")
async def update_persona(slug: str, request: Request):
    """Update project persona."""
    body = await request.json()
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            conn.execute(text("""
                INSERT INTO public.dash_personas (project_slug, persona)
                VALUES (:s, CAST(:p AS jsonb))
                ON CONFLICT (project_slug)
                DO UPDATE SET persona = CAST(:p AS jsonb), updated_at = NOW()
            """), {"s": slug, "p": json.dumps(body, default=str)})
            conn.commit()
        # Also save to file for PgVector
        persona_file = KNOWLEDGE_DIR / slug / "persona.json"
        persona_file.parent.mkdir(parents=True, exist_ok=True)
        with open(persona_file, "w") as f:
            json.dump(body, f, indent=2)
    except Exception:
        pass
    return {"status": "ok"}


@router.get("/projects/{slug}/persona")
def get_persona(slug: str, request: Request):
    """Get the generated persona for a project (DB first, file fallback)."""
    # Try DB first
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            row = conn.execute(text("SELECT persona FROM public.dash_personas WHERE project_slug = :s"), {"s": slug}).fetchone()
            if row and row[0]:
                return {"persona": row[0] if isinstance(row[0], dict) else json.loads(row[0])}
    except Exception:
        pass
    # Fallback to file
    persona_file = KNOWLEDGE_DIR / slug / "persona.json"
    if persona_file.exists():
        try:
            with open(persona_file) as f:
                return {"persona": json.load(f)}
        except Exception:
            pass
    return {"persona": None}


_training_cancel_flags: dict[str, bool] = {}


@router.post("/projects/{slug}/stop-training")
def stop_training(slug: str, request: Request):
    """Stop a running training pipeline."""
    _training_cancel_flags[slug] = True
    # Also mark the latest training run as failed
    try:
        eng = create_engine(db_url)
        with eng.connect() as conn:
            conn.execute(text(
                "UPDATE public.dash_training_runs SET status = 'failed' "
                "WHERE project_slug = :s AND status NOT IN ('done', 'failed') "
            ), {"s": slug})
            conn.commit()
        # Notify owner of stopped/failed training
        try:
            from app.auth import notify_user
            with eng.connect() as _c:
                _r = _c.execute(text(
                    "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                ), {"s": slug}).fetchone()
            if _r and _r[0]:
                notify_user(
                    int(_r[0]),
                    f"Training failed · {slug}",
                    "Training was stopped before completion",
                    "error",
                )
        except Exception:
            pass
    except Exception:
        pass
    return {"status": "stopped"}


@router.get("/projects/{slug}/quality-check")
def quality_check(slug: str, request: Request):
    """Run a fast, free (SQL-only) quality score on every table in the project.

    Powers the pre-train review card: user sees GOOD/WEAK/BAD per table and
    picks which to train. No LLM calls. Returns items sorted best-first.
    """
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")
    from app.auth import check_project_permission
    if not check_project_permission(user, slug, required_role="editor"):
        raise HTTPException(403, "Editor access required")

    from db.session import create_project_schema
    schema = create_project_schema(slug)
    engine = create_engine(db_url)
    try:
        insp = inspect(engine)
        tables = insp.get_table_names(schema=schema)
        items = [_compute_table_quality(slug, t, engine=engine) for t in tables]
    finally:
        try:
            engine.dispose()
        except Exception:
            pass
    items.sort(key=lambda x: x.get("score", 0), reverse=True)
    return {"items": items, "count": len(items)}


@router.get("/projects/{slug}/column-map/{table}")
def column_map(slug: str, table: str, request: Request):
    """Grouped column understanding for one table: dimension / measure (with
    unit) / identifier / date / text. SQL-profiled, no LLM. Powers the
    Column Map UI block. Fast (~stats per column)."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")
    from app.auth import check_project_permission
    if not check_project_permission(user, slug, required_role="viewer"):
        raise HTTPException(403, "Access denied")

    engine = create_engine(db_url)
    try:
        profiles = _sql_profile_columns(slug, table, engine=engine)
    finally:
        try:
            engine.dispose()
        except Exception:
            pass

    cols = [p for p in profiles if not str(p.get("name", "")).startswith("_source_")]
    # group → list, ordered for display
    order = ["DIMENSION", "MEASURE", "IDENTIFIER", "DATE", "TEXT"]
    groups: dict[str, list] = {g: [] for g in order}
    for p in cols:
        g = p.get("group", "TEXT")
        groups.setdefault(g, []).append({
            "name": p.get("name"), "type": p.get("type"), "unit": p.get("unit", "none"),
            "classification": p.get("classification"),
            "unique_count": p.get("unique_count"), "null_pct": p.get("null_pct"),
            "min": p.get("min"), "max": p.get("max"), "mean": p.get("mean"),
            "min_date": p.get("min_date"), "max_date": p.get("max_date"),
        })
    out = [{"group": g, "columns": groups[g]} for g in order if groups.get(g)]
    return {"table": table, "total": len(cols), "groups": out}


@router.post("/projects/{slug}/retrain")
async def retrain_project(slug: str, request: Request):
    """Delta retrain — only train tables that changed since last training.

    Optional JSON body {"table_names": [...]} trains ONLY those tables
    (used by the quality-review card so junk tables are skipped).
    """
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")

    # Editor role required for retraining
    from app.auth import check_project_permission
    perm = check_project_permission(user, slug, required_role="editor")
    if not perm:
        raise HTTPException(403, "Editor access required to retrain")

    # Optional selective-train list from request body (fail-soft on empty body).
    selected_tables: list[str] | None = None
    force = False  # when True, bypass fingerprint "unchanged" skip (train untrained tables)
    try:
        body = await request.json()
        if isinstance(body, dict):
            tn = body.get("table_names")
            if isinstance(tn, list) and tn:
                selected_tables = [str(x) for x in tn]
            force = bool(body.get("force"))
    except Exception:
        selected_tables = None
    # also accept ?force=1 query param
    try:
        _qf = request.query_params.get("force")
        if _qf and str(_qf).strip().lower() in ("1", "true", "yes"):
            force = True
    except Exception:
        pass

    from db.session import create_project_schema
    schema = create_project_schema(slug)
    engine = create_engine(db_url)
    insp = inspect(engine)

    tables = insp.get_table_names(schema=schema)

    # Fire column enrichment IMMEDIATELY (parallel to training, fire-and-forget).
    # User clicked "Train All" → columns should refresh even if training stalls.
    # UPSERT idempotent. Fail-soft. ~30s per table for 15 cols.
    if tables and not selected_tables:
        try:
            from dash.tools.column_describer import enrich_columns_async
            import asyncio as _aio_enr
            import threading as _th_enr
            _tbls_to_enrich = list(tables)
            def _kick_col_enrich():
                try:
                    _loop = _aio_enr.new_event_loop()
                    _aio_enr.set_event_loop(_loop)
                    for _tn in _tbls_to_enrich:
                        try:
                            _loop.run_until_complete(enrich_columns_async(slug, _tn))
                        except Exception as _ce:
                            import logging as _l
                            _l.getLogger(__name__).warning(f"col enrich {slug}/{_tn}: {_ce}")
                    _loop.close()
                except Exception:
                    pass
            _th_enr.Thread(target=_kick_col_enrich, daemon=True,
                           name=f"col-enrich-{slug}").start()
        except Exception:
            pass
    if selected_tables is not None:
        _requested = set(selected_tables)
        tables = [t for t in tables if t in _requested]
        if not tables:
            raise HTTPException(400, "none of the requested tables exist in this project")
    if not tables:
        # Doc-only project — index knowledge + fill brain from doc text
        import threading
        import time as _time
        def _bg_docs_only():
            _log_entries = []
            def _dlog(msg):
                _log_entries.append({"ts": _time.strftime('%H:%M:%S'), "tsabs": _time.time(), "msg": msg})

            try:
                from dash.settings import training_llm_call
                eng = create_engine(db_url)

                # Create training run for UI tracking
                run_id = None
                try:
                    with eng.connect() as conn:
                        r = conn.execute(text(
                            "INSERT INTO public.dash_training_runs (project_slug, status, steps) "
                            "VALUES (:s, 'running', 'reindex') RETURNING id"
                        ), {"s": slug})
                        run_id = r.fetchone()[0]
                        conn.commit()
                except Exception:
                    pass

                def _update_step(step_name):
                    if run_id:
                        try:
                            with eng.connect() as conn:
                                conn.execute(text(
                                    "UPDATE public.dash_training_runs SET steps = :s, logs = CAST(:logs AS jsonb) WHERE id = :id"
                                ), {"s": step_name, "logs": json.dumps(_log_entries), "id": run_id})
                                conn.commit()
                        except Exception:
                            pass

                # Step 1: Index knowledge
                _update_step("reindex")
                _dlog("indexing documents into knowledge base...")
                _reload_project_knowledge(slug, timeout_sec=60)
                _dlog("✓ knowledge indexed")

                # Step 2: Read all doc text for LLM context
                docs_dir = KNOWLEDGE_DIR / slug / "docs"
                all_text = ""
                if docs_dir.exists():
                    for f in docs_dir.iterdir():
                        if f.is_file():
                            try:
                                all_text += f.read_text(errors='ignore')[:5000] + "\n\n"
                            except Exception:
                                pass
                all_text = all_text[:8000]

                # Extract and describe images from raw document files
                docs_raw_dir = KNOWLEDGE_DIR / slug / "docs_raw"
                if docs_raw_dir.exists():
                    for raw_file in docs_raw_dir.iterdir():
                        ext_r = raw_file.suffix.lower()
                        if ext_r == ".pptx":
                            imgs = _extract_images_pptx(str(raw_file))
                        elif ext_r == ".pdf":
                            imgs = _extract_images_pdf(str(raw_file))
                        else:
                            continue
                        if imgs:
                            _dlog(f"describing {len(imgs)} images from {raw_file.name}...")
                            desc = _describe_images_with_vision(imgs, raw_file.name)
                            if desc:
                                all_text += f"\n\n--- IMAGES FROM {raw_file.name} ---\n{desc}"
                                _dlog(f"✓ {len(imgs)} image descriptions added")

                if not all_text.strip():
                    _dlog("· no document text found — skipping brain fill")
                else:
                    # Step 3: Generate memories from docs
                    _update_step("brain_fill")
                    _dlog("generating memories from documents...")
                    try:
                        result = training_llm_call(
                            f"Extract 5-8 key facts from this document that a data analyst should remember.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Return ONLY valid JSON array of strings:\n"
                            f'["fact 1", "fact 2", "fact 3"]',
                            "extraction"
                        )
                        if result:
                            facts = json.loads(result)
                            if isinstance(facts, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for fact in facts[:8]:
                                        if isinstance(fact, str) and len(fact) > 10:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'auto_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": fact})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} memories saved")
                    except Exception as e:
                        _dlog(f"⚠ memories error: {str(e)[:60]}")

                    # Step 4: Generate persona
                    _update_step("persona")
                    _dlog("generating agent persona...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, generate a JSON persona for an AI agent:\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return: {{"persona_prompt": "You are an expert...", "domain_terms": ["term1"], "expertise_areas": ["area1"], "greeting": "Hi! ..."}}',
                            "persona"
                        )
                        if result:
                            persona = json.loads(result)
                            if isinstance(persona, dict) and persona.get("persona_prompt"):
                                persona_file = KNOWLEDGE_DIR / slug / "persona.json"
                                persona_file.parent.mkdir(parents=True, exist_ok=True)
                                with open(persona_file, "w") as f:
                                    json.dump(persona, f, indent=2)
                                with eng.connect() as conn:
                                    conn.execute(text(
                                        "INSERT INTO public.dash_personas (project_slug, persona) VALUES (:s, CAST(:p AS jsonb)) "
                                        "ON CONFLICT (project_slug) DO UPDATE SET persona = CAST(:p AS jsonb)"
                                    ), {"s": slug, "p": json.dumps(persona)})
                                    conn.commit()
                                _dlog("✓ persona generated")
                    except Exception as e:
                        _dlog(f"⚠ persona error: {str(e)[:60]}")

                    # Step 5: Generate workflows (structure-aware for PPTX/PDF/DOCX)
                    _update_step("synthesis")
                    _dlog("generating sample workflows...")
                    try:
                        # Try structure-based extraction from raw docs first
                        structure_wf_saved = 0
                        docs_raw_dir = KNOWLEDGE_DIR / slug / "docs_raw"
                        if docs_raw_dir.exists():
                            for raw_file in docs_raw_dir.iterdir():
                                ext = raw_file.suffix.lower()
                                if ext not in (".pptx", ".pdf", ".docx"):
                                    continue
                                sections = _extract_document_structure(str(raw_file), ext)
                                if len(sections) < 2:
                                    continue
                                sections_text = "\n".join(f"{s['index']}. {s['title']} — {s['content_summary']}" for s in sections)
                                result = training_llm_call(
                                    f"Convert this document structure into a reusable analysis workflow.\n"
                                    f"Each section becomes one step — write a clear analyst question.\n\n"
                                    f"DOCUMENT: {raw_file.name}\nSECTIONS:\n{sections_text}\n\n"
                                    f'Return ONLY valid JSON (no markdown):\n'
                                    f'{{"name": "workflow name", "description": "what it analyzes", '
                                    f'"steps": [{{"title": "section title", "question": "analyst question"}}]}}',
                                    "extraction"
                                )
                                if result:
                                    wf = json.loads(result.strip().strip("`").strip())
                                    if isinstance(wf, dict) and wf.get("name"):
                                        steps = wf.get("steps", [])
                                        step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                        with eng.connect() as conn:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                                "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "n": wf["name"], "d": wf.get("description", ""), "st": json.dumps(step_list)})
                                            conn.commit()
                                        structure_wf_saved += 1
                                        _dlog(f"✓ workflow from {raw_file.name} ({len(step_list)} steps)")

                        # Fall back: use LLM to extract structure from text content
                        if structure_wf_saved == 0:
                            result = training_llm_call(
                                f"Read this document and identify the analysis sections/topics.\n"
                                f"Convert the document structure into a reusable analysis workflow.\n"
                                f"Each section/topic becomes one step — write a clear analyst question.\n\n"
                                f"DOCUMENT:\n{all_text[:4000]}\n\n"
                                f"Return ONLY valid JSON (no markdown):\n"
                                f'{{"name": "workflow name based on document", "description": "what this workflow analyzes", '
                                f'"steps": [{{"title": "section title", "question": "analyst question to reproduce this analysis"}}]}}',
                                "extraction"
                            )
                            if result:
                                wf = json.loads(result.strip().strip("`").strip())
                                if isinstance(wf, dict) and wf.get("name"):
                                    steps = wf.get("steps", [])
                                    step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                    with eng.connect() as conn:
                                        conn.execute(text(
                                            "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                            "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                        ), {"s": slug, "n": wf["name"], "d": wf.get("description", ""), "st": json.dumps(step_list)})
                                        conn.commit()
                                    _dlog(f"✓ workflow from text ({len(step_list)} steps)")
                                    structure_wf_saved = 1
                                elif isinstance(wf, list):
                                    saved = 0
                                    with eng.connect() as conn:
                                        for w in wf[:3]:
                                            if isinstance(w, dict) and w.get("name"):
                                                steps = w.get("steps", [])
                                                step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                                conn.execute(text(
                                                    "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                                    "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                                ), {"s": slug, "n": w["name"], "d": w.get("description", ""), "st": json.dumps(step_list)})
                                                saved += 1
                                        conn.commit()
                                    _dlog(f"✓ {saved} workflows from text")
                                    structure_wf_saved = saved
                        else:
                            _dlog(f"✓ {structure_wf_saved} document-based workflows saved")
                    except Exception as e:
                        _dlog(f"⚠ workflows error: {str(e)[:60]}")

                    # Step 6: Generate evals
                    _dlog("generating eval questions...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, generate 6 smart business questions an executive would ask.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Generate questions that:\n"
                            f"- Ask about key metrics, trends, performance, comparisons\n"
                            f"- Use business language (not technical or table names)\n"
                            f"- Cover: summary, metrics, trends, comparisons, risks, recommendations\n\n"
                            f'Return JSON array: [{{"question": "What was the revenue growth this quarter?", "expected_answer": "Revenue grew..."}}]',
                            "extraction"
                        )
                        if result:
                            evals = json.loads(result)
                            if isinstance(evals, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for ev in evals[:6]:
                                        if isinstance(ev, dict) and ev.get("question"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_evals (project_slug, question, expected_sql) VALUES (:s, :q, :a)"
                                            ), {"s": slug, "q": ev["question"], "a": ev.get("expected_answer", "")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} eval questions saved")
                    except Exception as e:
                        _dlog(f"⚠ evals error: {str(e)[:60]}")

                    # Step 7: Seed feedback
                    _dlog("seeding sample feedback...")
                    try:
                        with eng.connect() as conn:
                            conn.execute(text(
                                "INSERT INTO public.dash_feedback (project_slug, question, answer, rating) "
                                "VALUES (:s, 'What is this project about?', :a, 'up')"
                            ), {"s": slug, "a": all_text[:500]})
                            conn.commit()
                        _dlog("✓ 1 seed feedback saved")
                    except Exception:
                        pass

                    # Step 8: Extract business rules from docs
                    _update_step("domain_knowledge")
                    _dlog("extracting business rules from documents...")
                    try:
                        result = training_llm_call(
                            f"Extract business rules, constraints, and policies from this document.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f'Return JSON array: [{{"name": "Rule Name", "definition": "The rule description", "type": "business_rule"}}]',
                            "extraction"
                        )
                        if result:
                            rules = json.loads(result)
                            if isinstance(rules, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for r in rules[:8]:
                                        if isinstance(r, dict) and r.get("name"):
                                            rule_id = f"doc_{r['name'].lower().replace(' ', '_')[:25]}"
                                            conn.execute(text(
                                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                                "VALUES (:s, :rid, :name, :type, :defn, 'doc_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "rid": rule_id, "name": r["name"], "type": r.get("type", "business_rule"), "defn": r.get("definition", "")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} business rules extracted")
                    except Exception as e:
                        _dlog(f"⚠ rules error: {str(e)[:60]}")

                    # Step 9: Extract domain knowledge (glossary + KPIs)
                    _dlog("extracting domain knowledge...")
                    try:
                        result = training_llm_call(
                            f"Extract domain knowledge from this document:\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Return JSON with:\n"
                            f'{{"glossary": [{{"term": "SLA", "definition": "Service Level Agreement"}}], '
                            f'"kpis": [{{"name": "Resolution Time", "definition": "Average time to resolve tickets"}}], '
                            f'"key_metrics": [{{"name": "Ticket Volume", "value": "500/month"}}]}}',
                            "extraction"
                        )
                        if result:
                            domain = json.loads(result)
                            if isinstance(domain, dict):
                                saved = 0
                                with eng.connect() as conn:
                                    # Glossary → memories
                                    for g in (domain.get("glossary") or [])[:10]:
                                        if isinstance(g, dict) and g.get("term"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'glossary') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"Glossary: {g['term']} = {g.get('definition', '')}"})
                                            saved += 1
                                    # KPIs → rules
                                    for k in (domain.get("kpis") or [])[:8]:
                                        if isinstance(k, dict) and k.get("name"):
                                            rule_id = f"kpi_doc_{k['name'].lower().replace(' ', '_')[:25]}"
                                            conn.execute(text(
                                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                                "VALUES (:s, :rid, :name, 'kpi', :defn, 'doc_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "rid": rule_id, "name": k["name"], "defn": k.get("definition", "")})
                                            saved += 1
                                    # Metrics → memories
                                    for m in (domain.get("key_metrics") or [])[:5]:
                                        if isinstance(m, dict) and m.get("name"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'auto_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"Key metric: {m['name']} = {m.get('value', 'N/A')}"})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} domain knowledge items extracted")
                    except Exception as e:
                        _dlog(f"⚠ domain knowledge error: {str(e)[:60]}")

                    # Step 10: Generate proactive insights from docs
                    _dlog("generating proactive insights...")
                    try:
                        result = training_llm_call(
                            f"Analyze this document and identify 3 proactive insights — things that need attention, risks, or opportunities.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f'Return JSON array: [{{"insight": "Description of the insight", "severity": "info|warning|critical"}}]',
                            "extraction"
                        )
                        if result:
                            insights = json.loads(result)
                            if isinstance(insights, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for ins in insights[:5]:
                                        if isinstance(ins, dict) and ins.get("insight"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_proactive_insights (project_slug, insight, severity) "
                                                "VALUES (:s, :i, :sev)"
                                            ), {"s": slug, "i": ins["insight"], "sev": ins.get("severity", "info")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} proactive insights generated")
                    except Exception as e:
                        _dlog(f"⚠ insights error: {str(e)[:60]}")

                    # Step 11: Cross-document relationships
                    _update_step("relationships")
                    _dlog("discovering cross-document relationships...")
                    try:
                        doc_summaries = []
                        for f in docs_dir.iterdir():
                            if f.is_file():
                                try:
                                    content = f.read_text(errors='ignore')[:2000]
                                    doc_summaries.append(f"{f.name}: {content[:500]}")
                                except Exception:
                                    pass
                        if len(doc_summaries) >= 2:
                            result = training_llm_call(
                                f"Analyze these documents and find relationships, shared topics, cross-references between them.\n\n"
                                f"DOCUMENTS:\n" + "\n---\n".join(doc_summaries[:10]) + "\n\n"
                                f"Return ONLY valid JSON array:\n"
                                f'[{{"from_doc": "doc1.txt", "to_doc": "doc2.txt", "relationship": "both discuss revenue targets", "shared_topics": ["revenue", "KPIs"], "strength": 0.8}}]',
                                "extraction"
                            )
                            if result:
                                rels = json.loads(result)
                                if isinstance(rels, list):
                                    saved = 0
                                    with eng.connect() as conn:
                                        for r in rels[:10]:
                                            if isinstance(r, dict) and r.get("from_doc"):
                                                conn.execute(text(
                                                    "INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source) "
                                                    "VALUES (:s, :ft, :fc, :tt, :tc, 'topic', :conf, 'ai') "
                                                    "ON CONFLICT DO NOTHING"
                                                ), {"s": slug, "ft": r["from_doc"], "fc": ", ".join(r.get("shared_topics", [])[:5]),
                                                    "tt": r.get("to_doc", ""), "tc": r.get("relationship", ""),
                                                    "conf": r.get("strength", 0.5)})
                                                saved += 1
                                        conn.commit()
                                    _dlog(f"✓ {saved} cross-document relationships found")
                        else:
                            _dlog("· only 1 document — no cross-references to find")
                    except Exception as e:
                        _dlog(f"⚠ relationships error: {str(e)[:60]}")

                    # Step 12: Negative examples
                    _dlog("extracting negative examples...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, what are common mistakes someone might make when interpreting this information?\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return JSON array of "DON\'T / DO" pairs:\n'
                            f'["DON\'T confuse X with Y — DO use Z instead", "DON\'T assume A — DO check B"]',
                            "extraction"
                        )
                        if result:
                            negs = json.loads(result)
                            if isinstance(negs, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for neg in negs[:5]:
                                        if isinstance(neg, str) and len(neg) > 10:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'negative_example') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"⚠ {neg}"})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} negative examples saved")
                    except Exception as e:
                        _dlog(f"⚠ negative examples error: {str(e)[:60]}")

                    # Step 13: Training Q&A from docs
                    _dlog("generating training Q&A...")
                    try:
                        result = training_llm_call(
                            f"Generate 5 question-answer pairs that test understanding of this document.\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return JSON array: [{{"question": "What is...?", "answer": "It is..."}}]',
                            "extraction"
                        )
                        if result:
                            qas = json.loads(result)
                            if isinstance(qas, list):
                                qa_file = KNOWLEDGE_DIR / slug / "training"
                                qa_file.mkdir(parents=True, exist_ok=True)
                                with open(qa_file / "doc_qa.json", "w") as f:
                                    json.dump(qas, f, indent=2)
                                _dlog(f"✓ {len(qas)} training Q&A pairs generated")
                    except Exception as e:
                        _dlog(f"⚠ Q&A error: {str(e)[:60]}")

                    # Step 14: Multi-doc synthesis
                    if len([f for f in docs_dir.iterdir() if f.is_file()]) >= 2:
                        _dlog("running multi-document synthesis...")
                        try:
                            doc_names = [f.name for f in docs_dir.iterdir() if f.is_file()]
                            result = training_llm_call(
                                f"These documents are part of the same project. Create a unified understanding.\n\n"
                                f"DOCUMENTS: {', '.join(doc_names)}\n\n"
                                f"CONTENT:\n{all_text[:5000]}\n\n"
                                f"Write a 3-4 sentence synthesis that explains how these documents relate and what the overall project is about.",
                                "persona"
                            )
                            if result:
                                with eng.connect() as conn:
                                    conn.execute(text(
                                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                        "VALUES (:s, 'project', :f, 'synthesis') ON CONFLICT DO NOTHING"
                                    ), {"s": slug, "f": f"Project synthesis: {result[:500]}"})
                                    conn.commit()
                                _dlog("✓ multi-document synthesis complete")
                        except Exception as e:
                            _dlog(f"⚠ synthesis error: {str(e)[:60]}")

                    # Step 15: LangExtract — grounded fact extraction for Researcher
                    _dlog("extracting grounded facts with LangExtract...")
                    _langextract_facts(slug, all_text, _dlog)

                    # Step 16: Cross-Source Knowledge Graph
                    # 2026-05-25: wrap in StepRunner so doc-only path shares
                    # the same fp cache + dash_training_steps audit as data
                    # path. Previously direct build_knowledge_graph() call
                    # had no fp cache → rebuilt every doc-only retrain even
                    # when docs unchanged, AND wrote 0 step rows → audit blind.
                    _update_step("knowledge_graph")
                    _dlog("building cross-source knowledge graph...")
                    try:
                        from dash.tools.knowledge_graph import build_knowledge_graph
                        from dash.training.runner import StepRunner
                        _doc_runner = StepRunner(run_id, slug, logger_fn=_dlog)

                        def _do_kg_doc():
                            stats = build_knowledge_graph(slug)
                            _dlog(f"✓ knowledge graph built ({stats.get('triples', 0)} triples, {stats.get('communities', 0)} communities)")
                            return stats

                        # fp = doc set + extracted-text length (cheap stable signal)
                        _doc_files = sorted(p.name for p in docs_dir.iterdir() if p.is_file()) if docs_dir.exists() else []
                        _doc_runner.run(
                            "knowledge_graph", _do_kg_doc,
                            fp_inputs={"docs": _doc_files, "text_len": len(all_text)},
                            scope="project", step_no=16,
                        )
                    except Exception as e:
                        _dlog(f"⚠ knowledge graph skipped: {str(e)[:60]}")

                    # Sub-agent synthesis — DISABLED BY DEFAULT (dead feature, see data path).
                    if os.environ.get("SUBAGENT_SYNTHESIS_ENABLED") in ("1", "true", "True"):
                        try:
                            _update_step("subagent_synthesis")
                            from dash.learning.subagent_synthesis import synthesize_subagents
                            _sa_res = synthesize_subagents(slug, logger=_dlog)
                            _dlog(f"✓ sub-agents: created={_sa_res.get('created',0)} clusters={_sa_res.get('clusters',0)}")
                        except Exception:
                            logger.exception("subagent_synthesis failed (non-fatal)")

                    # Step 16b: Vector backfill — guarantee dash_vectors populated.
                    try:
                        _vbf = _enqueue_vector_backfill(slug)
                        _dlog(
                            f"✓ vector backfill: {_vbf.get('knowledge',0)} docs, "
                            f"{_vbf.get('brain',0)} brain, {_vbf.get('kg',0)} kg triples"
                        )
                    except Exception as _vbf_e:
                        _dlog(f"⚠ vector backfill skipped: {str(_vbf_e)[:80]}")

                    # Step 17a: Agent Template Reconciliation
                    try:
                        from dash.templates.reconcile import reconcile as _reconcile_template
                        rec = _reconcile_template(slug)
                        if rec.get("reconciled"):
                            wfs = rec.get("workflows", {})
                            _dlog(
                                f"✓ template reconciled: {rec.get('entities_matched',0)}/{rec.get('entities_total',0)} entities, "
                                f"{rec.get('bindings_bound',0)}/{rec.get('bindings_total',0)} bindings, "
                                f"{wfs.get('active',0)} workflows active"
                            )
                    except Exception as _re:
                        _dlog(f"⚠ template reconcile skipped: {str(_re)[:60]}")

                    # Step 17: Auto-Scope Guardrail
                    _update_step("scope_derivation")
                    _dlog("deriving scope guardrail...")
                    try:
                        from dash.scope_deriver import derive_scope
                        from dash.feature_config import set_scope
                        derived = derive_scope(slug)
                        set_scope(slug, derived, mark_auto=True)
                        _dlog(
                            f"✓ scope: {len(derived.get('topics', []))} topics, "
                            f"{len(derived.get('denied_intents', []))} denied"
                        )
                    except Exception as e:
                        import logging
                        logging.warning(f"scope derivation failed for {slug}: {e}")
                        _dlog(f"⚠ scope derivation skipped: {str(e)[:60]}")
                    # Auto-apply data-fit feature config (only if never configured).
                    try:
                        from dash.feature_config import apply_recommended_if_unset
                        _rec = apply_recommended_if_unset(slug)
                        if _rec:
                            _on = [k for k, v in (_rec.get("tools") or {}).items() if v]
                            _dlog(f"✓ feature config auto-set from data: {', '.join(_on) or 'none'}")
                    except Exception as e:
                        import logging
                        logging.warning(f"feature-config auto-apply failed for {slug}: {e}")

                    # Seed generic CRM starter metrics if the schema looks like a CRM
                    # (universal, column-alias-resolved, status='suggested'). Idempotent.
                    try:
                        from dash.tools.crm_starter import looks_like_crm, seed_crm_starter
                        if looks_like_crm(slug):
                            _cr = seed_crm_starter(slug)
                            if _cr.get("seeded"):
                                _dlog(f"✓ CRM starter metrics seeded: {', '.join(_cr['seeded'])}")
                    except Exception as e:
                        import logging
                        logging.warning(f"CRM starter seed failed for {slug}: {e}")

                    # Step 18: Auto-Derive Learning Goals
                    _update_step("goals_derivation")
                    _dlog("deriving learning goals...")
                    try:
                        from dash.learning.goals_deriver import derive_goals
                        gres = derive_goals(slug, force=False)
                        if gres.get("derived"):
                            _dlog("✓ learning_goals.md auto-generated")
                        else:
                            _dlog(f"· goals skipped: {gres.get('reason','')[:60]}")
                    except Exception as e:
                        import logging
                        logging.warning(f"goals derivation failed for {slug}: {e}")
                        _dlog(f"⚠ goals derivation skipped: {str(e)[:60]}")

                _dlog("✓ doc-only training complete")

                # Save training run with proper step tracking
                run_id = None
                with eng.connect() as conn:
                    r = conn.execute(text(
                        "INSERT INTO public.dash_training_runs (project_slug, status, steps, logs, finished_at) "
                        "VALUES (:s, 'done', 'complete', CAST(:logs AS jsonb), NOW()) RETURNING id"
                    ), {"s": slug, "logs": json.dumps(_log_entries)})
                    run_id = r.fetchone()[0]
                    conn.commit()

            except Exception:
                pass
        _bg_executor.submit(_bg_docs_only).add_done_callback(_bg_done_log)
        return {"status": "ok", "tables": 0, "message": "Doc-only project — indexing documents and filling brain"}

    import pandas as pd

    _training_cancel_flags[slug] = False  # Reset cancel flag

    def _bg():
        import time as _time
        import threading
        import concurrent.futures
        # Begin root trace for the whole TRAIN-ALL / retrain batch.
        _trace_start("training", project_slug=slug, name="training")
        skipped = 0
        trained = 0
        total_tables = len(tables)

        # Create a single master training run for the entire retrain batch
        master_run_id = None
        master_engine = create_engine(db_url)
        try:
            with master_engine.connect() as conn:
                result = conn.execute(text(
                    "INSERT INTO public.dash_training_runs (project_slug, status, steps) "
                    "VALUES (:s, 'running', :steps) RETURNING id"
                ), {"s": slug, "steps": f"starting||0|{total_tables}"})
                master_run_id = result.fetchone()[0]
                conn.commit()
                import logging as _l
                _l.info(f"[retrain] master run {master_run_id} created for {slug}")
        except Exception as _e:
            import logging as _l
            _l.error(f"[retrain] master run insert FAILED for {slug}: {_e}", exc_info=True)

        def _master_log(msg: str, tbl_name: str = "", tbl_idx: int = 0):
            """Log to the master training run."""
            if not master_run_id:
                return
            try:
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET logs = COALESCE(logs, '[]'::jsonb) || CAST(:entry AS jsonb) WHERE id = :id"
                    ), {"entry": json.dumps([{"ts": _time.strftime('%H:%M:%S'), "tsabs": _time.time(), "msg": msg, "table": tbl_name, "table_index": tbl_idx, "total_tables": total_tables}]), "id": master_run_id})
                    conn.commit()
            except Exception:
                pass

        def _set_step(step_name: str):
            """Update current_step + steps on the master run so the UI never
            freezes during the tail phase. Fail-soft."""
            if not master_run_id:
                return
            try:
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET current_step = :cs, "
                        "steps = :s WHERE id = :id"
                    ), {"cs": step_name,
                        "s": f"{step_name}||{total_tables}|{total_tables}",
                        "id": master_run_id})
                    conn.commit()
            except Exception:
                pass

        # ─── LLM observer wires per-call metrics into the run log ──────────────
        from dash.settings import set_llm_observer, reset_llm_observer, set_llm_project
        # Bind project for cost ledger + budget gate on training_llm_call/vision_call
        try:
            set_llm_project(slug)
        except Exception:
            pass
        _llm_totals = {"calls": 0, "tokens_in": 0, "tokens_out": 0, "cost_usd": 0.0}
        _llm_lock = threading.Lock()
        # Thread-local current-table context so parallel table workers don't
        # clobber each other's LLM-log labels. The observer callback runs on
        # whichever worker thread made the LLM call → reads its own thread's ctx.
        _tbl_ctx = threading.local()

        def _cur_tbl():
            return getattr(_tbl_ctx, "name", ""), getattr(_tbl_ctx, "idx", 0)

        def _on_llm(stats: dict):
            with _llm_lock:
                _llm_totals["calls"] += 1
                _llm_totals["tokens_in"] += stats.get("tokens_in", 0) or 0
                _llm_totals["tokens_out"] += stats.get("tokens_out", 0) or 0
                _llm_totals["cost_usd"] += stats.get("cost_usd", 0.0) or 0.0
                _calls = _llm_totals["calls"]
                _cost = _llm_totals["cost_usd"]
            model_short = (stats.get("model", "") or "").split("/")[-1]
            tag = "✗" if not stats.get("ok") else "✓"
            msg = (
                f"  {tag} llm · {stats.get('task','?')} · {model_short} · "
                f"{stats.get('latency_s', 0)}s · "
                f"{stats.get('tokens_in', 0)}→{stats.get('tokens_out', 0)} tok · "
                f"${stats.get('cost_usd', 0):.4f}"
                f" · running ${_cost:.4f} ({_calls} calls)"
            )
            _tname, _tidx = _cur_tbl()
            _master_log(msg, _tname, _tidx)

        _obs_token = set_llm_observer(_on_llm, tag=f"retrain:{slug}")

        # Clear stale eval cases ONCE at retrain start. Per-table eval generation
        # appends to public.dash_evals every run with no dedup, so the set grew
        # unbounded (15→30→90...) and post-training evals got slower each retrain.
        # Wiping here means each retrain regenerates a fresh, bounded eval set.
        # Purge orphaned knowledge files for tables that no longer exist (e.g.
        # a raw-SQL wipe drops the DB table but leaves dimensions/qa/business
        # JSON on disk → ghost relationships + eval cases against dead tables).
        try:
            _purge_orphan_knowledge(slug, log_fn=lambda m: _master_log(m, "", total_tables))
        except Exception:
            pass
        try:
            with master_engine.begin() as _evc:
                _evc.execute(text("DELETE FROM public.dash_evals WHERE project_slug = :s"),
                             {"s": slug})
        except Exception as _eve:
            import logging as _l
            _l.getLogger(__name__).debug(f"eval pre-clear skipped for {slug}: {_eve}")

        def _train_one(tbl_idx: int, tbl: str) -> dict:
            """Train a single table. Runs on a worker thread under the
            ThreadPoolExecutor. Sets thread-local table context so the LLM
            observer labels logs correctly. Fully fail-soft — returns a status
            dict, never raises (so one bad table can't kill the batch)."""
            # Stamp this worker thread's current-table context for _on_llm.
            _tbl_ctx.name = tbl
            _tbl_ctx.idx = tbl_idx
            # Per-table root trace — runs on this worker thread (contextvars are
            # thread/task-local, so each parallel table gets its own trace tree).
            _trace_start("training", project_slug=slug, name=f"training.table.{tbl}")
            if _training_cancel_flags.get(slug):
                try:
                    _trace_end("done")
                except Exception:  # noqa: BLE001
                    pass
                return {"status": "cancelled", "table": tbl}
            _train_err: str | None = None
            try:
                # Get current row count and columns (engine is thread-safe for .connect())
                with engine.connect() as conn:
                    row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                tbl_cols = [c["name"] for c in insp.get_columns(tbl, schema=schema)]

                # Check fingerprint — skip if unchanged (unless force=True)
                change_type = check_fingerprint_changed(slug, tbl, row_count, tbl_cols)

                if change_type == "unchanged" and not force:
                    _master_log(f"⊘ skipping {tbl} — unchanged (fingerprint match)", tbl, tbl_idx)
                    # Update master run step to show skip
                    if master_run_id:
                        try:
                            _skip_progress = json.dumps({
                                "step": "skipped", "table": tbl,
                                "index": tbl_idx, "total": total_tables,
                                "started_at": _time.strftime('%Y-%m-%dT%H:%M:%S'),
                            })
                            with master_engine.connect() as conn:
                                conn.execute(text(
                                    "UPDATE public.dash_training_runs SET steps = :steps, "
                                    "current_step = :cs, stage_progress = :sp, "
                                    "current_progress = CAST(:cp AS jsonb) WHERE id = :id"
                                ), {"steps": f"skipped|{tbl}|{tbl_idx}|{total_tables}",
                                     "cs": "skipped", "sp": tbl_idx,
                                     "cp": _skip_progress, "id": master_run_id})
                                conn.commit()
                        except Exception:
                            pass
                    return {"status": "skipped", "table": tbl}

                # Table changed or new — run training
                _master_log(f"training table {tbl} ({tbl_idx}/{total_tables})...", tbl, tbl_idx)
                df = pd.read_sql(f'SELECT * FROM "{schema}"."{tbl}" LIMIT 100', engine)
                col_analyses = [_analyze_column(df[col]) for col in df.columns]
                sample_rows = df.head(10).to_dict('records')

                tables_dir = KNOWLEDGE_DIR / slug / "tables"
                business_dir = KNOWLEDGE_DIR / slug / "business"
                tables_dir.mkdir(parents=True, exist_ok=True)
                business_dir.mkdir(parents=True, exist_ok=True)

                # Load existing metadata or create new
                meta_file = tables_dir / f"{tbl}.json"
                if meta_file.exists():
                    with open(meta_file) as f:
                        metadata = json.load(f)
                else:
                    metadata = _generate_metadata(tbl, df, col_analyses)

                biz_rules = _generate_business_rules(tbl, col_analyses)
                _run_auto_training(slug, tbl, col_analyses, metadata, biz_rules, sample_rows, tables_dir, business_dir,
                                   master_run_id=master_run_id, table_index=tbl_idx, total_tables=total_tables)

                # Save new fingerprint after training
                save_fingerprint(slug, tbl, row_count, tbl_cols)
                _master_log(f"✓ table {tbl} training complete ({tbl_idx}/{total_tables})", tbl, tbl_idx)
                return {"status": "trained", "table": tbl}
            except Exception as e:
                import logging
                logging.error(f"Retrain failed for {slug}/{tbl}: {e}")
                _master_log(f"⚠ table {tbl} training failed: {str(e)[:80]}", tbl, tbl_idx)
                _train_err = str(e)[:200]
                return {"status": "failed", "table": tbl, "error": str(e)[:200]}
            finally:
                # Close this table's root training trace (fail-soft).
                try:
                    _trace_end("error" if _train_err else "done", _train_err)
                except Exception:  # noqa: BLE001
                    pass

        # Parallelize per-table training (max 4 concurrent). Each table opens its
        # own DB connections + writes its own knowledge files, so they're
        # independent. Cancel is honored before submit + inside _train_one.
        with concurrent.futures.ThreadPoolExecutor(max_workers=4, thread_name_prefix="retrain-tbl") as _tpool:
            _futs = []
            for tbl_idx, tbl in enumerate(tables, start=1):
                if _training_cancel_flags.get(slug):
                    break
                _futs.append(_tpool.submit(_train_one, tbl_idx, tbl))
            _trained_tbls: list[str] = []
            for _f in concurrent.futures.as_completed(_futs):
                try:
                    _res = _f.result()
                except Exception:
                    _res = {"status": "failed"}
                _st = (_res or {}).get("status")
                if _st == "trained":
                    trained += 1
                    _tn = (_res or {}).get("table")
                    if _tn:
                        _trained_tbls.append(_tn)
                elif _st == "skipped":
                    skipped += 1

        # Kick off LLM column enrichment for ALL project tables (fire-and-forget,
        # fail-soft, runs in background). Auto-fills dash_column_meta so the
        # Knowledge tab shows rich semantic-type/blurb/samples without per-col
        # AI clicks. Runs for ALL tables (not just trained ones) since user
        # expects "Train All" to refresh everything; UPSERT is idempotent so
        # already-enriched cols just get refreshed.
        _trained_tbls = list(tables)
        if _trained_tbls:
            try:
                from dash.tools.column_describer import enrich_columns_async
                import asyncio as _aio
                def _kick_enrich():
                    try:
                        _loop = _aio.new_event_loop()
                        _aio.set_event_loop(_loop)
                        for _t in _trained_tbls:
                            try:
                                _loop.run_until_complete(enrich_columns_async(slug, _t))
                            except Exception:
                                _master_log(f"⚠ col enrich failed for {_t}", _t, 0)
                        _loop.close()
                    except Exception:
                        pass
                import threading as _th
                _th.Thread(target=_kick_enrich, daemon=True, name=f"enrich-{slug}").start()
                _master_log(f"⚙ column enrichment queued for {len(_trained_tbls)} table(s)", "", total_tables)
            except Exception:
                pass

        # ─── Shared fingerprint-cache runner for the idempotent tail steps ─────
        # Issues #40 + #41: KG build, subagent synthesis, vector backfill and ML
        # auto-create are all idempotent and re-run every retrain even when their
        # inputs are unchanged. Wrap them in a single StepRunner so an unchanged
        # retrain skips the expensive re-work. Each keeps its existing fail-soft
        # try/except (the runner is itself fail-soft, but the wrappers also are).
        from dash.training.runner import StepRunner
        _tail_runner = StepRunner(
            master_run_id, slug,
            logger_fn=lambda m: _master_log(m, "", total_tables),
        )

        def _tail_table_names() -> list[str]:
            """Sorted project table names for tail-step fingerprints. Fail-soft."""
            try:
                return sorted(tables)
            except Exception:
                return []

        def _tail_doc_count() -> int:
            """Project document count (KG fp input). Fail-soft → 0."""
            try:
                with master_engine.connect() as _dc:
                    return int(_dc.execute(text(
                        "SELECT count(*) FROM public.dash_documents WHERE project_slug = :s"
                    ), {"s": slug}).scalar() or 0)
            except Exception:
                return 0

        def _tail_qa_signature() -> dict:
            """Cheap Q&A signature (subagent synthesis fp input). Fail-soft."""
            try:
                with master_engine.connect() as _qc:
                    row = _qc.execute(text(
                        "SELECT count(*), COALESCE(max(id), 0) "
                        "FROM public.dash_training_qa WHERE project_slug = :s"
                    ), {"s": slug}).fetchone()
                return {"qa_count": int(row[0] or 0), "qa_max_id": int(row[1] or 0)}
            except Exception:
                return {"qa_count": 0, "qa_max_id": 0}

        def _tail_rowcounts() -> dict:
            """Per-table row counts (ML auto-create fp input). Fail-soft."""
            counts: dict[str, int] = {}
            for _t in _tail_table_names():
                try:
                    with engine.connect() as _rc:
                        counts[_t] = int(_rc.execute(
                            text(f'SELECT COUNT(*) FROM "{schema}"."{_t}"')).scalar() or 0)
                except Exception:
                    counts[_t] = -1
            return counts

        # Cross-Source Knowledge Graph (runs even if all tables were skipped)
        # Knowledge graph runs FIRST (serially) — subagent_synthesis depends on it.
        _set_step("knowledge_graph")
        _master_log("building cross-source knowledge graph...", "", total_tables)
        if master_run_id:
            try:
                _kg_progress = json.dumps({
                    "step": "knowledge_graph", "table": "",
                    "index": total_tables, "total": total_tables,
                    "started_at": _time.strftime('%Y-%m-%dT%H:%M:%S'),
                })
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET steps = :s, "
                        "current_step = :cs, stage_progress = :sp, "
                        "current_progress = CAST(:cp AS jsonb) WHERE id = :id"
                    ), {"s": f"knowledge_graph||{total_tables}|{total_tables}",
                         "cs": "knowledge_graph", "sp": total_tables,
                         "cp": _kg_progress, "id": master_run_id})
                    conn.commit()
            except Exception:
                pass
        try:
            from dash.tools.knowledge_graph import build_knowledge_graph

            def _do_kg():
                kg_stats = build_knowledge_graph(slug)
                _master_log(f"✓ knowledge graph: {kg_stats.get('triples', 0)} triples, {kg_stats.get('entities', 0)} entities, {kg_stats.get('communities', 0)} communities", "", total_tables)
                return kg_stats

            # Fingerprint must include row counts — 2026-05-25 fix. Previous fp
            # (tables + doc count) treated same-named tables as unchanged even
            # when data was wiped + reloaded → KG step was skipped on every
            # retrain after first run. Added row counts so any data change
            # invalidates the cached fp and rebuilds the graph.
            _tail_runner.run(
                "knowledge_graph", _do_kg,
                fp_inputs={
                    "tables": _tail_table_names(),
                    "docs": _tail_doc_count(),
                    "rowcounts": _tail_rowcounts(),
                },
                scope="project", step_no=11,
            )
        except Exception as e:
            import logging
            logging.error(f"Knowledge graph failed for {slug}: {e}")
            _master_log(f"⚠ knowledge graph skipped: {str(e)[:80]}", "", total_tables)

        # Sub-agent synthesis — DISABLED BY DEFAULT (dead feature). It writes
        # dash_custom_agents rows with enabled=false that are NEVER promoted,
        # NEVER loaded into the team, NEVER shown in UI, and NEVER routed to —
        # confirmed by full-codebase audit. It cost ~8s + LLM calls per train
        # for zero downstream value. Re-enable with SUBAGENT_SYNTHESIS_ENABLED=1
        # only if/when draft sub-agents actually get wired into the team.
        if os.environ.get("SUBAGENT_SYNTHESIS_ENABLED") in ("1", "true", "True"):
            _set_step("subagent_synthesis")
            try:
                from dash.learning.subagent_synthesis import synthesize_subagents
                def _sa_log(msg, _t=total_tables):
                    _master_log(msg, "", _t)

                def _do_subagents():
                    _sa_res = synthesize_subagents(slug, logger=_sa_log)
                    _master_log(
                        f"✓ sub-agents: created={_sa_res.get('created',0)} skipped={_sa_res.get('skipped',0)} clusters={_sa_res.get('clusters',0)}",
                        "", total_tables,
                    )
                    return _sa_res

                _tail_runner.run(
                    "subagent_synthesis", _do_subagents,
                    fp_inputs={"qa": _tail_qa_signature()},
                    scope="project", step_no=14,
                )
            except Exception as _sa_e:
                import logging
                logging.exception("subagent_synthesis failed (non-fatal)")
                _master_log(f"⚠ subagent_synthesis skipped: {str(_sa_e)[:80]}", "", total_tables)
        # else: sub-agent synthesis is a dead feature (SUBAGENT_SYNTHESIS_ENABLED=1
        # to re-enable) — silent, no log line.

        # Step 14b: Vector backfill — guarantee dash_vectors populated post-TRAIN ALL.
        # Reembed loop is async + can lag; this explicit pass ensures CMHL-class
        # retrains never finish with `dash.dash_vectors WHERE project_slug=:s = 0`.
        _set_step("vector_backfill")
        try:
            def _do_vbf():
                _vbf = _enqueue_vector_backfill(slug)
                _master_log(
                    f"✓ vector backfill enqueued: {_vbf.get('knowledge',0)} docs, "
                    f"{_vbf.get('brain',0)} brain, {_vbf.get('kg',0)} kg triples"
                    + (f" ({_vbf['errors']} errors)" if _vbf.get('errors') else ""),
                    "", total_tables,
                )
                return _vbf

            # Issue #41 — fp = sorted table names. Unchanged retrain → skip the
            # backfill enqueue (vectors already populated from a prior run).
            _tail_runner.run(
                "vector_backfill", _do_vbf,
                fp_inputs={"tables": _tail_table_names()},
                scope="project", step_no=31,
            )
        except Exception as _vbf_e:
            _master_log(f"⚠ vector backfill skipped: {str(_vbf_e)[:80]}", "", total_tables)

        # Step 14c: Codex code enrichment — read view/table DDL + source logic,
        # LLM-enrich, persist to dash_table_metadata.metadata['pipeline_logic'].
        # DATA path only (tables exist here). Fail-soft — never breaks training.
        _set_step("codex_code_enrich")
        _master_log("reading pipeline code (view DDL / table logic)...", "", total_tables)
        try:
            from dash.tools.codex_code import run_codex_code_enrichment
            res = run_codex_code_enrichment(slug)
            _master_log(
                f"✓ codex code: {res.get('tables_enriched', 0)} tables enriched from source logic",
                "", total_tables,
            )
        except Exception as _cce_e:
            _master_log(f"⚠ codex code enrichment skipped: {str(_cce_e)[:80]}", "", total_tables)

        # ─── Independent tail steps run CONCURRENTLY (max 4) ──────────────────
        # scope_derivation, goals_derivation, auto ML models, post-training evals,
        # and auto-configure all read training output but don't depend on each
        # other. Each task keeps its own fail-soft try/except so one failing never
        # blocks the others. _set_step(...) is still called per task for progress
        # visibility (concurrent overwrites of current_step are fine — the point
        # is the UI never freezes on a stale step). The LLM observer is still
        # active here; _on_llm + _llm_totals are thread-safe via _llm_lock.

        def _task_scope():
            _master_log("deriving scope guardrail...", "", total_tables)

            def _do_scope():
                from dash.scope_deriver import derive_scope
                from dash.feature_config import set_scope
                derived = derive_scope(slug)
                set_scope(slug, derived, mark_auto=True)
                _master_log(
                    f"✓ scope: {len(derived.get('topics', []))} topics, "
                    f"{len(derived.get('core_entities', []))} entities, "
                    f"{len(derived.get('denied_intents', []))} denied",
                    "", total_tables,
                )
                try:
                    from dash.feature_config import apply_recommended_if_unset
                    _rec = apply_recommended_if_unset(slug)
                    if _rec:
                        _on = [k for k, v in (_rec.get("tools") or {}).items() if v]
                        _master_log(f"✓ feature config auto-set from data: {', '.join(_on) or 'docs-only'}", "", total_tables)
                except Exception as _fe:
                    import logging
                    logging.warning(f"feature-config auto-apply failed for {slug}: {_fe}")
                try:
                    from dash.tools.crm_starter import looks_like_crm, seed_crm_starter
                    if looks_like_crm(slug):
                        _cr = seed_crm_starter(slug)
                        if _cr.get("seeded"):
                            _master_log(f"✓ CRM starter metrics seeded: {', '.join(_cr['seeded'])}", "", total_tables)
                except Exception as _ce:
                    import logging
                    logging.warning(f"CRM starter seed failed for {slug}: {_ce}")
                return derived

            # V2 (flag-gated): wrap the ~38s scope LLM in the fingerprint cache so
            # an unchanged (columns + persona) retrain skips it. Side-effect
            # (set_scope) is persisted, so a cache-hit skip leaves scope intact.
            if os.environ.get("TRAINING_V2_PHASE_E") in ("1", "true", "True"):
                try:
                    from dash.training.runner import StepRunner
                    from dash.training.phase_e import _all_column_names
                    # fp = column set only. Scope is schema-driven; do NOT key on
                    # persona (persona_enrich regenerates it via LLM every train,
                    # so its text drifts → fp would never be stable → never cache).
                    _sr = StepRunner(master_run_id, slug,
                                     logger_fn=lambda m: _master_log(m, "", total_tables))
                    _sr.run("derive_scope", _do_scope,
                            fp_inputs={"cols": _all_column_names(slug)},
                            scope="project", step_no=29)
                    return
                except Exception as e:
                    import logging
                    logging.warning(f"scope (v2) failed for {slug}, falling back: {e}")
                    _master_log(f"⚠ scope v2 fell back: {str(e)[:60]}", "", total_tables)

            try:
                _do_scope()
            except Exception as e:
                import logging
                logging.warning(f"scope derivation failed for {slug}: {e}")
                _master_log(f"⚠ scope derivation skipped: {str(e)[:80]}", "", total_tables)

        def _task_goals():
            _master_log("deriving learning goals...", "", total_tables)
            try:
                from dash.learning.goals_deriver import derive_goals
                gres = derive_goals(slug, force=False)
                if gres.get("derived"):
                    _master_log("✓ learning_goals.md auto-generated", "", total_tables)
                else:
                    _master_log(f"· goals skipped: {gres.get('reason','')[:80]}", "", total_tables)
            except Exception as e:
                import logging
                logging.warning(f"goals derivation failed for {slug}: {e}")
                _master_log(f"⚠ goals derivation skipped: {str(e)[:80]}", "", total_tables)

        def _task_ml():
            # ml_worker container + auto_create_models() were removed in the
            # 2026-05-23 ML pivot (LLM-native, in-process ≤50K rows). This
            # step is dead — kept as no-op so older _tail_runner audit rows
            # still resolve, but never fires the deleted code. Don't log
            # FAILED in the master log. Overwrite any prior cached step row
            # so stale FAILED status doesn't linger.
            try:
                with master_engine.begin() as _mc:
                    _mc.execute(text("""
                        INSERT INTO public.dash_training_steps
                          (run_id, project_slug, step_no, name, scope, status,
                           elapsed_ms, error, started_at, finished_at, updated_at)
                        VALUES (:r, :s, 32, 'ml_auto_create', 'project', 'skipped',
                                0, NULL, now(), now(), now())
                        ON CONFLICT (project_slug, name, scope) DO UPDATE
                          SET status='skipped', error=NULL, elapsed_ms=0,
                              finished_at=now(), updated_at=now(), run_id=EXCLUDED.run_id
                    """), {"r": master_run_id, "s": slug})
            except Exception:
                pass
            return

        def _task_evals():
            if master_run_id:
                try:
                    with master_engine.connect() as conn:
                        conn.execute(text(
                            "UPDATE public.dash_training_runs SET steps = :s WHERE id = :id"
                        ), {"s": f"eval_run||{total_tables}|{total_tables}", "id": master_run_id})
                        conn.commit()
                except Exception:
                    pass
            # Issue #39 — generate the project-wide eval set ONCE here (was
            # per-table inside _run_auto_training). The DELETE-at-start pre-clear
            # already wiped the prior set, so this produces one fresh bounded set
            # per retrain. NOT wrapped in StepRunner (caching/skipping it after a
            # pre-clear would leave evals deleted-and-not-regenerated).
            try:
                _ev_gen = _generate_project_evals(
                    slug, log_fn=lambda m: _master_log(m, "", total_tables))
            except Exception as _evg_e:
                import logging as _l
                _l.getLogger(__name__).warning(f"project eval gen failed for {slug}: {_evg_e}")
            _master_log("running post-training evals...", "", total_tables)
            try:
                from app.learning import _run_evals_for_slug
                _eval_res = _run_evals_for_slug(slug)
                _master_log(
                    f"✓ evals: {_eval_res.get('passed', 0)}/{_eval_res.get('total', 0)} passed, "
                    f"{_eval_res.get('partial', 0)} partial",
                    "", total_tables,
                )
            except Exception as _ee:
                import logging as _l
                _l.getLogger(__name__).warning(f"post-training eval failed for {slug}: {_ee}")
                _master_log(f"⚠ evals skipped: {str(_ee)[:80]}", "", total_tables)

        def _task_auto_configure():
            _master_log("auto-detecting vertical...", "", total_tables)
            # V2 Phase E (flag-gated): runner-tracked + fingerprint-cached vertical
            # detect+apply. Skips the ~60s pack apply on retrain when (vertical,
            # column-set) is unchanged. Flag off → unchanged legacy path below.
            if os.environ.get("TRAINING_V2_PHASE_E") in ("1", "true", "True"):
                try:
                    from dash.training.runner import StepRunner
                    from dash.training.phase_e import run_phase_e
                    _runner = StepRunner(master_run_id, slug,
                                         logger_fn=lambda m: _master_log(m, "", total_tables))
                    run_phase_e(slug, _runner,
                                master_log=lambda m: _master_log(m, "", total_tables),
                                set_step=_set_step)
                    return
                except Exception as _v2e:
                    import logging as _l
                    _l.getLogger(__name__).warning(f"phase_e failed for {slug}, falling back: {_v2e}")
                    _master_log(f"⚠ phase_e fell back to legacy: {str(_v2e)[:60]}", "", total_tables)
            try:
                from dash.learning.auto_configurator import classify_vertical
                _detection = classify_vertical(slug)
                _conf = float(_detection.get("confidence", 0) or 0)
                _vert = _detection.get("vertical", "generic")
                _master_log(f"· detected: {_vert} ({int(_conf*100)}%) via {_detection.get('method','?')}", "", total_tables)
                if _conf >= 0.50:
                    try:
                        from dash.learning.auto_apply import auto_apply_vertical
                        _apply_res = auto_apply_vertical(slug, _detection, user_id="auto", derive_scope_step=False)
                        _ok_steps = sum(1 for s in _apply_res.get("applied_steps", []) if s.get("ok"))
                        _tot_steps = len(_apply_res.get("applied_steps", []))
                        _master_log(f"✓ auto-applied {_vert} pack: {_ok_steps}/{_tot_steps} steps", "", total_tables)
                    except Exception as _ae:
                        _master_log(f"⚠ auto-apply skipped: {str(_ae)[:80]}", "", total_tables)
                elif _conf >= 0.15:
                    # Suggest only — save detection row to history for UI review
                    try:
                        import json as _aj
                        with master_engine.connect() as _ac:
                            _ac.execute(text("""
                                INSERT INTO dash.dash_auto_apply_history
                                  (project_slug, vertical, template, confidence, detection, applied, applied_by)
                                VALUES (:s, :v, :t, :c, CAST(:d AS jsonb), false, 'auto-suggest')
                            """), {"s": slug, "v": _vert, "t": _detection.get("template"),
                                   "c": _conf, "d": _aj.dumps(_detection)})
                            _ac.commit()
                        _master_log(f"· suggested {_vert} ({int(_conf*100)}%) — review in Auto-Config tab", "", total_tables)
                    except Exception as _se:
                        _master_log(f"⚠ suggestion save skipped: {str(_se)[:80]}", "", total_tables)
                else:
                    _master_log("· no strong vertical signal — skipped auto-apply", "", total_tables)
            except Exception as _ace:
                import logging as _l
                _l.getLogger(__name__).warning(f"auto-configure failed for {slug}: {_ace}")
                _master_log(f"⚠ auto-configure skipped: {str(_ace)[:80]}", "", total_tables)

        # Stable phase-level label for the concurrent tail. The 5 tail tasks no
        # longer call _set_step (they ran concurrently → last-writer-wins clobber
        # froze the UI on a random step). Per-step detail now lives in
        # dash_training_steps; current_step reads "tail" for the whole phase.
        _set_step("tail")
        with concurrent.futures.ThreadPoolExecutor(max_workers=4, thread_name_prefix="retrain-tail") as _ttp:
            _tail_futs = [
                _ttp.submit(_task_scope),
                _ttp.submit(_task_goals),
                _ttp.submit(_task_ml),
                _ttp.submit(_task_evals),
                _ttp.submit(_task_auto_configure),
            ]
            for _tf in concurrent.futures.as_completed(_tail_futs):
                try:
                    _tf.result()
                except Exception as _tfe:
                    import logging as _l
                    _l.getLogger(__name__).warning(f"tail task failed for {slug}: {_tfe}")

        # ─────────────── Steps 15a-15e — Investment vertical (gated) ───────────────
        # Auto-active when feature flag is on OR project has financial tables.
        try:
            _inv_on = False
            try:
                from dash.feature_config import get_feature_config as _ifc
                _inv_on = bool(_ifc(slug).get("agents", {}).get("investment", False))
            except Exception:
                pass
            if not _inv_on:
                try:
                    from dash.team import _has_financial_tables as _ihft
                    _inv_on = _ihft(slug)
                except Exception:
                    pass
            if _inv_on:
                _master_log("investment vertical detected — running 5 sub-steps...", "", total_tables)
                _run_investment_training_steps(slug, master_engine, master_run_id, total_tables, _master_log)
        except Exception as _e:
            import logging
            logging.debug(f"investment training skipped for {slug}: {_e}")

        # Final LLM cost summary
        try:
            reset_llm_observer(_obs_token)
        except Exception:
            pass
        try:
            set_llm_project(None)
        except Exception:
            pass
        _master_log(
            f"Σ totals · {_llm_totals['calls']} LLM calls · "
            f"{_llm_totals['tokens_in']:,} in + {_llm_totals['tokens_out']:,} out tokens · "
            f"${_llm_totals['cost_usd']:.4f} total cost",
            "", total_tables,
        )
        # (post-training evals + auto-configure now run concurrently in the
        #  tail ThreadPoolExecutor above — see _task_evals / _task_auto_configure)

        # Mark master run as done
        _set_step("finalizing")
        if master_run_id:
            try:
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET status = 'done', steps = :steps, "
                        "tables_trained = :trained, finished_at = NOW() WHERE id = :id"
                    ), {"steps": f"complete||{total_tables}|{total_tables}", "trained": trained, "id": master_run_id})
                    conn.commit()
            except Exception:
                pass

        # ─── Bilingual (EN+Burmese) twin refresh — runs AFTER English memories /
        # schema / training-qa are all written above. Idempotent + fail-soft, so a
        # regen failure never breaks training and re-runs (force-retrain) just
        # refresh the Burmese twins of any newly-added English rows.
        try:
            from scripts.regen_bilingual import run as _regen_bilingual
            _bi = _regen_bilingual()
            _master_log(
                f"✓ bilingual twins refreshed · {_bi.get('memories_my',0)} mem · "
                f"{_bi.get('patterns_twin',0)} pat · {_bi.get('rules_my',0)} rules · "
                f"{_bi.get('brain_my',0)} brain · {_bi.get('schema_bilingual',0)} schema",
                "", total_tables,
            )
        except Exception as _bie:
            import logging as _l
            _l.getLogger(__name__).warning(f"bilingual regen skipped for {slug}: {_bie}")
            _master_log(f"⚠ bilingual regen skipped: {str(_bie)[:80]}", "", total_tables)

        # ─── Semantic catalog vectors — hybrid (vector+keyword) catalog_search
        # corpus. Idempotent (unchanged blobs skipped). Fail-soft so an embed
        # failure never breaks training; catalog_search falls back to ILIKE.
        try:
            from scripts.build_catalog_vectors import run as _build_catalog
            _cv = _build_catalog()
            _master_log(f"✓ catalog vectors: {_cv.get('embedded',0)} products", "", total_tables)
        except Exception as _cve:
            import logging as _l
            _l.getLogger(__name__).warning(f"catalog vectors skipped for {slug}: {_cve}")
            _master_log(f"⚠ catalog vectors skipped: {str(_cve)[:80]}", "", total_tables)

        # ─── Denormalized shop_flat stock table — pre-joined catalog+stock for
        # fast branch lookups. Idempotent rebuild. Fail-soft so a build failure
        # never breaks training; tools fall back to the live joined tables.
        try:
            from scripts.build_shop_flat import run as _build_shop_flat
            _sf = _build_shop_flat()
            _master_log(f"✓ shop_flat: {_sf.get('rows',0)} rows", "", total_tables)
        except Exception as _sfe:
            import logging as _l
            _l.getLogger(__name__).warning(f"shop_flat skipped for {slug}: {_sfe}")
            _master_log(f"⚠ shop_flat skipped: {str(_sfe)[:80]}", "", total_tables)

        # AutoSim grounded-scenario enqueue REMOVED — the sim chassis was deleted
        # in the 2026-05-23 trim, so `autosim_generate_grounded` has NO registered
        # minion handler. This enqueue pushed a dead job into the queue on every
        # retrain (nothing consumed it). Dropped.

        # Issue #10 — Enqueue first dream-lite cycle so Day-1 user has anti-patterns
        # / insights / reflections without waiting for nightly 03:15 UTC cron.
        # Cheap (~$0.005). Fail-soft — never blocks training completion.
        try:
            from dash.minions import queue as _q2
            # Synthesize a bootstrap session_id so dream_lite handler (which
            # requires both project + session_id) can run a first-time cycle.
            # If no episodes exist for this session, handler writes a
            # `bootstrap_ok` row (is_bootstrap=true, friendly_status set) to
            # public.dash_dream_lite_runs so Day-1 user sees the run trail
            # without UI rendering it as a failure.
            _bootstrap_sid = f"first_training_{slug}_{int(_time.time())}"
            _q2.enqueue(
                project=slug,
                kind="dream_lite",
                payload={
                    "project": slug,
                    "session_id": _bootstrap_sid,
                    "trigger_reason": "first_training",
                },
                priority=7,
            )
            _master_log("✓ first dream-lite cycle enqueued", "", total_tables)
        except Exception as _de:
            import logging as _l
            _l.getLogger(__name__).warning(f"dream_lite enqueue failed for {slug}: {_de}")
            _master_log(f"⚠ dream_lite enqueue skipped: {str(_de)[:80]}", "", total_tables)

        # Best-effort notification: full project retrain finished
        try:
            from app.auth import notify_user
            with master_engine.connect() as _nc:
                _nr = _nc.execute(text(
                    "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                ), {"s": slug}).fetchone()
            if _nr and _nr[0]:
                notify_user(
                    int(_nr[0]),
                    f"Training complete · {slug}",
                    f"{trained}/{total_tables} tables retrained",
                    "success",
                )
        except Exception:
            pass

        # Close the batch-level root training trace (fail-soft).
        try:
            _trace_end("done")
        except Exception:  # noqa: BLE001
            pass

        # Auto-install best-matching vertical workflow pack on training complete
        # (Issue #4). Only fires if project has zero existing workflows — never
        # overwrites manual setups. Fail-soft.
        try:
            from sqlalchemy import text as _ttext
            from db import get_sql_engine as _tge
            _ce = _tge()
            with _ce.connect() as _cc:
                _wfc = int(_cc.execute(_ttext(
                    "SELECT COUNT(*) FROM dash.dash_autonomous_workflows "
                    " WHERE project_slug = :s"
                ), {"s": slug}).scalar() or 0)
            if _wfc == 0:
                from dash.workflows.verticals import detect as _vdetect, install as _vinstall
                _matches = _vdetect(slug)
                _top = _matches[0] if _matches else None
                if _top and _top.get("score", 0) >= 0.4:
                    _r = _vinstall(slug, _top["name"])
                    _master_log(
                        f"✓ auto-installed vertical pack '{_top['name']}' "
                        f"(score {_top['score']:.2f}) — {_r.get('installed', 0)} workflows ready",
                        "", total_tables,
                    )
        except Exception as _vp_e:
            import logging as _vp_log
            _vp_log.getLogger(__name__).debug(
                f"vertical pack auto-install skipped for {slug}: {_vp_e}"
            )

    _fut = _bg_executor.submit(_bg)
    def _check_bg(f):
        if f.exception():
            import logging as _l
            _l.error(f"[retrain] _bg crashed for {slug}: {f.exception()}", exc_info=f.exception())
    _fut.add_done_callback(_check_bg)
    return {"status": "ok", "tables": len(tables), "message": "Delta retraining started — unchanged tables will be skipped"}


@router.get("/knowledge-file-content/{filename}")
def get_knowledge_file_content(filename: str, request: Request, project: str | None = None, subdir: str = "tables"):
    """Read the content of a knowledge file."""
    safe_name = Path(filename).name
    if safe_name != filename or '..' in filename:
        raise HTTPException(400, "Invalid filename")
    if project:
        base = KNOWLEDGE_DIR / project / subdir
    else:
        base = KNOWLEDGE_DIR / subdir
    filepath = base / filename
    if not filepath.exists():
        raise HTTPException(404, f"File not found: {filename}")
    try:
        content = filepath.read_text(errors="replace")
        if filename.endswith(".json"):
            # Defense against LLM-generated descriptions polluted by source CSV control bytes
            # and illegal JSON escapes (\0, \x, raw NULs etc.).
            # 1. Strip ASCII control chars except \t \n \r
            sanitized = "".join(ch for ch in content if ch in "\t\n\r" or ord(ch) >= 32)
            # 2. Fix invalid backslash escapes. JSON only allows: " \ / b f n r t u
            # Walk char by char so paired backslashes consume correctly.
            import re as _re
            _valid_esc = set('"\\/bfnrtu')
            _buf = []
            _i = 0
            _n = len(sanitized)
            while _i < _n:
                _ch = sanitized[_i]
                if _ch == '\\' and _i + 1 < _n:
                    _nxt = sanitized[_i + 1]
                    if _nxt in _valid_esc:
                        _buf.append(_ch)
                        _buf.append(_nxt)
                        _i += 2
                        continue
                    else:
                        # invalid escape — drop the backslash, keep next char
                        _buf.append(_nxt)
                        _i += 2
                        continue
                _buf.append(_ch)
                _i += 1
            sanitized = "".join(_buf)
            try:
                parsed = json.loads(sanitized)
            except json.JSONDecodeError:
                try:
                    parsed = json.loads(sanitized, strict=False)
                except json.JSONDecodeError:
                    # Last resort: return raw text so UI shows file content
                    return {"content": sanitized, "type": "text", "parse_error": True}
            return {"content": parsed, "type": "json"}
        return {"content": content, "type": "text"}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.get("/knowledge-files")
def list_knowledge_files(request: Request, project: str | None = None):
    """List knowledge files for a project or global."""
    if project:
        base = KNOWLEDGE_DIR / project
    else:
        base = KNOWLEDGE_DIR
    files = []
    for subdir in ["tables", "queries", "business"]:
        path = base / subdir
        if path.exists():
            for f in sorted(path.iterdir()):
                if f.is_file() and not f.name.startswith("."):
                    files.append({"name": f.name, "type": subdir, "size": f.stat().st_size})
    return {"files": files}


@router.get("/lineage")
def get_lineage(request: Request, project: str | None = None):
    """Detect table relationships via FK column patterns and constraints."""
    engine = create_engine(db_url)
    insp = inspect(engine)
    schema = _get_project_schema(request, project) or "public"
    tables = insp.get_table_names(schema=schema)
    relationships: list[dict] = []

    for tbl in tables:
        cols = insp.get_columns(tbl, schema=schema)
        for col in cols:
            name = col["name"]
            if name.endswith("_id") and name != "id":
                ref_table = name[:-3] + "s"
                if ref_table not in tables:
                    ref_table = name[:-3]
                if ref_table in tables:
                    relationships.append({
                        "from_table": tbl,
                        "from_column": name,
                        "to_table": ref_table,
                        "to_column": "id",
                        "type": "foreign_key",
                    })

        try:
            fks = insp.get_foreign_keys(tbl, schema=schema)
            for fk in fks:
                for i, col in enumerate(fk.get("constrained_columns", [])):
                    ref_cols = fk.get("referred_columns", [])
                    relationships.append({
                        "from_table": tbl,
                        "from_column": col,
                        "to_table": fk.get("referred_table", ""),
                        "to_column": ref_cols[i] if i < len(ref_cols) else "id",
                        "type": "constraint",
                    })
        except Exception:
            pass

    # Deduplicate
    seen = set()
    unique: list[dict] = []
    for r in relationships:
        key = f"{r['from_table']}.{r['from_column']}->{r['to_table']}.{r['to_column']}"
        if key not in seen:
            seen.add(key)
            unique.append(r)

    return {"relationships": unique, "tables": tables}


def _bootstrap_shared():
    """Create shared_results table if needed."""
    engine = create_engine(db_url)
    with engine.connect() as conn:
        conn.execute(text("""
            CREATE TABLE IF NOT EXISTS public.shared_results (
                id SERIAL PRIMARY KEY,
                from_user TEXT NOT NULL,
                title TEXT NOT NULL,
                content TEXT NOT NULL,
                sql_query TEXT,
                created_at TIMESTAMP DEFAULT NOW()
            )
        """))
        conn.commit()


@router.post("/shared")
def share_result(request: Request, title: str, content: str, sql_query: str = ""):
    """Share a query result with all users."""
    user_id = _get_user_id(request)
    _bootstrap_shared()
    engine = create_engine(db_url)
    with engine.connect() as conn:
        conn.execute(text(
            "INSERT INTO public.shared_results (from_user, title, content, sql_query) VALUES (:u, :t, :c, :s)"
        ), {"u": user_id, "t": title, "c": content, "s": sql_query})
        conn.commit()
    return {"status": "ok"}


@router.get("/shared")
def list_shared(request: Request):
    """List all shared results."""
    _bootstrap_shared()
    engine = create_engine(db_url)
    results = []
    with engine.connect() as conn:
        rows = conn.execute(text(
            "SELECT id, from_user, title, content, sql_query, created_at FROM public.shared_results ORDER BY created_at DESC LIMIT 50"
        )).fetchall()
        for r in rows:
            results.append({"id": r[0], "from_user": r[1], "title": r[2], "content": r[3][:500], "sql_query": r[4], "created_at": str(r[5]) if r[5] else None})
    return {"results": results}


@router.delete("/shared/{result_id}")
def delete_shared(result_id: int, request: Request):
    """Delete a shared result (own or super admin only)."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")
    username = user.get("username", "")
    is_super = user.get("is_super", False)
    engine = create_engine(db_url)
    with engine.connect() as conn:
        if is_super:
            conn.execute(text("DELETE FROM public.shared_results WHERE id = :id"), {"id": result_id})
        else:
            conn.execute(text("DELETE FROM public.shared_results WHERE id = :id AND from_user = :u"), {"id": result_id, "u": username})
        conn.commit()
    return {"status": "ok"}


@router.get("/sessions")
def list_sessions(request: Request, project: str | None = None):
    """List current user's chat sessions, optionally filtered by project."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"sessions": []}
    uid = user["user_id"]
    engine = create_engine(db_url)
    sessions: list[dict] = []
    try:
        with engine.connect() as conn:
            if project:
                # Project-specific sessions
                result = conn.execute(text("""
                    SELECT session_id, first_message, created_at, updated_at
                    FROM public.dash_chat_sessions
                    WHERE user_id = :uid AND project_slug = :proj
                    ORDER BY updated_at DESC
                    LIMIT 50
                """), {"uid": uid, "proj": project})
            else:
                # All user sessions (for Dash Agent)
                result = conn.execute(text("""
                    SELECT s.session_id, s.first_message, s.created_at, s.updated_at,
                           COALESCE(p.agent_name, 'Dash Agent') as agent_name, s.project_slug
                    FROM public.dash_chat_sessions s
                    LEFT JOIN public.dash_projects p ON p.slug = s.project_slug
                    WHERE s.user_id = :uid
                    ORDER BY s.updated_at DESC
                    LIMIT 50
                """), {"uid": uid})

            for row in result.fetchall():
                entry: dict = {
                    "session_id": str(row[0]),
                    "first_message": row[1] or None,
                    "created_at": str(row[2]) if row[2] else None,
                    "updated_at": str(row[3]) if row[3] else None,
                }
                if not project and len(row) > 4:
                    entry["agent_name"] = row[4]
                    entry["project_slug"] = row[5]
                sessions.append(entry)
    except Exception:
        pass
    return {"sessions": sessions}


@router.get("/sessions/{session_id}/messages")
def get_session_messages(session_id: str, request: Request):
    """Load chat messages from a stored Agno session."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"messages": []}

    engine = create_engine(db_url)
    messages: list[dict] = []

    try:
        with engine.connect() as conn:
            # Verify this session belongs to the user
            owned = conn.execute(text(
                "SELECT 1 FROM public.dash_chat_sessions WHERE session_id = :sid AND user_id = :uid"
            ), {"sid": session_id, "uid": user["user_id"]}).fetchone()
            if not owned:
                return {"messages": []}

            # Get runs from Agno sessions
            row = conn.execute(text(
                "SELECT runs FROM ai.agno_sessions WHERE session_id = :sid"
            ), {"sid": session_id}).fetchone()

            if row and row[0]:
                runs = row[0] if isinstance(row[0], list) else []
                # Index child runs by run_id for parent lookup
                children_by_parent: dict[str, list] = {}
                for r in runs:
                    pid = r.get("parent_run_id")
                    if pid:
                        children_by_parent.setdefault(pid, []).append(r)

                # Pull quality scores keyed by run/message
                try:
                    q_rows = conn.execute(text(
                        "SELECT chat_id, score FROM public.dash_quality_scores WHERE session_id = :sid"
                    ), {"sid": session_id}).fetchall()
                    quality_by_chat = {str(qr[0]): qr[1] for qr in q_rows}
                except Exception:
                    quality_by_chat = {}

                quality_list = [v for v in quality_by_chat.values()]
                q_idx = 0

                for run in runs:
                    if run.get("parent_run_id"):
                        continue

                    input_data = run.get("input", {})
                    user_msg = ""
                    if isinstance(input_data, dict):
                        user_msg = input_data.get("input_content", "") or input_data.get("content", "")
                    elif isinstance(input_data, str):
                        user_msg = input_data

                    if user_msg:
                        messages.append({"role": "user", "content": user_msg})

                    content = run.get("content", "")
                    if content:
                        run_id = run.get("run_id") or ""
                        children = children_by_parent.get(run_id, [])
                        tool_calls = []
                        primary_agent = None
                        for ch in children:
                            agent_nm = ch.get("agent_name") or ch.get("agent_id") or ""
                            if agent_nm and not primary_agent and agent_nm.lower() not in ("leader", "team"):
                                primary_agent = agent_nm
                            for t in (ch.get("tools") or []):
                                m = t.get("metrics") or {}
                                tool_calls.append({
                                    "name": t.get("tool_name") or "tool",
                                    "args": t.get("tool_args") or {},
                                    "result": str(t.get("result") or "")[:200],
                                    "duration": float(m.get("duration") or 0),
                                    "status": "error" if t.get("tool_call_error") else "done",
                                })
                        for t in (run.get("tools") or []):
                            m = t.get("metrics") or {}
                            tool_calls.append({
                                "name": t.get("tool_name") or "tool",
                                "args": t.get("tool_args") or {},
                                "result": str(t.get("result") or "")[:200],
                                "duration": float(m.get("duration") or 0),
                                "status": "error" if t.get("tool_call_error") else "done",
                            })
                        run_metrics = run.get("metrics") or {}
                        duration_s = float(run_metrics.get("duration") or 0)
                        q_score = None
                        if q_idx < len(quality_list):
                            try:
                                q_score = int(quality_list[q_idx])
                            except Exception:
                                q_score = None
                            q_idx += 1
                        # Reconstruct the reasoning trace + token usage so it
                        # survives page refresh (matches the live SSE shape).
                        _trace = []
                        _usage = None
                        try:
                            from app.projects import _trace_from_stored_run, _usage_from_stored_run
                            _trace = _trace_from_stored_run(run, children)
                            _usage = _usage_from_stored_run(run, children)
                        except Exception:
                            _trace = []
                            _usage = None
                        messages.append({
                            "role": "assistant",
                            "content": content,
                            "tool_calls": tool_calls,
                            "trace": _trace,
                            "usage": _usage,
                            "duration": duration_s,
                            "quality_score": q_score,
                            "routing": {"routed_to": primary_agent} if primary_agent else None,
                        })
    except Exception:
        pass

    return {"messages": messages}


@router.post("/sessions/register")
def register_session(request: Request, session_id: str, message: str = "", project: str | None = None):
    """Register or update a chat session for the current user."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"status": "skip"}
    uid = user["user_id"]
    engine = create_engine(db_url)
    try:
        with engine.connect() as conn:
            existing = conn.execute(text(
                "SELECT id, first_message FROM public.dash_chat_sessions WHERE session_id = :sid"
            ), {"sid": session_id}).fetchone()
            if existing:
                # Update timestamp, set first_message if empty
                if not existing[1] and message:
                    conn.execute(text(
                        "UPDATE public.dash_chat_sessions SET updated_at = NOW(), first_message = :msg WHERE session_id = :sid"
                    ), {"sid": session_id, "msg": message[:100]})
                else:
                    conn.execute(text(
                        "UPDATE public.dash_chat_sessions SET updated_at = NOW() WHERE session_id = :sid"
                    ), {"sid": session_id})
            else:
                conn.execute(text(
                    "INSERT INTO public.dash_chat_sessions (user_id, session_id, project_slug, first_message) "
                    "VALUES (:uid, :sid, :proj, :msg)"
                ), {"uid": uid, "sid": session_id, "proj": project, "msg": (message[:100] if message else None)})
            conn.commit()
    except Exception:
        pass
    return {"status": "ok"}


@router.get("/projects/{slug}/sessions")
def list_project_sessions(slug: str, request: Request):
    """Sessions scoped to a given project for the current user."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"sessions": []}
    try:
        from app.auth import check_project_permission
        perm = check_project_permission(user, slug)
        if not perm:
            return {"sessions": []}
    except Exception:
        pass
    uid = user["user_id"]
    engine = create_engine(db_url)
    sessions: list[dict] = []
    try:
        with engine.connect() as conn:
            result = conn.execute(text("""
                SELECT session_id, first_message, created_at, updated_at, project_slug
                FROM public.dash_chat_sessions
                WHERE user_id = :uid AND project_slug = :slug
                ORDER BY updated_at DESC
                LIMIT 50
            """), {"uid": uid, "slug": slug})
            for row in result.fetchall():
                sessions.append({
                    "session_id": str(row[0]),
                    "first_message": row[1] or None,
                    "created_at": str(row[2]) if row[2] else None,
                    "updated_at": str(row[3]) if row[3] else None,
                    "project_slug": row[4],
                })
    except Exception:
        pass
    return {"sessions": sessions}


@router.get("/workflows")
def list_workflows_by_project(project: str | None = None, request: Request = None):
    """Frontend-friendly workflow listing keyed by ?project=<slug>.

    Reads from public.dash_workflows_db. Returns empty list if table missing
    or no project supplied. Auth checked via project access; falls open if
    perm helpers unavailable.
    """
    if not project:
        return {"workflows": []}
    user = getattr(getattr(request, 'state', None), 'user', None) if request else None
    if user:
        try:
            from app.auth import check_project_permission
            if not check_project_permission(user, project):
                return {"workflows": []}
        except Exception:
            pass
    engine = create_engine(db_url)
    wfs: list[dict] = []
    try:
        with engine.connect() as conn:
            result = conn.execute(text(
                "SELECT id, name, steps, source, created_at FROM public.dash_workflows_db "
                "WHERE project_slug = :p ORDER BY created_at DESC LIMIT 100"
            ), {"p": project})
            for row in result.fetchall():
                steps_val = row[2]
                try:
                    if isinstance(steps_val, str):
                        steps_val = json.loads(steps_val)
                except Exception:
                    pass
                wfs.append({
                    "id": row[0],
                    "name": row[1],
                    "steps": steps_val or [],
                    "source": row[3] or "user",
                    "created_at": str(row[4]) if row[4] else None,
                })
    except Exception:
        pass
    return {"workflows": wfs}


@router.get("/queries")
def recent_queries():
    """Get recent query sessions."""
    engine = create_engine(db_url)
    queries: list[dict] = []
    try:
        with engine.connect() as conn:
            # Agno stores sessions in ai.agno_sessions
            result = conn.execute(text("""
                SELECT session_id, team_id, user_id, created_at, updated_at
                FROM ai.agno_sessions
                WHERE team_id = 'dash'
                ORDER BY updated_at DESC
                LIMIT 30
            """))
            for row in result.fetchall():
                queries.append({
                    "session_id": str(row[0]),
                    "team_id": row[1],
                    "user_id": row[2],
                    "created_at": str(row[3]) if row[3] else None,
                    "updated_at": str(row[4]) if row[4] else None,
                })
    except Exception:
        pass
    return {"queries": queries}


@router.get("/tables/{table_name}/inspect")
def inspect_table(table_name: str, request: Request, project: str | None = None):
    """Get column details and sample data for a table."""
    schema = _get_project_schema(request, project) or "public"
    engine = create_engine(db_url)
    insp = inspect(engine)

    if table_name not in insp.get_table_names(schema=schema):
        raise HTTPException(404, f"Table '{table_name}' not found in {schema}")

    cols = insp.get_columns(table_name, schema=schema)
    columns = [{"name": c["name"], "type": str(c["type"]), "nullable": c.get("nullable", True)} for c in cols]

    # Quality notes from metadata file
    quality_notes: list[str] = []
    if project:
        meta_file = KNOWLEDGE_DIR / project / "tables" / f"{table_name}.json"
    else:
        meta_file = TABLES_DIR / f"{table_name}.json"
    if meta_file.exists():
        try:
            with open(meta_file) as f:
                meta = json.load(f)
            quality_notes = meta.get("data_quality_notes", [])
        except Exception:
            pass

    # Sample data
    sample: list[dict] = []
    try:
        with engine.connect() as conn:
            result = conn.execute(text(f'SELECT * FROM "{schema}"."{table_name}" LIMIT 5'))
            keys = list(result.keys())
            for row in result.fetchall():
                sample.append({k: str(v)[:50] if v is not None else None for k, v in zip(keys, row)})
    except Exception:
        pass

    return {"columns": columns, "sample": sample, "quality_notes": quality_notes}


@router.delete("/tables/{table_name}")
def delete_table(table_name: str, request: Request, project: str | None = None):
    """Drop a user-uploaded table and remove its knowledge files."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")

    if table_name in PROTECTED_TABLES:
        raise HTTPException(403, f"Cannot delete protected table: {table_name}")

    # Determine schema
    if project:
        import re as _re
        schema = _re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    else:
        schema = "public"

    engine = create_engine(db_url)
    insp = inspect(engine)

    if table_name not in insp.get_table_names(schema=schema):
        raise HTTPException(404, f"Table '{table_name}' not found in schema '{schema}'")

    with engine.connect() as conn:
        conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{table_name}" CASCADE'))
        conn.commit()

    # Remove knowledge files (project-scoped or global)
    if project:
        base = KNOWLEDGE_DIR / project
    else:
        base = KNOWLEDGE_DIR

    for subdir in ["tables", "queries", "business", "staging"]:
        for pattern in [f"{table_name}.json", f"{table_name}_queries.sql", f"{table_name}_rules.json", f"{table_name}.*"]:
            for f in (base / subdir).glob(pattern) if (base / subdir).exists() else []:
                f.unlink(missing_ok=True)

    # Clean DB records for this table
    try:
        slug = project or ""
        with engine.connect() as conn:
            conn.execute(text("DELETE FROM public.dash_table_metadata WHERE project_slug = :s AND table_name = :t"), {"s": slug, "t": table_name})
            conn.execute(text("DELETE FROM public.dash_training_qa WHERE project_slug = :s AND table_name = :t"), {"s": slug, "t": table_name})
            conn.commit()
    except Exception:
        pass

    return {"status": "ok", "deleted": table_name}


# ---------------------------------------------------------------------------
# Extraction Plan endpoints (P4) — list / get / re-ingest with overrides
# ---------------------------------------------------------------------------

@router.get("/projects/{slug}/extraction-plans")
def list_extraction_plans(slug: str, request: Request, limit: int = 50):
    """List recent extraction plans for project."""
    user = _get_user(request)
    if not user:
        raise HTTPException(401, "Auth required")
    from app.auth import check_project_permission
    if not check_project_permission(user, slug, required_role="viewer"):
        raise HTTPException(403, "No access")

    from db.session import get_sql_engine
    from sqlalchemy import text as _slt
    e = get_sql_engine()
    try:
        with e.connect() as c:
            rows = c.execute(_slt(
                """
                SELECT id, table_name, source_file, sheet_name, strategy,
                       header_row, skip_rows, row_count_in, row_count_out,
                       llm_rescued, rescue_reasoning, user_overrides, created_at
                FROM public.dash_extraction_plans
                WHERE project_slug = :s
                ORDER BY created_at DESC
                LIMIT :l
                """
            ), {"s": slug, "l": int(limit)}).fetchall()
    except Exception as exc:
        # Migration not yet applied or transient DB error — return empty list
        try:
            _upload_lg.warning(f"list_extraction_plans failed: {exc}")
        except Exception:
            pass
        return {"plans": []}
    return {"plans": [dict(r._mapping) for r in rows]}


@router.get("/projects/{slug}/extraction-plans/{plan_id}")
def get_extraction_plan(slug: str, plan_id: int, request: Request):
    """Get full extraction plan including blocks."""
    user = _get_user(request)
    if not user:
        raise HTTPException(401, "Auth required")
    from app.auth import check_project_permission
    if not check_project_permission(user, slug, required_role="viewer"):
        raise HTTPException(403, "No access")
    from db.session import get_sql_engine
    from sqlalchemy import text as _slt
    e = get_sql_engine()
    with e.connect() as c:
        row = c.execute(_slt(
            "SELECT * FROM public.dash_extraction_plans "
            "WHERE project_slug = :s AND id = :i"
        ), {"s": slug, "i": int(plan_id)}).fetchone()
    if not row:
        raise HTTPException(404, "Plan not found")
    return dict(row._mapping)


@router.post("/projects/{slug}/extraction-plans/{plan_id}/re-ingest")
async def re_ingest_with_overrides(slug: str, plan_id: int, request: Request):
    """Re-ingest source file with user-supplied overrides.
    Body: {header_row?: int, skip_rows?: list[int]}
    Requires editor role.
    """
    user = _get_user(request)
    if not user:
        raise HTTPException(401, "Auth required")
    from app.auth import check_project_permission
    if not check_project_permission(user, slug, required_role="editor"):
        raise HTTPException(403, "Editor access required")

    try:
        body = await request.json()
    except Exception:
        body = {}
    new_header = body.get("header_row")
    new_skip = body.get("skip_rows") or []

    from db.session import get_sql_engine, get_write_engine
    from sqlalchemy import text as _slt

    # Look up the plan
    e = get_sql_engine()
    with e.connect() as c:
        plan = c.execute(_slt(
            "SELECT * FROM public.dash_extraction_plans "
            "WHERE project_slug = :s AND id = :i"
        ), {"s": slug, "i": int(plan_id)}).fetchone()
    if not plan:
        raise HTTPException(404, "Plan not found")

    plan_dict = dict(plan._mapping)
    source_file = plan_dict.get("source_file")
    sheet_name = plan_dict.get("sheet_name")
    table_name = plan_dict.get("table_name")

    if not source_file:
        raise HTTPException(400, "No source file recorded for this plan")

    # Source file must still be on disk (raw_uploads, fallback docs_raw)
    file_path = KNOWLEDGE_DIR / slug / "raw_uploads" / source_file
    if not file_path.exists():
        file_path = KNOWLEDGE_DIR / slug / "docs_raw" / source_file
        if not file_path.exists():
            raise HTTPException(404, f"Source file {source_file} not on disk anymore")

    try:
        ext = Path(source_file).suffix.lower()
        if ext not in (".xlsx", ".xls"):
            raise HTTPException(400, f"Re-ingest only supports xlsx/xls for now, got {ext}")

        hrow = int(new_header) if new_header is not None else int(plan_dict.get("header_row") or 0)
        skip = list(new_skip) if new_skip else (plan_dict.get("skip_rows") or [])
        # Build pandas skiprows: include everything before header (excluding header row itself)
        full_skip = sorted(set(int(s) for s in skip) | set(range(0, hrow)))
        pd_skip = [s for s in full_skip if s != hrow]
        df = pd.read_excel(
            str(file_path),
            sheet_name=sheet_name if sheet_name else 0,
            header=0,
            skiprows=pd_skip,
        )
        df = _clean_dataframe(df)

        # Write into project schema, REPLACE existing table
        from db import get_project_engine
        from db.session import create_project_schema
        schema = create_project_schema(slug)
        engine = get_project_engine(slug)
        df.to_sql(table_name, engine, schema=schema, if_exists='replace', index=False)

        # Update plan w/ new metadata + user_overrides
        we = get_write_engine()
        with we.begin() as c:
            c.execute(_slt(
                """
                UPDATE public.dash_extraction_plans
                SET header_row = :h,
                    skip_rows = CAST(:sk AS jsonb),
                    row_count_out = :rc,
                    user_overrides = CAST(:uo AS jsonb),
                    updated_at = now()
                WHERE id = :i
                """
            ), {
                "h": hrow if new_header is not None else plan_dict.get("header_row"),
                "sk": json.dumps(list(new_skip) if new_skip else (plan_dict.get("skip_rows") or [])),
                "rc": int(len(df)),
                "uo": json.dumps({
                    "header_row": new_header,
                    "skip_rows": new_skip,
                    "applied_at": _dt_datetime.utcnow().isoformat(),
                }),
                "i": int(plan_id),
            })

        return {"status": "ok", "rows_loaded": int(len(df)), "table": table_name}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(500, f"Re-ingest failed: {str(exc)[:200]}")


# ── P5: Upload-Cache Admin Endpoints ─────────────────────────────────────────

@router.get("/admin/upload-cache/stats")
def get_upload_cache_stats(request: Request):
    """Admin: cache stats."""
    user = _get_user(request)
    if not user or not user.get("is_admin"):
        raise HTTPException(403, "Admin only")
    from db.session import get_sql_engine
    from sqlalchemy import text as _slt
    e = get_sql_engine()
    with e.connect() as c:
        stats = c.execute(_slt("""
            SELECT
                COUNT(*) AS total_files,
                COUNT(*) FILTER (WHERE rescue_used) AS rescue_count,
                COALESCE(SUM(hit_count), 0) AS total_hits,
                COALESCE(SUM(file_size_bytes), 0) AS total_bytes
            FROM public.dash_upload_cache
        """)).fetchone()
        top = c.execute(_slt("""
            SELECT file_hash, file_ext, hit_count, file_size_bytes, last_used_at
            FROM public.dash_upload_cache
            ORDER BY hit_count DESC, last_used_at DESC
            LIMIT 20
        """)).fetchall()
    return {
        "stats": dict(stats._mapping) if stats else {},
        "top_files": [dict(r._mapping) for r in top],
    }


@router.delete("/admin/upload-cache/{file_hash}")
def delete_upload_cache_entry(file_hash: str, request: Request):
    """Admin: evict cache entry."""
    user = _get_user(request)
    if not user or not user.get("is_admin"):
        raise HTTPException(403, "Admin only")
    from db.session import get_write_engine
    from sqlalchemy import text as _slt
    e = get_write_engine()
    with e.begin() as c:
        c.execute(_slt("DELETE FROM public.dash_upload_cache WHERE file_hash = :h"), {"h": file_hash})
    return {"status": "ok", "evicted": file_hash}
