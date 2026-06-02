"""Dash-OS Phase 2A — FileGenerationTools.

Wraps PDF/PPTX/CSV/XLSX/DOCX/JSON/MD generation as Agno @tool functions so
any agent can produce downloadable files mid-conversation.

All outputs registered in dash.dash_generated_files; downloadable via
GET /api/reporter/files/{file_id}.

Gating: tool functions always run. Behind EXPERIMENTAL_AGI=1 the Reporter
agent is registered in the default team (Phase 8); otherwise must be
invoked directly.
"""
from __future__ import annotations

import io
import json as _json
import logging
import os
import secrets
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

def _root() -> Path:
    return Path(os.getenv("DASH_GENERATED_ROOT", "/app/knowledge/generated"))


# ── ContextVar pull (fail-soft) ──────────────────────────────────────────
def _ctx() -> Dict[str, Any]:
    try:
        from dash.agentic.hooks import (
            current_project_slug,
            current_user_id,
            current_run_id,
            current_agent_name,
        )
        return {
            "project_slug": current_project_slug.get(),
            "user_id": current_user_id.get(),
            "run_id": current_run_id.get(),
            "agent_name": current_agent_name.get(),
        }
    except Exception:
        return {"project_slug": None, "user_id": None, "run_id": None, "agent_name": None}


def _get_engine():
    try:
        from db.session import get_sql_engine
        return get_sql_engine()
    except Exception:
        try:
            from db import get_sql_engine  # type: ignore
            return get_sql_engine()
        except Exception:
            return None


def _save_and_register(
    file_type: str,
    filename: str,
    data: bytes,
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    file_id = "gen_" + secrets.token_hex(4)
    ctx = _ctx()
    scope = ctx["project_slug"] or "_global"
    dirpath = _root() / scope
    dirpath.mkdir(parents=True, exist_ok=True)
    ext = filename.split(".")[-1] if "." in filename else file_type
    storage_path = dirpath / f"{file_id}.{ext}"
    storage_path.write_bytes(data)

    eng = _get_engine()
    if eng is not None:
        try:
            from sqlalchemy import text
            with eng.begin() as conn:
                conn.execute(
                    text(
                        """
                        INSERT INTO dash.dash_generated_files
                          (id, project_slug, user_id, agent_name, run_id,
                           file_type, filename, storage_path, size_bytes, metadata)
                        VALUES
                          (:id, :ps, :uid, :ag, :rid,
                           :ft, :fn, :sp, :sz, CAST(:md AS jsonb))
                        ON CONFLICT (id) DO NOTHING
                        """
                    ),
                    {
                        "id": file_id,
                        "ps": ctx["project_slug"],
                        "uid": ctx["user_id"],
                        "ag": ctx["agent_name"],
                        "rid": ctx["run_id"],
                        "ft": file_type,
                        "fn": filename,
                        "sp": str(storage_path),
                        "sz": len(data),
                        "md": _json.dumps(metadata or {}),
                    },
                )
        except Exception as e:  # pragma: no cover
            logger.warning("dash_generated_files insert failed: %s", e)

    return {
        "ok": True,
        "file_id": file_id,
        "filename": filename,
        "file_type": file_type,
        "size_bytes": len(data),
        "download_url": f"/api/reporter/files/{file_id}",
    }


# ── Tool wrappers (graceful imports) ────────────────────────────────────
def _try_tool(fn):
    """Wrap as Agno @tool if available, else return raw fn."""
    try:
        from agno.tools import tool
        return tool(fn)
    except Exception:
        return fn


def _make_csv(rows: List[Dict[str, Any]], filename: Optional[str] = None) -> Dict[str, Any]:
    try:
        import pandas as pd
    except Exception as e:
        return {"ok": False, "error": f"missing dep: pandas ({e})"}
    df = pd.DataFrame(rows)
    buf = io.StringIO()
    df.to_csv(buf, index=False)
    return _save_and_register("csv", filename or "data.csv", buf.getvalue().encode("utf-8"))


def _make_json(payload: Any, filename: Optional[str] = None) -> Dict[str, Any]:
    body = _json.dumps(payload, indent=2, default=str).encode("utf-8")
    return _save_and_register("json", filename or "data.json", body)


def _make_md(title: str, body: str, filename: Optional[str] = None) -> Dict[str, Any]:
    text = f"# {title}\n\n{body}\n"
    return _save_and_register("md", filename or f"{title[:40].replace(' ', '_')}.md", text.encode("utf-8"))


def _make_xlsx(sheets: Dict[str, List[Dict[str, Any]]], filename: Optional[str] = None) -> Dict[str, Any]:
    try:
        import pandas as pd
        import xlsxwriter  # noqa: F401
    except Exception as e:
        return {"ok": False, "error": f"missing dep: pandas/xlsxwriter ({e})"}
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        for sheet_name, rows in sheets.items():
            pd.DataFrame(rows).to_excel(writer, sheet_name=sheet_name[:31], index=False)
    return _save_and_register("xlsx", filename or "workbook.xlsx", buf.getvalue())


def _make_pdf(title: str, sections: List[Dict[str, Any]], filename: Optional[str] = None) -> Dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    except Exception as e:
        # fallback: HTML pretending to be PDF
        html = f"<html><head><title>{title}</title></head><body><h1>{title}</h1>"
        for s in sections:
            html += f"<h2>{s.get('heading','')}</h2><p>{s.get('body','')}</p>"
        html += "</body></html>"
        return _save_and_register("pdf", filename or "report.pdf", html.encode("utf-8"), {"fallback": "html", "reason": str(e)})
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for s in sections:
        if s.get("heading"):
            story.append(Paragraph(s["heading"], styles["Heading2"]))
        if s.get("body"):
            story.append(Paragraph(s["body"], styles["BodyText"]))
        story.append(Spacer(1, 8))
    doc.build(story)
    return _save_and_register("pdf", filename or "report.pdf", buf.getvalue())


def _make_pptx(
    title: str,
    slides: List[Dict[str, Any]],
    theme: str = "midnight_executive",
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        return {"ok": False, "error": f"missing dep: python-pptx ({e})"}
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title
    for slide in slides:
        content_layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = str(slide.get("title", ""))
        body = s.placeholders[1].text_frame
        body.text = str(slide.get("content", ""))
    buf = io.BytesIO()
    prs.save(buf)
    return _save_and_register(
        "pptx", filename or "deck.pptx", buf.getvalue(), {"theme": theme, "slides": len(slides)}
    )


def _make_docx(title: str, sections: List[Dict[str, Any]], filename: Optional[str] = None) -> Dict[str, Any]:
    try:
        from docx import Document
    except Exception as e:
        return {"ok": False, "error": f"missing dep: python-docx ({e})"}
    doc = Document()
    doc.add_heading(title, level=0)
    for s in sections:
        if s.get("heading"):
            doc.add_heading(s["heading"], level=1)
        if s.get("body"):
            doc.add_paragraph(s["body"])
    buf = io.BytesIO()
    doc.save(buf)
    return _save_and_register("docx", filename or "report.docx", buf.getvalue())


def _list_generated_files(project_slug: Optional[str] = None, limit: int = 20) -> Dict[str, Any]:
    eng = _get_engine()
    if eng is None:
        return {"ok": False, "error": "db unavailable", "files": []}
    try:
        from sqlalchemy import text
        with eng.connect() as conn:
            rows = conn.execute(
                text(
                    """
                    SELECT id, project_slug, file_type, filename, size_bytes, created_at
                    FROM dash.dash_generated_files
                    WHERE (:ps IS NULL OR project_slug = :ps)
                      AND (expires_at IS NULL OR expires_at > now())
                    ORDER BY created_at DESC
                    LIMIT :lim
                    """
                ),
                {"ps": project_slug, "lim": limit},
            ).mappings().all()
        return {"ok": True, "files": [dict(r) for r in rows]}
    except Exception as e:
        return {"ok": False, "error": str(e), "files": []}


# Public exported Agno tools
make_csv = _try_tool(_make_csv)
make_json = _try_tool(_make_json)
make_md = _try_tool(_make_md)
make_xlsx = _try_tool(_make_xlsx)
make_pdf = _try_tool(_make_pdf)
make_pptx = _try_tool(_make_pptx)
make_docx = _try_tool(_make_docx)
list_generated_files = _try_tool(_list_generated_files)
