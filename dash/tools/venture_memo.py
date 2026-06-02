"""
Venture IC memo export — PDF + PPTX.

PDF: reportlab platypus, 1-page exec summary + DCF table + sensitivity grid.
PPTX: python-pptx, 6-slide deck (cover / thesis / financials / sensitivity / risks / verdict).

Pulls deal + scenarios from dash.dash_venture_* tables. Fail-soft.
"""
from __future__ import annotations

import io
import json
import logging
from datetime import datetime
from typing import Optional

from sqlalchemy import text

logger = logging.getLogger(__name__)

# Supply Chain Sentry — Sprint 3. Graceful skip when backend not deployed.
try:
    from dash.tools.supply_tools import summarize_supply_for_memo  # type: ignore
except Exception:
    summarize_supply_for_memo = None  # type: ignore


# ---------- helpers ----------

def _engine():
    try:
        from db.session import get_sql_engine
        return get_sql_engine()
    except Exception as e:
        raise RuntimeError(f"db unavailable: {e}")


def _fetch_deal(deal_id: str) -> dict:
    """Fetch deal + most-recent scenarios. Returns {deal, scenarios[]}."""
    eng = _engine()
    with eng.connect() as cx:
        d = cx.execute(text("""
            SELECT id, project_slug, name, stage, sector, geography,
                   ask_amount, pre_money, post_money, status, created_at
            FROM dash.dash_venture_deals WHERE id = :i
        """), {"i": deal_id}).fetchone()
        if not d:
            return {}
        scens = cx.execute(text("""
            SELECT id, name, irr, moic, payback_yrs, npv, inputs, verdict, created_at
            FROM dash.dash_venture_scenarios
            WHERE deal_id = :d
            ORDER BY created_at DESC LIMIT 10
        """), {"d": deal_id}).fetchall()
        comps = cx.execute(text("""
            SELECT name, share_pct, moat, source
            FROM dash.dash_venture_competitors WHERE deal_id = :d LIMIT 20
        """), {"d": deal_id}).fetchall()

    return {
        "deal": {
            "id": str(d[0]), "project_slug": d[1], "name": d[2], "stage": d[3],
            "sector": d[4], "geography": d[5], "ask_amount": float(d[6]) if d[6] else None,
            "pre_money": float(d[7]) if d[7] else None,
            "post_money": float(d[8]) if d[8] else None,
            "status": d[9], "created_at": d[10].isoformat() if d[10] else None,
        },
        "scenarios": [
            {
                "id": str(s[0]), "name": s[1],
                "irr": float(s[2]) if s[2] is not None else None,
                "moic": float(s[3]) if s[3] is not None else None,
                "payback_yrs": float(s[4]) if s[4] is not None else None,
                "npv": float(s[5]) if s[5] is not None else None,
                "inputs": s[6] if isinstance(s[6], dict) else (json.loads(s[6]) if s[6] else {}),
                "verdict": s[7],
                "created_at": s[8].isoformat() if s[8] else None,
            } for s in scens
        ],
        "competitors": [
            {"name": c[0], "share_pct": float(c[1]) if c[1] is not None else None,
             "moat": c[2], "source": c[3]}
            for c in comps
        ],
    }


def _fmt_money(n):
    if n is None: return "—"
    a = abs(n)
    if a >= 1e9: return f"{n/1e9:.2f}B"
    if a >= 1e6: return f"{n/1e6:.2f}M"
    if a >= 1e3: return f"{n/1e3:.1f}K"
    return f"{n:.0f}"


def _fmt_pct(n):
    if n is None: return "—"
    return f"{n*100:.1f}%"


def _verdict_color_hex(v: Optional[str]) -> str:
    v = (v or "hold").lower()
    if v == "go": return "#2c7a3f"
    if v == "pass": return "#c0392b"
    return "#c96342"  # hold / coral


def _render_sensitivity_heatmap(scenarios: list[dict]) -> bytes:
    """Render NPV sensitivity heatmap (WACC x Terminal Growth) as PNG bytes.

    Returns empty bytes if <2 scenarios with inputs.wacc or inputs.cashflows,
    or on any failure (matplotlib unavailable, etc).
    """
    try:
        # gate: need >=2 scenarios with wacc or cashflows in inputs
        eligible = [
            s for s in (scenarios or [])
            if isinstance(s.get("inputs"), dict)
            and (s["inputs"].get("wacc") is not None or s["inputs"].get("cashflows"))
        ]
        if len(eligible) < 2:
            return b""

        base = eligible[0]
        inputs = base.get("inputs") or {}
        cashflows = inputs.get("cashflows")
        if not cashflows or not isinstance(cashflows, list):
            return b""

        from dash.tools.venture_tools import sensitivity_grid

        wacc_range = [0.08, 0.10, 0.12, 0.14, 0.16]
        growth_range = [0.01, 0.02, 0.03, 0.04, 0.05]

        result = sensitivity_grid(cashflows, wacc_range, growth_range)
        if not result or not result.get("ok"):
            return b""
        grid = result.get("grid") or []
        if not grid:
            return b""

        # Replace None entries (invalid g>=w combos) w/ NaN so imshow renders gray
        import numpy as _np
        grid_arr = _np.array(
            [[(v if v is not None else _np.nan) for v in row] for row in grid],
            dtype=float,
        )

        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(6, 4), dpi=100)
        im = ax.imshow(grid_arr, cmap="RdYlGn", aspect="auto")

        ax.set_xticks(range(len(growth_range)))
        ax.set_xticklabels([f"{g*100:.0f}%" for g in growth_range])
        ax.set_yticks(range(len(wacc_range)))
        ax.set_yticklabels([f"{w*100:.0f}%" for w in wacc_range])
        ax.set_xlabel("Terminal Growth")
        ax.set_ylabel("WACC")
        ax.set_title("Sensitivity: NPV vs WACC × Terminal Growth")

        # annotate cells
        for i in range(len(wacc_range)):
            for j in range(len(growth_range)):
                try:
                    val = grid[i][j]
                except (IndexError, TypeError):
                    continue
                ax.text(j, i, _fmt_money(val), ha="center", va="center",
                        fontsize=8, color="#1a1614")

        fig.colorbar(im, ax=ax, label="NPV")
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", dpi=100)
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception as e:
        logger.debug(f"sensitivity heatmap render failed: {e}")
        return b""


# ---------- PDF ----------

def generate_pdf(deal_id: str) -> bytes:
    """Generate 1-page IC memo PDF. Returns bytes."""
    data = _fetch_deal(deal_id)
    if not data:
        raise ValueError(f"deal {deal_id} not found")

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
        )
    except Exception as e:
        raise RuntimeError(f"reportlab unavailable: {e}")

    deal = data["deal"]
    scens = data["scenarios"]
    base = scens[0] if scens else {}

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=18*mm, rightMargin=18*mm,
                            topMargin=18*mm, bottomMargin=18*mm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=styles["Heading1"],
                        fontSize=18, textColor=colors.HexColor("#1a1614"),
                        spaceAfter=4)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"],
                        fontSize=11, textColor=colors.HexColor("#c96342"),
                        spaceAfter=4, spaceBefore=8)
    body = ParagraphStyle("body", parent=styles["BodyText"],
                          fontSize=9.5, textColor=colors.HexColor("#2c2a26"),
                          leading=13)
    muted = ParagraphStyle("muted", parent=body, textColor=colors.HexColor("#6b6557"),
                           fontSize=8.5)

    story = []

    # Header
    story.append(Paragraph(f"IC Memo — {deal['name']}", h1))
    story.append(Paragraph(
        f"{deal.get('stage','').replace('_',' ').title()} · "
        f"{deal.get('sector','—')} · {deal.get('geography','—')} · "
        f"Ask {_fmt_money(deal.get('ask_amount'))} · "
        f"Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
        muted))
    story.append(Spacer(1, 8))

    # Verdict + KPIs
    verdict = (base.get("verdict") or "hold").upper()
    vcolor = colors.HexColor(_verdict_color_hex(base.get("verdict")))
    kpi_data = [
        [Paragraph(f"<b>VERDICT</b><br/><font color='{_verdict_color_hex(base.get('verdict'))}' size='14'><b>{verdict}</b></font>", body),
         Paragraph(f"<b>IRR</b><br/><font size='14'>{_fmt_pct(base.get('irr'))}</font>", body),
         Paragraph(f"<b>MOIC</b><br/><font size='14'>{base.get('moic',0):.2f}×</font>" if base.get("moic") is not None else "<b>MOIC</b><br/>—", body),
         Paragraph(f"<b>NPV</b><br/><font size='14'>{_fmt_money(base.get('npv'))}</font>", body),
         Paragraph(f"<b>PAYBACK</b><br/><font size='14'>{base.get('payback_yrs','—')} yrs</font>", body)],
    ]
    t = Table(kpi_data, colWidths=[34*mm, 34*mm, 34*mm, 34*mm, 34*mm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#f7f6f3")),
        ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#d6cfbe")),
        ("INNERGRID", (0,0), (-1,-1), 0.5, colors.HexColor("#d6cfbe")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("PADDING", (0,0), (-1,-1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 10))

    # Investment thesis
    story.append(Paragraph("Investment thesis", h2))
    thesis = (
        f"{deal['name']} is a {deal.get('stage','').replace('_',' ')} "
        f"opportunity in {deal.get('sector','—')} ({deal.get('geography','—')}). "
        f"Ask of {_fmt_money(deal.get('ask_amount'))} at "
        f"pre-money {_fmt_money(deal.get('pre_money'))}."
    )
    story.append(Paragraph(thesis, body))
    story.append(Spacer(1, 4))

    # Financial model
    story.append(Paragraph("Financial model", h2))
    if base:
        inputs = base.get("inputs") or {}
        story.append(Paragraph(
            f"Scenario: <b>{base.get('name','base')}</b> · "
            f"WACC {_fmt_pct(inputs.get('wacc') or inputs.get('WACC'))} · "
            f"Terminal growth {_fmt_pct(inputs.get('terminal_growth') or inputs.get('tg'))}",
            body))

    # Scenarios table
    if scens:
        rows = [["Scenario", "IRR", "MOIC", "Payback", "NPV", "Verdict"]]
        for s in scens[:5]:
            rows.append([
                s.get("name","—"),
                _fmt_pct(s.get("irr")),
                f"{s.get('moic',0):.2f}×" if s.get("moic") is not None else "—",
                f"{s.get('payback_yrs','—')} yrs",
                _fmt_money(s.get("npv")),
                (s.get("verdict") or "hold").upper(),
            ])
        st = Table(rows, colWidths=[40*mm, 25*mm, 22*mm, 25*mm, 30*mm, 25*mm])
        st.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1a1614")),
            ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("FONTSIZE", (0,0), (-1,-1), 9),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("ALIGN", (1,0), (-1,-1), "RIGHT"),
            ("ALIGN", (0,0), (0,-1), "LEFT"),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#f7f6f3")]),
            ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#d6cfbe")),
            ("INNERGRID", (0,0), (-1,-1), 0.25, colors.HexColor("#e8e3d6")),
            ("PADDING", (0,0), (-1,-1), 5),
        ]))
        story.append(Spacer(1, 4))
        story.append(st)
        story.append(Spacer(1, 8))

    # Sensitivity heatmap
    try:
        img_bytes = _render_sensitivity_heatmap(scens)
        if img_bytes:
            from reportlab.platypus import Image as RLImage
            story.append(Paragraph("Sensitivity analysis", h2))
            img = RLImage(io.BytesIO(img_bytes), width=160*mm, height=107*mm, kind='proportional')
            story.append(img)
            story.append(Spacer(1, 8))
    except Exception as e:
        logger.debug(f"pdf sensitivity embed failed: {e}")

    # Competitors
    if data["competitors"]:
        story.append(Paragraph("Competitive landscape", h2))
        rows = [["Name", "Share %", "Moat", "Source"]]
        for c in data["competitors"][:10]:
            rows.append([
                c.get("name","—"),
                _fmt_pct((c.get("share_pct") or 0)/100) if c.get("share_pct") else "—",
                (c.get("moat") or "—")[:40],
                (c.get("source") or "—")[:20],
            ])
        ct = Table(rows, colWidths=[50*mm, 25*mm, 60*mm, 30*mm])
        ct.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#e8e3d6")),
            ("FONTSIZE", (0,0), (-1,-1), 8.5),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("VALIGN", (0,0), (-1,-1), "TOP"),
            ("PADDING", (0,0), (-1,-1), 4),
            ("INNERGRID", (0,0), (-1,-1), 0.25, colors.HexColor("#e8e3d6")),
        ]))
        story.append(ct)
        story.append(Spacer(1, 8))

    # Market (Market Sentinel — sprint 2). Graceful skip if no data.
    try:
        from dash.tools.market_tools import summarize_market_for_memo
        msum = summarize_market_for_memo(deal['id'])
    except Exception as _e:
        logger.debug(f"market summary skipped: {_e}")
        msum = None
    if msum and msum.get("ok") and not msum.get("empty"):
        story.append(Paragraph("Market", h2))
        # TAM/SAM/SOM card
        if msum.get("tam_usd") or msum.get("sam_usd") or msum.get("som_usd"):
            tam_data = [[
                Paragraph(f"<b>TAM</b><br/><font size='12'>{_fmt_money(msum.get('tam_usd'))}</font>", body),
                Paragraph(f"<b>SAM</b><br/><font size='12'>{_fmt_money(msum.get('sam_usd'))}</font>", body),
                Paragraph(f"<b>SOM</b><br/><font size='12'>{_fmt_money(msum.get('som_usd'))}</font>", body),
                Paragraph(f"<b>Method</b><br/>{msum.get('methodology') or '—'}", body),
            ]]
            tt = Table(tam_data, colWidths=[42*mm, 42*mm, 42*mm, 42*mm])
            tt.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#f0ebde")),
                ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#c96342")),
                ("INNERGRID", (0,0), (-1,-1), 0.3, colors.HexColor("#d6cfbe")),
                ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                ("PADDING", (0,0), (-1,-1), 6),
            ]))
            story.append(tt)
            assumptions = msum.get("assumptions") or {}
            conf = assumptions.get("confidence")
            if conf is not None:
                story.append(Paragraph(
                    f"Methodology: {msum.get('methodology') or '—'} · "
                    f"Confidence: {conf:.2f}", muted))
            story.append(Spacer(1, 6))
        # Top 3 competitors.
        comps = msum.get("top_competitors") or []
        if comps:
            mc_rows = [["Competitor", "Share %", "Geography"]]
            for c in comps[:3]:
                mc_rows.append([
                    c.get("name", "—"),
                    f"{c.get('share_pct'):.1f}%" if c.get("share_pct") is not None else "—",
                    c.get("geography") or "—",
                ])
            mct = Table(mc_rows, colWidths=[70*mm, 35*mm, 60*mm])
            mct.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#e8e3d6")),
                ("FONTSIZE", (0,0), (-1,-1), 8.5),
                ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                ("PADDING", (0,0), (-1,-1), 4),
                ("INNERGRID", (0,0), (-1,-1), 0.25, colors.HexColor("#e8e3d6")),
            ]))
            story.append(mct)
            story.append(Spacer(1, 4))
        # Recent signals.
        sigs = msum.get("key_signals") or []
        if sigs:
            story.append(Paragraph("Recent market signals", muted))
            for s in sigs[:5]:
                story.append(Paragraph(
                    f"• <b>{s.get('signal_type', '?')}</b> — "
                    f"{(s.get('title') or '')[:140]}", body))
        story.append(Spacer(1, 8))

    # Supply Chain Risk (Supply Chain Sentry — sprint 3). Graceful skip if no data.
    ssum = None
    if summarize_supply_for_memo is not None:
        try:
            ssum = summarize_supply_for_memo(deal.get('project_slug') or deal['id'])
        except Exception as _e:
            logger.debug(f"supply summary skipped: {_e}")
            ssum = None
    if ssum and not ssum.get("empty"):
        story.append(Paragraph("Supply Chain Risk", h2))
        rollup = ssum.get("rollup") or {}
        if rollup:
            sr_data = [[
                Paragraph(f"<b>GREEN</b><br/><font size='14' color='#2c7a3f'><b>{rollup.get('green', 0)}</b></font>", body),
                Paragraph(f"<b>YELLOW</b><br/><font size='14' color='#c96342'><b>{rollup.get('yellow', 0)}</b></font>", body),
                Paragraph(f"<b>RED</b><br/><font size='14' color='#c0392b'><b>{rollup.get('red', 0)}</b></font>", body),
            ]]
            srt = Table(sr_data, colWidths=[56*mm, 56*mm, 56*mm])
            srt.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#f7f6f3")),
                ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#d6cfbe")),
                ("INNERGRID", (0,0), (-1,-1), 0.3, colors.HexColor("#d6cfbe")),
                ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                ("PADDING", (0,0), (-1,-1), 8),
            ]))
            story.append(srt)
            story.append(Spacer(1, 6))
        top_risks = ssum.get("top_risks") or []
        if top_risks:
            tr_rows = [["Supplier", "Country", "Score"]]
            for s in top_risks[:3]:
                tr_rows.append([
                    s.get("name", "—"),
                    s.get("country", "—"),
                    f"{s.get('score'):.1f}" if s.get("score") is not None else "—",
                ])
            trt = Table(tr_rows, colWidths=[80*mm, 50*mm, 38*mm])
            trt.setStyle(TableStyle([
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#e8e3d6")),
                ("FONTSIZE", (0,0), (-1,-1), 8.5),
                ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                ("PADDING", (0,0), (-1,-1), 4),
                ("INNERGRID", (0,0), (-1,-1), 0.25, colors.HexColor("#e8e3d6")),
            ]))
            story.append(trt)
            story.append(Spacer(1, 4))
        recent = ssum.get("recent_events") or []
        if recent:
            story.append(Paragraph("Recent critical events", muted))
            for ev in recent[:3]:
                story.append(Paragraph(
                    f"• <b>{(ev.get('severity') or '').upper()}</b> — "
                    f"{(ev.get('title') or ev.get('event_type') or '')[:140]}", body))
        alt = ssum.get("key_alt_rec")
        if alt:
            story.append(Spacer(1, 4))
            story.append(Paragraph(
                f"<b>Alt supplier rec:</b> {alt.get('sku', '—')} → "
                f"<b>{alt.get('name', '—')}</b> "
                f"(switch {_fmt_money(alt.get('switching_cost_usd'))}, "
                f"lead Δ {alt.get('lead_time_delta_days', '—')}d)",
                body))
        story.append(Spacer(1, 8))

    # Recommendation
    story.append(Paragraph("Recommendation", h2))
    rec_text = {
        "go": "Proceed to IC vote. Draft term sheet with reserve for follow-on.",
        "hold": "Continue diligence. Surface 3 outstanding risks before re-scoring.",
        "pass": "Decline. Archive in pipeline w/ trigger conditions for revisit.",
    }.get((base.get("verdict") or "hold").lower(), "Continue diligence.")
    story.append(Paragraph(rec_text, body))

    # Footer
    story.append(Spacer(1, 12))
    story.append(Paragraph(
        f"Confidential — internal use only · Deal ID {deal['id'][:8]} · "
        f"Project {deal.get('project_slug','')}", muted))

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()


# ---------- PPTX ----------

def generate_pptx(deal_id: str) -> bytes:
    """6-slide IC deck. Returns bytes."""
    data = _fetch_deal(deal_id)
    if not data:
        raise ValueError(f"deal {deal_id} not found")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as e:
        raise RuntimeError(f"python-pptx unavailable: {e}")

    deal = data["deal"]
    scens = data["scenarios"]
    base = scens[0] if scens else {}

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BG_DARK = RGBColor(0x1A, 0x16, 0x14)
    FG_LIGHT = RGBColor(0xE8, 0xE3, 0xD6)
    CORAL = RGBColor(0xC9, 0x63, 0x42)
    INK = RGBColor(0x2C, 0x2A, 0x26)
    GREEN = RGBColor(0x2C, 0x7A, 0x3F)
    RED = RGBColor(0xC0, 0x39, 0x2B)
    MUTED = RGBColor(0x6B, 0x65, 0x57)

    blank = prs.slide_layouts[6]

    def _set_bg(slide, color):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = color

    def _add_text(slide, text, left, top, width, height, *,
                  size=18, color=INK, bold=False, align=None):
        from pptx.enum.text import PP_ALIGN
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align == "center": p.alignment = PP_ALIGN.CENTER
        if align == "right":  p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def _verdict_rgb(v):
        v = (v or "hold").lower()
        if v == "go": return GREEN
        if v == "pass": return RED
        return CORAL

    # --- Slide 1: Cover ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG_DARK)
    _add_text(s, "IC MEMO", Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
              size=14, color=CORAL, bold=True)
    _add_text(s, deal["name"], Inches(0.8), Inches(1.2), Inches(12), Inches(1.5),
              size=44, color=FG_LIGHT, bold=True)
    _add_text(s,
              f"{deal.get('stage','').replace('_',' ').title()} · "
              f"{deal.get('sector','—')} · {deal.get('geography','—')}",
              Inches(0.8), Inches(3.0), Inches(12), Inches(0.5),
              size=20, color=FG_LIGHT)
    _add_text(s, f"Ask {_fmt_money(deal.get('ask_amount'))}",
              Inches(0.8), Inches(3.8), Inches(12), Inches(0.6),
              size=32, color=CORAL, bold=True)
    _add_text(s, f"Confidential · {datetime.utcnow().strftime('%Y-%m-%d')}",
              Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
              size=10, color=MUTED)

    # --- Slide 2: Verdict + KPIs ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, FG_LIGHT)
    _add_text(s, "Verdict", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
              size=14, color=CORAL, bold=True)
    verdict = (base.get("verdict") or "hold").upper()
    _add_text(s, verdict, Inches(0.6), Inches(1.0), Inches(12), Inches(1.5),
              size=72, color=_verdict_rgb(base.get("verdict")), bold=True)

    # KPI strip
    kpis = [
        ("IRR", _fmt_pct(base.get("irr"))),
        ("MOIC", f"{base.get('moic',0):.2f}×" if base.get("moic") is not None else "—"),
        ("NPV", _fmt_money(base.get("npv"))),
        ("Payback", f"{base.get('payback_yrs','—')} yrs"),
    ]
    x0 = 0.6
    w = 3.0
    gap = 0.15
    for i, (label, val) in enumerate(kpis):
        left = Inches(x0 + i*(w + gap))
        _add_text(s, label, left, Inches(4.5), Inches(w), Inches(0.4),
                  size=11, color=MUTED, bold=True)
        _add_text(s, val, left, Inches(4.9), Inches(w), Inches(1.1),
                  size=36, color=INK, bold=True)

    # --- Slide 3: Investment thesis ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, FG_LIGHT)
    _add_text(s, "Investment thesis", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
              size=14, color=CORAL, bold=True)
    thesis = (
        f"{deal['name']} — {deal.get('stage','').replace('_',' ')} round in "
        f"{deal.get('sector','—')}, geography {deal.get('geography','—')}.\n\n"
        f"Ask: {_fmt_money(deal.get('ask_amount'))}\n"
        f"Pre-money: {_fmt_money(deal.get('pre_money'))}\n"
        f"Post-money: {_fmt_money(deal.get('post_money'))}\n"
    )
    _add_text(s, thesis, Inches(0.6), Inches(1.2), Inches(12), Inches(5),
              size=18, color=INK)

    # --- Slide 4: Scenarios ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, FG_LIGHT)
    _add_text(s, "Scenarios", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
              size=14, color=CORAL, bold=True)
    if scens:
        rows = len(scens[:6]) + 1
        cols = 6
        tbl = s.shapes.add_table(rows, cols,
                                  Inches(0.6), Inches(1.2),
                                  Inches(12), Inches(0.5 * rows + 0.5)).table
        headers = ["Scenario", "IRR", "MOIC", "Payback", "NPV", "Verdict"]
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = h
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(12)
                    r.font.color.rgb = FG_LIGHT
            c.fill.solid(); c.fill.fore_color.rgb = BG_DARK
        for i, sc in enumerate(scens[:6], start=1):
            vals = [
                sc.get("name","—"),
                _fmt_pct(sc.get("irr")),
                f"{sc.get('moic',0):.2f}×" if sc.get("moic") is not None else "—",
                f"{sc.get('payback_yrs','—')} yrs",
                _fmt_money(sc.get("npv")),
                (sc.get("verdict") or "hold").upper(),
            ]
            for j, v in enumerate(vals):
                c = tbl.cell(i, j)
                c.text = str(v)
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(11)
                        r.font.color.rgb = INK
    else:
        _add_text(s, "No scenarios recorded yet.", Inches(0.6), Inches(1.5), Inches(12), Inches(1),
                  size=18, color=MUTED)

    # --- Slide 4b: Sensitivity analysis (only if heatmap renders) ---
    try:
        sens_bytes = _render_sensitivity_heatmap(scens)
        if sens_bytes:
            s = prs.slides.add_slide(blank)
            _set_bg(s, FG_LIGHT)
            _add_text(s, "Sensitivity analysis", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                      size=14, color=CORAL, bold=True)
            s.shapes.add_picture(io.BytesIO(sens_bytes), Inches(1.5), Inches(1.2), width=Inches(10))
    except Exception as e:
        logger.debug(f"pptx sensitivity slide failed: {e}")

    # --- Slide 5: Risks ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, FG_LIGHT)
    _add_text(s, "Key risks", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
              size=14, color=CORAL, bold=True)
    risks = [
        "Regulatory — local permits + compliance lead-time",
        "Adoption — market readiness + customer acquisition velocity",
        "Competition — incumbent moat + capital intensity",
        "Execution — team experience + delivery track record",
    ]
    body = "\n\n".join(f"• {r}" for r in risks)
    _add_text(s, body, Inches(0.6), Inches(1.2), Inches(12), Inches(5),
              size=20, color=INK)

    # --- Slide 5b: Market (Market Sentinel — sprint 2). Skip if empty. ---
    try:
        from dash.tools.market_tools import summarize_market_for_memo
        msum = summarize_market_for_memo(deal['id'])
    except Exception as _e:
        logger.debug(f"pptx market summary skipped: {_e}")
        msum = None
    if msum and msum.get("ok") and not msum.get("empty"):
        s = prs.slides.add_slide(blank)
        _set_bg(s, FG_LIGHT)
        _add_text(s, "Market", Inches(0.6), Inches(0.4), Inches(12),
                  Inches(0.6), size=14, color=CORAL, bold=True)
        # TAM/SAM/SOM big numbers row.
        cards = [
            ("TAM", _fmt_money(msum.get("tam_usd"))),
            ("SAM", _fmt_money(msum.get("sam_usd"))),
            ("SOM", _fmt_money(msum.get("som_usd"))),
        ]
        x0 = 0.6; w = 4.0; gap = 0.15
        for i, (lab, val) in enumerate(cards):
            left = Inches(x0 + i*(w + gap))
            _add_text(s, lab, left, Inches(1.1), Inches(w), Inches(0.4),
                      size=12, color=MUTED, bold=True)
            _add_text(s, val, left, Inches(1.5), Inches(w), Inches(1.1),
                      size=36, color=INK, bold=True)
        assumptions = msum.get("assumptions") or {}
        meth = msum.get("methodology") or "—"
        conf = assumptions.get("confidence")
        meta_line = f"Method: {meth}"
        if conf is not None:
            meta_line += f"  ·  Confidence: {conf:.2f}"
        _add_text(s, meta_line, Inches(0.6), Inches(2.8), Inches(12),
                  Inches(0.4), size=12, color=MUTED)
        # Top competitors.
        comps = msum.get("top_competitors") or []
        if comps:
            _add_text(s, "Top competitors", Inches(0.6), Inches(3.4),
                      Inches(12), Inches(0.4), size=13, color=CORAL, bold=True)
            comp_lines = []
            for c in comps[:3]:
                name = c.get("name", "—")
                sp = c.get("share_pct")
                share_str = f"{sp:.1f}%" if sp is not None else "—"
                geo = c.get("geography") or "—"
                comp_lines.append(f"• {name}  —  {share_str}  ({geo})")
            comp_text = "\n".join(comp_lines)
            _add_text(s, comp_text, Inches(0.6), Inches(3.85), Inches(12),
                      Inches(2), size=16, color=INK)

    # --- Slide 5c: Supply Risk (Supply Chain Sentry — sprint 3). Skip if empty. ---
    ssum2 = None
    if summarize_supply_for_memo is not None:
        try:
            ssum2 = summarize_supply_for_memo(deal.get('project_slug') or deal['id'])
        except Exception as _e:
            logger.debug(f"pptx supply summary skipped: {_e}")
            ssum2 = None
    if ssum2 and not ssum2.get("empty"):
        s = prs.slides.add_slide(blank)
        _set_bg(s, FG_LIGHT)
        _add_text(s, "Supply Risk", Inches(0.6), Inches(0.4), Inches(12),
                  Inches(0.6), size=14, color=CORAL, bold=True)
        rollup = ssum2.get("rollup") or {}
        tiles = [
            ("GREEN", str(rollup.get("green", 0)), GREEN),
            ("YELLOW", str(rollup.get("yellow", 0)), CORAL),
            ("RED", str(rollup.get("red", 0)), RED),
        ]
        x0 = 0.6; w = 4.0; gap = 0.15
        for i, (lab, val, col) in enumerate(tiles):
            left = Inches(x0 + i*(w + gap))
            _add_text(s, lab, left, Inches(1.1), Inches(w), Inches(0.4),
                      size=12, color=MUTED, bold=True)
            _add_text(s, val, left, Inches(1.5), Inches(w), Inches(1.4),
                      size=48, color=col, bold=True)
        # Top risks (left col) + alt rec callout (right col)
        risks = ssum2.get("top_risks") or []
        if risks:
            _add_text(s, "Top at-risk suppliers", Inches(0.6), Inches(3.5),
                      Inches(6), Inches(0.4), size=13, color=CORAL, bold=True)
            risk_lines = []
            for r in risks[:3]:
                nm = r.get("name", "—")
                ctry = r.get("country", "—")
                sc = r.get("score")
                sc_str = f"{sc:.1f}" if sc is not None else "—"
                risk_lines.append(f"• {nm} ({ctry}) — score {sc_str}")
            _add_text(s, "\n".join(risk_lines), Inches(0.6), Inches(3.95),
                      Inches(6), Inches(2.5), size=14, color=INK)
        recent = ssum2.get("recent_events") or []
        if recent:
            _add_text(s, "Recent events", Inches(7), Inches(3.5),
                      Inches(6), Inches(0.4), size=13, color=CORAL, bold=True)
            ev_lines = []
            for ev in recent[:3]:
                sev = (ev.get("severity") or "").upper()
                title = (ev.get("title") or ev.get("event_type") or "")[:60]
                ev_lines.append(f"• [{sev}] {title}")
            _add_text(s, "\n".join(ev_lines), Inches(7), Inches(3.95),
                      Inches(6), Inches(2.5), size=14, color=INK)
        alt = ssum2.get("key_alt_rec")
        if alt:
            alt_text = (
                f"Alt rec: {alt.get('sku', '—')} → "
                f"{alt.get('name', '—')}  |  "
                f"switch ${alt.get('switching_cost_usd', 0):,.0f}  ·  "
                f"lead Δ {alt.get('lead_time_delta_days', '—')}d"
            )
            _add_text(s, alt_text, Inches(0.6), Inches(6.4), Inches(12),
                      Inches(0.5), size=14, color=CORAL, bold=True)

    # --- Slide 6: Recommendation ---
    s = prs.slides.add_slide(blank)
    _set_bg(s, BG_DARK)
    _add_text(s, "Recommendation", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
              size=14, color=CORAL, bold=True)
    rec = {
        "go": "PROCEED TO IC VOTE.\nDraft term sheet. Reserve 15% follow-on.",
        "hold": "CONTINUE DILIGENCE.\nSurface 3 outstanding risks before re-scoring.",
        "pass": "DECLINE.\nArchive in pipeline with trigger conditions for revisit.",
    }.get((base.get("verdict") or "hold").lower(), "CONTINUE DILIGENCE.")
    _add_text(s, rec, Inches(0.6), Inches(2.0), Inches(12), Inches(3),
              size=36, color=FG_LIGHT, bold=True)
    _add_text(s, f"Deal ID {deal['id'][:8]} · {deal.get('project_slug','')}",
              Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
              size=10, color=MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
