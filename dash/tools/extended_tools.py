"""Tools lifted from former Demo-OS agents. No agent wrappers, no HITL."""
from __future__ import annotations

import ast
import json
import logging
import operator as op
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

import httpx
from agno.tools import tool

log = logging.getLogger(__name__)
logger = logging.getLogger(__name__)


# --- doc/web ---

_HTTP_TIMEOUT = 20.0
_MAX_BYTES = 2_000_000  # 2 MB cap per fetch


@tool(name="fetch_llms_txt",
      description="Fetch and parse the /llms.txt index for a docs site (e.g. https://docs.agno.com). "
                  "Returns a list of section headings + links the site exposes for LLM consumption. "
                  "Args: url (str) — base URL without trailing /llms.txt")
def fetch_llms_txt(url: str) -> dict[str, Any]:
    base = url.rstrip("/")
    target = base if base.endswith("/llms.txt") else f"{base}/llms.txt"
    try:
        with httpx.Client(timeout=_HTTP_TIMEOUT, follow_redirects=True) as cli:
            r = cli.get(target)
            r.raise_for_status()
            text = r.text[:_MAX_BYTES]
    except Exception as e:
        return {"ok": False, "url": target, "error": str(e)}

    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for line in text.splitlines():
        s = line.strip()
        if s.startswith("#"):
            if current:
                sections.append(current)
            current = {"heading": s.lstrip("# ").strip(), "links": []}
        elif s.startswith("- ") and current is not None:
            current["links"].append(s[2:].strip())
    if current:
        sections.append(current)

    return {"ok": True, "url": target, "sections": sections, "raw_chars": len(text)}


@tool(name="web_search",
      description="Run a web search (Tavily/Brave/Perplexity fallback chain) for docs not in llms.txt. "
                  "Args: query (str)")
def web_search(query: str) -> dict[str, Any]:
    try:
        from dash.learning.web_search import search as _search
        resp = _search(query, max_results=5)
        return {
            "ok": resp.error is None,
            "query": query,
            "source": resp.source_type,
            "summary": getattr(resp, "summary", None),
            "results": [
                {"title": r.title, "url": r.url, "snippet": r.snippet}
                for r in (resp.results or [])
            ],
            "error": resp.error,
        }
    except Exception as e:
        return {"ok": False, "query": query, "error": f"web_search unavailable: {e}"}


@tool(name="parse_doc_url",
      description="Fetch a documentation URL and return markdown. Uses pymupdf4llm for PDFs, "
                  "BeautifulSoup for HTML. Args: url (str)")
def parse_doc_url(url: str) -> dict[str, Any]:
    try:
        with httpx.Client(timeout=_HTTP_TIMEOUT, follow_redirects=True) as cli:
            r = cli.get(url)
            r.raise_for_status()
            ctype = (r.headers.get("content-type") or "").lower()
            content = r.content[:_MAX_BYTES]
    except Exception as e:
        return {"ok": False, "url": url, "error": str(e)}

    if "pdf" in ctype or url.lower().endswith(".pdf"):
        try:
            import io
            import pymupdf4llm  # type: ignore
            import pymupdf  # type: ignore
            doc = pymupdf.open(stream=io.BytesIO(content), filetype="pdf")
            md = pymupdf4llm.to_markdown(doc)
            return {"ok": True, "url": url, "kind": "pdf", "markdown": md[:50_000]}
        except Exception as e:
            return {"ok": False, "url": url, "kind": "pdf", "error": f"pdf parse failed: {e}"}

    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(content, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        title = (soup.title.string.strip() if soup.title and soup.title.string else "")
        text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
        return {"ok": True, "url": url, "kind": "html", "title": title, "markdown": text[:50_000]}
    except Exception as e:
        return {"ok": False, "url": url, "kind": "html", "error": f"html parse failed: {e}"}


# --- reports ---

def _reports_dir(project_slug: Optional[str]) -> Path:
    base = Path(os.getenv("KNOWLEDGE_DIR", "knowledge"))
    slug = project_slug or "_global"
    out = base / slug / "reports"
    out.mkdir(parents=True, exist_ok=True)
    return out


def _ts() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")


@tool
def generate_pdf(title: str, sections: List[Dict[str, Any]], project_slug: Optional[str] = None) -> Dict[str, Any]:
    """sections=[{heading, body}]. Tries app/export.py, falls back to reportlab."""
    out_path = _reports_dir(project_slug) / f"{_ts()}_{title[:40].replace(' ', '_')}.pdf"
    # Try existing exporter.
    try:
        from app import export as _exp  # type: ignore
        for fn_name in ("generate_pdf", "make_pdf", "export_pdf"):
            fn = getattr(_exp, fn_name, None)
            if callable(fn):
                try:
                    fn(title=title, sections=sections, output_path=str(out_path))
                    if out_path.exists():
                        return {"ok": True, "file_path": str(out_path)}
                except Exception:
                    pass
    except Exception:
        pass

    # Fallback: reportlab.
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        doc = SimpleDocTemplate(str(out_path), pagesize=letter)
        styles = getSampleStyleSheet()
        flow = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
        for s in sections or []:
            flow.append(Paragraph(s.get("heading", ""), styles["Heading2"]))
            flow.append(Paragraph(s.get("body", ""), styles["BodyText"]))
            flow.append(Spacer(1, 12))
        doc.build(flow)
        return {"ok": True, "file_path": str(out_path)}
    except Exception as e:
        logger.exception("generate_pdf failed: %s", e)
        return {"ok": False, "error": str(e)}


@tool
def generate_pptx(title: str, slides: List[Dict[str, Any]], project_slug: Optional[str] = None) -> Dict[str, Any]:
    """slides=[{title, content}]. Tries app/export.py slide_agent, falls back to python-pptx."""
    out_path = _reports_dir(project_slug) / f"{_ts()}_{title[:40].replace(' ', '_')}.pptx"
    try:
        from app import export as _exp  # type: ignore
        for fn_name in ("slide_agent", "generate_pptx", "make_pptx", "export_pptx"):
            fn = getattr(_exp, fn_name, None)
            if callable(fn):
                try:
                    fn(title=title, slides=slides, output_path=str(out_path))
                    if out_path.exists():
                        return {"ok": True, "file_path": str(out_path)}
                except Exception:
                    pass
    except Exception:
        pass

    try:
        from pptx import Presentation
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        body_layout = prs.slide_layouts[1]
        s0 = prs.slides.add_slide(title_layout)
        s0.shapes.title.text = title
        for sl in slides or []:
            s = prs.slides.add_slide(body_layout)
            s.shapes.title.text = sl.get("title", "")
            try:
                s.placeholders[1].text = sl.get("content", "")
            except Exception:
                pass
        prs.save(str(out_path))
        return {"ok": True, "file_path": str(out_path)}
    except Exception as e:
        logger.exception("generate_pptx failed: %s", e)
        return {"ok": False, "error": str(e)}


@tool
def generate_csv(rows: List[Dict[str, Any]], filename: Optional[str] = None, project_slug: Optional[str] = None) -> Dict[str, Any]:
    """rows=[{col:val}]. Uses pandas to_csv."""
    try:
        import pandas as pd
        name = (filename or f"{_ts()}_report.csv").replace(" ", "_")
        if not name.endswith(".csv"):
            name += ".csv"
        out_path = _reports_dir(project_slug) / name
        df = pd.DataFrame(rows or [])
        df.to_csv(out_path, index=False)
        return {"ok": True, "file_path": str(out_path), "rows": len(df)}
    except Exception as e:
        logger.exception("generate_csv failed: %s", e)
        return {"ok": False, "error": str(e)}


# --- Safe calculator: AST-based, optional numexpr accelerator. ---
_AST_OPS = {
    ast.Add: op.add, ast.Sub: op.sub, ast.Mult: op.mul, ast.Div: op.truediv,
    ast.Mod: op.mod, ast.Pow: op.pow, ast.FloorDiv: op.floordiv,
    ast.UAdd: op.pos, ast.USub: op.neg,
}


def _safe_eval(node):
    if isinstance(node, ast.Num):
        return node.n
    if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
        return node.value
    if isinstance(node, ast.BinOp) and type(node.op) in _AST_OPS:
        return _AST_OPS[type(node.op)](_safe_eval(node.left), _safe_eval(node.right))
    if isinstance(node, ast.UnaryOp) and type(node.op) in _AST_OPS:
        return _AST_OPS[type(node.op)](_safe_eval(node.operand))
    raise ValueError(f"unsupported expression node: {type(node).__name__}")


@tool
def calculator(expression: str) -> Dict[str, Any]:
    """Safe math eval. Tries numexpr if installed; otherwise AST walk."""
    try:
        try:
            import numexpr  # type: ignore
            val = float(numexpr.evaluate(expression).item())
            return {"ok": True, "result": val}
        except Exception:
            pass
        tree = ast.parse(expression, mode="eval")
        return {"ok": True, "result": _safe_eval(tree.body)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# --- pii ---

# PII patterns
_PII_PATTERNS: dict[str, re.Pattern[str]] = {
    "email": re.compile(r"[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}"),
    "phone": re.compile(r"\b(?:\+?\d{1,3}[\s\-.]?)?\(?\d{3}\)?[\s\-.]?\d{3}[\s\-.]?\d{4}\b"),
    "ssn":   re.compile(r"\b\d{3}-\d{2}-\d{4}\b"),
    "credit_card": re.compile(r"\b(?:\d[ \-]?){13,16}\b"),
}


@tool(name="scan_for_pii",
      description="Scan text for PII patterns (email, phone, SSN, credit card). "
                  "Returns counts per category and the first match (redacted). Args: text (str)")
def scan_for_pii(text: str) -> dict[str, Any]:
    hits: dict[str, dict[str, Any]] = {}
    for label, pat in _PII_PATTERNS.items():
        matches = pat.findall(text or "")
        if matches:
            sample = matches[0]
            # Redact middle chars in the sample
            redacted = sample[:2] + "***" + sample[-2:] if len(sample) > 4 else "***"
            hits[label] = {"count": len(matches), "sample_redacted": redacted}
    return {"pii_found": bool(hits), "categories": hits}


# --- schedules ---

# 5-field cron regex fallback (minute hour day month dow). Permissive enough
# for `* / , -` and numeric ranges; rejects obvious garbage.
_CRON_FIELD = r"(\*|\d+(-\d+)?(/\d+)?(,\d+(-\d+)?(/\d+)?)*|\*/\d+)"
_CRON_RE = re.compile(rf"^\s*{_CRON_FIELD}(\s+{_CRON_FIELD}){{4}}\s*$")


def _valid_cron(expr: str) -> bool:
    try:
        from crontab import CronSlices  # type: ignore
        return bool(CronSlices.is_valid(expr))
    except Exception:
        return bool(_CRON_RE.match(expr or ""))


def _engine():
    from db.session import get_sql_engine
    return get_sql_engine()


@tool
def create_schedule(name: str, prompt: str, cron: str, project_slug: str) -> Dict[str, Any]:
    if not _valid_cron(cron):
        return {"ok": False, "error": f"invalid cron expression: {cron!r}"}
    try:
        from sqlalchemy import text
        eng = _engine()
        with eng.begin() as conn:
            row = conn.execute(
                text(
                    "INSERT INTO public.dash_schedules "
                    "(name, prompt, cron, project_slug, enabled, created_at, metadata) "
                    "VALUES (:n, :p, :c, :s, true, :ts, CAST(:m AS jsonb)) RETURNING id"
                ),
                {
                    "n": name, "p": prompt, "c": cron, "s": project_slug,
                    "ts": datetime.now(timezone.utc),
                    "m": json.dumps({"source": "scheduler_agent"}),
                },
            ).fetchone()
            return {"ok": True, "schedule_id": int(row[0]), "cron": cron}
    except Exception as e:
        logger.exception("create_schedule failed: %s", e)
        return {"ok": False, "error": str(e)}


@tool
def list_schedules(project_slug: str) -> Dict[str, Any]:
    try:
        from sqlalchemy import text
        eng = _engine()
        with eng.connect() as conn:
            rows = conn.execute(
                text(
                    "SELECT id, name, cron, enabled, created_at FROM public.dash_schedules "
                    "WHERE project_slug = :s ORDER BY id DESC LIMIT 100"
                ),
                {"s": project_slug},
            ).fetchall()
        return {
            "ok": True,
            "schedules": [
                {"id": int(r[0]), "name": r[1], "cron": r[2], "enabled": bool(r[3]),
                 "created_at": r[4].isoformat() if r[4] else None}
                for r in rows
            ],
        }
    except Exception as e:
        logger.exception("list_schedules failed: %s", e)
        return {"ok": False, "error": str(e)}


@tool
def delete_schedule(schedule_id: int) -> Dict[str, Any]:
    try:
        from sqlalchemy import text
        eng = _engine()
        with eng.begin() as conn:
            res = conn.execute(
                text("DELETE FROM public.dash_schedules WHERE id = :i"),
                {"i": int(schedule_id)},
            )
            return {"ok": True, "deleted": res.rowcount}
    except Exception as e:
        logger.exception("delete_schedule failed: %s", e)
        return {"ok": False, "error": str(e)}


@tool
def enable_schedule(schedule_id: int, enabled: bool) -> Dict[str, Any]:
    try:
        from sqlalchemy import text
        eng = _engine()
        with eng.begin() as conn:
            res = conn.execute(
                text("UPDATE public.dash_schedules SET enabled = :e WHERE id = :i"),
                {"e": bool(enabled), "i": int(schedule_id)},
            )
            return {"ok": True, "updated": res.rowcount, "enabled": bool(enabled)}
    except Exception as e:
        logger.exception("enable_schedule failed: %s", e)
        return {"ok": False, "error": str(e)}
