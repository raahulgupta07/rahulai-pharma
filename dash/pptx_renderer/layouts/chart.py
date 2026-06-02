"""Chart layout (mirrors build.js layoutChart).

Delegates chart drawing to dash.pptx_renderer.charts.add_chart (other agent).
This file's job: header, action_line, and a placeholder chart frame.
"""
from __future__ import annotations

from typing import TYPE_CHECKING

from ._common import (
    PALETTE,
    add_text_box, apply_theme_background, header, hex_to_rgbcolor,
    set_speaker_notes, subtitle_color, _theme_font,
)

if TYPE_CHECKING:  # pragma: no cover
    from pptx.slide import Slide
    from dash.pptx_renderer.themes import Theme


# Chart frame geometry — bumped down to clear title + action_line zones above.
# eyebrow:     0.30 h=0.22  (set in header())
# title:       0.55 h=0.55  (set in header(), shrinks to fit)
# action_line: 1.18 h=0.25
# chart:       1.55 h=3.45
# source:      5.05 h=0.25
CHART_FRAME = {"x": 0.5, "y": 1.55, "w": 9.0, "h": 3.45}


def render(slide: "Slide", data: dict, theme: "Theme") -> None:
    """Render a chart slide.

    data keys: title, eyebrow?, chart_type, x?: [str], y?: [num]
               OR series?: [{name, data}], action_line?, source?,
               chart_data?: {labels, series} (build.js-compat)
    """
    data = data or {}
    apply_theme_background(slide, theme)
    header(slide, data.get("eyebrow") or "", data.get("title") or "", theme)

    body_font = _theme_font(theme, "body_font", "Calibri")
    muted_hex = getattr(theme, "muted", None) or PALETTE["muted"]
    gray_hex = getattr(theme, "gray", None) or PALETTE["gray"]
    rose_hex = getattr(theme, "rose", None) or PALETTE["rose"]
    action_color = subtitle_color(theme)

    if data.get("action_line"):
        add_text_box(
            slide, 0.5, 1.18, 9, 0.25,
            str(data["action_line"]),
            font=body_font, size=11, italic=True,
            color=action_color,
        )

    # Delegate chart drawing — other agent owns dash.pptx_renderer.charts.
    # Wrap in try/except so any chart-side failure renders inline error text
    # instead of killing the whole slide/deck render.
    chart_drawn = False
    try:
        from dash.pptx_renderer.charts import add_chart  # type: ignore
    except ImportError:
        # charts module not yet shipped — render placeholder below
        add_chart = None  # type: ignore

    if add_chart is not None:
        try:
            add_chart(slide, data, theme, frame=CHART_FRAME)
            chart_drawn = True
        except Exception as exc:  # pragma: no cover - defensive
            add_text_box(
                slide, 0.5, 2.5, 9, 0.5,
                f"Chart render error: {exc}",
                font=body_font, size=10,
                color=hex_to_rgbcolor(rose_hex),
                align="center",
            )
            chart_drawn = True

    if not chart_drawn:
        # Friendly placeholder so deck still renders without charts module.
        add_text_box(
            slide,
            CHART_FRAME["x"], CHART_FRAME["y"] + (CHART_FRAME["h"] - 0.5) / 2,
            CHART_FRAME["w"], 0.5,
            "(chart pending: dash.pptx_renderer.charts not loaded)",
            font=body_font, size=12, italic=True,
            color=hex_to_rgbcolor(gray_hex),
            align="center",
        )

    # Source / footnote
    source = data.get("source")
    if source:
        add_text_box(
            slide, 0.5, 5.0, 9, 0.25,
            f"Source: {source}",
            font=body_font, size=8, italic=True,
            color=hex_to_rgbcolor(muted_hex),
        )

    set_speaker_notes(slide, data.get("speaker_notes"))
