"""Shared helpers for pptx layout renderers.

Mirrors helpers in dash/render_js/build.js (header / footer / sectionDivider /
gridCards). Uses python-pptx 1.x API. Coordinates in inches (matches build.js
SLIDE_W=10, SLIDE_H=5.625).
"""
from __future__ import annotations

from typing import TYPE_CHECKING, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

if TYPE_CHECKING:  # pragma: no cover
    from pptx.slide import Slide

    from dash.pptx_renderer.themes import Theme


# Slide geometry (matches build.js)
SLIDE_W = 10.0
SLIDE_H = 5.625


# ---- Palette tokens (mirror build.js C dict) -----------------------------------
PALETTE = {
    "ink":         "0B1220",
    "navy":        "0F172A",
    "slate":       "1E293B",
    "steel":       "334155",
    "gray":        "64748B",
    "muted":       "94A3B8",
    "mutedLight":  "CBD5E1",
    "divider":     "E2E8F0",
    "paper":       "F8FAFC",
    "white":       "FFFFFF",
    "teal":        "14B8A6",
    "tealDark":    "0D9488",
    "amber":       "F59E0B",
    "rose":        "F43F5E",
    "emerald":     "10B981",
    "cyan":        "06B6D4",
    "violet":      "8B5CF6",
    "sky":         "0EA5E9",
}


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def hex_to_rgbcolor(hex_str: str) -> RGBColor:
    """Convert hex string to RGBColor. Accepts '#abc123' or 'abc123'."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    h = hex_str.lstrip("#").strip()
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return RGBColor.from_string(h.upper())
    except Exception:
        return RGBColor(0, 0, 0)


def _theme_color(theme, attr: str, fallback_hex: str) -> RGBColor:
    """Pull a hex string off theme (if present) and convert to RGBColor."""
    val = None
    if theme is not None:
        val = getattr(theme, attr, None)
    return hex_to_rgbcolor(val or fallback_hex)


def _theme_font(theme, attr: str, fallback: str) -> str:
    val = None
    if theme is not None:
        val = getattr(theme, attr, None)
    return val or fallback


# ---------------------------------------------------------------------------
# Luminance + theme-aware text color helpers
# ---------------------------------------------------------------------------

def _hex_luminance(hex_str: str) -> float:
    """Return perceived luminance 0..1 for a hex color.

    Empty/invalid → 1.0 (assume light, so callers fall back to ink text).
    """
    if not hex_str:
        return 1.0
    h = hex_str.lstrip("#").strip()
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return 1.0
    try:
        r = int(h[0:2], 16) / 255.0
        g = int(h[2:4], 16) / 255.0
        b = int(h[4:6], 16) / 255.0
    except Exception:
        return 1.0
    # Rec. 601 luma
    return 0.299 * r + 0.587 * g + 0.114 * b


def bg_is_dark(theme) -> bool:
    """True when theme background luminance < 0.5 (dark theme)."""
    if theme is None:
        return True
    bg = getattr(theme, "bg", None)
    return _hex_luminance(bg) < 0.5


def title_color(theme) -> RGBColor:
    """Theme-aware title color: white on dark bg, theme.ink on light bg."""
    if bg_is_dark(theme):
        return hex_to_rgbcolor(getattr(theme, "white", None) or PALETTE["white"])
    return hex_to_rgbcolor(getattr(theme, "ink", None) or PALETTE["ink"])


def _lighten_hex(hex_str: str, factor: float = 0.45) -> str:
    """Blend a hex color toward white by `factor` (0=unchanged, 1=white).
    Used to brighten dim secondary text so it stays readable on dark cards."""
    h = (hex_str or "").lstrip("#")
    if len(h) != 6:
        return hex_str
    try:
        r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
    except Exception:
        return hex_str
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"{r:02X}{g:02X}{b:02X}"


def subtitle_color(theme) -> RGBColor:
    """Theme-aware subtitle color. On dark themes, brighten ink_soft toward white
    (raw ink_soft is too dim against dark bg_alt cards). On light themes use ink_soft."""
    if bg_is_dark(theme):
        base = getattr(theme, "ink_soft", None) or PALETTE["mutedLight"]
        return hex_to_rgbcolor(_lighten_hex(base, 0.45))
    return hex_to_rgbcolor(
        getattr(theme, "ink_soft", None)
        or getattr(theme, "muted", None)
        or PALETTE["muted"]
    )


# ---------------------------------------------------------------------------
# Background helper
# ---------------------------------------------------------------------------

def set_slide_background(slide: "Slide", hex_color: str) -> None:
    """Force a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def apply_theme_background(slide: "Slide", theme) -> None:
    """Set slide background to theme.bg (or fall back to PALETTE white).

    Theme-aware shortcut used by every layout's render() so non-cover/section
    slides don't render with the default white PPT background (which makes
    white titles invisible on dark themes like midnight_executive).
    """
    bg_hex = None
    if theme is not None:
        bg_hex = getattr(theme, "bg", None)
    set_slide_background(slide, bg_hex or PALETTE["white"])


# ---------------------------------------------------------------------------
# Text box helper
# ---------------------------------------------------------------------------

_ALIGN_MAP = {
    "left":   PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right":  PP_ALIGN.RIGHT,
}


def responsive_size(text: str, base_size: int, min_size: int = 14,
                    long_threshold: int = 30, very_long_threshold: int = 60,
                    huge_threshold: int = 90) -> int:
    """Shrink font size based on text length.
    base_size is the size for short text (<long_threshold chars).
    Returns at least min_size."""
    n = len(text or "")
    if n <= long_threshold:
        return base_size
    if n <= very_long_threshold:
        return max(min_size, int(base_size * 0.72))
    if n <= huge_threshold:
        return max(min_size, int(base_size * 0.55))
    return max(min_size, int(base_size * 0.42))


def add_text_box(
    slide: "Slide",
    x: float, y: float, w: float, h: float,
    text: str,
    *,
    font: Optional[str] = None,
    size: int = 11,
    color: Optional[RGBColor] = None,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    valign: Optional[str] = None,
    char_spacing: Optional[float] = None,
    auto_shrink: bool = False,
) -> None:
    """Drop a single-paragraph text box on the slide.

    Coordinates in inches.
    auto_shrink: enable PowerPoint's "shrink text on overflow" autofit.
    """
    from pptx.enum.text import MSO_AUTO_SIZE  # local import
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    if auto_shrink:
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE_OR_SHRINK_TEXT
        except Exception:
            pass
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if font:
        run.font.name = font
    if color is not None:
        run.font.color.rgb = color
    # char_spacing on python-pptx requires manipulating XML; skip silently.
    _ = char_spacing


# ---------------------------------------------------------------------------
# Filled rectangle helper
# ---------------------------------------------------------------------------

def add_rect(
    slide: "Slide",
    x: float, y: float, w: float, h: float,
    fill_hex: str,
    line_hex: Optional[str] = None,
    line_width_pt: float = 0.0,
) -> None:
    """Add a filled rectangle. Use None/empty line_hex to remove line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgbcolor(fill_hex)
    line = shape.line
    if line_hex:
        line.color.rgb = hex_to_rgbcolor(line_hex)
        if line_width_pt:
            line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    # Strip default text frame margins so shape stays clean
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Inches(0)
        tf.margin_top = tf.margin_bottom = Inches(0)
    return shape


# ---------------------------------------------------------------------------
# header() / footer() — match build.js geometry
# ---------------------------------------------------------------------------

def header(slide: "Slide", eyebrow: str, title: str, theme) -> None:
    """Top-of-slide eyebrow + title block (matches build.js header()).

    Title uses theme-aware `title_color(theme)` so it stays readable on both
    dark and light theme backgrounds.
    """
    title_font = _theme_font(theme, "title_font", "Georgia")
    body_font = _theme_font(theme, "body_font", "Calibri")
    accent = _theme_color(theme, "accent", PALETTE["teal"])
    t_color = title_color(theme)

    if eyebrow:
        add_text_box(
            slide, 0.5, 0.35, 9, 0.28,
            str(eyebrow).upper(),
            font=body_font, size=10, color=accent, bold=True,
            align="left",
        )
    if title:
        # Dynamic size: shrink long titles to avoid overflow into chart frame.
        title_size = responsive_size(
            str(title), base_size=24, min_size=14,
            long_threshold=40, very_long_threshold=80, huge_threshold=120,
        )
        add_text_box(
            slide, 0.5, 0.62, 9, 0.55,
            title,
            font=title_font, size=title_size, color=t_color, bold=True,
            align="left",
            auto_shrink=True,
        )


def footer(slide: "Slide", brand: str, page_no: int, total: int, theme) -> None:
    """Bottom-of-slide brand + page number (matches build.js footer())."""
    body_font = _theme_font(theme, "body_font", "Calibri")
    muted = _theme_color(theme, "muted", PALETTE["muted"])

    add_text_box(
        slide, 0.5, 5.30, 6, 0.25,
        str(brand or ""),
        font=body_font, size=9, color=muted,
        align="left",
    )
    if total:
        page_text = f"{page_no} / {total}"
    else:
        page_text = str(page_no)
    add_text_box(
        slide, 9.0, 5.30, 0.5, 0.25,
        page_text,
        font=body_font, size=9, color=muted,
        align="right",
    )


# ---------------------------------------------------------------------------
# Section divider (used by section_divider.py + cover/closing)
# ---------------------------------------------------------------------------

def section_divider(
    slide: "Slide",
    part: str, title: str, subtitle: str,
    theme,
) -> None:
    """Section divider with theme-aware background + accent rail.

    Uses theme.bg (not theme.ink) so light themes get a cream/white background
    and dark themes keep their dark background. Title/subtitle colors derive
    from luminance check on bg.
    """
    bg_hex = getattr(theme, "bg", None) or PALETTE["ink"]
    accent_hex = getattr(theme, "accent", None) or PALETTE["teal"]

    set_slide_background(slide, bg_hex)
    add_rect(slide, 0, 0, 0.18, SLIDE_H, accent_hex)

    title_font = _theme_font(theme, "title_font", "Georgia")
    body_font = _theme_font(theme, "body_font", "Calibri")

    t_color = title_color(theme)
    sub_color = subtitle_color(theme)

    if part:
        add_text_box(
            slide, 0.6, 1.2, 8.8, 0.4,
            str(part).upper(),
            font=body_font, size=12, color=hex_to_rgbcolor(accent_hex),
            bold=True,
        )
    add_text_box(
        slide, 0.6, 1.8, 8.8, 1.6,
        title or "",
        font=title_font, size=54, color=t_color,
        bold=True,
    )
    if subtitle:
        add_text_box(
            slide, 0.6, 3.55, 8.8, 0.6,
            subtitle,
            font=body_font, size=16, color=sub_color,
            italic=True,
        )


# ---------------------------------------------------------------------------
# Speaker notes helper
# ---------------------------------------------------------------------------

def set_speaker_notes(slide: "Slide", text: str) -> None:
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(text)
    except Exception:
        pass
