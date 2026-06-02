"""Table layout (mirrors build.js layoutTable)."""
from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.util import Inches, Pt

from ._common import (
    PALETTE,
    add_text_box, apply_theme_background, bg_is_dark, header, hex_to_rgbcolor,
    set_speaker_notes, subtitle_color, title_color, _theme_font,
)
from dash.pptx_renderer.markdown_runs import render_markdown_text

if TYPE_CHECKING:  # pragma: no cover
    from pptx.slide import Slide
    from dash.pptx_renderer.themes import Theme


def render(slide: "Slide", data: dict, theme: "Theme") -> None:
    """Render a table slide.

    data keys: title, eyebrow?, headers/columns: [str], rows: [[str]],
               action_line?, note?
    """
    data = data or {}
    apply_theme_background(slide, theme)
    header(slide, data.get("eyebrow") or "", data.get("title") or "", theme)

    cols = data.get("headers") or data.get("columns") or []
    rows = data.get("rows") or []
    if not isinstance(cols, list) or not cols:
        return
    if not isinstance(rows, list):
        rows = []

    body_font = _theme_font(theme, "body_font", "Calibri")
    mono_font = _theme_font(theme, "mono_font", "JetBrains Mono")
    # Theme-aware fills
    dark = bg_is_dark(theme)
    # Header row fill: slate works on both themes (dark slate stripe reads
    # fine on dark bg via theme.bg_alt and on light cream bg via slate).
    header_fill_hex = (
        getattr(theme, "slate", None) or PALETTE["slate"]
    ) if not dark else (
        getattr(theme, "bg_alt", None) or getattr(theme, "slate", None) or PALETTE["slate"]
    )
    # Banded body rows: alternate bg / bg_alt so banding stays visible on dark
    band_a_hex = getattr(theme, "bg_alt", None) or getattr(theme, "paper", None) or PALETTE["paper"]
    band_b_hex = getattr(theme, "bg", None) or getattr(theme, "white", None) or PALETTE["white"]
    # Header text: always bright white for contrast on slate fill
    header_text_hex = getattr(theme, "white", None) or PALETTE["white"]
    # Body cell text: theme-aware (ink on light, white on dark)
    body_text_hex = (
        getattr(theme, "white", None) or PALETTE["white"]
    ) if dark else (
        getattr(theme, "ink", None) or PALETTE["ink"]
    )
    muted_hex = getattr(theme, "muted", None) or PALETTE["muted"]
    action_color = subtitle_color(theme)

    if data.get("action_line"):
        add_text_box(
            slide, 0.5, 1.18, 9, 0.25,
            str(data["action_line"]),
            font=body_font, size=11, italic=True,
            color=action_color,
        )

    n_cols = len(cols)
    n_rows = 1 + len(rows)  # +1 for header

    # Place table
    x, y, w = 0.5, 1.4, 9.0
    row_h = 0.32
    # Cap rows so table fits roughly within 3.6 vertical inches
    max_body_rows = max(0, int((5.0 - y) / row_h) - 1)
    if len(rows) > max_body_rows:
        rows = rows[:max_body_rows]
        n_rows = 1 + len(rows)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows),
    )
    tbl = table_shape.table

    # Header row
    for ci, col_name in enumerate(cols):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgbcolor(header_fill_hex)
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        header_text = "" if col_name is None else str(col_name)
        render_markdown_text(
            tf,
            header_text,
            base_font=body_font,
            base_size=11,
            base_color=header_text_hex,
            mono_font=mono_font,
        )
        # Force bold on all header runs (headers are visually bold regardless
        # of markdown). render_markdown_text bolds only **wrapped** spans.
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.bold = True

    # Body rows w/ banded fill
    for ri, row in enumerate(rows):
        banded = band_a_hex if (ri % 2 == 0) else band_b_hex
        for ci in range(n_cols):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgbcolor(banded)
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            val = row[ci] if isinstance(row, (list, tuple)) and ci < len(row) else ""
            cell_text = "" if val is None else str(val)
            render_markdown_text(
                tf,
                cell_text,
                base_font=body_font,
                base_size=10,
                base_color=body_text_hex,
                mono_font=mono_font,
            )

    # Optional footnote
    if data.get("note"):
        add_text_box(
            slide, 0.5, 5.0, 9, 0.25,
            str(data["note"]),
            font=body_font, size=9, italic=True,
            color=hex_to_rgbcolor(muted_hex),
        )

    set_speaker_notes(slide, data.get("speaker_notes"))
