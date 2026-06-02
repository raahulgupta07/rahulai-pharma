"""Top-level orchestrator for the native python-pptx deck renderer.

Wires theme + layouts + charts together to produce a complete ``.pptx`` file
without invoking the Node sidecar/subprocess. Callers opt in via the
``PPTX_ENGINE=native`` env switch in ``dash.tools.render_pptxgenjs``.

Slide geometry mirrors ``dash/render_js/build.js``:
    SLIDE_W = 10.0 in   SLIDE_H = 5.625 in   (16:9)
"""
from __future__ import annotations

import io
import logging
import os
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches


log = logging.getLogger(__name__)


# Slide geometry — must match build.js / layouts/_common.py
_SLIDE_W_IN = 10.0
_SLIDE_H_IN = 5.625

# Blank layout index in the default python-pptx template
_BLANK_LAYOUT_IDX = 6


__all__ = ["render_to_path", "render_to_bytes"]


def _new_presentation(theme: Any) -> Presentation:
    """Build an empty 16:9 ``Presentation`` with the theme applied."""
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)

    # apply_theme is owned by the themes agent; we just call it.
    try:
        from dash.pptx_renderer.themes import apply_theme  # type: ignore
        apply_theme(prs, theme)
    except Exception as e:  # pragma: no cover - defensive
        log.warning("apply_theme failed (continuing without theme cascade): %s", e)

    return prs


def _blank_layout(prs: Presentation) -> Any:
    """Return the blank slide layout, falling back to last layout if missing."""
    layouts = prs.slide_layouts
    try:
        return layouts[_BLANK_LAYOUT_IDX]
    except (IndexError, KeyError):
        # python-pptx default template has 11 layouts; index 6 = Blank.
        # Fall back gracefully if a custom template was loaded.
        return layouts[len(layouts) - 1]


def _render_verified_badge(slide: Any) -> None:
    """Draw a small coral ✓ verified badge in the bottom-right corner.

    Phase 2 — truth-grounded slides. Coral #c96342, 9pt, sits just above the
    footer baseline so it doesn't collide with brand / page-number.
    """
    try:
        from dash.pptx_renderer.layouts._common import add_text_box, hex_to_rgbcolor  # type: ignore
        # Bottom-right, just above the standard footer (which lives at y=5.30).
        add_text_box(
            slide,
            7.4, 5.05, 2.45, 0.22,
            "✓ verified vs pinned metric",
            font="Calibri", size=9,
            color=hex_to_rgbcolor("c96342"),
            bold=True, align="right",
        )
    except Exception as e:  # pragma: no cover - never break a render
        log.warning("verified badge render failed: %s", e)


def _render_slides(prs: Presentation, spec: Dict[str, Any], theme: Any) -> None:
    """Iterate spec['slides'] and delegate each to layouts.render_layout."""
    from dash.pptx_renderer.layouts import render_layout  # type: ignore

    slides = spec.get("slides") or []
    if not isinstance(slides, list):
        raise ValueError("spec['slides'] must be a list")

    blank = _blank_layout(prs)
    for idx, slide_spec in enumerate(slides):
        if not isinstance(slide_spec, dict):
            log.warning("skipping non-dict slide at index %d", idx)
            continue
        layout_name = slide_spec.get("layout") or "content_grid"
        slide = prs.slides.add_slide(blank)
        try:
            render_layout(layout_name, slide, slide_spec, theme)
        except Exception as e:
            log.exception(
                "render_layout failed (layout=%s, idx=%d): %s",
                layout_name, idx, e,
            )
            # Continue with remaining slides rather than aborting the deck.
        # Phase 2 — overlay verified badge for truth-grounded slides
        if slide_spec.get("verified"):
            _render_verified_badge(slide)


def _resolve_theme(theme_name: str) -> Any:
    from dash.pptx_renderer.themes import get_theme  # type: ignore
    return get_theme(theme_name)


def render_to_path(
    spec: Dict[str, Any],
    output_path: str,
    *,
    theme_name: str = "midnight_executive",
    theme_override: Any = None,
) -> str:
    """Render a deck spec to ``output_path``. Returns the absolute path.

    theme_override: pre-built Theme object (e.g. from AI stylist with
                    accent/palette overrides). When provided, wins over
                    theme_name lookup.
    """
    if not isinstance(spec, dict):
        raise ValueError("spec must be a dict")
    if not output_path:
        raise ValueError("output_path must be a non-empty path")

    out_abs = os.path.abspath(output_path)
    out_dir = os.path.dirname(out_abs)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    theme = theme_override or _resolve_theme(theme_name)
    prs = _new_presentation(theme)
    _render_slides(prs, spec, theme)
    prs.save(out_abs)
    return out_abs


def render_to_bytes(
    spec: Dict[str, Any],
    *,
    theme_name: str = "midnight_executive",
) -> bytes:
    """Render a deck spec and return the ``.pptx`` bytes (no disk I/O)."""
    if not isinstance(spec, dict):
        raise ValueError("spec must be a dict")
    theme = _resolve_theme(theme_name)
    prs = _new_presentation(theme)
    _render_slides(prs, spec, theme)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
