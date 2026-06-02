"""Chart-rendering helpers for python-pptx PPTX renderer.

Replicates the 8 chart types from `dash/render_js/build.js` (pptxgenjs) using
python-pptx native charts where supported, and raw XML for candlestick.

Public API: `add_chart(slide, chart_type, x, y, w, h, *, ...) -> chart | None`
Registry:   `CHART_REGISTRY` — dispatch by chart_type string (+ aliases).
"""
from __future__ import annotations

import logging
from typing import Any

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.slide import Slide
from pptx.util import Inches

log = logging.getLogger(__name__)

# Color-blind safe default palette (no leading '#').
DEFAULT_PALETTE = ["c96342", "3a8dff", "10b981", "f59e0b", "9b6dff"]


def _norm_hex(c: str | None) -> str | None:
    """Strip leading '#' / spaces. Return None if invalid."""
    if not c:
        return None
    s = str(c).strip().lstrip("#")
    if len(s) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in s):
        return s.lower()
    return None


def _palette_for(series: list[dict], theme: Any) -> list[str]:
    """Resolve final hex color per series. Series color > theme.palette > DEFAULT."""
    theme_palette: list[str] = []
    if theme is not None:
        try:
            theme_palette = [_norm_hex(c) or "" for c in (getattr(theme, "palette", []) or [])]
            theme_palette = [c for c in theme_palette if c]
        except Exception:
            theme_palette = []
    base = theme_palette or DEFAULT_PALETTE
    out: list[str] = []
    for i, s in enumerate(series):
        explicit = _norm_hex(s.get("color") if isinstance(s, dict) else None)
        out.append(explicit or base[i % len(base)])
    return out


def _apply_series_color(chart_series, hex_color: str) -> None:
    """Best-effort fill color on a chart series."""
    try:
        fill = chart_series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color)
    except Exception as exc:  # legacy chart types may not support .format.fill
        log.debug("series color apply failed (%s): %s", hex_color, exc)


def _apply_common(chart, *, title: str, show_legend: bool, show_data_labels: bool) -> None:
    try:
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
    except Exception:
        pass
    try:
        chart.has_legend = bool(show_legend)
        if show_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
    except Exception:
        pass
    if show_data_labels:
        try:
            for s in chart.series:
                s.data_labels.show_value = True
        except Exception:
            pass


def _build_category_data(categories: list[str], series: list[dict]) -> CategoryChartData:
    cd = CategoryChartData()
    cd.categories = [str(c) for c in (categories or [])]
    for s in series or []:
        name = str(s.get("name", "Series"))
        values = [float(v) if v is not None else 0.0 for v in (s.get("data") or [])]
        cd.add_series(name, values)
    return cd


def _add_category_chart(
    slide: Slide,
    xl_type: XL_CHART_TYPE,
    x: float, y: float, w: float, h: float,
    *,
    categories: list[str],
    series: list[dict],
    theme: Any,
    title: str,
    show_legend: bool,
    show_data_labels: bool,
):
    cd = _build_category_data(categories, series)
    gframe = slide.shapes.add_chart(xl_type, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gframe.chart
    colors = _palette_for(series, theme)
    try:
        for i, cs in enumerate(chart.series):
            if i < len(colors):
                _apply_series_color(cs, colors[i])
    except Exception as exc:
        log.debug("palette apply failed: %s", exc)
    _apply_common(chart, title=title, show_legend=show_legend, show_data_labels=show_data_labels)
    return chart


def _add_bar(slide, x, y, w, h, **kw):
    return _add_category_chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, **kw)


def _add_line(slide, x, y, w, h, **kw):
    return _add_category_chart(slide, XL_CHART_TYPE.LINE, x, y, w, h, **kw)


def _add_pie(slide, x, y, w, h, **kw):
    # Pie: collapse to single series if multiple were given.
    series = kw.get("series") or []
    if len(series) > 1:
        kw["series"] = [series[0]]
    chart = _add_category_chart(slide, XL_CHART_TYPE.PIE, x, y, w, h, **kw)
    # Color each data point individually for pie.
    try:
        colors = _palette_for([{"color": None} for _ in (kw.get("categories") or [])], kw.get("theme"))
        if chart.series:
            pts = chart.series[0].points
            for i, pt in enumerate(pts):
                if i < len(colors):
                    fill = pt.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(colors[i])
    except Exception as exc:
        log.debug("pie slice color apply failed: %s", exc)
    return chart


def _add_area(slide, x, y, w, h, **kw):
    return _add_category_chart(slide, XL_CHART_TYPE.AREA, x, y, w, h, **kw)


def _add_grouped_bar(slide, x, y, w, h, **kw):
    return _add_category_chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, **kw)


def _add_stacked_bar(slide, x, y, w, h, **kw):
    return _add_category_chart(slide, XL_CHART_TYPE.BAR_STACKED, x, y, w, h, **kw)


def _add_scatter(slide, x, y, w, h, *, categories, series, theme, title, show_legend, show_data_labels):
    """XY scatter. X = numeric index of category, Y = series.data."""
    cd = XyChartData()
    for s in series or []:
        name = str(s.get("name", "Series"))
        sub = cd.add_series(name)
        data = s.get("data") or []
        for i, v in enumerate(data):
            try:
                yv = float(v) if v is not None else 0.0
            except Exception:
                yv = 0.0
            # Use category position as x; if categories[i] is numeric, prefer that.
            xv: float = float(i)
            if categories and i < len(categories):
                try:
                    xv = float(categories[i])
                except Exception:
                    xv = float(i)
            sub.add_data_point(xv, yv)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    chart = gframe.chart
    colors = _palette_for(series, theme)
    try:
        for i, cs in enumerate(chart.series):
            if i < len(colors):
                _apply_series_color(cs, colors[i])
    except Exception:
        pass
    _apply_common(chart, title=title, show_legend=show_legend, show_data_labels=show_data_labels)
    return chart


def _add_candlestick(slide, x, y, w, h, *, categories, series, theme, title, show_legend,
                     show_data_labels, options=None):
    """Candlestick fallback: render a line chart of CLOSE prices.

    python-pptx 1.0.2 has no native stockChart support. Building raw OOXML
    via lxml is brittle (charts in pptx live in their own part with bespoke
    rels + embedded xlsx workbook). We fall back to a line chart of close
    prices so the deck never breaks. Document this clearly in the report.
    """
    options = options or {}
    ohlc = options.get("ohlc") or []
    if ohlc:
        cats = [str(p.get("date", i)) for i, p in enumerate(ohlc)]
        closes = [float(p.get("close", 0) or 0) for p in ohlc]
        series = [{"name": "Close", "data": closes, "color": None}]
        categories = cats
    log.warning("candlestick chart requested — falling back to line chart of close prices")
    return _add_line(
        slide, x, y, w, h,
        categories=categories, series=series, theme=theme,
        title=title or "Candlestick (close)", show_legend=show_legend,
        show_data_labels=show_data_labels,
    )


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------

CHART_REGISTRY: dict[str, Any] = {
    "bar": _add_bar,
    "line": _add_line,
    "pie": _add_pie,
    "scatter": _add_scatter,
    "area": _add_area,
    "grouped_bar": _add_grouped_bar,
    "stacked_bar": _add_stacked_bar,
    "candlestick": _add_candlestick,
    # aliases
    "column": _add_bar,
    "donut": _add_pie,
    "doughnut": _add_pie,
}


def _dispatch_chart(
    slide: Slide,
    chart_type: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    categories: list[str],
    series: list[dict],
    theme: Any = None,
    title: str = "",
    show_legend: bool = True,
    show_data_labels: bool = False,
    options: dict | None = None,
):
    """Core dispatch: route to registered chart fn, apply common safety."""
    key = (chart_type or "bar").strip().lower()
    fn = CHART_REGISTRY.get(key, _add_bar)
    kwargs: dict[str, Any] = dict(
        categories=categories or [],
        series=series or [],
        theme=theme,
        title=title or "",
        show_legend=show_legend,
        show_data_labels=show_data_labels,
    )
    if fn is _add_candlestick:
        kwargs["options"] = options or {}
    try:
        return fn(slide, x, y, w, h, **kwargs)
    except Exception as exc:
        log.exception("add_chart(%s) failed: %s", key, exc)
        return None


def add_chart(slide, *args, **kwargs):
    """Flexible chart adder. Supports two call styles:

    Style A (positional, original):
        add_chart(slide, chart_type, x, y, w, h, *,
                  categories, series, theme=None, title="",
                  show_legend=True, show_data_labels=False, options=None)

    Style B (data-dict + frame, used by layouts/chart.py):
        add_chart(slide, data: dict, theme=None, *, frame: dict = None)
      where `data` contains keys: chart_type, categories (or x),
      series (or y), title, action_line, etc.

    Returns the python-pptx Chart object or None.
    """
    # Detect call style: Style B if first positional arg is a dict
    if args and isinstance(args[0], dict):
        data = args[0]
        theme = args[1] if len(args) > 1 else kwargs.get("theme")
        frame = kwargs.get("frame") or {"x": 0.5, "y": 1.1, "w": 9.0, "h": 3.8}
        chart_type = data.get("chart_type", "bar")

        # Support build.js-compat shape: chart_data: {labels, series}
        chart_data = data.get("chart_data") or {}
        categories = (
            data.get("categories")
            or data.get("x")
            or chart_data.get("labels")
            or []
        )
        series_in = (
            data.get("series")
            or data.get("y")
            or chart_data.get("series")
            or []
        )

        # Normalize: if series is a flat list of numbers, wrap as single series.
        if series_in and not isinstance(series_in[0], dict):
            series = [{"name": data.get("title", "Series"), "data": list(series_in)}]
        else:
            series = list(series_in)

        return _dispatch_chart(
            slide,
            chart_type,
            float(frame.get("x", 0.5)),
            float(frame.get("y", 1.1)),
            float(frame.get("w", 9.0)),
            float(frame.get("h", 3.8)),
            categories=categories,
            series=series,
            theme=theme,
            # Slide layout's header() already prints title at top of slide.
            # Use chart_title (default empty) to avoid duplicate inside chart.
            title=data.get("chart_title", "") or "",
            show_legend=data.get("show_legend", True),
            show_data_labels=data.get("show_data_labels", False),
            options=data.get("options"),
        )

    # Style A: positional + kwargs (original signature)
    chart_type = args[0] if args else kwargs.pop("chart_type", "bar")
    x = args[1] if len(args) > 1 else kwargs.pop("x")
    y = args[2] if len(args) > 2 else kwargs.pop("y")
    w = args[3] if len(args) > 3 else kwargs.pop("w")
    h = args[4] if len(args) > 4 else kwargs.pop("h")
    # Strip frame if accidentally passed in Style A
    kwargs.pop("frame", None)
    return _dispatch_chart(slide, chart_type, x, y, w, h, **kwargs)
