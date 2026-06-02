"""Markdown → python-pptx text runs helper.

Parses a tiny subset of markdown (bold, italic, code) into a flat list of
run-dicts and applies them to a python-pptx TextFrame so that the rendered
slide preserves inline formatting.

This module is intentionally minimal — links, headers, and lists are NOT
supported here. Layout-level helpers handle bullets/headings.
"""
from __future__ import annotations

import re
from typing import Any, List, Optional

from pptx.dml.color import RGBColor
from pptx.util import Pt


__all__ = ["parse_markdown", "apply_runs", "render_markdown_text"]


# ---- Parser ---------------------------------------------------------------

# Tokens (longest first for correct alternation):
#   \\X   → escape (X is literal)
#   **    → bold delim
#   `     → code delim
#   *     → italic delim (also _)
#   _     → italic delim
_TOKEN_RE = re.compile(r"(\\.|\*\*|`|\*|_)")


def parse_markdown(text: str) -> List[dict]:
    """Parse a small subset of markdown to a flat list of run-dicts.

    Supports:
      **bold**           → bold
      *italic* / _italic_→ italic
      `code`             → monospace
      \\* \\_ \\`        → literal escape

    Returns: ``[{"text": str, "bold": bool, "italic": bool, "code": bool}, ...]``
    """
    if text is None:
        return []
    s = str(text)
    if not s:
        return []

    # State flags
    bold = False
    italic = False
    code = False

    runs: List[dict] = []
    buf: List[str] = []

    def flush() -> None:
        if not buf:
            return
        runs.append(
            {
                "text": "".join(buf),
                "bold": bold,
                "italic": italic,
                "code": code,
            }
        )
        buf.clear()

    pos = 0
    for m in _TOKEN_RE.finditer(s):
        start = m.start()
        if start > pos:
            buf.append(s[pos:start])
        tok = m.group(0)
        pos = m.end()

        # Escape: \\X → literal X
        if tok.startswith("\\") and len(tok) == 2:
            buf.append(tok[1])
            continue

        # Inside code, only ` closes it; everything else is literal.
        if code:
            if tok == "`":
                flush()
                code = False
            else:
                buf.append(tok)
            continue

        if tok == "`":
            flush()
            code = True
        elif tok == "**":
            flush()
            bold = not bold
        elif tok == "*" or tok == "_":
            flush()
            italic = not italic
        else:
            buf.append(tok)

    if pos < len(s):
        buf.append(s[pos:])
    flush()

    # Filter empty-text runs (artifacts of consecutive toggles).
    return [r for r in runs if r["text"]]


# ---- Applier --------------------------------------------------------------

def _set_color(font: Any, color_hex: Optional[str]) -> None:
    if not color_hex:
        return
    hx = color_hex.lstrip("#")
    if len(hx) != 6:
        return
    try:
        r = int(hx[0:2], 16)
        g = int(hx[2:4], 16)
        b = int(hx[4:6], 16)
        font.color.rgb = RGBColor(r, g, b)
    except (ValueError, AttributeError):
        pass


def apply_runs(
    text_frame: Any,
    runs: List[dict],
    *,
    base_font: str = "Inter",
    base_size: int = 14,
    base_color: Optional[str] = None,
    mono_font: str = "JetBrains Mono",
) -> None:
    """Apply parsed run-dicts to a python-pptx TextFrame.

    First run replaces ``text_frame.text``; subsequent runs are added via
    ``paragraph.add_run()``. Sets ``font.name``, ``font.bold``, ``font.italic``,
    ``font.size``, and (optionally) ``font.color.rgb`` per run.

    Code runs use ``mono_font`` and ``base_size - 1`` for size.
    """
    if text_frame is None:
        return
    if not runs:
        text_frame.text = ""
        return

    # Use clear() to reset to a single empty paragraph with NO runs.
    # Note: text_frame.text = "..." adds a default run that can interfere with
    # subsequent paragraph.add_run() calls — runs get silently dropped or
    # overwritten because the default run's XML state isn't fully initialized.
    try:
        text_frame.clear()
    except Exception:
        # Fallback for TextFrame impls without clear()
        text_frame.text = ""
    paragraph = text_frame.paragraphs[0]

    def _style(r_obj: Any, spec: dict) -> None:
        if r_obj is None:
            return
        font = r_obj.font
        is_code = bool(spec.get("code"))
        font.name = mono_font if is_code else base_font
        font.bold = bool(spec.get("bold"))
        font.italic = bool(spec.get("italic"))
        size = max(1, int(base_size) - 1) if is_code else int(base_size)
        try:
            font.size = Pt(size)
        except Exception:
            pass
        _set_color(font, base_color)

    for i, spec in enumerate(runs):
        text = spec.get("text", "")
        if i == 0 and paragraph.runs:
            # First run: reuse the empty Run that clear() may have left behind.
            r_obj = paragraph.runs[0]
            r_obj.text = text
        else:
            # Always add_run() then set .text — python-pptx canonical pattern.
            r_obj = paragraph.add_run()
            r_obj.text = text
        _style(r_obj, spec)


def render_markdown_text(text_frame: Any, text: str, **kwargs: Any) -> None:
    """Shortcut: parse markdown ``text`` and apply runs to ``text_frame``."""
    runs = parse_markdown(text)
    apply_runs(text_frame, runs, **kwargs)
